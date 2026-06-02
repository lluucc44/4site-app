"""
generar_pptx_4site.py  v4
━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
Diseño basado en ppttipo.pptx (referencia oficial 4SITE).
Colores y posiciones extraídos directamente del archivo de referencia.

Paleta exacta:
  #0D0D0D  fondo panel portada
  #BB944F  dorado acento
  #27AE60  verde score excelente
  #F39C12  amarillo score bueno
  #E74C3C  rojo score bajo
  #262626  card score (gris muy oscuro)
  #1A1A1A  footer / header fondo
  #F5F5F5  fondo KPI cards
  #6B6B6B  texto secundario
  #1A1A1A  texto primario (no negro puro)
  #FFFFFF  texto sobre oscuro

Fuente: Montserrat (bold para títulos, regular para cuerpo)
"""
from __future__ import annotations
import json, os, subprocess, tempfile, datetime, re
from io import BytesIO

_THIS_DIR  = os.path.dirname(os.path.abspath(__file__))
_JS_SCRIPT = os.path.join(_THIS_DIR, "pptx_4site_builder.js")

def _find_asset(names):
    for d in [_THIS_DIR, os.path.join(_THIS_DIR,"assets"), os.path.join(_THIS_DIR,"static")]:
        for n in names:
            p = os.path.join(d, n)
            if os.path.exists(p): return p
    return ""

def _find_node():
    for c in ["node","/usr/bin/node","/usr/local/bin/node"]:
        try:
            r = subprocess.run([c,"--version"], capture_output=True, timeout=5)
            if r.returncode == 0: return c
        except: pass
    return ""

def _node_ok():
    return bool(_find_node()) and os.path.exists(_JS_SCRIPT)

def _safe_list(v):
    if v is None: return []
    if isinstance(v, list): return v
    return list(v)

# ══════════════════════════════════════════════════════════════
def generar_pptx_por_tier(
    tier_key="free", ubicacion="", score=0, desglose=None,
    analisis="", competidores=None, idioma="es", lat=0.0, lng=0.0,
    modo="validar", tipo_negocio="", recomendaciones=None,
    demografia=None, contexto=None, ticket_data=None,
    trafico_data=None, mercado_data=None, forecast_data=None,
    roi_data=None, **kwargs,
):
    tipo_display = tipo_negocio or ""
    if not tipo_display and recomendaciones:
        tipo_display = (_safe_list(recomendaciones)[0].get("nombre","") if recomendaciones else "")

    bg_path   = _find_asset(["bg_calle.jpg","bg_ciudad.jpg","bg_ciudad.png"])
    logo_path = _find_asset(["logo_4site.png","logo_4site.jpg"])

    payload = {
        "tier_key":tier_key,"ubicacion":ubicacion,"score":float(score),
        "desglose":desglose or {},"analisis":analisis or "",
        "competidores":_safe_list(competidores),"idioma":idioma,
        "lat":float(lat or 0),"lng":float(lng or 0),"modo":modo,
        "tipo_negocio":tipo_display,"recomendaciones":_safe_list(recomendaciones),
        "demografia":demografia or {},"contexto":contexto or {},
        "ticket_data":ticket_data,"trafico_data":trafico_data,
        "mercado_data":mercado_data,"forecast_data":forecast_data,
        "roi_data":roi_data,"bg_path":bg_path,"logo_path":logo_path,
    }

    if _node_ok():
        result = _with_node(payload)
        if result: return result
    return _build_pptx(payload)


def _with_node(payload):
    node = _find_node(); out_path = None
    try:
        with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as tf:
            out_path = tf.name
        payload["out_path"] = out_path
        r = subprocess.run([node, _JS_SCRIPT],
                           input=json.dumps(payload, ensure_ascii=False, default=str),
                           capture_output=True, text=True, timeout=90, cwd=_THIS_DIR)
        if r.returncode != 0: return None
        if not os.path.exists(out_path) or os.path.getsize(out_path) < 1000: return None
        with open(out_path,"rb") as f: buf = BytesIO(f.read())
        buf.seek(0); return buf
    except: return None
    finally:
        if out_path:
            try: os.unlink(out_path)
            except: pass


# ══════════════════════════════════════════════════════════════
# CONSTRUCTOR PYTHON-PPTX — DISEÑO REFERENCIA
# ══════════════════════════════════════════════════════════════
def _build_pptx(payload: dict) -> BytesIO | None:
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN
        from pptx.oxml.ns import qn
        from lxml import etree
    except ImportError:
        return None

    # ── Paleta exacta del ppttipo.pptx ────────────────────────
    C = {
        "dorado":  "BB944F",
        "verde":   "27AE60",
        "amarillo":"F39C12",
        "rojo":    "E74C3C",
        "blanco":  "FFFFFF",
        "oscuro":  "1A1A1A",  # texto primario / headers
        "gris_t":  "6B6B6B",  # texto secundario
        "gris_c":  "F5F5F5",  # fondo cards
        "card_sc": "262626",  # card score (gris muy oscuro)
        "footer":  "1A1A1A",  # footer bg
        "panel":   "1A1A1A",  # panel portada (gris oscuro ppttipo)
    }
    MN = "Montserrat"   # fuente principal

    def rgb(h): h=h.lstrip("#"); return RGBColor(int(h[0:2],16),int(h[2:4],16),int(h[4:6],16))

    def bg(sl, col):
        sl.background.fill.solid()
        sl.background.fill.fore_color.rgb = rgb(col)

    def rect(sl, x, y, w, h, fill, lc=None, lw=0.5, alpha=None):
        s = sl.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
        s.fill.solid(); s.fill.fore_color.rgb = rgb(fill)
        if alpha is not None:
            try:
                sp=s._element; spPr=sp.find(qn("p:spPr"))
                sf=spPr.find(".//" + qn("a:solidFill"))
                if sf is not None:
                    sr=sf.find(qn("a:srgbClr"))
                    if sr is None:
                        sr=etree.SubElement(sf,qn("a:srgbClr")); sr.set("val",fill.lstrip("#"))
                    al=etree.SubElement(sr,qn("a:alpha")); al.set("val",str(int(alpha*1000)))
            except: pass
        if lc: s.line.color.rgb=rgb(lc); s.line.width=Pt(lw)
        else: s.line.fill.background()
        return s

    def txt(sl, text, x, y, w, h, sz=9, bold=False, col=None, al=None, it=False):
        col = col or C["oscuro"]
        tb = sl.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
        tf = tb.text_frame; tf.word_wrap = True
        p = tf.paragraphs[0]; p.text = str(text)
        if al: p.alignment = al
        r = p.runs[0]
        r.font.size = Pt(sz); r.font.bold = bold; r.font.italic = it
        r.font.color.rgb = rgb(col)
        try: r.font.name = MN
        except: pass

    def img(sl, path, x, y, w, h):
        try: sl.shapes.add_picture(path, Inches(x), Inches(y), Inches(w), Inches(h))
        except: pass

    # ── Header estándar (slides 2+) ───────────────────────────
    def header(sl, title, sub=""):
        """Header gris oscuro con acento dorado izquierdo — como ppttipo"""
        rect(sl, 0, 0, 13.33, 0.88, C["oscuro"])
        rect(sl, 0, 0, 0.2,   0.88, C["dorado"])
        txt(sl, title, 0.35, 0.1, 11.0, 0.55, sz=18, bold=True, col=C["blanco"])
        if sub:
            txt(sl, sub, 0.35, 0.64, 11.5, 0.22, sz=8, col=C["dorado"])

    # ── Footer estándar ───────────────────────────────────────
    def footer(sl, page_num=None):
        """Footer negro con texto gris — exacto al ppttipo"""
        rect(sl, 0, 7.22, 13.33, 0.28, C["footer"])
        ft = f"4SITE · Don't guess. Foresee. · 4site.mx"
        if page_num:
            ft += f"  ·  {page_num}"
        txt(sl, ft, 0.3, 7.23, 12.7, 0.25, sz=7.5, col=C["gris_t"])

    # ── KPI card ─────────────────────────────────────────────
    def kpi(sl, x, y, w, val, label, col_val=None):
        """Card blanca con valor grande y label pequeño — como slide 3"""
        col_val = col_val or C["oscuro"]
        rect(sl, x, y, w, 1.0, C["gris_c"])
        txt(sl, str(val), x, y+0.12, w, 0.55, sz=22, bold=True,
            col=col_val, al=PP_ALIGN.CENTER)
        txt(sl, label, x, y+0.67, w, 0.3, sz=8.5, col=C["gris_t"],
            al=PP_ALIGN.CENTER)

    # ── Barra horizontal desglose ─────────────────────────────
    def barra(sl, x_bg, y_bg, total_w, pct, label, val_txt, row_y_label):
        """Barra dorada sobre fondo gris claro — exacta al ppttipo slide 2"""
        rect(sl, x_bg, y_bg, total_w, 0.22, C["gris_c"])
        bar_w = max(0.04, total_w * min(pct, 100) / 100)
        rect(sl, x_bg, y_bg, bar_w, 0.22, C["dorado"])
        txt(sl, label, 4.1, row_y_label, 2.4, 0.32, sz=9, col=C["oscuro"])
        txt(sl, str(val_txt), 12.2, row_y_label, 0.8, 0.32, sz=9,
            bold=True, col=C["oscuro"])

    # ── Caja de interpretación dorada ─────────────────────────
    def interp(sl, y, text, h=0.55):
        rect(sl, 0.35, y, 12.6, h, "FFF8EE", C["dorado"], 0.6)
        rect(sl, 0.35, y, 0.06, h, C["dorado"])
        txt(sl, text, 0.52, y+0.07, 12.3, h-0.1,
            sz=9, col="3D3000", it=True)

    # ── Score color ───────────────────────────────────────────
    def sc_col(s): return C["verde"] if s>=75 else C["amarillo"] if s>=50 else C["rojo"]
    def sc_lbl(s): return ("EXCELENTE" if s>=80 else "BUENO" if s>=65 else
                            "REGULAR" if s>=50 else "BAJO" if s>=35 else "MUY BAJO")
    def sc_desc(s):
        if s>=75: return "Ubicación con condiciones muy favorables. Alta probabilidad de éxito."
        if s>=65: return "Buena ubicación. Condiciones favorables con algunos riesgos menores."
        if s>=50: return "Ubicación viable con precauciones. Evalúa los factores de riesgo."
        if s>=35: return "Condiciones difíciles. Se recomienda análisis más profundo antes de proceder."
        return "Ubicación con factores de riesgo importantes. Se recomienda no proceder."

    # ── Datos ─────────────────────────────────────────────────
    sc   = float(payload.get("score", 0))
    tk   = payload.get("tier_key", "free")
    ub   = payload.get("ubicacion", "")
    an   = payload.get("analisis", "")
    tp   = payload.get("tipo_negocio", "")
    dsg  = payload.get("desglose", {}) or {}
    dem  = payload.get("demografia", {}) or {}
    ctx  = payload.get("contexto", {}) or {}
    comp = payload.get("competidores", []) or []
    tkt  = payload.get("ticket_data") or {}
    traf = payload.get("trafico_data") or {}
    merc = payload.get("mercado_data") or {}
    fore = payload.get("forecast_data") or {}
    roi  = payload.get("roi_data") or {}
    bg_p = payload.get("bg_path", "")
    logo_p = payload.get("logo_path", "")
    _api = payload.get("google_api_key","") or os.getenv("GOOGLE_API_KEY","")
    lat  = payload.get("lat", 0)
    lng_ = payload.get("lng", 0)

    tl   = {"free":"GRATIS","basic":"BÁSICO · $99","pro":"PRO · $299",
            "premium":"PREMIUM · $999"}.get(tk, tk.upper())
    fecha = datetime.datetime.now().strftime("%d de %b de %Y")
    isB  = tk in ("basic","basico")
    isP  = tk == "pro"
    isPr = tk in ("premium","premium2")
    total = {"free":2,"basic":5,"basico":5,"pro":7,
             "premium":10,"premium2":10}.get(tk, 2)

    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    sn = 0; tmp_files = []

    # ══════════════════════════════════════════════════════════
    # SLIDE 1 — PORTADA (réplica exacta del ppttipo)
    # ══════════════════════════════════════════════════════════
    sn += 1
    sl = prs.slides.add_slide(blank)
    # NO hay fondo sólido — la foto ES el fondo, como en ppttipo
    # bg(sl, ...) removido para que la foto sea el único fondo

    # Foto de fondo a pantalla completa — composite PIL con panel derecho
    if bg_p and os.path.exists(bg_p):
        try:
            from PIL import Image as _PIL, ImageDraw as _IDraw
            im = _PIL.open(bg_p).convert("RGB")
            tw, th = 1333, 750  # proporciones de la slide (13.33 × 7.5)
            ir = im.width / im.height; pr = tw / th
            nw,nh = (int(th*ir),th) if ir>pr else (tw,int(tw/ir))
            im2 = im.resize((nw,nh), _PIL.LANCZOS if hasattr(_PIL,'LANCZOS') else _PIL.ANTIALIAS)
            lx,ly = (nw-tw)//2, (nh-th)//2
            im3 = im2.crop((lx,ly,lx+tw,ly+th)).convert("RGBA")

            # Panel oscuro derecho exacto del ppttipo
            # x=7.441/13.333*1333=743, y=1.613/7.5*750=161, w=4.313/13.333*1333=431, h=4.296/7.5*750=429
            overlay = _PIL.new("RGBA", (tw, th), (0,0,0,0))
            draw = _IDraw.Draw(overlay)
            draw.rectangle([743, 161, 743+431, 161+429], fill=(26, 26, 26, 195))  # ~76% opacidad

            im_final = _PIL.alpha_composite(im3, overlay).convert("RGB")
            _tf = tempfile.NamedTemporaryFile(suffix=".jpg", delete=False)
            im_final.save(_tf.name, "JPEG", quality=93); _tf.close()
            tmp_files.append(_tf.name)
            # Posición exacta del ppttipo: y=-1.45
            sl.shapes.add_picture(_tf.name, Inches(0), Inches(-1.45),
                                  Inches(13.333), Inches(8.952))
        except: pass

    # Logo arriba izquierda (x=-0.135 y=-0.825 w=3.6 h=1.178 como ref)
    if logo_p and os.path.exists(logo_p):
        try: sl.shapes.add_picture(logo_p, Inches(-0.135), Inches(-0.825), Inches(3.6), Inches(1.178))
        except:
            txt(sl, "4SITE", 0.25, 0.1, 2.5, 0.65, sz=30, bold=True, col=C["blanco"])
    else:
        # Texto logo estilizado sin archivo
        txt(sl, "4", 0.2, 0.1, 0.6, 0.65, sz=34, bold=True, col=C["dorado"])
        txt(sl, "SITE", 0.72, 0.12, 2.0, 0.6, sz=30, bold=True, col=C["blanco"])

    # Contenido del panel derecho — posiciones exactas del ppttipo
    PX = 7.904  # x base del panel

    # Badge tier (x=7.904 y=1.68 w=1.264 h=0.28)
    rect(sl, PX, 1.68, 1.264, 0.28, C["dorado"])
    txt(sl, tl, PX, 1.68, 1.264, 0.28, sz=8, bold=True, col=C["oscuro"],
        al=PP_ALIGN.CENTER)

    # Línea dorada (x=7.956 y=2.123 w=0.806 — línea fina)
    rect(sl, 7.956, 2.12, 0.806, 0.025, C["dorado"])

    # Etiqueta reporte
    txt(sl, "REPORTE DE ANALISIS DE UBICACION", PX, 2.18, 3.356, 0.26,
        sz=7, bold=True, col=C["gris_t"])

    # Dirección (x=7.904 y=2.42 w=3.356 h=1.3)
    txt(sl, ub, PX, 2.42, 3.356, 1.3, sz=16, bold=True, col=C["blanco"])

    # Tipo de negocio (x=7.904 y=3.72)
    if tp:
        txt(sl, tp, PX, 3.72, 3.356, 0.4, sz=10.5, col=C["dorado"])

    # Score / 100 (x=7.904 y=4.56 w=2.0 h=0.52)
    txt(sl, f"{int(sc)}/100", PX, 4.56, 2.0, 0.52,
        sz=28, bold=True, col=sc_col(sc))

    # EXCELENTE / BUENO / etc. (x=7.904 y=5.06)
    txt(sl, sc_lbl(sc), PX, 5.06, 2.0, 0.3,
        sz=14, bold=True, col=sc_col(sc))

    # Fecha y firma
    txt(sl, fecha, PX, 5.42, 3.356, 0.25, sz=8, col=C["gris_t"])
    txt(sl, "4SITE · 4site.mx", PX, 5.68, 3.356, 0.22, sz=8, col=C["dorado"])

    # Footer
    txt(sl, "4SITE - Don't guess. Foresee. · 4site.mx", 0.0, 7.22,
        13.33, 0.25, sz=7, col=C["gris_t"], al=PP_ALIGN.CENTER)

    # ══════════════════════════════════════════════════════════
    # SLIDE 2 — SCORE + DESGLOSE (Basic+)
    # Réplica exacta slide 2 de ppttipo
    # ══════════════════════════════════════════════════════════
    if isB or isP or isPr:
        sn += 1
        sl = prs.slides.add_slide(blank)
        bg(sl, C["blanco"])
        header(sl, "Análisis del Score de Viabilidad",
               f"{ub[:60]}  ·  {tp}" if tp else ub[:70])

        # Panel izquierdo gris claro (x=0.3 y=1.05 w=3.5 h=5.8)
        rect(sl, 0.3, 1.05, 3.5, 5.8, C["gris_c"])

        # Card score oscura (x=1.038 y=1.488 w=2.024 h=2.024)
        rect(sl, 1.038, 1.488, 2.024, 2.024, C["card_sc"])

        # Número score grande (x=0.95 y=1.95)
        txt(sl, str(int(sc)), 0.95, 1.95, 2.2, 0.85,
            sz=44, bold=True, col=sc_col(sc), al=PP_ALIGN.CENTER)

        # SCORE / 100
        txt(sl, "SCORE / 100", 0.95, 2.8, 2.2, 0.3,
            sz=9, col=C["gris_t"], al=PP_ALIGN.CENTER)

        # Badge verde EXCELENTE (x=1.332 y=3.15 w=1.435 h=0.3)
        rect(sl, 1.332, 3.15, 1.435, 0.3, sc_col(sc))
        txt(sl, sc_lbl(sc), 1.332, 3.15, 1.435, 0.3,
            sz=9, bold=True, col=C["blanco"], al=PP_ALIGN.CENTER)

        # Label grande EXCELENTE abajo (x=0.3 y=3.75)
        txt(sl, sc_lbl(sc), 0.3, 3.75, 3.5, 0.5,
            sz=16, bold=True, col=sc_col(sc), al=PP_ALIGN.CENTER)

        # Descripción del score
        txt(sl, sc_desc(sc), 0.4, 4.3, 3.3, 0.9,
            sz=9, col=C["gris_t"], al=PP_ALIGN.CENTER)

        # ── Desglose por categoría (lado derecho) ──
        txt(sl, "Desglose por categoría", 4.1, 1.05, 8.8, 0.4,
            sz=12, bold=True, col=C["oscuro"])

        # Barras de desglose — posiciones exactas del ppttipo
        BAR_X, BAR_W, BAR_TOTAL = 6.6, 5.5, 5.5

        rows = [
            ("Densidad",         dsg.get("densidad", 0),        1.55, 1.61),
            ("Calidad",          dsg.get("calidad", 0),         2.07, 2.13),
            ("Consolidacion",    dsg.get("consolidacion", 0),   2.59, 2.65),
            ("Num Competidores", min(dsg.get("num_competidores", 0), 20)*5, 3.11, 3.17),
            ("Rating Promedio",  min((dsg.get("rating_promedio") or 0)*20, 100), 3.63, 3.69),
        ]
        for lbl, pct, y_lbl, y_bar in rows:
            val_raw = dsg.get(lbl.lower().replace(" ","_"),
                     dsg.get(lbl.replace(" ","_").lower(), 0))
            # Para campos especiales, usar el valor directo
            if "num" in lbl.lower():
                val_raw = dsg.get("num_competidores", 0)
                val_display = str(val_raw)
                pct_bar = min(val_raw * 5, 100)
            elif "rating" in lbl.lower():
                val_raw = dsg.get("rating_promedio", 0) or 0
                try: val_display = str(round(float(val_raw), 1))
                except: val_display = str(val_raw)
                try: pct_bar = min(float(val_raw) * 20, 100)
                except: pct_bar = 0
            else:
                val_display = str(int(pct))
                pct_bar = pct

            rect(sl, BAR_X, y_bar, BAR_TOTAL, 0.22, C["gris_c"])
            bw = max(0.04, BAR_TOTAL * pct_bar / 100)
            rect(sl, BAR_X, y_bar, bw, 0.22, C["dorado"])
            txt(sl, lbl, 4.1, y_lbl, 2.4, 0.32, sz=9, col=C["oscuro"])
            txt(sl, val_display, 12.2, y_lbl, 0.8, 0.32, sz=9, bold=True, col=C["oscuro"])

        # Análisis vial al pie (si existe)
        av = ctx.get("analisis_vial", {})
        if av.get("badge"):
            adj = av.get("ajuste_score", 0)
            ac = C["verde"] if adj>0 else C["rojo"] if adj<-5 else C["amarillo"]
            rect(sl, 4.1, 5.0, 8.8, 0.6, "F0F0F0", C["dorado"], 0.5)
            txt(sl, f"Vialidad: {av.get('badge','')}  |  Impacto score: {'+'if adj>=0 else ''}{adj} pts",
                4.2, 5.07, 7.5, 0.25, sz=9, bold=True, col=ac)
            txt(sl, av.get("descripcion","")[:120],
                4.2, 5.35, 8.5, 0.22, sz=8, col=C["gris_t"], it=True)

        footer(sl, sn)

    # ══════════════════════════════════════════════════════════
    # SLIDE 3 — COMPETENCIA (Basic+)
    # Réplica exacta slide 3 de ppttipo
    # ══════════════════════════════════════════════════════════
    if isB or isP or isPr:
        sn += 1
        sl = prs.slides.add_slide(blank)
        bg(sl, C["blanco"])

        rats  = [float(c.get("rating",0)) for c in comp if c.get("rating")]
        avg_r = round(sum(rats)/len(rats), 1) if rats else 0
        tk_c  = tkt.get("ticket_competencia", 0) if tkt else 0
        rz    = tkt.get("rango_zona", "N/D") if tkt else "N/D"

        header(sl, "Análisis de Competencia",
               f"{len(comp)} establecimientos identificados en el área de influencia")

        # 4 KPI cards (x=0.3 y=1.05 w=2.4 h=1.0) — posiciones exactas
        kpi(sl, 0.3,  1.05, 2.4, len(comp),      "Competidores cercanos", C["oscuro"])
        kpi(sl, 2.85, 1.05, 2.4, f"{avg_r} ★",   "Rating promedio",      C["dorado"])
        kpi(sl, 5.4,  1.05, 2.4, f"${int(tk_c):,}" if tk_c else "N/D", "Ticket promedio zona", C["dorado"])
        kpi(sl, 7.95, 1.05, 2.4, rz,              "Nivel de precio",      C["oscuro"])

        # Tabla de competidores (x=0.3 y=2.25)
        comps = comp[:12]
        TY = 2.25
        # Header de tabla (fondo oscuro)
        COL_W = [0.45, 3.3, 1.8, 1.15, 1.1, 1.05, 0.8]  # #, Establecimiento, Tipo, Rating, Reseñas, Precio, Nivel
        HEADERS = ["#", "Establecimiento", "Tipo", "Rating", "Reseñas", "Precio", "Nivel"]
        rect(sl, 0.3, TY, 9.65, 0.42, C["oscuro"])
        cx = 0.35
        for h_txt, cw in zip(HEADERS, COL_W):
            txt(sl, h_txt, cx, TY+0.1, cw, 0.26, sz=8.5, bold=True, col=C["blanco"])
            cx += cw

        for ri, c in enumerate(comps):
            ry = TY + 0.42 + ri * 0.3
            row_bg = C["blanco"] if ri%2==0 else C["gris_c"]
            rect(sl, 0.3, ry, 9.65, 0.3, row_bg, "EEEEEE", 0.25)

            nombre = (c.get("displayName",{}).get("text","") or
                      c.get("name", c.get("nombre","—")))[:32]
            tipo_c = (c.get("primaryTypeDisplayName",{}).get("text","") or
                      c.get("primaryType","—"))[:16]
            rating = c.get("rating", "—")
            reviews = c.get("userRatingCount", "—")
            pm = {1:"$",2:"$$",3:"$$$",4:"$$$$",
                  "PRICE_LEVEL_INEXPENSIVE":"$","PRICE_LEVEL_MODERATE":"$$",
                  "PRICE_LEVEL_EXPENSIVE":"$$$","PRICE_LEVEL_VERY_EXPENSIVE":"$$$$"}
            precio = pm.get(c.get("priceLevel", c.get("price_level","")),"N/D")
            try:
                rv = float(rating)
                rc = C["verde"] if rv>=4.3 else C["amarillo"] if rv>=3.5 else C["rojo"]
                nv = "Alto" if rv>=4.3 else "Med." if rv>=3.5 else "Bajo"
            except: rc, nv = C["gris_t"], "—"

            vals = [str(ri+1), nombre, tipo_c, f"{rating} ★", str(reviews), precio, nv]
            cols_c = [C["gris_t"], C["oscuro"], C["gris_t"], rc, C["gris_t"], C["dorado"], rc]
            cols_b = [False, False, False, True, False, True, True]
            cx = 0.35
            for val, cw, col_v, bold_v in zip(vals, COL_W, cols_c, cols_b):
                aln = PP_ALIGN.CENTER if val in [str(ri+1),f"{rating} ★",str(reviews),precio,nv] else None
                txt(sl, val, cx, ry+0.06, cw, 0.2, sz=8.5, col=col_v, bold=bold_v, al=aln)
                cx += cw

        # Interpretación
        iy = TY + 0.42 + len(comps)*0.3 + 0.12
        top_r = sum(1 for r in rats if r>=4.3)
        weak_r = sum(1 for r in rats if r<3.5)
        tpl_c = (tp or "").lower()

        if len(comp)==0:
            itxt = (f"Sin competidores directos en 500m — esto es una oportunidad excepcional para '{tp or 'tu negocio'}'. "
                    f"Tendrás todo el mercado disponible sin presión de precio ni calidad. "
                    f"El riesgo es que la ausencia de competidores podría indicar baja demanda en la zona: "
                    f"valida con observación directa del tráfico peatonal antes de invertir.")
        elif len(comp)<=3:
            itxt = (f"Solo {len(comp)} competidor(es) en 500m — saturación baja, ideal para '{tp or 'tu negocio'}'. "
                    f"Con rating promedio de {avg_r}★, {'el mercado está insatisfecho: los clientes buscan una mejor opción y tú puedes ser esa alternativa.' if avg_r<3.8 else 'hay competidores establecidos pero con espacio para diferenciarte.'} "
                    f"{'Enfoca desde el primer día en superar ese ' + str(avg_r) + '★ — un negocio nuevo con 4.5★ en Google captura mercado rápidamente.' if rats else 'La poca competencia te da tiempo para posicionarte antes de que lleguen más jugadores.'}")
        elif len(comp)<=7:
            itxt = (f"{len(comp)} negocios similares en 500m (competencia moderada). Rating promedio {avg_r}★ — "
                    f"{'hay ' + str(weak_r) + ' competidor(es) con rating bajo (<3.5★): son tu oportunidad directa de capturar sus clientes insatisfechos con un servicio claramente superior.' if weak_r>0 else 'el estándar de la zona es medio-alto.'} "
                    f"{'Con ' + str(top_r) + ' negocios fuertes (≥4.3★), necesitas diferenciarte en: precio, especialización o experiencia del cliente — no puedes competir siendo igual.' if top_r>=2 else 'La mayoría de competidores son vulnerables: con buena atención y presencia digital puedes posicionarte como el referente.'}")
        else:
            itxt = (f"Zona con alta saturación: {len(comp)} competidores, {top_r} de ellos con ≥4.3★. "
                    f"Para '{tp or 'tu negocio'}' esto significa que el mercado existe y tiene demanda probada, pero la diferenciación es crítica. "
                    f"Opciones viables: (1) nicho dentro del segmento que nadie atiende, (2) precio agresivo en apertura + fidelización, "
                    f"(3) ubicación dentro del mismo radio con mayor visibilidad que los actuales.")
        if iy < 6.5:
            interp(sl, iy, itxt, h=0.65 if len(itxt)>200 else 0.55)

        footer(sl, sn)

    # ══════════════════════════════════════════════════════════
    # SLIDE 4 — DEMOGRAFÍA (Basic+)
    # ══════════════════════════════════════════════════════════
    if isB or isP or isPr:
        sn += 1
        sl = prs.slides.add_slide(blank)
        bg(sl, C["blanco"])
        fuente = dem.get("fuente","INEGI Censo 2020 · CONAPO")
        header(sl, "Perfil Demográfico del Área — 500m", f"Fuente: {fuente}")

        yr  = datetime.datetime.now().year
        pob = dem.get("poblacion_actual", dem.get("poblacion_estimada", 0))
        viv = dem.get("viviendas_actual", dem.get("viviendas_habitadas", 0))
        nse = dem.get("nse_predominante", "C")
        ing = dem.get("ingreso_actual", dem.get("ingreso_promedio_mensual", 0))
        gst = dem.get("gasto_actual",   dem.get("gasto_promedio_mensual", 0))
        tasa= dem.get("tasa_crecimiento_pct", 0.55)
        nc  = {"A":"7B1FA2","A/B":"9C27B0","B":"1976D2","B/C+":"0097A7","C+":"2E7D32",
               "C":"558B2F","C/D+":"F57F17","D+":"E65100","D/E":"B71C1C"}.get(nse,"1A76D2")

        # 6 KPI cards en una fila
        kpis6 = [(f"{int(pob):,}",  f"Población {yr}",   "1A76D2"),
                 (f"{int(viv):,}",  "Viviendas",          "1A76D2"),
                 (nse,              "NSE Predominante",    nc),
                 (f"${int(ing):,}", "Ingreso prom/mes",   "27AE60"),
                 (f"${int(gst):,}", "Gasto prom/mes",     "F39C12"),
                 (f"{float(tasa):.2f}%", "Crec. anual",   C["dorado"])]
        for i,(v,l,c) in enumerate(kpis6):
            kpi(sl, 0.3+i*2.17, 1.0, 1.98, v, l, c)

        # Barra de edad
        txt(sl, "Distribución de Edad Estimada del Entorno:", 0.3, 2.12,
            9, 0.3, sz=10, bold=True, col=C["oscuro"])
        grupos = [("0–14",0.18,"5EEAD4"),("15–29",0.22,"14B8A6"),
                  ("30–44",0.24,"0D9488"),("45–59",0.20,"0F766E"),("60+",0.16,"134E4A")]
        bx = 0.3
        for g, p, c in grupos:
            bw = 12.7*p; rect(sl, bx, 2.5, bw, 0.58, c)
            if bw > 0.9:
                txt(sl, f"{g}\n{int(p*100)}%", bx+0.05, 2.56, bw-0.1, 0.48,
                    sz=8.5, bold=True, col=C["blanco"], al=PP_ALIGN.CENTER)
            bx += bw

        # NSE implicación
        nse_d = {"A":"NSE A: >$85K/mes. Consumo premium sin restricción de precio.",
                 "A/B":"NSE A/B: >$45K/mes. Perfil urbano sofisticado, paga por calidad y experiencia.",
                 "B":"NSE B: >$32K/mes. Valora calidad, fidelizable con propuesta diferenciada.",
                 "B/C+":"NSE B/C+: ~$22K/mes. Relación precio-valor determinante.",
                 "C+":"NSE C+: ~$14K/mes. Leal si percibe buen trato y precio justo.",
                 "C":"NSE C: ~$10K/mes. Mercado masivo, competencia por precio.",
                 "C/D+":"NSE C/D+: ~$7K/mes. Ticket bajo, alta rotación necesaria.",
                 "D+":"NSE D+: ~$4.5K/mes. Volumen alto, margen limitado.",
                 "D/E":"NSE D/E: Segmento de subsistencia."}.get(nse,"")
        txt(sl, f"NSE {nse} — Implicación directa:", 0.3, 3.2, 9, 0.28,
            sz=10, bold=True, col=C["oscuro"])
        interp(sl, 3.54, nse_d, h=0.5)

        # Qué significa para el negocio
        txt(sl, f"¿Qué significa para \"{tp or 'tu negocio'}\"?",
            0.3, 4.16, 9, 0.28, sz=10, bold=True, col=C["oscuro"])
        nse_alto = nse in ["A","A/B","B"]
        nse_med  = nse in ["B/C+","C+","C"]
        tpl = (tp or "").lower()
        pct_gasto_cat = 0.15  # % ingreso destinado a categoría genérica

        if "caf" in tpl:
            pct_gasto_cat = 0.08
            if nse_alto:
                di = (f"El perfil NSE {nse} con ${int(ing):,}/mes de ingreso es el cliente ideal de una cafetería premium. "
                      f"Estas personas consumen café fuera de casa en promedio 4-5 veces por semana y valoran la calidad del origen, "
                      f"el ambiente y el barista sobre el precio. Un ticket de $150-250 MXN es natural para este segmento — "
                      f"no necesitas justificar el precio, necesitas justificar la experiencia. "
                      f"Con {int(pob):,} habitantes en 500m, tu mercado potencial de café diario es ~{int(pob*0.12):,} personas.")
            else:
                di = (f"NSE {nse} con ${int(ing):,}/mes — el cliente valora la cafetería como un gusto accesible, no un lujo. "
                      f"Ticket objetivo $80-130 MXN con propuesta de calidad visible. "
                      f"Los programas de fidelización (sellitos, 10a gratis, suscripción mensual) funcionan muy bien en este segmento "
                      f"porque el cliente es sensible al precio pero leal si recibe buen trato. "
                      f"Horario matutino 7-10h concentra el 60% de ventas — prioriza esa ventana.")
        elif "restaurante" in tpl or "comida" in tpl:
            pct_gasto_cat = 0.18
            gasto_food = int(ing * pct_gasto_cat)
            if nse_alto:
                di = (f"NSE {nse} — este segmento destina ~${gasto_food:,}/mes en alimentos fuera de casa por persona. "
                      f"El concepto gastronómico medio-alto es perfectamente viable: ambiente cuidado, ingredientes de calidad y "
                      f"experiencia de servicio son los diferenciadores que este cliente paga. "
                      f"Con {int(viv):,} viviendas en 500m, tu mercado de comida es de ~{int(viv*2.8*0.18):,} comensales potenciales.")
            else:
                di = (f"NSE {nse} con ${int(ing):,}/mes — el gasto en food service es ~${int(ing*0.18):,}/mes/persona. "
                      f"El menú ejecutivo ($80-150 MXN) + opciones para llevar maximizan rotación y ticket promedio. "
                      f"El cliente de NSE {nse} prioriza: precio justo, ración generosa y rapidez. "
                      f"Considera desayunos y comidas como anclas del día — la cena es secundaria en este segmento.")
        elif "gym" in tpl or "fitness" in tpl:
            if nse_alto:
                di = (f"NSE {nse} con ${int(ing):,}/mes — membresías de $800-2,000 MXN/mes son completamente viables. "
                      f"Este segmento ya tiene el hábito de ir al gym y está dispuesto a pagar por amenidades premium: "
                      f"clases especializadas, equipamiento nuevo, vestidores dignos y app de reservas. "
                      f"El {int((dem.get('distribucion_edad',{}).get('18-35',30) + dem.get('distribucion_edad',{}).get('36-55',25)))}% "
                      f"de adultos activos en la zona son tu mercado directo — ~{int(pob*0.55*0.12):,} personas.")
            else:
                di = (f"NSE {nse} con ${int(ing):,}/mes — el precio de membresía es el factor #1 de decisión. "
                      f"Rango viable: $400-700 MXN/mes con contratos de 3-6 meses. "
                      f"Las clases grupales son clave: reducen el churn porque el cliente se engancha con la comunidad, "
                      f"no solo con el equipo. Horarios de 6-8am y 6-8pm concentran el 70% de asistencia.")
        elif any(x in tpl for x in ["acab","piso","lamina","persiana","cortina","duela"]):
            di = (f"NSE {nse} con ${int(ing):,}/mes — el cliente de acabados del hogar tiene ticket promedio de compra "
                  f"de $3,000-15,000 MXN por proyecto. Con {int(viv):,} viviendas en 500m, el mercado anual de remodelación "
                  f"es estimado en ${int(viv*0.08*8000):,} MXN (8% de hogares remodelan/año, ticket $8K promedio). "
                  f"{'El NSE alto indica predisposición a marcas premium y diseño — el showroom y las muestras físicas son críticos.' if nse_alto else 'NSE medio: relación precio-calidad es determinante. El cliente compara 3+ opciones antes de decidir.'} "
                  f"Las mujeres lideran el 70% de estas decisiones de compra (ENIGH 2022) — adapta tu comunicación visual.")
        elif "empan" in tpl or "antoj" in tpl or "quesad" in tpl:
            di = (f"NSE {nse} con ${int(ing):,}/mes — el negocio de antojitos tiene alta demanda en todos los NSE. "
                  f"Con {int(pob):,} hab en 500m y un ticket de $40-100 MXN, tu mercado diario potencial "
                  f"es ~{int(pob*0.06):,} clientes. "
                  f"{'NSE alto: el cliente busca calidad artesanal, ingredientes visibles y presentación cuidada.' if nse_alto else 'NSE medio: precio accesible, ración generosa y rapidez son las 3 claves del éxito.'} "
                  f"El horario crítico: 7-10am y 1-3pm concentran el 75% de las ventas.")
        else:
            porc = int(float(tasa)*100)
            di = (f"El entorno de {int(pob):,} habitantes con NSE {nse} e ingreso ${int(ing):,}/mes "
                  f"{'tiene alta capacidad de gasto y disposición a pagar por propuestas de valor diferenciadas.' if nse_alto else 'requiere una propuesta precio-valor muy competitiva para captar y retener clientes.' if nse_med else 'exige precios muy accesibles y alta rotación de clientes para ser rentable.'} "
                  f"La densidad de {int(dem.get('densidad_actual',0) if dem else 0):,} hab/km² indica "
                  f"{'zona urbana densa — excelente para negocios de captación por tráfico.' if (dem.get('densidad_actual',0) if dem else 0)>10000 else 'densidad media — el marketing activo es necesario para atraer clientes del radio completo.'} "
                  f"Con {float(tasa):.2f}% de crecimiento anual, la zona "
                  f"{'está en expansión — demanda creciente en los próximos 5 años.' if float(tasa)>0.7 else 'es estable — mercado consolidado con base de clientes predecible.'}")
        interp(sl, 4.5, di, h=0.82 if len(di)>280 else 0.72)
        footer(sl, sn)

    # ══════════════════════════════════════════════════════════
    # SLIDE 5 — TICKET (Pro+)
    # ══════════════════════════════════════════════════════════
    if (isP or isPr) and tkt:
        sn += 1
        sl = prs.slides.add_slide(blank)
        bg(sl, C["blanco"])
        header(sl, "Ticket Promedio — Análisis de Precio de la Zona",
               "Fuente: price_level Google Places + rangos ENIGH 2022 por categoría")

        tk_c = tkt.get("ticket_competencia", 0)
        rz   = tkt.get("rango_zona", "—")
        tk_r = tkt.get("ticket_recomendado", 0)
        pos  = tkt.get("posicionamiento", "")

        kpi(sl, 0.3,  1.0, 3.9, f"${int(tk_c):,} MXN", "Ticket promedio\ncompetencia zona", C["dorado"])
        kpi(sl, 4.45, 1.0, 3.9, rz, "Nivel de precio del área", "1A76D2")
        kpi(sl, 8.6,  1.0, 4.35, f"${int(tk_r):,} MXN", "Ticket recomendado\npara tu negocio", C["verde"])

        itk = tkt.get("interpretacion_prem", tkt.get("interpretacion_pro",""))
        if itk: interp(sl, 2.12, itk, h=0.62)

        txt(sl, "Estrategia de Posicionamiento:", 0.3, 2.86, 9, 0.28,
            sz=10, bold=True, col=C["oscuro"])
        pm2 = {"premium": (C["dorado"], f"PREMIUM: Cobrar por encima de ${int(tk_c):,} MXN. Tu ubicación y NSE lo justifican. Enfoca comunicación en calidad y experiencia."),
               "competitivo": ("1A76D2", f"COMPETITIVA: Igualar ${int(tk_c):,} MXN. Gana por conveniencia y servicio. Compensa con programa de fidelización."),
               "accesible":   (C["verde"], f"ACCESIBLE: Por debajo de ${int(tk_c):,} MXN para penetrar rápido. Construye base de clientes en primeros 90 días, luego ajusta.")}
        pc2, pt2 = pm2.get(pos, ("1A76D2", f"Calibra tu precio según tu propuesta de valor vs competencia actual (${int(tk_c):,} MXN promedio zona)."))
        rect(sl, 0.3, 3.2, 12.7, 0.52, "F0F8FF", pc2, 0.6)
        txt(sl, pt2, 0.46, 3.24, 12.42, 0.44, sz=9.5, col=C["oscuro"])

        det = tkt.get("detalle_competidores",[])[:6]
        if det:
            txt(sl, "Detalle de ticket por competidor:", 0.3, 3.84, 5, 0.26,
                sz=9.5, bold=True, col=C["gris_t"])
            for di2, d in enumerate(det):
                dx = 0.3+(di2%3)*4.3; dy = 4.16+(di2//3)*0.56
                rect(sl, dx, dy, 4.1, 0.5, C["gris_c"], "EBEBEB", 0.3)
                txt(sl, (d.get("nombre",""))[:30], dx+0.1, dy+0.05, 2.8, 0.2,
                    sz=8, bold=True, col=C["oscuro"])
                txt(sl, f"${int(d.get('ticket_estimado',0)):,}",
                    dx+2.9, dy+0.05, 1.1, 0.2, sz=9, bold=True,
                    col=C["dorado"], al=PP_ALIGN.RIGHT)
                txt(sl, d.get("nivel",""), dx+0.1, dy+0.28, 2.8, 0.18,
                    sz=7.5, col=C["gris_t"])

        diff = int(tk_r) - int(tk_c)
        nse_d = dem.get("nse_predominante","C") if dem else "C"
        ing_d = dem.get("ingreso_actual", dem.get("ingreso_promedio_mensual",0)) if dem else 0
        tpl_t = (tp or "").lower()
        if diff > 0:
            dt_txt = (f"Tu ticket recomendado ${int(tk_r):,} MXN está ${diff:,} por ENCIMA del promedio zona (${int(tk_c):,} MXN). "
                      f"Esto es viable porque el NSE {nse_d} del área, con ingreso promedio ${int(ing_d):,}/mes, "
                      f"tiene capacidad de pago para propuestas premium. "
                      f"Para sostener este ticket: comunica claramente qué te hace diferente "
                      f"({'calidad del producto, origen, ambiente y experiencia de servicio' if 'caf' in tpl_t else 'especialización, calidad de materiales y garantía' if 'acab' in tpl_t or 'piso' in tpl_t else 'propuesta de valor clara y atención diferenciada'}). "
                      f"Un ticket ${diff:,} más alto que la competencia requiere que el cliente lo perciba justificado desde el primer contacto.")
        elif diff < 0:
            dt_txt = (f"Tu ticket recomendado ${int(tk_r):,} MXN está ${abs(diff):,} por DEBAJO del promedio zona (${int(tk_c):,} MXN). "
                      f"Esta estrategia de precio bajo es ideal para la etapa de apertura: captura clientes rápidamente, "
                      f"construye tu base de reseñas y tu rating en Google. "
                      f"Una vez que tengas 50+ reseñas con ≥4.3★, considera subir gradualmente el precio hacia ${int(tk_c):,} MXN. "
                      f"Compensa el menor margen con mayor volumen de clientes y frecuencia de visita.")
        else:
            dt_txt = (f"Tu ticket recomendado ${int(tk_r):,} MXN está en línea con el promedio de la zona. "
                      f"Esto significa que debes competir por atributos diferentes al precio: "
                      f"calidad del servicio, velocidad de atención, ambiente, presencia digital y fidelización. "
                      f"Con NSE {nse_d} e ingreso ${int(ing_d):,}/mes, el cliente puede pagar más si percibe valor — "
                      f"considera subir ${int(tk_c*0.1):,} MXN una vez establecido.")
        interp(sl, 5.3, dt_txt, h=0.75 if len(dt_txt)>250 else 0.55)
        footer(sl, sn)

    # ══════════════════════════════════════════════════════════
    # SLIDE 6 — TRÁFICO + MERCADO (Pro+)
    # ══════════════════════════════════════════════════════════
    if (isP or isPr) and traf:
        sn += 1
        sl = prs.slides.add_slide(blank)
        bg(sl, C["blanco"])
        zt = ctx.get("tipo_zona","mixto")
        header(sl, "Análisis de Tráfico · Mercado Potencial",
               f"Zona '{zt}' · Basado en densidad y tipo de negocio")

        kpi(sl, 0.3,  0.98, 3.0, traf.get("nivel_general","—"), "Flujo general",   "1A76D2")
        kpi(sl, 3.5,  0.98, 3.0, traf.get("dia_pico","—"),      "Mejor día",        C["verde"])
        kpi(sl, 6.7,  0.98, 3.0, traf.get("dia_bajo","—"),      "Día más bajo",     C["amarillo"])

        txt(sl, "Horas pico para tu tipo de negocio:", 0.3, 2.1,
            9, 0.28, sz=10, bold=True, col=C["oscuro"])
        picos = traf.get("horas_pico",[])[:5]
        for pi, p in enumerate(picos):
            px = 0.3+pi*1.95; fl = p.get("flujo",0)
            pc3 = C["verde"] if fl>=80 else C["amarillo"] if fl>=60 else "888888"
            rect(sl, px, 2.44, 1.75, 0.78, pc3, None, 0)
            txt(sl, p.get("hora_str",""), px, 2.5, 1.75, 0.35,
                sz=12, bold=True, col=C["blanco"], al=PP_ALIGN.CENTER)
            txt(sl, f"{p.get('nivel','')} · {fl}%", px, 2.82, 1.75, 0.3,
                sz=8, col=C["blanco"], al=PP_ALIGN.CENTER)

        # Gráfica horaria
        txt(sl, "Perfil de flujo horario (0h–23h):", 0.3, 3.34,
            5, 0.24, sz=8.5, bold=True, col=C["gris_t"])
        horario = traf.get("trafico_horario",[0]*24)
        bw_e = 12.7/24
        rect(sl, 0.3, 3.62, 12.7, 0.85, "F0F0F0", "E0E0E0", 0.3)
        for hi, val in enumerate(horario):
            nrm = max(0,min(100,val)); bh = 0.8*nrm/100
            bc = "3A3A3A" if nrm>=70 else C["dorado"] if nrm>=45 else "CCCCCC"
            if bh>0.01:
                rect(sl, 0.3+hi*bw_e+0.01, 3.62+0.8-bh, bw_e-0.03, bh, bc)
        for h in [0,6,12,18,23]:
            txt(sl, f"{h}h", 0.28+h*bw_e, 4.5, 0.5, 0.18, sz=7, col="999999")

        # Mercado potencial
        if merc:
            txt(sl, "Mercado Potencial:", 0.3, 4.78, 4, 0.28,
                sz=10, bold=True, col=C["oscuro"])
            mt  = merc.get("mercado_total",0); mc_v = merc.get("mercado_capturable",0)
            cpd = merc.get("clientes_potenciales_dia",0)
            pen = merc.get("penetracion_estimada_pct",15)
            kpi(sl, 0.3,  5.12, 3.0, f"${int(mt):,}",  "Mercado total área/mes",  "1A76D2")
            kpi(sl, 3.5,  5.12, 3.0, f"${int(mc_v):,}","Mercado capturable/mes",  C["verde"])
            kpi(sl, 6.7,  5.12, 3.0, f"{int(cpd)}/día","Clientes potenciales/día", C["dorado"])
            kpi(sl, 9.9,  5.12, 3.1, f"{pen}%",         "Penetración estimada",    C["amarillo"])

        picoH = ", ".join([p.get("hora_str","") for p in picos[:3]]) or "—"
        niv = traf.get("nivel_general","")
        tpl_tr = (tp or "").lower()
        niv_l = str(niv).lower()

        if "alto" in niv_l:
            niv_txt = "Flujo alto en la zona — ventaja directa para captar clientes sin necesidad de marketing intenso desde el inicio."
        elif "medio" in niv_l or "moderado" in niv_l:
            niv_txt = "Flujo moderado — necesitas estrategias activas: señalización visible, Google My Business optimizado y redes sociales con geolocalización para captar el tráfico existente."
        else:
            niv_txt = "Flujo bajo — el marketing local proactivo es indispensable: volanteo en 500m, alianzas con negocios vecinos y campañas geolocalizadas para crear el hábito de visita."

        if "caf" in tpl_tr:
            hor_txt = f"Para cafetería, los picos {picoH} concentran el 70% de las ventas del día. Maximiza staffing en esas ventanas y ofrece el menú más rápido posible — el cliente matutino no puede esperar más de 3 minutos."
        elif "gym" in tpl_tr or "fitness" in tpl_tr:
            hor_txt = f"Para gimnasio, los picos {picoH} son los horarios críticos de tu negocio. Programa tus clases grupales más populares exactamente en esas franjas — llenan el espacio y reducen churn de membresías."
        elif "restaurante" in tpl_tr or "comida" in tpl_tr:
            hor_txt = f"Para restaurante, los picos {picoH} son tus ventanas de ingreso principal. El servicio en esas horas debe ser impecable — el 80% de las reseñas de Google se escriben después de las horas pico."
        else:
            hor_txt = f"Los picos de tráfico en {picoH} son tus ventanas clave. Concentra tu máximo personal, tus promociones y tu marketing activo en esas franjas horarias."

        ti = f"{niv_txt} {hor_txt}"
        if merc:
            ti += f" Meta realista mes 3: ${int(mc_v*pen/100):,} MXN (={pen}% del mercado capturable de ${int(mc_v):,}/mes)."
        interp(sl, 6.2, ti, h=0.65 if len(ti)>240 else 0.52)
        footer(sl, sn)

    # ══════════════════════════════════════════════════════════
    # SLIDE 7 — FORECAST + ROI (Pro+)
    # ══════════════════════════════════════════════════════════
    if (isP or isPr) and fore:
        sn += 1
        sl = prs.slides.add_slide(blank)
        bg(sl, "3A3A3A")  # gris oscuro (ppttipo)
        rect(sl, 0, 0, 13.33, 0.88, C["dorado"])
        rect(sl, 0, 0, 0.2,   0.88, C["dorado"])
        txt(sl, "Proyección de Ventas · ROI Proyectado",
            0.35, 0.1, 10, 0.55, sz=18, bold=True, col=C["blanco"])
        txt(sl, f"Negocio: {tp}",
            0.35, 0.64, 10, 0.22, sz=8, col="444444")

        txt(sl, "Escenarios de Ventas Mensuales:", 0.3, 0.98,
            9, 0.28, sz=10, bold=True, col=C["blanco"])
        scens = [("Pesimista", fore.get("pesimista",0), C["rojo"],    "Arranque lento — penetración < 10%"),
                 ("Realista",  fore.get("realista",0),  C["amarillo"],"Escenario base — penetración ~15%, mes 3"),
                 ("Optimista", fore.get("optimista",0), C["verde"],   "Penetración alta — buen marketing desde día 1")]
        max_v = max(s[1] for s in scens) or 1
        for i,(nm,v,c,desc) in enumerate(scens):
            sy = 1.32+i*0.88
            rect(sl, 0.3, sy, 12.7, 0.78, "2E2E2E", "3A3A3A", 0.3)
            txt(sl, nm,   0.45, sy+0.12, 1.8, 0.28, sz=10, bold=True, col=c)
            txt(sl, desc, 0.45, sy+0.46, 1.8, 0.24, sz=8, col="777777", it=True)
            blen = max(0.3, 8.5*v/max_v)
            rect(sl, 2.5, sy+0.25, 8.5, 0.18, "3A3A3A")
            rect(sl, 2.5, sy+0.25, blen, 0.18, c)
            txt(sl, f"${int(v):,} MXN/mes",
                2.5+blen+0.15, sy+0.15, 4, 0.3, sz=12, bold=True, col=c)

        nota = fore.get("nota","")
        if nota: interp(sl, 3.98, nota, h=0.5)

        if roi:
            rect(sl, 0.3, 4.6, 12.7, 0.03, "555555")
            txt(sl, "Retorno sobre Inversión (ROI):", 0.3, 4.72, 9, 0.28,
                sz=10, bold=True, col=C["blanco"])
            inv_min = roi.get("inversion_min",0)
            inv_max = roi.get("inversion_max",0)
            ret_m   = roi.get("retorno_mensual", fore.get("realista",0))
            mr      = roi.get("meses_recuperacion_realista","—")
            roi_a   = roi.get("roi_anual_pct",
                              round(ret_m*12/inv_min*100) if inv_min and ret_m else "—")
            for i,(v,l,c) in enumerate([
                (f"${int(inv_min):,}", "Inversión mínima", "888888"),
                (f"${int(inv_max):,}", "Inversión máxima", "AAAAAA"),
                (f"${int(ret_m):,}/mes","Retorno mensual",  C["dorado"]),
                (f"{mr} meses",        "Recuperación",      C["verde"]),
                (f"{roi_a}%",          "ROI anual",         C["verde"]),
            ]):
                kpi(sl, 0.3+i*2.57, 5.06, 2.38, v, l, c)
            tpl_r = (tp or "").lower()
            try: roi_a_num = float(roi_a)
            except: roi_a_num = 0
            if roi_a_num > 100:
                roi_calificacion = f"ROI {roi_a}% es excepcional — supera ampliamente el promedio de inversión retail en México (60-80% anual). Significa que por cada $1 invertido, recuperas $1 + ${roi_a/100:.2f} al final del primer año."
            elif roi_a_num > 50:
                roi_calificacion = f"ROI {roi_a}% es sólido y competitivo para el sector. Es una inversión viable si se ejecuta el plan de negocio correctamente."
            elif roi_a_num > 0:
                roi_calificacion = f"ROI {roi_a}% es marginal — suficiente para no perder, pero requiere optimizar costos operativos para mejorar márgenes. Prioriza controlar tu costo de nómina y renta."
            else:
                roi_calificacion = f"ROI negativo — la inversión no se recupera en el primer año bajo el escenario base. Considera reducir la inversión inicial o buscar una ubicación con menor competencia."

            ri_t = (f"Con inversión de ${int(inv_min):,}–${int(inv_max):,} MXN y ventas realistas de ${int(ret_m):,}/mes, "
                    f"recuperas tu capital en ~{mr} meses. {roi_calificacion} "
                    f"El escenario pesimista (${int(fore.get('pesimista',0) if fore else 0):,}/mes) aún debe cubrir tus costos fijos — "
                    f"si no, reduce la inversión inicial o negocia condiciones de renta más favorables.")
            interp(sl, 6.12, ri_t, h=0.72 if len(ri_t)>250 else 0.55)
        footer(sl, sn)

    # ── Competencia y Demografía también en FREE tier ──────────
    if tk == "free" and comp:
        sn += 1; sl = prs.slides.add_slide(blank); bg(sl, C["blanco"])
        rats_f = [float(c.get("rating",0)) for c in comp if c.get("rating")]
        avg_f = round(sum(rats_f)/len(rats_f),1) if rats_f else 0
        header(sl, "Competencia Cercana - 500m", f"{len(comp)} negocios similares")
        kpi(sl, 0.3,  0.98, 4.5, len(comp),     "Competidores cercanos", "1A76D2")
        kpi(sl, 5.0,  0.98, 4.5, f"{avg_f}*",   "Rating promedio zona", C["dorado"])
        kpi(sl, 9.65, 0.98, 3.3, "Baja" if len(comp)<=4 else "Media" if len(comp)<=8 else "Alta",
            "Saturacion", C["verde"] if len(comp)<=4 else C["amarillo"] if len(comp)<=8 else C["rojo"])
        comps_f = comp[:10]
        TY_f = 2.08
        rect(sl, 0.3, TY_f, 12.7, 0.35, C["oscuro"])
        for _lbl, _cx, _cw in [("#",0.32,0.28),("Nombre",0.62,4.2),("Rating",4.87,1.0),("Precio",5.92,0.9),("Nivel",6.87,1.4)]:
            txt(sl, _lbl, _cx, TY_f+0.07, _cw, 0.22, sz=8.5, bold=True, col=C["blanco"])
        for _ri, _c in enumerate(comps_f):
            _ry = TY_f + 0.35 + _ri*0.3
            rect(sl, 0.3, _ry, 12.7, 0.28, C["blanco"] if _ri%2==0 else C["gris_c"], "EEEEEE", 0.2)
            _nm = (_c.get("displayName",{}).get("text","") or _c.get("name",_c.get("nombre","—")))[:35]
            _rt = _c.get("rating","—")
            _pm = {1:"$",2:"$$",3:"$$$",4:"$$$$","PRICE_LEVEL_INEXPENSIVE":"$","PRICE_LEVEL_MODERATE":"$$","PRICE_LEVEL_EXPENSIVE":"$$$","PRICE_LEVEL_VERY_EXPENSIVE":"$$$$"}
            _pr = _pm.get(_c.get("priceLevel",_c.get("price_level","")),"—")
            try:
                _rv=float(_rt); _rc=C["verde"] if _rv>=4.3 else C["amarillo"] if _rv>=3.5 else C["rojo"]
                _nv="Fuerte" if _rv>=4.3 else "Moderado" if _rv>=3.5 else "Debil"
            except: _rc,_nv=C["gris_t"],"—"
            txt(sl, str(_ri+1), 0.32, _ry+0.05, 0.28, 0.2, sz=8, col=C["gris_t"], al=PP_ALIGN.CENTER)
            txt(sl, _nm, 0.62, _ry+0.05, 4.2, 0.2, sz=8.5, col=C["oscuro"])
            txt(sl, f"* {_rt}", 4.87, _ry+0.05, 1.0, 0.2, sz=8.5, bold=True, col=_rc, al=PP_ALIGN.CENTER)
            txt(sl, _pr, 5.92, _ry+0.05, 0.9, 0.2, sz=8.5, bold=True, col=C["dorado"], al=PP_ALIGN.CENTER)
            txt(sl, _nv, 6.87, _ry+0.05, 1.4, 0.2, sz=8, bold=True, col=_rc, al=PP_ALIGN.CENTER)
        _iy_f = TY_f + 0.35 + len(comps_f)*0.3 + 0.1
        if len(comp)<=3: _itf=f"Solo {len(comp)} competidor(es) en 500m. Zona de baja saturacion — ideal para entrar primero."
        elif len(comp)<=7: _itf=f"{len(comp)} negocios en 500m, rating {avg_f}*. {'Mercado insatisfecho — entra con calidad superior.' if avg_f<3.8 else 'Estandar medio-alto — diferenciacion necesaria.'}"
        else: _itf=f"{len(comp)} competidores — zona saturada. Diferenciacion critica para sobrevivir."
        if _iy_f < 6.5: interp(sl, _iy_f, _itf, h=0.55)
        footer(sl, sn)

    if tk == "free" and dem:
        sn += 1; sl = prs.slides.add_slide(blank); bg(sl, C["blanco"])
        header(sl, "Perfil Demografico del Area", f"Fuente: {dem.get('fuente','INEGI 2020')}")
        _pob_f=dem.get("poblacion_actual",0); _nse_f=dem.get("nse_predominante","C")
        _ing_f=dem.get("ingreso_actual",0); _tsa_f=dem.get("tasa_crecimiento_pct",0.55)
        _nc_f={"A":"7B1FA2","A/B":"9C27B0","B":"1976D2","B/C+":"0097A7","C+":"2E7D32","C":"558B2F","C/D+":"F57F17","D+":"E65100","D/E":"B71C1C"}.get(_nse_f,"1A76D2")
        kpi(sl,0.3,0.98,3.0,f"{int(_pob_f):,}","Poblacion area","1A76D2")
        kpi(sl,3.5,0.98,3.0,_nse_f,"NSE Predominante",_nc_f)
        kpi(sl,6.7,0.98,3.0,f"${int(_ing_f):,}","Ingreso prom/mes",C["verde"])
        kpi(sl,9.9,0.98,3.1,f"{float(_tsa_f):.2f}%","Crec. anual",C["dorado"])
        txt(sl,"Que significa para tu negocio:",0.3,2.12,9,0.28,sz=10,bold=True,col=C["oscuro"])
        _nse_alto_f=_nse_f in["A","A/B","B"]
        _tpl_f=(tp or "").lower()
        if "caf" in _tpl_f: _di_f=f"NSE {_nse_f}, ${int(_ing_f):,}/mes. {'Ticket $150-250 MXN natural para este segmento. Clientes con habito de cafe premium establecido.' if _nse_alto_f else 'Calidad accesible + fidelizacion. La frecuencia de visita es tu palanca de ingresos.'}"
        elif "restaurante" in _tpl_f or "comida" in _tpl_f: _di_f=f"Gasto food service ~${int(_ing_f*0.18):,}/mes/persona. {'Concepto medio-alto viable.' if _nse_alto_f else 'Menu ejecutivo + para llevar maximizan rotacion.'}"
        else: _di_f=f"Zona {int(_pob_f):,} hab, NSE {_nse_f}, ingreso ${int(_ing_f):,}/mes. {'Alta capacidad de gasto.' if _nse_alto_f else 'Propuesta precio-valor competitiva necesaria.'} Crecimiento {float(_tsa_f):.2f}%/ano."
        interp(sl, 2.48, _di_f, h=0.72)
        footer(sl, sn)

    # ══════════════════════════════════════════════════════════
    # RECOMENDACIÓN FINAL
    # Réplica exacta slide 4 del ppttipo
    # ══════════════════════════════════════════════════════════
    sn += 1
    sl = prs.slides.add_slide(blank)
    bg(sl, "3D3D3D")  # gris oscuro (ppttipo tx1 75% lum)

    # Barra dorada izquierda vertical (x=0 y=0 w=0.2 h=7.5) — como ppttipo slide 4
    rect(sl, 0, 0, 0.2, 7.5, C["dorado"])

    # Título grande (x=0.35 y=0.2 w=10.0 h=0.65)
    txt(sl, "RECOMENDACIÓN FINAL", 0.35, 0.2, 10.0, 0.65,
        sz=22, bold=True, col=C["blanco"])

    # Sub-header dirección + score (x=0.35 y=0.8)
    txt(sl, f"{ub[:70]}  ·  Score: {int(sc)}/100",
        0.35, 0.8, 12.0, 0.35, sz=10, col=C["dorado"])

    # Línea separadora dorada (x=0.35 y=1.28 w=6.0 h=0.04)
    rect(sl, 0.35, 1.28, 6.0, 0.04, C["dorado"])

    # Texto de análisis (x=0.35 y=1.4 w=12.55 h=4.2)
    y = 1.4; ML, MW = 0.35, 12.55; MAX_Y = 5.45
    cur_sl = sl

    def ensure(needed):
        nonlocal cur_sl, y, sn
        if y + needed > MAX_Y:
            sn += 1
            cur_sl = prs.slides.add_slide(blank)
            bg(cur_sl, "3D3D3D")
            rect(cur_sl, 0, 0, 0.2, 7.5, C["dorado"])
            txt(cur_sl, "RECOMENDACIÓN FINAL  (continuación)",
                0.35, 0.2, 10.0, 0.5, sz=18, bold=True, col=C["blanco"])
            txt(cur_sl, f"{ub[:60]}  ·  Score: {int(sc)}/100",
                0.35, 0.75, 12.0, 0.3, sz=9, col=C["dorado"])
            rect(cur_sl, 0.35, 1.1, 6.0, 0.04, C["dorado"])
            y = 1.2

    for line in (an or "").split("\n"):
        line = line.rstrip()
        if not line.strip(): y += 0.1; continue

        # VEREDICTO
        if ("VEREDICTO" in line or
                (("PROCEDER" in line or "PRECAUCIÓN" in line or "PRECAUCION" in line.upper())
                 and "🎯" in line)):
            is_p = "PROCEDER ✅" in line or ("PROCEDER" in line and "NO PROCEDER" not in line)
            is_pr = "PRECAUCIÓN" in line or "PRECAUCION" in line.upper()
            vc = C["verde"] if is_p else C["amarillo"] if is_pr else C["rojo"]
            cl = re.sub(r"\*\*","",line).strip()
            ensure(0.68)
            rect(cur_sl, ML, y, MW, 0.6, vc, None, 0)
            txt(cur_sl, cl, ML+0.15, y+0.1, MW-0.25, 0.44,
                sz=11, bold=True, col=C["blanco"])
            y += 0.70; continue

        # Encabezado bold
        if re.match(r"^\s*\*\*",line) or re.match(r"^\s*[✅⚠️📊💰🎫👥🛣️🔍💡📈🎯]",line):
            cl = re.sub(r"\*\*","",line).strip()
            ensure(0.42)
            rect(cur_sl, ML, y+0.05, 0.06, 0.28, C["dorado"])
            txt(cur_sl, cl, ML+0.14, y, MW-0.16, 0.4,
                sz=10.5, bold=True, col=C["blanco"])
            y += 0.42; continue

        # Bullet
        if re.match(r"^\s*[-•]\s",line):
            cl = re.sub(r"^\s*[-•]\s+","",line)
            cl = re.sub(r"\*\*(.+?)\*\*",r"\1",cl).strip()
            ensure(0.38)
            rect(cur_sl, ML+0.22, y+0.14, 0.08, 0.08, C["dorado"])
            txt(cur_sl, cl, ML+0.38, y, MW-0.42, 0.36,
                sz=9.5, col="DDDDDD")
            y += 0.38; continue

        # Nivel confianza
        if "confianza" in line.lower():
            cl = re.sub(r"\*\*","",line).strip()
            ensure(0.34); y += 0.06
            txt(cur_sl, cl, ML, y, MW, 0.3, sz=8.5, it=True, col="888888")
            y += 0.34; continue

        # Texto normal
        cl = re.sub(r"\*\*(.+?)\*\*",r"\1",line).strip()
        if not cl: continue
        ensure(0.36)
        txt(cur_sl, cl, ML, y, MW, 0.34, sz=9.5, col="CCCCCC")
        y += 0.36

    # ── Caja de tips al pie (x=0.35 y=5.55 — réplica exacta slide 4) ──
    tips_y = 5.55
    if y < tips_y:
        # Usar las posiciones exactas del ppttipo
        rect(cur_sl, 0.35, tips_y, 12.55, 1.65, "3A3A3A")
        txt(cur_sl, "La ubicación es el punto de partida — el éxito también depende de:",
            0.5, tips_y+0.1, 12.2, 0.38, sz=9.5, bold=True, col=C["dorado"])
        tips = [
            (0.42,      "🎫 Ticket promedio",    "Calibra tu precio al NSE y competencia de la zona."),
            (4.553,     "🤝 Atención al cliente", "Tu rating en Google es tu segundo score de viabilidad."),
            (8.687,     "📱 Presencia digital",  "Google Maps, Instagram y WhatsApp Business son esenciales."),
        ]
        for tx, title, desc in tips:
            txt(cur_sl, title, tx, tips_y+0.55, 4.0, 0.32, sz=9, bold=True, col=C["dorado"])
            txt(cur_sl, desc, tx,  tips_y+0.88, 4.0, 0.55, sz=8, col=C["gris_t"])

    # Footer final
    txt(cur_sl, "4SITE · Don't guess. Foresee. · 4site.mx", 0.0, 7.18,
        13.33, 0.28, sz=7.5, col=C["gris_t"], al=PP_ALIGN.CENTER)

    # ══════════════════════════════════════════════════════════
    # Guardar
    # ══════════════════════════════════════════════════════════
    buf = BytesIO(); prs.save(buf); buf.seek(0)
    for f in tmp_files:
        try: os.unlink(f)
        except: pass
    return buf
