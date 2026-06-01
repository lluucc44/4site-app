import streamlit as st
import googlemaps
import os
from dotenv import load_dotenv
from anthropic import Anthropic
import json
from io import BytesIO
from reportlab.lib.pagesizes import letter
from reportlab.lib import colors
from reportlab.lib.units import inch
from reportlab.platypus import SimpleDocTemplate, Table, TableStyle, Paragraph, Spacer, Image as RLImage, PageBreak
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib.enums import TA_CENTER, TA_LEFT, TA_RIGHT
from reportlab.pdfgen import canvas
# PPTX — nuevo sistema de reportes (reemplaza PDF)
try:
    from generar_pptx_4site import generar_pptx_por_tier as _generar_pptx_por_tier
    PPTX_OK = True
except ImportError:
    PPTX_OK = False
    _generar_pptx_por_tier = None
from reportlab.pdfbase import pdfmetrics
from reportlab.pdfbase.ttfonts import TTFont
import requests
import re
import math
import hashlib
import datetime

load_dotenv()  # Carga el .env local (no afecta Streamlit Cloud)

# Módulos 4SITE — datos enriquecidos
try:
    from modulos_4site import (
        obtener_datos_inegi, clasificar_densidad,
        generar_reporte_trafico, calcular_mercado_potencial,
        generar_forecast, calcular_roi, DIAS_SEMANA,
        interpretar_demografia_negocio, calcular_ticket_competencia
    )
    MODULOS_OK = True
except ImportError:
    MODULOS_OK = False

try:
    from graficas_4site import (
        grafica_score_gauge, grafica_desglose_score,
        grafica_trafico_horario, grafica_trafico_semanal,
        grafica_forecast, grafica_mercado_donut,
        grafica_roi_dashboard, grafica_demografia,
        grafica_dashboard_premium, grafica_comparativa,
        grafica_edad_genero
    )
    GRAFICAS_OK = True
except ImportError:
    GRAFICAS_OK = False

try:
    from mapas_4site import (
        crear_mapa_competidores, crear_heatmap_competencia,
        crear_mapa_isocronas, crear_mapa_canibalizacion,
        crear_mapa_vialidades,
        render_mapa_con_interpretacion
    )
    MAPAS_OK = True
except ImportError:
    MAPAS_OK = False

# ORS API key para isócronas (opcional — gratuito en openrouteservice.org)
try:
    ORS_API_KEY = st.secrets.get("ORS_API_KEY", None)
except:
    ORS_API_KEY = os.getenv("ORS_API_KEY", None)

# ============================================
# CONFIGURACIÓN
# ============================================

st.set_page_config(
    page_title="4SITE - Location Analysis",
    page_icon="📍",
    layout="wide"
)

# ── CSS completo 4SITE embebido ──
def _cargar_css():
    st.markdown("""<style>
@import url('https://fonts.googleapis.com/css2?family=Montserrat:ital,wght@0,300;0,400;0,500;0,700;1,300&display=swap');

/* Tipografía global */
html, body, [class*="css"], .stApp, .stMarkdown, p, span, div, label, input, select, textarea, button {
    font-family: 'Montserrat', sans-serif !important;
}

/* Fondo principal — blanco limpio como en redes 4SITE */
.stApp { background-color: #FFFFFF !important; }
.main .block-container { 
    background-color: #FFFFFF !important; 
    padding: 2rem 2.5rem !important;
    max-width: 1100px;
}

/* Títulos */
h1,h2,h3,h4,h5 {
    font-family: 'Montserrat', sans-serif !important;
    font-weight: 700 !important;
    color: #0D0D0D !important;
}

/* Sidebar negro */
[data-testid="stSidebar"] {
    background-color: #0D0D0D !important;
    border-right: none !important;
}
[data-testid="stSidebar"] .stMarkdown,
[data-testid="stSidebar"] p,
[data-testid="stSidebar"] span,
[data-testid="stSidebar"] label,
[data-testid="stSidebar"] div {
    color: #CCCCCC !important;
    font-family: 'Montserrat', sans-serif !important;
}
[data-testid="stSidebar"] h1,
[data-testid="stSidebar"] h2,
[data-testid="stSidebar"] h3 {
    color: #FFFFFF !important;
}
[data-testid="stSidebar"] .stSelectbox label,
[data-testid="stSidebar"] .stTextInput label {
    color: #888888 !important;
    font-size: 10px !important;
    letter-spacing: 0.1em !important;
    text-transform: uppercase !important;
}
/* Radio sidebar */
[data-testid="stSidebar"] [data-testid="stRadio"] label { color: #CCCCCC !important; }

/* Botón primario — dorado */
div.stButton > button[kind="primary"],
div.stButton > button[data-testid="baseButton-primary"] {
    background: linear-gradient(135deg, #BB944F, #D4A85A) !important;
    color: #0D0D0D !important;
    border: none !important;
    font-weight: 700 !important;
    border-radius: 10px !important;
    padding: 0.5rem 1.5rem !important;
    font-size: 14px !important;
    box-shadow: 0 2px 8px rgba(187,148,79,0.3) !important;
    transition: all 0.2s !important;
}
div.stButton > button[kind="primary"]:hover {
    background: linear-gradient(135deg, #D4A85A, #E8BC6A) !important;
    box-shadow: 0 4px 14px rgba(187,148,79,0.4) !important;
}

/* Botón secundario — negro */
div.stButton > button:not([kind="primary"]) {
    background-color: #0D0D0D !important;
    color: #FFFFFF !important;
    border: none !important;
    font-weight: 500 !important;
    border-radius: 10px !important;
    font-family: 'Montserrat', sans-serif !important;
}
div.stButton > button:not([kind="primary"]):hover {
    background-color: #2A2A2A !important;
}

/* Tabs */
.stTabs [data-baseweb="tab-list"] { gap: 4px; border-bottom: 1px solid #EBEBEB !important; }
.stTabs [data-baseweb="tab"] {
    font-family: 'Montserrat', sans-serif !important;
    font-weight: 500 !important;
    color: #6B6B6B !important;
    border-radius: 8px 8px 0 0 !important;
    padding: 8px 16px !important;
}
.stTabs [aria-selected="true"] {
    color: #BB944F !important;
    border-bottom: 3px solid #BB944F !important;
    font-weight: 700 !important;
}

/* Inputs */
.stTextInput input, .stSelectbox select {
    font-family: 'Montserrat', sans-serif !important;
    border-radius: 8px !important;
    border: 1.5px solid #D5D5D5 !important;
    background: #FAFAFA !important;
    color: #0D0D0D !important;
}
.stTextInput input:focus { border-color: #BB944F !important; box-shadow: 0 0 0 2px rgba(187,148,79,0.15) !important; }

/* Métricas */
[data-testid="stMetricValue"] { font-weight: 700 !important; font-family: 'Montserrat', sans-serif !important; color: #0D0D0D !important; font-size: 1.6rem !important; }
[data-testid="stMetricLabel"] { font-family: 'Montserrat', sans-serif !important; color: #6B6B6B !important; font-size: 11px !important; font-weight: 500 !important; text-transform: uppercase !important; letter-spacing: 0.06em !important; }
[data-testid="stMetricDelta"] { font-family: 'Montserrat', sans-serif !important; }

/* Cards / contenedores */
[data-testid="stExpander"] {
    border: 1px solid #EBEBEB !important;
    border-radius: 10px !important;
}
.streamlit-expanderHeader {
    font-family: 'Montserrat', sans-serif !important;
    font-weight: 600 !important;
    color: #0D0D0D !important;
}

/* Caption y texto pequeño */
.stCaption, [data-testid="stCaptionContainer"] {
    font-family: 'Montserrat', sans-serif !important;
    color: #6B6B6B !important;
    font-size: 12px !important;
}

/* Divider */
hr { border-color: #BB944F !important; opacity: 0.25; margin: 1rem 0; }

/* Info/Warning/Success boxes */
.stAlert { border-radius: 10px !important; font-family: 'Montserrat', sans-serif !important; }
[data-testid="stInfoBox"] { border-left: 3px solid #BB944F !important; background: #FFF9F0 !important; }

/* Radio buttons */
.stRadio [data-testid="stMarkdownContainer"] p { font-family: 'Montserrat', sans-serif !important; }
/* Radio button color — dorado 4SITE */
.stRadio [data-baseweb="radio"] [data-checked="true"] div:first-child {
    background-color: #BB944F !important;
    border-color: #BB944F !important;
}
.stRadio [data-baseweb="radio"] div:first-child {
    border-color: #BB944F !important;
}
/* Evitar el rosa/pink del theme por default */
[data-testid="stWidgetLabel"] { color: #0D0D0D !important; font-family: 'Montserrat', sans-serif !important; }

/* Selectbox */
[data-baseweb="select"] [data-testid="stMarkdownContainer"] { font-family: 'Montserrat', sans-serif !important; }

/* Scrollbar */
::-webkit-scrollbar { width: 6px; }
::-webkit-scrollbar-track { background: #F5F5F5; }
::-webkit-scrollbar-thumb { background: #BB944F; border-radius: 3px; }

</style>""", unsafe_allow_html=True)

_cargar_css()

GA_MEASUREMENT_ID = "G-WW1XFBQ1PN"
ga_script = f"""
<script async src="https://www.googletagmanager.com/gtag/js?id={GA_MEASUREMENT_ID}"></script>
<script>
  window.dataLayer = window.dataLayer || [];
  function gtag(){{dataLayer.push(arguments);}}
  gtag('js', new Date());
  gtag('config', '{GA_MEASUREMENT_ID}');
</script>
"""
st.markdown(ga_script, unsafe_allow_html=True)

# ============================================
# META PIXEL (browser fallback)
# ============================================
try:
    META_PIXEL_ID = st.secrets.get("META_PIXEL_ID", "1436219981860379")
except:
    META_PIXEL_ID = os.getenv("META_PIXEL_ID", "1436219981860379")

st.markdown(f"""
<script>
!function(f,b,e,v,n,t,s)
{{if(f.fbq)return;n=f.fbq=function(){{n.callMethod?
n.callMethod.apply(n,arguments):n.queue.push(arguments)}};
if(!f._fbq)f._fbq=n;n.push=n;n.loaded=!0;n.version='2.0';
n.queue=[];t=b.createElement(e);t.async=!0;
t.src=v;s=b.getElementsByTagName(e)[0];
s.parentNode.insertBefore(t,s)}}(window, document,'script',
'https://connect.facebook.net/en_US/fbevents.js');
fbq('init', '{META_PIXEL_ID}');
fbq('track', 'PageView');
</script>
<noscript><img height="1" width="1" style="display:none"
src="https://www.facebook.com/tr?id={META_PIXEL_ID}&ev=PageView&noscript=1"/></noscript>
""", unsafe_allow_html=True)

# ============================================
# META CONVERSIONS API (server-side)
# ============================================
try:
    META_CAPI_TOKEN = st.secrets.get("META_CAPI_TOKEN", "")
except:
    META_CAPI_TOKEN = os.getenv("META_CAPI_TOKEN", "")

META_TEST_CODE = "TEST87947"  # Quitar cuando termines pruebas

def _meta_capi(event_name, value=None, currency="MXN", content_name=None, event_id=None):
    """Envía evento a Meta Conversions API desde el servidor."""
    if not META_CAPI_TOKEN:
        return
    import time, uuid
    payload = {
        "data": [{
            "event_name": event_name,
            "event_time": int(time.time()),
            "event_id": event_id or str(uuid.uuid4()),
            "action_source": "website",
            "event_source_url": "https://4site-app-mx.streamlit.app",
            "custom_data": {}
        }]
    }
    if value is not None:
        payload["data"][0]["custom_data"]["value"] = value
        payload["data"][0]["custom_data"]["currency"] = currency
    if content_name:
        payload["data"][0]["custom_data"]["content_name"] = content_name
    if META_TEST_CODE:
        payload["test_event_code"] = META_TEST_CODE
    try:
        requests.post(
            "https://graph.facebook.com/v18.0/" + META_PIXEL_ID + "/events",
            params={"access_token": META_CAPI_TOKEN},
            json=payload,
            timeout=3
        )
    except Exception:
        pass

# PageView al cargar (una vez por sesión)
if "meta_pageview_sent" not in st.session_state:
    _meta_capi("PageView")
    st.session_state["meta_pageview_sent"] = True

try:
    GOOGLE_API_KEY = st.secrets["GOOGLE_API_KEY"]
    CLAUDE_API_KEY = st.secrets["CLAUDE_API_KEY"]
except:
    GOOGLE_API_KEY = os.getenv("GOOGLE_API_KEY", "")
    CLAUDE_API_KEY = os.getenv("CLAUDE_API_KEY", "")

# Google Sheet donde el webhook guarda los códigos
SHEET_ID = "107pcvu1mgINU5AqJhiHLeBMQF3U8FqNhDTxVvPqNSsg"
SHEET_CSV_URL = f"https://docs.google.com/spreadsheets/d/{SHEET_ID}/export?format=csv"

def cargar_codigos_sheet():
    """Lee códigos activos desde Google Sheets. Retorna dict {codigo: tier}"""
    try:
        import pandas as pd
        df = pd.read_csv(SHEET_CSV_URL)
        df.columns = [c.strip().lower() for c in df.columns]
        codigos = {}
        for _, row in df.iterrows():
            cod   = str(row.get("codigo", "")).strip().upper()
            tier  = str(row.get("tier", "")).strip().lower()
            estado = str(row.get("estado", "")).strip().lower()
            if cod and tier and estado == "activo":
                codigos[cod] = tier
        return codigos
    except Exception as e:
        return {}

gmaps  = googlemaps.Client(key=GOOGLE_API_KEY)
claude = Anthropic(api_key=CLAUDE_API_KEY)

# ============================================
# SISTEMA DE TIERS
# ============================================

TIERS = {
    "free": {
        "nombre": "Gratis",
        "precio": "$0",
        "color": "#9E9E9E",
        "competidores_mostrados": 3,   # filas en tabla
        "top_negocios": 3,              # en modo recomendador
        "muestra_demografia": False,
        "muestra_badges": False,
        "muestra_score_desglose": False,
        "muestra_fotos_comp": False,
        "muestra_mapa_comp": False,
        "muestra_forecast": False,
        "muestra_mercado": False,
        "muestra_roi": False,
        "muestra_comparativa": False,
        "pdf_paginas": "Análisis básico",
        "watermark": True,
    },
    "basic": {
        "nombre": "Básico",
        "precio": "$99 MXN",  # ✅ precio actualizado
        "color": "#0D0D0D",
        "competidores_mostrados": 10,
        "top_negocios": 5,
        "muestra_demografia": True,
        "muestra_badges": True,
        "muestra_score_desglose": True,
        "muestra_fotos_comp": False,    # Fase 2
        "muestra_mapa_comp": False,     # Fase 2
        "muestra_forecast": False,
        "muestra_mercado": False,
        "muestra_roi": False,
        "muestra_comparativa": False,
        "pdf_paginas": "Análisis completo",
        "watermark": False,
    },
    "pro": {
        "nombre": "PRO",
        "precio": "$499 MXN",
        "color": "#BB944F",
        "competidores_mostrados": 20,
        "top_negocios": 10,
        "muestra_demografia": True,
        "muestra_badges": True,
        "muestra_score_desglose": True,
        "muestra_fotos_comp": True,     # Fase 2
        "muestra_mapa_comp": True,      # Fase 2
        "muestra_forecast": True,       # Fase 2
        "muestra_mercado": True,        # Fase 2
        "muestra_roi": False,
        "muestra_comparativa": False,
        "pdf_paginas": "Análisis PRO",
        "watermark": False,
    },
    "premium": {
        "nombre": "PREMIUM",
        "precio": "$999 MXN",
        "color": "#FFD700",
        "competidores_mostrados": 20,
        "top_negocios": 15,
        "muestra_demografia": True,
        "muestra_badges": True,
        "muestra_score_desglose": True,
        "muestra_fotos_comp": True,
        "muestra_mapa_comp": True,
        "muestra_forecast": True,
        "muestra_mercado": True,
        "muestra_roi": True,
        "muestra_comparativa": True,
        "pdf_paginas": "Análisis PREMIUM",
        "watermark": False,
    },
    "premium2": {
        "nombre": "🔍 ZONA ELITE",
        "precio": "$1,500 MXN",
        "color": "#BB944F",
        "competidores_mostrados": 20,
        "top_negocios": 15,
        "muestra_demografia": True,
        "muestra_badges": True,
        "muestra_score_desglose": True,
        "muestra_fotos_comp": True,
        "muestra_mapa_comp": True,
        "muestra_forecast": True,
        "muestra_mercado": True,
        "muestra_roi": True,
        "muestra_comparativa": False,
        "pdf_paginas": "Reporte Zona Potencial",
        "watermark": False,
    }
}

# ============================================
# SISTEMA DE USOS — límite por tier
# ============================================
# Básico  → 1 uso  (código de un solo análisis)
# PRO     → 2 usos
# Premium → 3 usos
# Los usos se validan en Google Sheets (estado activo/usado).
# El archivo JSON es solo respaldo local para la sesión.

MAX_USOS_POR_TIER = {
    "basic":    1,
    "basico":   1,
    "pro":      2,
    "premium":  3,
    "premium2": 3,
    "free":     0,
}
ARCHIVO_USOS_JSON = os.path.join(os.path.dirname(os.path.abspath(__file__)), "codigos_usos.json")

def _max_usos_para_codigo(codigo: str) -> int:
    """Infiere el máximo de usos permitidos según el prefijo del código."""
    codigo = codigo.strip().upper()
    if "4SITE-99-"  in codigo or codigo.startswith("4SITE-99-"):  return 1   # básico
    if "4SITE-499-" in codigo or codigo.startswith("4SITE-499-"): return 2   # pro
    if "4SITE-999-" in codigo or codigo.startswith("4SITE-999-"): return 3   # premium
    # Fallback: leer tier del Sheet si no podemos inferirlo del prefijo
    return 1

def _cargar_conteo_usos() -> dict:
    """Devuelve dict {CODIGO: usos_totales} desde archivo JSON."""
    try:
        if os.path.exists(ARCHIVO_USOS_JSON):
            with open(ARCHIVO_USOS_JSON, "r") as f:
                return json.load(f)
    except Exception:
        pass
    return {}

def _guardar_conteo_usos(conteo: dict):
    """Persiste el conteo de usos en archivo JSON."""
    try:
        with open(ARCHIVO_USOS_JSON, "w") as f:
            json.dump(conteo, f)
    except Exception:
        pass

def usos_restantes(codigo: str) -> int:
    """Retorna cuántos usos le quedan al código según su tier."""
    codigo = codigo.strip().upper()
    es_test = codigo.endswith("-TEST") or codigo == "TEST"
    if es_test:
        return 99  # ilimitado

    max_usos = _max_usos_para_codigo(codigo)
    conteo = _cargar_conteo_usos()
    total_usado = conteo.get(codigo, 0)
    return max(max_usos - total_usado, 0)

def codigo_agotado(codigo: str) -> bool:
    """True si el código ya no tiene usos disponibles."""
    return usos_restantes(codigo) <= 0

def _registrar_uso(codigo: str):
    """Incrementa el contador de usos del código (archivo + session_state)."""
    codigo = codigo.strip().upper()
    es_test = codigo.endswith("-TEST") or codigo == "TEST"
    if es_test:
        return

    # Actualizar archivo
    conteo = _cargar_conteo_usos()
    conteo[codigo] = conteo.get(codigo, 0) + 1
    _guardar_conteo_usos(conteo)

    # Actualizar sesión actual
    if "usos_sesion_actual" not in st.session_state:
        st.session_state["usos_sesion_actual"] = {}
    st.session_state["usos_sesion_actual"][codigo] = (
        st.session_state["usos_sesion_actual"].get(codigo, 0) + 1
    )

def validar_codigo(codigo_ingresado):
    """
    Valida código contra Google Sheets y retorna el tier correspondiente.
    Retorna "agotado" si el código superó los MAX_USOS_POR_CODIGO análisis.
    Códigos de prueba (-TEST) nunca se agotan.
    """
    codigo = codigo_ingresado.strip().upper()
    if not codigo:
        return "free"

    es_test = codigo.endswith("-TEST") or codigo == "TEST"

    # Verificar usos disponibles (solo códigos reales)
    if not es_test and codigo_agotado(codigo):
        return "agotado"

    # Códigos de prueba hardcodeados como fallback
    if codigo == "BASIC-TEST":
        return "basic"
    if codigo == "PRO-TEST":
        return "pro"
    if codigo == "PREM-TEST":
        return "premium"
    if codigo == "ZONA-TEST":
        return "premium2"

    # Leer de Google Sheets
    codigos_sheet = cargar_codigos_sheet()
    if codigo in codigos_sheet:
        tier = codigos_sheet[codigo]
        # Normalizar nombre del tier
        if tier in ["basico", "basic"]:
            return "basic"
        if tier in ["pro"]:
            return "pro"
        if tier in ["premium"]:
            return "premium"
        if tier in ["premium2", "zona_elite", "zona elite"]:
            return "premium2"

    return None  # Código inválido

def activar_codigo(codigo):
    """
    Registra un uso del código:
    1. Localmente en session_state (respaldo inmediato)
    2. En Google Sheets vía Apps Script (persistente entre sesiones)
    """
    _registrar_uso(codigo)

    # Marcar como 'usado' en Google Sheets cuando se agoten los usos
    es_test = codigo.strip().upper().endswith("-TEST") or codigo.strip().upper() == "TEST"
    if not es_test:
        try:
            _restantes = usos_restantes(codigo)  # ya decrementado por _registrar_uso
            _url    = st.secrets.get("APPS_SCRIPT_URL", "")
            _secret = st.secrets.get("APPS_SCRIPT_SECRET", "")
            if _url and _secret and _restantes <= 0:
                # Solo marca 'usado' en Sheets cuando ya no quedan usos
                requests.get(
                    _url,
                    params={"action": "marcar_usado", "codigo": codigo.strip().upper(), "secret": _secret},
                    timeout=5
                )
        except Exception:
            pass  # Si falla la llamada, el conteo local sigue activo como respaldo

# ============================================
# TIPOS DE NEGOCIO
# ============================================

TIPOS_NEGOCIO = {
    "cafe_premium":       {"nombre": "☕ Cafetería Premium/Specialty",   "keywords": ["cafe","coffee","specialty","cafetería"],               "inversion_min": 300000,  "inversion_max": 600000,  "competencia_optima": 2, "descripcion": "Café de especialidad, third wave, ambiente trendy"},
    "cafe_casual":        {"nombre": "☕ Cafetería Casual/Cadena",        "keywords": ["cafe","coffee","starbucks","italian coffee"],          "inversion_min": 400000,  "inversion_max": 800000,  "competencia_optima": 3, "descripcion": "Café estilo cadena, servicio rápido"},
    "restaurante_casual": {"nombre": "🍕 Restaurante Casual",             "keywords": ["restaurant","restaurante","comida","food"],            "inversion_min": 500000,  "inversion_max": 1000000, "competencia_optima": 4, "descripcion": "Restaurante familiar, menú variado"},
    "restaurante_fino":   {"nombre": "🍽️ Restaurante Fino/Gourmet",      "keywords": ["restaurant","fine dining","gourmet"],                  "inversion_min": 1000000, "inversion_max": 3000000, "competencia_optima": 1, "descripcion": "Alta cocina, experiencia premium"},
    "comida_rapida":      {"nombre": "🌮 Comida Rápida/Fast Food",        "keywords": ["fast food","taco","burger","torta","comida rápida"],   "inversion_min": 200000,  "inversion_max": 500000,  "competencia_optima": 5, "descripcion": "Servicio rápido, alto volumen"},
    "gimnasio_boutique":  {"nombre": "🏋️ Gimnasio Boutique/Crossfit",    "keywords": ["gym","gimnasio","crossfit","fitness boutique"],        "inversion_min": 800000,  "inversion_max": 1500000, "competencia_optima": 1, "descripcion": "Fitness especializado, clases grupales"},
    "gimnasio_regular":   {"nombre": "🏋️ Gimnasio Regular/Cadena",       "keywords": ["gym","gimnasio","fitness","sportium"],                 "inversion_min": 1500000, "inversion_max": 3000000, "competencia_optima": 2, "descripcion": "Gym completo, equipamiento variado"},
    "farmacia":           {"nombre": "💊 Farmacia",                        "keywords": ["pharmacy","farmacia","drogueria"],                     "inversion_min": 400000,  "inversion_max": 800000,  "competencia_optima": 2, "descripcion": "Venta de medicamentos y productos de salud"},
    "tienda_conveniencia":{"nombre": "🏪 Tienda de Conveniencia",          "keywords": ["convenience store","oxxo","7-eleven","minisuper"],     "inversion_min": 300000,  "inversion_max": 600000,  "competencia_optima": 2, "descripcion": "Productos básicos 24/7"},
    "panaderia":          {"nombre": "🥖 Panadería/Repostería",            "keywords": ["bakery","panaderia","pastry","reposteria"],            "inversion_min": 250000,  "inversion_max": 500000,  "competencia_optima": 3, "descripcion": "Pan artesanal, pasteles, café"},
    "bar":                {"nombre": "🍺 Bar/Cantina",                     "keywords": ["bar","pub","cantina","cerveceria"],                    "inversion_min": 600000,  "inversion_max": 1200000, "competencia_optima": 3, "descripcion": "Bebidas alcohólicas, ambiente social"},
    "yoga_wellness":      {"nombre": "🧘 Yoga/Wellness Studio",            "keywords": ["yoga","wellness","spa","pilates"],                     "inversion_min": 400000,  "inversion_max": 800000,  "competencia_optima": 1, "descripcion": "Clases de yoga, meditación, bienestar"},
    "guarderia":          {"nombre": "👶 Guardería/Kinder",                "keywords": ["daycare","guarderia","kinder","preescolar"],           "inversion_min": 500000,  "inversion_max": 1000000, "competencia_optima": 2, "descripcion": "Cuidado infantil, educación temprana"},
    "libreria":           {"nombre": "📚 Librería/Papelería",              "keywords": ["bookstore","libreria","papeleria","stationery"],       "inversion_min": 300000,  "inversion_max": 600000,  "competencia_optima": 2, "descripcion": "Libros, útiles escolares, regalos"},
    "servicios":          {"nombre": "💇 Servicios (Peluquería/Tintorería)","keywords": ["salon","peluqueria","barberia","tintoreria"],         "inversion_min": 200000,  "inversion_max": 400000,  "competencia_optima": 3, "descripcion": "Servicios personales del día a día"},
    "ferreteria":         {"nombre": "🔧 Ferretería / Materiales Eléctricos",      "keywords": ["ferreteria","hardware","electricidad","materiales","herramientas","plomeria"], "inversion_min": 350000, "inversion_max": 800000, "competencia_optima": 2, "descripcion": "Herramientas, material eléctrico, plomería, construcción"},
    "acabados_hogar":     {"nombre": "🪟 Acabados del Hogar (Pisos/Persianas/Laminado)", "keywords": ["persiana","piso","laminado","alfombra","lambrin","decoracion","revestimiento","cortina","tapiz","duela"], "inversion_min": 400000, "inversion_max": 900000, "competencia_optima": 2, "descripcion": "Persianas, pisos laminados, alfombras, lambrín, decoración de interiores"},
    "nutricion":          {"nombre": "🥗 Consultorio de Nutrición",                 "keywords": ["nutriologo","nutricionista","dietista","nutricion","dieta","nutri"], "inversion_min": 80000,  "inversion_max": 250000,  "competencia_optima": 1, "descripcion": "Consulta nutricional, planes de dieta, control de peso"},
    "medicina_estetica":  {"nombre": "💉 Medicina Estética / Spa Médico",           "keywords": ["medicina estetica","med estetica","estetica medica","clinica estetica","spa medico","botox","relleno","laser","depilacion laser","dermatologia","rejuvenecimiento","tratamiento facial","liposuccion","mesoterapia","hilos","rinomodelacion","carboxiterapia","clinica spa","multiestetic","euroestetica"], "inversion_min": 500000, "inversion_max": 2000000, "competencia_optima": 1, "descripcion": "Tratamientos estéticos médicos, botox, rellenos, laser, spa médico"},
    "empanadas_antojitos":{"nombre": "🫔 Empanadas / Antojitos",                    "keywords": ["empanada","antojito","taco","quesadilla","gordita","tlayuda","sope"], "inversion_min": 80000,  "inversion_max": 250000,  "competencia_optima": 4, "descripcion": "Antojitos mexicanos, empanadas, comida rápida de alto volumen y ticket bajo-medio"},
}

TYPE_MAPPING = {
    "cafe_premium":       ["cafe", "restaurant"],
    "cafe_casual":        ["cafe", "restaurant"],
    "restaurante_casual": ["restaurant", "meal_delivery", "meal_takeaway"],
    "restaurante_fino":   ["restaurant"],
    "comida_rapida":      ["restaurant", "meal_takeaway"],
    "gimnasio_boutique":  ["gym"],
    "gimnasio_regular":   ["gym", "sports_complex"],
    "farmacia":           ["pharmacy", "drugstore"],
    "tienda_conveniencia":["convenience_store", "supermarket"],
    "panaderia":          ["bakery", "cafe", "restaurant"],
    "bar":                ["bar", "night_club"],
    "yoga_wellness":      ["gym", "spa", "beauty_salon"],
    "guarderia":          ["school", "primary_school"],
    "libreria":           ["book_store"],
    "servicios":          ["beauty_salon", "hair_care", "laundry", "spa"],
    "ferreteria":         ["hardware_store", "home_goods_store", "electrician", "plumber"],
    "acabados_hogar":     ["home_goods_store", "furniture_store", "interior_designer", "flooring_store"],
    "nutricion":          ["doctor", "health", "physiotherapist"],
    "medicina_estetica":  ["doctor", "beauty_salon", "spa", "health"],
    "empanadas_antojitos":["restaurant", "meal_takeaway", "bakery"],
}

# ============================================
# TICKET — CONFIABILIDAD Y RANGOS REALES
# ============================================
# Google Places price_level es poco confiable para servicios de salud/estética
# (salones de belleza reportan "bajo" aunque un botox cueste $5,000 MXN).
# Para estos tipos forzamos un rango real y mostramos advertencia.

TICKET_DATOS_REALES = {
    # tipos donde price_level de Google Places es representativo
    "confiable": [
        "cafe_premium","cafe_casual","restaurante_casual","restaurante_fino",
        "comida_rapida","panaderia","bar","tienda_conveniencia","libreria",
        "empanadas_antojitos",
    ],
    # tipos donde price_level NO es confiable — override con rangos de mercado MX
    "referencial": {
        "medicina_estetica": {
            "rango_txt":     "Servicios de alto valor — price_level Google Places no aplica",
            "ticket_min":    1500,
            "ticket_max":    8000,
            "ticket_tipico": 3500,
            "ejemplos":      "Limpieza facial $800–$1,500 · Botox $3,000–$6,000 · Rellenos $4,000–$12,000 · Láser $1,500–$4,000",
            "nota_ui":       (
                "⚠️ El dato de Google Places no refleja los precios reales de medicina estética — "
                "los negocios aparecen como 'sin precio' o 'económico' aunque sus tratamientos "
                "cuestan $1,500–$8,000+ MXN. El ticket mostrado es referencial basado en el mercado "
                "mexicano, no en los listings de Google."
            ),
        },
        "nutricion": {
            "rango_txt":     "Servicio profesional — price_level Google Places no aplica",
            "ticket_min":    400,
            "ticket_max":    1500,
            "ticket_tipico": 700,
            "ejemplos":      "Consulta inicial $600–$1,200 · Seguimiento $400–$800 · Paquete mensual $1,200–$2,500",
            "nota_ui":       (
                "⚠️ Los consultorios médicos y de nutrición rara vez publican precios en Google Places — "
                "el dato puede aparecer como $0 o 'económico'. El ticket mostrado es referencial "
                "basado en el mercado mexicano de consulta privada ($400–$1,500 MXN por visita)."
            ),
        },
        "gimnasio_boutique": {
            "rango_txt":     "Membresía/clases — price_level parcialmente confiable",
            "ticket_min":    500,
            "ticket_max":    2500,
            "ticket_tipico": 1200,
            "ejemplos":      "Clase suelta $200–$400 · Membresía mensual $800–$2,500 · Pack 10 clases $1,500–$3,500",
            "nota_ui":       (
                "ℹ️ El price_level de Google para gimnasios boutique puede subestimar el ticket real. "
                "Considera que una membresía mensual suele estar entre $800–$2,500 MXN, "
                "no el gasto por visita individual."
            ),
        },
        "yoga_wellness": {
            "rango_txt":     "Membresía/clases — price_level parcialmente confiable",
            "ticket_min":    300,
            "ticket_max":    1800,
            "ticket_tipico": 900,
            "ejemplos":      "Clase suelta $150–$350 · Membresía mensual $600–$1,800 · Retiro $2,500–$8,000",
            "nota_ui":       (
                "ℹ️ Los estudios de yoga/wellness suelen no publicar precios en Google. "
                "El ticket mostrado es referencial — la membresía mensual típica en México es $600–$1,800 MXN."
            ),
        },
        "acabados_hogar": {
            "rango_txt":     "Compra por proyecto — price_level no aplica",
            "ticket_min":    2000,
            "ticket_max":    30000,
            "ticket_tipico": 8000,
            "ejemplos":      "Cortinas/persianas $3,000–$15,000 · Piso laminado m² $400–$900 · Proyecto completo $15,000–$80,000",
            "nota_ui":       (
                "⚠️ Los negocios de acabados del hogar no tienen 'ticket por visita' — "
                "la compra es por proyecto. Google Places no tiene datos de precio confiables. "
                "El rango mostrado es estimado por proyecto típico en México."
            ),
        },
        "ferreteria": {
            "rango_txt":     "Compra variable por proyecto — referencial",
            "ticket_min":    200,
            "ticket_max":    5000,
            "ticket_tipico": 800,
            "ejemplos":      "Compra ocasional $200–$600 · Proyecto pequeño $1,500–$5,000 · Remodelación $10,000–$50,000",
            "nota_ui":       (
                "ℹ️ El ticket de ferretería varía enormemente según si el cliente compra por proyecto "
                "o en pequeñas visitas. Google Places no publica precios para este giro. "
                "El dato mostrado es referencial basado en comportamiento de compra típico en México."
            ),
        },
        "guarderia": {
            "rango_txt":     "Mensualidad — price_level no aplica",
            "ticket_min":    2500,
            "ticket_max":    8000,
            "ticket_tipico": 4000,
            "ejemplos":      "Guardería básica $2,500–$4,000/mes · Kinder privado $4,000–$8,000/mes · Tiempo completo",
            "nota_ui":       (
                "ℹ️ Las guarderías y kinders cobran mensualidad, no por visita. "
                "Google Places no refleja estos precios. El rango mostrado es la mensualidad "
                "típica de cuidado infantil privado en México."
            ),
        },
    }
}

def _get_ticket_nota(tipo_key: str) -> str | None:
    """Retorna la nota de confiabilidad del ticket para un tipo de negocio, o None si es confiable."""
    ref = TICKET_DATOS_REALES["referencial"]
    if tipo_key in ref:
        return ref[tipo_key]["nota_ui"]
    return None

def _get_ticket_override(tipo_key: str) -> dict | None:
    """Retorna el override de ticket para tipos con datos no confiables de Google Places."""
    ref = TICKET_DATOS_REALES["referencial"]
    if tipo_key in ref:
        return ref[tipo_key]
    return None

# ============================================
# DICCIONARIO BILINGÜE
# ============================================

TEXTOS = {
    "es": {
        "titulo": "4SITE",
        "subtitulo": "Don't guess. Foresee.",
        "selector_idioma": "Idioma / Language",
        "input_ubicacion": "📍 Ingresa la ubicación que quieres analizar",
        "placeholder_direccion": "Ej: Av. Insurgentes Sur 1458, CDMX",
        "boton_analizar": "🔍 Analizar Ubicación",
        "modo_analisis": "🎯 ¿Qué quieres hacer?",
        "modo_validar": "Validar mi idea de negocio",
        "modo_recomendar": "Descubrir qué negocio es mejor",
        "selecciona_tipo": "Selecciona el tipo de negocio:",
        "titulo_analisis": "📊 Análisis de Ubicación",
        "score_viabilidad": "Score de Viabilidad",
        "resumen": "Resumen Ejecutivo",
        "competencia_titulo": "Competencia Cercana (500m)",
        "tabla_nombre": "Nombre",
        "tabla_rating": "Rating",
        "tabla_reseñas": "Reseñas",
        "descargar_pdf": "📄 Descargar Reporte PDF",
        "top_recomendaciones": "🎯 TOP NEGOCIOS RECOMENDADOS",
        "mejor_opcion": "⭐ MEJOR OPCIÓN",
        "inversion_estimada": "💰 Inversión estimada",
        "evitar": "❌ Evitar estos negocios",
        "watermark": "Reporte Gratis · 4SITE · Actualiza para análisis completo",
        "ingresa_codigo": "🔑 ¿Tienes un código de acceso?",
        "placeholder_codigo": "Ej: BASIC-ABC123",
        "validar_codigo": "Aplicar código",
        "codigo_invalido": "❌ Código inválido. Verifica e intenta de nuevo.",
        "codigo_valido": "✅ Código válido —",
        "sin_codigo": "Continuar gratis",
        "tier_activo": "Plan activo:",
        "desglose_score": "📊 Desglose del Score",
        "densidad_comp": "Densidad competencia",
        "calidad_comp": "Calidad competencia",
        "consolidacion": "Consolidación mercado",
        "demografía": "📍 Datos del Área (500m)",
        "poblacion": "Población estimada",
        "nse": "NSE Predominante",
        "densidad_hab": "Densidad",
        "datos_estimados": "💡 Datos estimados basados en la zona.",
        "cta_titulo": "🚀 ¿Quieres el análisis completo?",
        "cta_intro": "Upgrades disponibles:",
        "footer_derechos": "© 2026 4SITE - Todos los derechos reservados",
        "footer_contacto": "hola@4site.mx",
        "coming_soon": "🔜 Próximamente",
        "pro_badge": "PRO",
        "premium_badge": "PREMIUM",
    },
    "en": {
        "titulo": "4SITE",
        "subtitulo": "Don't guess. Foresee.",
        "selector_idioma": "Language / Idioma",
        "input_ubicacion": "📍 Enter the location you want to analyze",
        "placeholder_direccion": "Ex: 1458 Insurgentes Sur Ave, Mexico City",
        "boton_analizar": "🔍 Analyze Location",
        "modo_analisis": "🎯 What do you want to do?",
        "modo_validar": "Validate my business idea",
        "modo_recomendar": "Discover which business is best",
        "selecciona_tipo": "Select business type:",
        "titulo_analisis": "📊 Location Analysis",
        "score_viabilidad": "Viability Score",
        "resumen": "Executive Summary",
        "competencia_titulo": "Nearby Competition (500m)",
        "tabla_nombre": "Name",
        "tabla_rating": "Rating",
        "tabla_reseñas": "Reviews",
        "descargar_pdf": "📄 Download PDF Report",
        "top_recomendaciones": "🎯 TOP RECOMMENDED BUSINESSES",
        "mejor_opcion": "⭐ BEST OPTION",
        "inversion_estimada": "💰 Estimated investment",
        "evitar": "❌ Avoid these businesses",
        "watermark": "Free Report · 4SITE · Upgrade for full analysis",
        "ingresa_codigo": "🔑 Do you have an access code?",
        "placeholder_codigo": "Ex: BASIC-ABC123",
        "validar_codigo": "Apply code",
        "codigo_invalido": "❌ Invalid code. Check and try again.",
        "codigo_valido": "✅ Valid code —",
        "sin_codigo": "Continue for free",
        "tier_activo": "Active plan:",
        "desglose_score": "📊 Score Breakdown",
        "densidad_comp": "Competition density",
        "calidad_comp": "Competition quality",
        "consolidacion": "Market consolidation",
        "demografía": "📍 Area Data (500m)",
        "poblacion": "Estimated population",
        "nse": "Predominant NSE",
        "densidad_hab": "Density",
        "datos_estimados": "💡 Estimated data based on the area.",
        "cta_titulo": "🚀 Want the complete analysis?",
        "cta_intro": "Available upgrades:",
        "footer_derechos": "© 2026 4SITE - All rights reserved",
        "footer_contacto": "hello@4site.mx",
        "coming_soon": "🔜 Coming soon",
        "pro_badge": "PRO",
        "premium_badge": "PREMIUM",
    }
}

# ============================================
# FUNCIONES AUXILIARES - GEOCODING
# ============================================

def geocodificar_direccion(direccion):
    try:
        result = gmaps.geocode(direccion)
        if result:
            loc = result[0]['geometry']['location']
            return loc['lat'], loc['lng']
    except Exception as e:
        st.error(f"Error al geocodificar: {e}")
    return None, None

def geocodificar_inversa(lat, lng):
    try:
        result = gmaps.reverse_geocode((lat, lng))
        if result:
            return result[0]['formatted_address']
    except:
        pass
    return f"{lat:.6f}, {lng:.6f}"

# ============================================
# MOTOR VIAL — Clasificación de tipo de vialidad
# ============================================

# Clasificación de negocios por dependencia de tráfico de paso
NEGOCIOS_PASO = {
    "alto":  ["cafe_premium","cafe_casual","comida_rapida","tienda_conveniencia",
              "panaderia","farmacia","servicios","empanadas_antojitos"],
    "medio": ["restaurante_casual","bar","libreria","nutricion"],
    "bajo":  ["restaurante_fino","gimnasio_boutique","gimnasio_regular",
              "yoga_wellness","guarderia","medicina_estetica","acabados_hogar","ferreteria"],
}

# Palabras clave de vialidades principales en México
KEYWORDS_VIALIDAD_PRINCIPAL = [
    "autopista","periférico","periferico","circuito","anillo","viaducto",
    "eje vial","eje central","boulevard","blvd","calzada","paseo",
    "insurgentes","reforma","constituyentes","revolución","revolucion",
    "tecnológico","tecnologico","universitaria","churubusco","tlalpan",
    "zaragoza","canal nacional","canal del norte","indios verdes",
    "lázaro cárdenas","lazaro cardenas","adolfo lópez mateos","adolfo lopez mateos",
]

KEYWORDS_AVENIDA_SECUNDARIA = [
    "avenida","av.","av ","avenue",
]

KEYWORDS_CALLE_LOCAL = [
    "calle","cjon","callejón","callejon","privada","cerrada","andador",
    "manzana","mza","priv.",
]



# ══════════════════════════════════════════════════════════════════════════
# DETECCIÓN DE TIPO DE PLAZA COMERCIAL
# ══════════════════════════════════════════════════════════════════════════

# Keywords que identifican MALL CERRADO (Galerías, Antara, etc.)
_KEYWORDS_MALL = [
    "galerias","galerías","antara","perisur","santa fe mall","multiplaza",
    "forum","liverpool","palacio de hierro","plaza satélite","satélite",
    "interlomas","town center","midtown","parque delta","coyoacán mall",
    "fashion mall","lifestyle center","gran plaza","power center",
    "altabrisa","peninsula","cumbres","galerias","magnocentro",
    "punto valle","arcos vallarta","andares","zapopan","san agustín",
    "sendero","villahermosa","metepec mall","galerías metepec",
    "galerías toluca","outlet",
]
# Keywords que identifican STRIP CENTER (plaza abierta)
_KEYWORDS_STRIP = [
    "plaza ","c.c.","cc ","local ","bodega ","bodega comercial",
    "plaza comercial","centro comercial barrio","plaza de barrio",
    "plaza express","strip","nave ","bodega aurora",
]
# Palabras clave genéricas de PLAZA (requieren preguntar al usuario)
_KEYWORDS_PLAZA_GENERICA = [
    "plaza","mall","local","centre","center","cc","c.c","centro comercial","galería",
]

def detectar_tipo_plaza(direccion: str) -> str | None:
    """
    Detecta si la dirección corresponde a una plaza comercial.
    Retorna: 'mall' | 'strip_center' | 'plaza_ambigua' | None
    La detección es case-insensitive.
    """
    d = (direccion or "").lower()
    # Primero revisar si es mall cerrado conocido
    if any(kw in d for kw in _KEYWORDS_MALL):
        return "mall"
    # Luego strip center
    if any(kw in d for kw in _KEYWORDS_STRIP):
        return "strip_center"
    # Genérica — puede ser cualquiera
    if any(kw in d for kw in _KEYWORDS_PLAZA_GENERICA):
        return "plaza_ambigua"
    return None


def clasificar_vialidad_por_nombre(direccion):
    """
    Clasifica el tipo de vialidad basándose en el nombre de la dirección.
    Retorna: 'principal', 'avenida', 'local'
    """
    dl = direccion.lower()

    # Primero verificar vialidades principales conocidas
    for kw in KEYWORDS_VIALIDAD_PRINCIPAL:
        if kw in dl:
            return "principal"

    # Luego avenidas secundarias
    for kw in KEYWORDS_AVENIDA_SECUNDARIA:
        if kw in dl:
            return "avenida"

    # Calles locales
    for kw in KEYWORDS_CALLE_LOCAL:
        if kw in dl:
            return "local"

    return "desconocida"


def buscar_vialidad_principal_cercana(lat, lng, radio_metros=150):
    """
    Busca si hay una vialidad principal en el radio especificado.
    Usa Google Roads API (snap to roads) o reverse geocoding de puntos cercanos.
    Retorna: {'distancia_m': X, 'nombre': 'Av. Tecnológico', 'tipo': 'principal'}
    """
    # Estrategia: hacer reverse geocoding en 4 puntos alrededor de la ubicación
    # a diferentes distancias para detectar vialidades cercanas
    import math

    puntos_busqueda = []
    for distancia in [50, 100, 150]:
        for angulo in [0, 90, 180, 270]:  # Norte, Este, Sur, Oeste
            rad = math.radians(angulo)
            # Convertir metros a grados aproximados
            dlat = (distancia / 111320) * math.cos(rad)
            dlng = (distancia / (111320 * math.cos(math.radians(lat)))) * math.sin(rad)
            puntos_busqueda.append((lat + dlat, lng + dlng, distancia))

    vialidades_encontradas = []

    for p_lat, p_lng, dist in puntos_busqueda:
        try:
            result = gmaps.reverse_geocode((p_lat, p_lng))
            if result:
                for comp in result[0].get('address_components', []):
                    tipos_comp = comp.get('types', [])
                    nombre = comp.get('long_name', '')
                    if 'route' in tipos_comp:
                        tipo_vial = clasificar_vialidad_por_nombre(nombre)
                        if tipo_vial in ['principal', 'avenida']:
                            vialidades_encontradas.append({
                                'nombre': nombre,
                                'tipo': tipo_vial,
                                'distancia_m': dist
                            })
        except:
            pass

    if not vialidades_encontradas:
        return None

    # Retornar la más cercana y más importante
    vialidades_encontradas.sort(
        key=lambda x: (0 if x['tipo'] == 'principal' else 1, x['distancia_m'])
    )
    return vialidades_encontradas[0]


def calcular_factor_vial(tipo_vialidad_propia, vialidad_cercana, tipo_negocio_key):
    """
    Calcula el ajuste de score basado en la vialidad.
    tipo_vialidad_propia: 'principal', 'avenida', 'local', 'desconocida'
    vialidad_cercana: dict con 'tipo', 'distancia_m', 'nombre' o None
    tipo_negocio_key: clave del tipo de negocio
    """
    # Determinar dependencia del negocio al tráfico de paso
    if tipo_negocio_key in NEGOCIOS_PASO["alto"]:
        dep = "alto"
    elif tipo_negocio_key in NEGOCIOS_PASO["medio"]:
        dep = "medio"
    else:
        dep = "bajo"

    ajuste = 0
    descripcion = ""

    # ── Análisis de la vialidad propia ──
    if tipo_vialidad_propia == "principal":
        if dep == "alto":
            ajuste += 25
            descripcion = "✅ Sobre vialidad principal — máximo tráfico de paso para este tipo de negocio"
        elif dep == "medio":
            ajuste += 12
            descripcion = "✅ Sobre vialidad principal — buena visibilidad y accesibilidad"
        else:  # bajo
            ajuste += 5
            descripcion = "ℹ️ Sobre vialidad principal — accesibilidad alta aunque no es crítica para este negocio"

    elif tipo_vialidad_propia == "avenida":
        if dep == "alto":
            ajuste += 12
            descripcion = "🟡 Sobre avenida secundaria — flujo moderado para negocio de paso"
        elif dep == "medio":
            ajuste += 8
            descripcion = "✅ Sobre avenida secundaria — buena accesibilidad"
        else:
            ajuste += 5
            descripcion = "✅ Sobre avenida — accesibilidad adecuada para este negocio"

    elif tipo_vialidad_propia == "local":
        if dep == "alto":
            ajuste -= 15
            descripcion = "⚠️ Calle local — flujo peatonal bajo, crítico para negocio de paso"
        elif dep == "medio":
            ajuste -= 5
            descripcion = "🟡 Calle local — acceso limitado, considera visibilidad"
        else:  # bajo — destino
            ajuste += 8
            descripcion = "✅ Calle local/residencial — ventaja para negocio de destino (menos ruido, más tranquilo)"

    # ── Bonus/Penalización por vialidad principal cercana ──
    if vialidad_cercana and tipo_vialidad_propia in ["local", "avenida"]:
        dist = vialidad_cercana.get('distancia_m', 999)
        nombre_vial = vialidad_cercana.get('nombre', '')
        tipo_cercana = vialidad_cercana.get('tipo', '')

        if tipo_cercana == "principal":
            if dist <= 80:
                if dep == "alto":
                    ajuste += 15
                    descripcion += f" | ✅ A {dist}m de {nombre_vial} — captación de tráfico alta"
                elif dep == "medio":
                    ajuste += 8
                    descripcion += f" | ✅ A {dist}m de {nombre_vial}"
            elif dist <= 150:
                if dep == "alto":
                    ajuste += 8
                    descripcion += f" | 🟡 A {dist}m de {nombre_vial} — captación de tráfico parcial"
                elif dep == "medio":
                    ajuste += 4
                    descripcion += f" | 🟡 A {dist}m de {nombre_vial}"
            else:
                if dep == "alto":
                    ajuste -= 5
                    descripcion += f" | ⚠️ {nombre_vial} a {dist}m — distancia reduce captación significativamente"

    return ajuste, descripcion, tipo_vialidad_propia


def analizar_vialidad(lat, lng, direccion, tipo_negocio_key="cafe_premium"):
    """
    Función principal — analiza la vialidad y retorna el análisis completo.
    """
    tipo_propio = clasificar_vialidad_por_nombre(direccion)
    vialidad_cercana = None

    # Solo buscar vialidad cercana si estamos en calle local o desconocida
    if tipo_propio in ["local", "desconocida"]:
        try:
            vialidad_cercana = buscar_vialidad_principal_cercana(lat, lng, radio_metros=150)
        except:
            pass

    ajuste, descripcion, tipo_final = calcular_factor_vial(
        tipo_propio, vialidad_cercana, tipo_negocio_key
    )

    # Badge para UI
    if tipo_propio == "principal":
        badge = "🛣️ Vialidad principal"
    elif tipo_propio == "avenida":
        badge = "🚦 Avenida secundaria"
    elif tipo_propio == "local":
        if vialidad_cercana and vialidad_cercana.get('distancia_m', 999) <= 100:
            badge = f"📍 Calle local (cerca de {vialidad_cercana.get('nombre','')})"
        else:
            badge = "🏘️ Calle local"
    else:
        badge = "📍 Vialidad detectada"

    return {
        "tipo_vialidad":     tipo_propio,
        "vialidad_cercana":  vialidad_cercana,
        "ajuste_score":      ajuste,
        "descripcion":       descripcion,
        "badge":             badge,
        "dependencia_paso":  "alto" if tipo_negocio_key in NEGOCIOS_PASO["alto"] else
                             "medio" if tipo_negocio_key in NEGOCIOS_PASO["medio"] else "bajo",
    }


# ============================================
# FUNCIONES DE CONTEXTO Y DEMOGRAFÍA
# ============================================

def detectar_contexto_ubicacion(lat, lng, direccion, tipo_negocio_key="cafe_premium"):
    contexto = {"tipo_zona": "residencial", "pois_cercanos": [], "trafico": "medio",
                "badges": [], "analisis_vial": None}
    dl = direccion.lower()

    # ── Análisis vial mejorado ──
    av = None
    try:
        av = analizar_vialidad(lat, lng, direccion, tipo_negocio_key)
    except Exception as e_vial:
        pass  # Se maneja abajo con fallback

    # Si el motor vial no pudo detectar o retornó tipo desconocido,
    # intentar con reverse geocoding de Google
    if av is None or av.get("tipo_vialidad") == "desconocida":
        try:
            result = gmaps.reverse_geocode((lat, lng))
            if result:
                for comp in result[0].get('address_components', []):
                    if 'route' in comp.get('types', []):
                        nombre_ruta = comp.get('long_name', '')
                        tipo_detectado = clasificar_vialidad_por_nombre(nombre_ruta)
                        if tipo_detectado != "desconocida":
                            ajuste, desc, _ = calcular_factor_vial(
                                tipo_detectado, None, tipo_negocio_key)
                            badge_vial = (f"🛣️ {nombre_ruta}" if tipo_detectado == "principal"
                                         else f"🚦 {nombre_ruta}" if tipo_detectado == "avenida"
                                         else f"🏘️ {nombre_ruta}")
                            dep = ("alto" if tipo_negocio_key in NEGOCIOS_PASO["alto"]
                                   else "medio" if tipo_negocio_key in NEGOCIOS_PASO["medio"]
                                   else "bajo")
                            av = {
                                "tipo_vialidad":    tipo_detectado,
                                "vialidad_cercana": None,
                                "ajuste_score":     ajuste,
                                "descripcion":      desc,
                                "badge":            badge_vial,
                                "dependencia_paso": dep,
                            }
                            break
        except:
            pass

    # Si aún no tenemos análisis vial, crear uno básico por nombre de dirección
    if av is None or av.get("tipo_vialidad") == "desconocida":
        tipo_fallback = clasificar_vialidad_por_nombre(direccion)
        if tipo_fallback == "desconocida":
            # Default: avenida si no se puede determinar
            tipo_fallback = "avenida"
        ajuste_fb, desc_fb, _ = calcular_factor_vial(tipo_fallback, None, tipo_negocio_key)
        dep_fb = ("alto" if tipo_negocio_key in NEGOCIOS_PASO["alto"]
                  else "medio" if tipo_negocio_key in NEGOCIOS_PASO["medio"] else "bajo")
        badge_fb = {"principal": "🛣️ Vialidad principal",
                    "avenida": "🚦 Avenida secundaria",
                    "local": "🏘️ Calle local"}.get(tipo_fallback, "📍 Vialidad")
        av = {
            "tipo_vialidad":    tipo_fallback,
            "vialidad_cercana": None,
            "ajuste_score":     ajuste_fb,
            "descripcion":      desc_fb if desc_fb else f"Vialidad tipo {tipo_fallback} detectada",
            "badge":            badge_fb,
            "dependencia_paso": dep_fb,
        }

    # Guardar en contexto y actualizar tipo_zona
    contexto["analisis_vial"] = av
    contexto["badges"].append(av["badge"])
    if av["tipo_vialidad"] == "principal":
        contexto.update({"tipo_zona": "paso", "trafico": "alto"})
    elif av["tipo_vialidad"] == "avenida":
        contexto.update({"tipo_zona": "comercial", "trafico": "medio"})
    elif av["tipo_vialidad"] == "local":
        vc = av.get("vialidad_cercana")
        if vc and vc.get("distancia_m", 999) <= 80:
            contexto.update({"tipo_zona": "comercial", "trafico": "medio"})
        else:
            contexto.update({"tipo_zona": "residencial", "trafico": "bajo"})

    url = "https://places.googleapis.com/v1/places:searchNearby"
    headers = {"Content-Type": "application/json", "X-Goog-Api-Key": GOOGLE_API_KEY,
               "X-Goog-FieldMask": "places.displayName,places.types"}
    data = {
        "includedTypes": ["gas_station","shopping_mall","hospital","school","university","transit_station","park"],
        "maxResultCount": 20,
        "locationRestriction": {"circle": {"center": {"latitude": lat, "longitude": lng}, "radius": 500.0}}
    }
    try:
        resp = requests.post(url, headers=headers, json=data)
        if resp.status_code == 200:
            for place in resp.json().get('places', []):
                tipos = place.get('types', [])
                nombre = place.get('displayName', {}).get('text', '')
                if 'gas_station' in tipos:
                    contexto["pois_cercanos"].append({"tipo": "gasolinera", "nombre": nombre})
                    contexto["badges"].append("⛽ Gasolinera cercana")
                    contexto["tipo_zona"] = "paso"
                elif 'shopping_mall' in tipos:
                    contexto["pois_cercanos"].append({"tipo": "plaza", "nombre": nombre})
                    contexto["badges"].append("🏬 Plaza comercial")
                    contexto["tipo_zona"] = "comercial"
                elif 'hospital' in tipos:
                    contexto["pois_cercanos"].append({"tipo": "hospital", "nombre": nombre})
                    contexto["badges"].append("🏥 Hospital cercano")
                elif 'school' in tipos or 'university' in tipos:
                    contexto["pois_cercanos"].append({"tipo": "escuela", "nombre": nombre})
                    contexto["badges"].append("🎓 Escuela/Universidad")
                elif 'transit_station' in tipos:
                    contexto["pois_cercanos"].append({"tipo": "transporte", "nombre": nombre})
                    contexto["badges"].append("🚌 Transporte público")
    except:
        pass
    contexto["badges"] = list(set(contexto["badges"]))
    return contexto


def obtener_demografia(lat, lng, direccion=""):
    """Wrapper — usa INEGI Engine si está disponible, fallback si no"""
    if MODULOS_OK and direccion:
        try:
            return obtener_datos_inegi(lat, lng, direccion, gmaps)
        except:
            pass
    # Fallback básico
    dem = {
        "poblacion_estimada": 0, "densidad_hab_km2": 8000,
        "nse_predominante": "C", "ingreso_promedio_mensual": 15000,
        "distribucion_edad": {"0-17": 25, "18-35": 35, "36-55": 25, "56+": 15},
        "viviendas_habitadas": 0, "gasto_promedio_mensual": 12000,
        "ingreso_actual": 15000, "gasto_actual": 12000,
        "poblacion_actual": 0, "viviendas_actual": 0,
        "densidad_actual": 8000, "personas_manzana": 80,
        "tasa_crecimiento_pct": 0.55, "zona_detectada": "default",
        "fuente": "Estimación basada en zona",
    }
    try:
        result = gmaps.reverse_geocode((lat, lng))
        if result:
            for comp in result[0].get('address_components', []):
                tipos = comp.get('types', [])
                name = comp.get('long_name', '').lower()
                if 'locality' in tipos or 'administrative_area_level_1' in tipos:
                    if any(x in name for x in ['polanco','condesa','roma','santa fe','lomas']):
                        dem.update({"nse_predominante": "A/B", "ingreso_promedio_mensual": 35000, "densidad_hab_km2": 15000})
                    elif 'ciudad de méxico' in name or 'cdmx' in name:
                        dem.update({"nse_predominante": "C+", "ingreso_promedio_mensual": 18000, "densidad_hab_km2": 12000})
                    elif any(x in name for x in ['guadalajara','monterrey','zapopan','san pedro']):
                        dem.update({"nse_predominante": "B/C+", "ingreso_promedio_mensual": 20000, "densidad_hab_km2": 10000})
                    elif any(x in name for x in ['querétaro','puebla','león','mérida']):
                        dem.update({"nse_predominante": "C", "ingreso_promedio_mensual": 15000, "densidad_hab_km2": 8000})
                    elif any(x in name for x in ['estado de méxico','edomex','toluca']):
                        dem.update({"nse_predominante": "C/D+", "ingreso_promedio_mensual": 12000, "densidad_hab_km2": 9000})
    except:
        pass
    area_km2 = 0.785
    dem["poblacion_estimada"]  = int(dem["densidad_hab_km2"] * area_km2)
    dem["poblacion_actual"]    = dem["poblacion_estimada"]
    dem["viviendas_habitadas"] = int(dem["poblacion_estimada"] / 3.5)
    dem["viviendas_actual"]    = dem["viviendas_habitadas"]
    dem["densidad_actual"]     = dem["densidad_hab_km2"]
    dem["personas_manzana"]    = int(dem["densidad_hab_km2"] * 0.01)
    dem["ingreso_actual"]      = dem["ingreso_promedio_mensual"]
    dem["gasto_actual"]        = int(dem["ingreso_promedio_mensual"] * 0.8)
    dem["gasto_promedio_mensual"] = dem["gasto_actual"]
    return dem


def formatear_densidad(densidad_hab_km2):
    personas_manzana = int(densidad_hab_km2 * 0.01)
    if densidad_hab_km2 < 2000:
        return {"clasificacion": "Baja 📉", "personas_manzana": personas_manzana, "descripcion": "Zona suburbana o de baja densidad"}
    elif densidad_hab_km2 < 8000:
        return {"clasificacion": "Media 📊", "personas_manzana": personas_manzana, "descripcion": "Zona residencial típica"}
    elif densidad_hab_km2 < 15000:
        return {"clasificacion": "Alta 📈", "personas_manzana": personas_manzana, "descripcion": "Zona urbana densamente poblada - Ideal para comercio"}
    else:
        return {"clasificacion": "Muy Alta 🔥", "personas_manzana": personas_manzana, "descripcion": "Centro urbano de muy alta densidad"}

# ============================================
# FUNCIONES DE COMPETENCIA Y SCORING
# ============================================

def buscar_competencia_por_tipo(lat, lng, tipo_key, radio=500.0, tipo_plaza=None):
    """
    Busca competidores. Ajusta radio y max_results según tipo de plaza:
    - mall:         radio 1000m, más resultados (competencia interna relevante)
    - strip_center: radio 600m
    - None:         radio estándar 500m
    """
    tipo = TIPOS_NEGOCIO[tipo_key]
    included_types = TYPE_MAPPING.get(tipo_key, ["restaurant", "cafe", "store"])
    keywords = tipo.get("keywords", [])

    # Ajustar radio según tipo de plaza
    if tipo_plaza == "mall":
        radio_efectivo = 1000.0
        max_results    = 20
    elif tipo_plaza == "strip_center":
        radio_efectivo = 600.0
        max_results    = 20
    else:
        radio_efectivo = float(radio)
        max_results    = 20

    url = "https://places.googleapis.com/v1/places:searchNearby"
    headers = {
        "Content-Type": "application/json",
        "X-Goog-Api-Key": GOOGLE_API_KEY,
        "X-Goog-FieldMask": "places.displayName,places.rating,places.userRatingCount,places.location,places.primaryTypeDisplayName,places.id,places.types,places.photos,places.priceLevel"
    }
    data = {
        "includedTypes": included_types,
        "maxResultCount": max_results,
        "locationRestriction": {"circle": {"center": {"latitude": lat, "longitude": lng}, "radius": radio_efectivo}}
    }
    try:
        resp = requests.post(url, headers=headers, json=data)
        if resp.status_code != 200:
            return []
        places = resp.json().get('places', [])
        filtered = []

        # Palabras que CLARAMENTE indican no-competencia para medicina estética
        _excluir_medicina_estetica = [
            "uñas","nails","nail art","manicure","pedicure",
            "barberia","barbería","barber",
            "peluqueria","peluquería","hair salon","hair studio","salon de cabello",
            "tintes","corte de cabello","blow dry","blow out",
            "tatuaje","tattoo","piercing",
            "bronceado","bronzeado","tanning",
        ]
        # Palabras que CONFIRMAN que es competencia de medicina estética
        _incluir_medicina_estetica = [
            "estetica","estética","estetico","estético",
            "botox","relleno","laser","láser","depilacion laser","depilación láser",
            "dermatolog","medicina","clinica","clínica","spa medic",
            "rejuvenec","lifting","mesoterapia","hilos","rinomodel",
            "carboxiter","liposuccion","liposucción","sculptra","coolsculpt",
            "multiestetic","euroestetica","euroestética","centros estetic",
            "tratamiento facial","peeling","dermapen","hydrafacial",
            "biorevitaliz","biorrevitaliz",
        ]

        for p in places:
            nombre = p.get('displayName', {}).get('text', '').lower()
            tipos  = p.get('types', [])

            if tipo_key == "medicina_estetica":
                # Paso 1: descartar los claramente no-médicos
                if any(excl in nombre for excl in _excluir_medicina_estetica):
                    continue
                # Paso 2: incluir si tiene keyword de medicina estética en el nombre
                _tiene_kw = any(kw.lower() in nombre for kw in keywords) or \
                            any(inc in nombre for inc in _incluir_medicina_estetica)
                # Paso 3: incluir si Google lo clasifica como spa/beauty_salon/doctor/health
                # (no excluir beauty_salon porque así clasifica Google a las clínicas estéticas)
                _es_tipo_valido = any(t in ["beauty_salon","spa","doctor","health","physiotherapist"] for t in tipos)
                if _tiene_kw or _es_tipo_valido:
                    filtered.append(p)

            elif any(kw.lower() in nombre for kw in keywords) or any(t in included_types for t in tipos):
                filtered.append(p)
        return filtered
    except:
        return []


def buscar_competencia_general(lat, lng):
    url = "https://places.googleapis.com/v1/places:searchNearby"
    headers = {
        "Content-Type": "application/json",
        "X-Goog-Api-Key": GOOGLE_API_KEY,
        "X-Goog-FieldMask": "places.displayName,places.rating,places.userRatingCount,places.location,places.primaryTypeDisplayName,places.id,places.photos"
    }
    data = {
        "includedTypes": ["restaurant","cafe","bar","gym","pharmacy","store","bakery","book_store"],
        "maxResultCount": 20,
        "locationRestriction": {"circle": {"center": {"latitude": lat, "longitude": lng}, "radius": 500.0}}
    }
    try:
        resp = requests.post(url, headers=headers, json=data)
        if resp.status_code == 200:
            return resp.json().get('places', [])
    except:
        pass
    return []


def calcular_score_competencia(competidores, tipo_key):
    """Score metodológico: densidad 50% + calidad 30% + consolidación 20%"""
    tipo  = TIPOS_NEGOCIO[tipo_key]
    n     = len(competidores)
    optimo = tipo["competencia_optima"]

    # Clasificar si es negocio de destino o de paso
    # Destino: el cliente busca específicamente el lugar (presencia de competidores = mercado existente)
    # Paso: el cliente entra por visibilidad/conveniencia (menos competidores = mejor)
    _negocios_destino = {
        "acabados_hogar", "ferreteria", "gimnasio_boutique",
        "gimnasio_regular", "yoga_wellness", "guarderia",
        "libreria", "restaurante_fino", "servicios",
        "bar", "farmacia",
    }
    es_destino = tipo_key in _negocios_destino

    # Densidad de competencia
    if es_destino:
        # Para negocios de destino: 0 competidores = posible falta de mercado
        if n == 0:           d = 45   # Penalizar: sin competidores, sin mercado demostrado
        elif n == 1:         d = 75   # Hay mercado pero poca competencia → oportunidad
        elif n <= optimo:    d = 90   # Óptimo: mercado confirmado + espacio
        elif n <= optimo+2:  d = 65
        elif n <= optimo+4:  d = 45
        else:                d = 25
    else:
        # Para negocios de paso (café, tienda, comida rápida):
        # menos competidores = mejor
        if n == 0:           d = 100
        elif n <= optimo:    d = 85
        elif n <= optimo+2:  d = 65
        elif n <= optimo+4:  d = 45
        else:                d = 25

    # Calidad (ratings)
    ratings = [c.get('rating', 0) for c in competidores if c.get('rating')]
    avg_r = sum(ratings)/len(ratings) if ratings else 0
    if avg_r == 0:
        q = 60 if es_destino else 100  # Sin competidores: neutro para destino
    elif avg_r < 3.5:    q = 80
    elif avg_r < 4.0:    q = 60
    elif avg_r < 4.3:    q = 40
    else:                q = 20

    # Consolidación (reseñas)
    reviews = [c.get('userRatingCount', 0) for c in competidores]
    avg_rv = sum(reviews)/len(reviews) if reviews else 0
    if avg_rv == 0:
        c_score = 60 if es_destino else 100  # Sin datos: neutro para destino
    elif avg_rv < 100:   c_score = 80
    elif avg_rv < 500:   c_score = 60
    elif avg_rv < 1500:  c_score = 40
    else:                c_score = 20

    score_final = int(d*0.50 + q*0.30 + c_score*0.20)
    return score_final, {"densidad": int(d), "calidad": int(q), "consolidacion": int(c_score),
                         "num_competidores": n, "rating_promedio": round(avg_r, 1) if avg_r else 0,
                         "reseñas_promedio": int(avg_rv) if avg_rv else 0}


def ajustar_score_por_contexto(score_base, tipo_key, contexto):
    score = score_base
    zona  = contexto["tipo_zona"]
    pois  = [p["tipo"] for p in contexto["pois_cercanos"]]

    # ── Ajuste por POIs y tipo de zona (lógica original) ──
    if "cafe" in tipo_key:
        if zona == "paso" or "gasolinera" in pois: score += 15
        elif zona == "comercial": score += 10
        elif "escuela" in pois:   score += 10
        elif zona == "residencial": score += 5
    elif tipo_key == "comida_rapida":
        if zona == "paso": score += 20
        elif "gasolinera" in pois: score += 15
        elif "escuela" in pois:    score += 10
    elif tipo_key == "restaurante_casual":
        if zona == "comercial" or "plaza" in pois: score += 15
        elif zona == "paso": score += 10
    elif tipo_key == "restaurante_fino":
        if zona == "comercial" or "plaza" in pois: score += 10
        elif zona == "paso": score -= 10
    elif "gimnasio" in tipo_key:
        if zona == "residencial": score += 15
        elif "plaza" in pois:    score += 10
        elif zona == "paso":     score -= 5
    elif tipo_key == "farmacia":
        if "hospital" in pois:   score += 20
        elif zona == "residencial": score += 10
        elif zona == "paso":     score += 5
    elif tipo_key == "bar":
        if zona == "comercial" or "plaza" in pois: score += 10
        elif zona == "paso":        score -= 15
        elif zona == "residencial": score -= 10
    elif tipo_key == "tienda_conveniencia":
        if zona == "paso" or "gasolinera" in pois: score += 15
        elif zona == "residencial": score += 10
    elif tipo_key == "guarderia":
        if zona == "residencial": score += 15
        elif zona == "comercial": score += 5
        elif zona == "paso":      score -= 10
    elif tipo_key == "nutricion":
        # Consultorio de nutrición — semi-destino, funciona bien en zonas residenciales o comerciales tranquilas
        if zona == "residencial": score += 12
        elif zona == "comercial": score += 8
        elif "hospital" in pois:  score += 10
        elif zona == "paso":      score -= 5   # Calle de paso = ambiente no ideal para consulta
    elif tipo_key == "medicina_estetica":
        # Medicina estética — negocio de destino puro, necesita calle tranquila
        if zona == "residencial": score += 15
        elif zona == "comercial": score += 5
        elif "plaza" in pois:     score += 8
        elif zona == "paso":      score -= 15  # Vialidad de alto tráfico = imagen no aspiracional
    elif tipo_key == "empanadas_antojitos":
        # Alto volumen, ticket bajo-medio — necesita tráfico de paso
        if zona == "paso": score += 20
        elif "gasolinera" in pois: score += 15
        elif "escuela" in pois:    score += 12
        elif zona == "comercial":  score += 10
        elif zona == "residencial": score += 3

    # ── Ajuste por análisis vial (nuevo motor) ──
    analisis_vial = contexto.get("analisis_vial")
    if analisis_vial:
        ajuste_vial = analisis_vial.get("ajuste_score", 0)
        score += ajuste_vial

    return max(0, min(100, int(score)))


def ajustar_score_por_demografia(score_base, tipo_key, dem):
    score = score_base
    nse   = dem["nse_predominante"]
    densidad = dem["densidad_hab_km2"]
    edad_joven = dem["distribucion_edad"]["18-35"]
    edad_ninos = dem["distribucion_edad"]["0-17"]
    if tipo_key in ["cafe_premium","restaurante_fino","gimnasio_boutique","yoga_wellness"]:
        if "A" in nse or "B" in nse: score += 15
        elif "C+" in nse:             score += 5
        elif "D" in nse or "E" in nse: score -= 20
    elif tipo_key in ["comida_rapida","tienda_conveniencia"]:
        if "C" in nse or "D" in nse:  score += 10
        elif "A" in nse or "B" in nse: score -= 5
    elif tipo_key in ["acabados_hogar", "ferreteria"]:
        # Acabados del hogar: mejor en NSE B/C+ con vivienda propia y remodelación
        if "B" in nse or "C+" in nse:  score += 12
        elif "A" in nse:               score += 5   # Muy alto = más online/boutique
        elif "C/" in nse or "D" in nse: score -= 15  # Bajo poder adquisitivo = no remodelación
    elif tipo_key == "nutricion":
        # NSE C+ requerido mínimo — servicio de salud preventiva con costo
        if "A" in nse or "B" in nse:   score += 15
        elif "C+" in nse:              score += 10
        elif "C" in nse and "C+" not in nse: score += 0   # C sin plus: neutro
        elif "D" in nse or "E" in nse: score -= 20  # NSE bajo = poco gasto preventivo
    elif tipo_key == "medicina_estetica":
        # NSE A/B obligatorio — ticket alto, cliente elige por reputación no por precio
        if "A" in nse:                 score += 20
        elif "B" in nse:               score += 15
        elif "C+" in nse:              score += 5
        elif "C" in nse and "C+" not in nse: score -= 10
        elif "D" in nse or "E" in nse: score -= 30  # Inviable en NSE bajo
    elif tipo_key == "empanadas_antojitos":
        # Alta rotación y ticket bajo-medio — funciona en C y D, penalizar A/B leve
        if "C" in nse or "D" in nse:   score += 12
        elif "E" in nse:               score += 5
        elif "A" in nse or "B" in nse: score -= 5   # Zona premium = ticket esperado más alto
    if densidad > 12000:
        if tipo_key != "restaurante_fino": score += 10
    elif densidad < 5000:
        score -= 10
    if tipo_key in ["cafe_premium","gimnasio_boutique","bar","comida_rapida"]:
        if edad_joven > 35:  score += 10
        elif edad_joven < 25: score -= 10
    elif tipo_key in ["medicina_estetica","yoga_wellness","nutricion"]:
        # Perfil principal 25-55: edad_joven (18-35) + adultos
        if edad_joven > 30:  score += 8
        elif edad_joven < 20: score -= 5
    elif tipo_key == "empanadas_antojitos":
        # Funciona bien con población joven y niños
        if edad_joven > 30:  score += 8
    elif tipo_key == "guarderia":
        if edad_ninos > 25:  score += 15
        elif edad_ninos < 15: score -= 15
    return max(0, min(100, int(score)))



def ajustar_score_por_plaza(score_base: int, tipo_key: str, tipo_plaza: str | None,
                              zona_mall: str = "no_especificado",
                              nivel_renta: str = "desconocido",
                              ancla: str = "sin_ancla") -> tuple[int, dict]:
    """
    Ajusta el score y genera contexto adicional según el tipo de plaza.

    Strip Center:
    - +10 si hay ancla fuerte (Walmart, Chedraui → tráfico garantizado)
    - +5 si hay OXXO / farmacias (tráfico de paso)
    - Sin penalización de ROI (renta similar a calle)

    Mall Cerrado:
    - +15 si zona food_court + negocio de alimentos (tráfico de corredor garantizado)
    - +10 si cerca de ancla principal (Liverpool, Cine → cola de tráfico)
    - +5  planta baja / entrada
    - -10 piso superior (menor tráfico orgánico)
    - -10 nivel de renta muy alto (impacto en ROI)
    - El score de competencia se recalcula internamente con mayor peso a competidores internos
    """
    score = score_base
    nota_plaza = ""
    ajuste_roi_pct = 0  # % de penalización al ROI por renta

    if tipo_plaza is None:
        return score, {"tipo_plaza": None, "nota": "", "ajuste_roi_pct": 0}

    if tipo_plaza == "strip_center":
        nota_plaza = "📍 Strip Center / Plaza abierta"
        ancla_map = {
            "walmart":   (+10, "Walmart/Bodega Aurrerá como ancla garantiza flujo constante de visitantes."),
            "chedraui":  (+10, "Chedraui como ancla genera tráfico recurrente y familias de compra semanal."),
            "soriana":   (+8,  "Soriana como ancla aporta flujo familiar constante al strip center."),
            "oxxo":      (+5,  "OXXO genera tráfico de paso permanente, beneficia a todos los locales."),
            "farmacias": (+5,  "Farmacia cadena como ancla = clientela recurrente y de confianza."),
            "banco":     (+3,  "Banco/cajero atrae clientes con capacidad de compra activa."),
            "otra":      (+5,  "Tienda ancla identificada — tráfico compartido beneficia el local."),
            "sin_ancla": (0,   "Sin ancla identificada. El tráfico dependerá de la visibilidad del local."),
        }
        adj, msg_ancla = ancla_map.get(ancla, (0, ""))
        score += adj
        nota_plaza += (chr(10) + msg_ancla) if msg_ancla else ""
        nota_plaza += chr(10) + "Visibilidad desde la calle y acceso vehicular son factores clave de exito."

    elif tipo_plaza == "mall":
        nota_plaza = "🏬 Mall / Plaza Cerrada"
        # Ajuste por zona
        zona_adj_map = {
            "food_court":          +15,
            "ancla_principal":     +10,
            "planta_baja_entrada": +5,
            "moda_accesorios":     0,
            "servicios":           0,
            "piso_superior":       -10,
            "no_especificado":     0,
        }
        adj_zona = zona_adj_map.get(zona_mall, 0)
        score += adj_zona

        zona_labels = {
            "food_court":          "Food court — corredor de alto tráfico, excelente para alimentos y bebidas.",
            "ancla_principal":     "Cerca de ancla principal — tráfico de cola garantizado.",
            "planta_baja_entrada": "Planta baja / entrada — máxima visibilidad en el mall.",
            "moda_accesorios":     "Zona de moda/accesorios — tráfico selectivo de compras.",
            "servicios":           "Zona de servicios — clientela recurrente y de destino.",
            "piso_superior":       "Piso superior — tráfico naturalmente menor, requiere señalización y marketing interno.",
            "no_especificado":     "Zona dentro del mall no especificada.",
        }
        nota_plaza += chr(10) + zona_labels.get(zona_mall, "")

        # Penalización por renta
        renta_adj_map = {
            "desconocido": (0,    0),
            "bajo":        (0,    0),
            "medio":       (-5,   15),   # (score_adj, ajuste_roi_pct)
            "alto":        (-8,   30),
            "muy_alto":    (-12,  50),
        }
        adj_renta, ajuste_roi_pct = renta_adj_map.get(nivel_renta, (0, 0))
        score += adj_renta
        if ajuste_roi_pct > 0:
            nota_plaza += chr(10) + f"⚠️ Renta elevada ({nivel_renta}): el ROI se reduce ~{ajuste_roi_pct}% vs local en calle equivalente."
        nota_plaza += chr(10) + "💡 En mall cerrado el trafico esta garantizado pero el costo de renta impacta significativamente el ROI."

    return max(0, min(100, int(score))), {
        "tipo_plaza":     tipo_plaza,
        "nota":           nota_plaza,
        "ajuste_roi_pct": ajuste_roi_pct,
        "zona_mall":      zona_mall,
        "ancla":          ancla,
    }


def nivel_score(score, idioma="es"):
    if score >= 80:   return ("ALTO", "#4CAF50") if idioma == "es" else ("HIGH", "#4CAF50")
    elif score >= 65: return ("MEDIO-ALTO", "#8BC34A") if idioma == "es" else ("MEDIUM-HIGH", "#8BC34A")
    elif score >= 50: return ("MEDIO", "#FFC107") if idioma == "es" else ("MEDIUM", "#FFC107")
    else:             return ("BAJO", "#F44336") if idioma == "es" else ("LOW", "#F44336")


def recomendar_tipos_negocio(lat, lng, direccion):
    # Usar cafe_casual como tipo base para detección vial inicial
    # Se re-detectará con el top1 después de calcular recomendaciones
    contexto  = detectar_contexto_ubicacion(lat, lng, direccion, "cafe_casual")
    demografia = obtener_demografia(lat, lng)
    recomendaciones = []
    for tipo_key, tipo_info in TIPOS_NEGOCIO.items():
        comps = buscar_competencia_por_tipo(lat, lng, tipo_key)
        score_base, desglose = calcular_score_competencia(comps, tipo_key)
        score_ctx  = ajustar_score_por_contexto(score_base, tipo_key, contexto)
        score_final = ajustar_score_por_demografia(score_ctx, tipo_key, demografia)
        recomendaciones.append({
            "tipo_key": tipo_key, "nombre": tipo_info["nombre"],
            "score": score_final, "desglose": desglose,
            "num_competidores": len(comps),
            "inversion_min": tipo_info["inversion_min"],
            "inversion_max": tipo_info["inversion_max"],
            "descripcion": tipo_info["descripcion"],
            "competidores": comps,
        })
    recomendaciones.sort(key=lambda x: x["score"], reverse=True)
    return recomendaciones, contexto, demografia

# ============================================
# ANÁLISIS CLAUDE
# ============================================

def generar_analisis_claude(ubicacion, competidores, idioma, modo, tipo_negocio=None, recomendaciones=None, tipo_plaza=None, info_plaza=None):
    if modo == "validar":
        tipo_info = TIPOS_NEGOCIO[tipo_negocio]
        prompt = f"""Eres analista experto en ubicaciones comerciales para {tipo_info['nombre']}.

UBICACIÓN: {ubicacion}
Descripción del negocio: {tipo_info['descripcion']}
Inversión: ${tipo_info['inversion_min']:,}–${tipo_info['inversion_max']:,} MXN
{"TIPO DE PLAZA: " + ("Mall / Plaza Cerrada" if tipo_plaza == "mall" else "Strip Center" if tipo_plaza == "strip_center" else "No aplica") if tipo_plaza else "Ubicación en calle o edificio estándar"}
{f"Zona en mall: {info_plaza.get('zona_mall','')}" if tipo_plaza == 'mall' and info_plaza else ""}
{f"Nota de plaza: {info_plaza.get('nota','')}" if info_plaza and info_plaza.get('nota') else ""}

Competencia en {500 if not tipo_plaza else (1000 if tipo_plaza == 'mall' else 600)}m ({len(competidores)} negocios similares):
{json.dumps([{'nombre': c.get('displayName',{}).get('text','N/A'), 'rating': c.get('rating','N/A'), 'reviews': c.get('userRatingCount',0)} for c in competidores], indent=2, ensure_ascii=False)}

{"INSTRUCCIÓN ESPECIAL MALL: Los competidores incluyen negocios del mismo mall. Esto tiene MAYOR peso que la competencia exterior. Analiza si la zona del mall (food court, moda, etc.) es favorable para este tipo de negocio y si la renta elevada impacta el ROI." if tipo_plaza == "mall" else ""}
{"INSTRUCCIÓN ESPECIAL STRIP CENTER: Considera la visibilidad desde la calle y la tienda ancla. El tráfico es mixto (paso + destino). La accesibilidad y señalización son factores críticos." if tipo_plaza == "strip_center" else ""}

Responde en {idioma} con este formato EXACTO:
**Resumen:** [2-3 líneas sobre la oportunidad en esta ubicación específica]

**✅ Ventajas:**
- [ventaja 1]
- [ventaja 2]
- [ventaja 3]

**⚠️ Riesgos:**
- [riesgo 1]
- [riesgo 2]
- [riesgo 3]

**🎯 Recomendación:** [TOMAR / RECHAZAR / ANALIZAR MÁS] — [1 línea de justificación]

Sé específico. No inventes datos."""
    else:
        top5 = recomendaciones[:5]
        prompt = f"""Eres analista experto en ubicaciones comerciales.

UBICACIÓN: {ubicacion}

TOP 5 negocios viables (por score):
{json.dumps([{'tipo': r['nombre'], 'score': r['score'], 'competidores': r['num_competidores']} for r in top5], indent=2, ensure_ascii=False)}

Responde en {idioma}:
**Análisis de la zona:**
[2 líneas sobre el perfil comercial de esta ubicación]

**Por qué el #1 lidera:**
[1-2 líneas específicas]

**Oportunidad clave:**
[1 línea de la mayor oportunidad detectada]

No inventes datos demográficos."""

    try:
        msg = claude.messages.create(
            model="claude-sonnet-4-20250514",
            max_tokens=1000,
            messages=[{"role": "user", "content": prompt}]
        )
        return msg.content[0].text
    except Exception as e:
        return f"Error al generar análisis: {e}"

# ============================================
# RECOMENDACIÓN FINAL — Veredicto integrado
# ============================================

def generar_recomendacion_final(
        ubicacion, score, desglose, competidores, demografia,
        contexto, idioma, tier_key, modo,
        tipo_negocio=None, recomendaciones=None,
        trafico_data=None, mercado_data=None,
        forecast_data=None, roi_data=None,
        ticket_data=None):
    """
    Genera el Recomendación Final final integrando TODOS los datos disponibles.
    El prompt varía según el tier — más datos = conclusión más precisa.
    """

    # ── Datos base (todos los tiers) ──
    tipo_info = TIPOS_NEGOCIO.get(tipo_negocio, {}) if tipo_negocio else {}
    tipo_nombre = tipo_info.get('nombre', '')

    # En modo recomendar, usar el #1
    if modo == "recomendar" and recomendaciones:
        top1 = recomendaciones[0]
        tipo_nombre   = top1['nombre']
        tipo_negocio  = top1['tipo_key']
        tipo_info     = TIPOS_NEGOCIO.get(tipo_negocio, {})
        competidores  = top1.get('competidores', competidores)

    nivel_txt, _ = nivel_score(score, idioma)
    n_comp        = len(competidores or [])
    ratings       = [c.get('rating',0) for c in (competidores or []) if c.get('rating')]
    avg_rating    = round(sum(ratings)/len(ratings), 1) if ratings else 0
    n_verdes      = sum(1 for r in ratings if r >= 4.3)

    nse           = demografia.get('nse_predominante', 'C') if demografia else 'C'
    ingreso       = demografia.get('ingreso_actual', demografia.get('ingreso_promedio_mensual', 0)) if demografia else 0
    densidad      = demografia.get('densidad_actual', demografia.get('densidad_hab_km2', 0)) if demografia else 0
    pob           = demografia.get('poblacion_actual', demografia.get('poblacion_estimada', 0)) if demografia else 0
    tasa_crec     = demografia.get('tasa_crecimiento_pct', 0) if demografia else 0
    tipo_zona     = contexto.get('tipo_zona', 'mixto') if contexto else 'mixto'
    badges        = contexto.get('badges', []) if contexto else []
    inversion_min = tipo_info.get('inversion_min', 0)
    inversion_max = tipo_info.get('inversion_max', 0)

    # ── Construir contexto de datos según tier ──
    datos_base = f"""
UBICACIÓN: {ubicacion}
TIPO DE NEGOCIO: {tipo_nombre}
INVERSIÓN REQUERIDA: ${inversion_min:,} – ${inversion_max:,} MXN

SCORE DE VIABILIDAD: {score}/100 ({nivel_txt})
- Densidad competencia: {desglose.get('densidad',0)}/100 (50%)
- Calidad competencia: {desglose.get('calidad',0)}/100 (30%)
- Consolidación mercado: {desglose.get('consolidacion',0)}/100 (20%)

COMPETENCIA EN 500m:
- Total competidores: {n_comp}
- Rating promedio: {avg_rating}/5
- Competidores excelentes (≥4.3★): {n_verdes}

DEMOGRAFÍA DEL ÁREA:
- Población estimada: {pob:,} hab
- NSE predominante: {nse}
- Ingreso promedio mensual: ${ingreso:,}
- Densidad urbana: {densidad:,} hab/km²
- Crecimiento anual zona: {tasa_crec:.2f}%
- Tipo de zona: {tipo_zona}
- Características: {', '.join(badges) if badges else 'No detectadas'}

ANÁLISIS DE VIALIDAD:
- Tipo de vialidad: {((contexto or {}).get('analisis_vial') or {}).get('badge', 'No detectada')}
- Impacto en score: {'+' if ((contexto or {}).get('analisis_vial') or {}).get('ajuste_score', 0) >= 0 else ''}{((contexto or {}).get('analisis_vial') or {}).get('ajuste_score', 0)} puntos
- Análisis: {((contexto or {}).get('analisis_vial') or {}).get('descripcion', 'No disponible')}
- Dependencia al tráfico de paso: {((contexto or {}).get('analisis_vial') or {}).get('dependencia_paso', 'medio')}"""

    # ── Plaza comercial — extraer de contexto si existe ──
    _info_plaza_ctx = (contexto or {}).get("info_plaza", {}) if contexto else {}
    _tipo_plaza_ctx = _info_plaza_ctx.get("tipo_plaza") if _info_plaza_ctx else None
    _nota_plaza_ctx = _info_plaza_ctx.get("nota", "") if _info_plaza_ctx else ""
    _zona_mall_ctx  = _info_plaza_ctx.get("zona_mall", "") if _info_plaza_ctx else ""
    _roi_adj_ctx    = _info_plaza_ctx.get("ajuste_roi_pct", 0) if _info_plaza_ctx else 0

    if _tipo_plaza_ctx:
        _plaza_label = "Mall / Plaza Cerrada" if _tipo_plaza_ctx == "mall" else "Strip Center / Plaza Abierta"
        datos_base += f"""

TIPO DE UBICACIÓN — PLAZA COMERCIAL:
- Tipo de plaza: {_plaza_label}
{"- Zona dentro del mall: " + _zona_mall_ctx if _zona_mall_ctx and _tipo_plaza_ctx == "mall" else ""}
{"- Penalización ROI por renta: -" + str(_roi_adj_ctx) + "% vs local en calle equivalente" if _roi_adj_ctx > 0 else ""}
- Notas de plaza: {_nota_plaza_ctx}
INSTRUCCIÓN: Incluye una sección específica sobre las implicaciones de estar en {"un mall cerrado (tráfico captivo, renta alta, competencia interna, food court vs zona servicios)" if _tipo_plaza_ctx == "mall" else "un strip center (visibilidad desde calle, impacto de tienda ancla, tráfico mixto)"} en tu recomendación final."""

    datos_pro = ""
    if trafico_data:
        picos = trafico_data.get('horas_pico', [])
        pico_str = ', '.join([p['hora_str'] for p in picos[:2]]) if picos else 'N/A'
        datos_pro += f"""
TRÁFICO ESTIMADO:
- Nivel general: {trafico_data.get('nivel_general','N/A')}
- Horas pico: {pico_str}
- Mejor día: {trafico_data.get('dia_pico','N/A')}
- Día más bajo: {trafico_data.get('dia_bajo','N/A')}"""

    if mercado_data:
        datos_pro += f"""
MERCADO POTENCIAL:
- Mercado total área/mes: ${mercado_data.get('mercado_total_mensual',0):,}
- Tu captura estimada/mes: ${mercado_data.get('mercado_captura_mensual',0):,}
- Tasa de captura: {mercado_data.get('factor_captura_pct',0):.0f}%
- Clientes/día estimados: {mercado_data.get('clientes_dia_estimados',0)}"""

    if forecast_data:
        esc = forecast_data.get('escenarios', {})
        datos_pro += f"""
FORECAST DE VENTAS (año 1):
- Pesimista: ${esc.get('pesimista',{}).get('total_anual',0):,}
- Base: ${esc.get('base',{}).get('total_anual',0):,}
- Optimista: ${esc.get('optimista',{}).get('total_anual',0):,}"""

    # Ticket data para PRO y PREMIUM
    if ticket_data:
        # Aplicar override de ticket si el tipo no es confiable en Google Places
        _tk_ov = _get_ticket_override(tipo_negocio or "")
        if _tk_ov:
            ticket_data = dict(ticket_data)
            ticket_data["ticket_competencia"] = _tk_ov["ticket_tipico"]
            ticket_data["rango_zona"]         = _tk_ov["rango_txt"]
            ticket_data["ticket_recomendado"] = int(_tk_ov["ticket_tipico"] * 1.1)
            datos_pro += f"""
TICKET PROMEDIO — REFERENCIA DE MERCADO (Google Places no confiable para este tipo):
- Ticket referencial zona: ${ticket_data['ticket_competencia']:,} MXN
- Rango real del mercado MX: {_tk_ov.get('ejemplos','')}
- Nota: {_tk_ov['nota_ui'][:150]}"""
        else:
            datos_pro += f"""
TICKET PROMEDIO DE COMPETENCIA (GOOGLE PLACES + ENIGH 2022):
- Ticket promedio zona: ${ticket_data.get('ticket_competencia',0):,} MXN
- Nivel de precio zona: {ticket_data.get('rango_zona','N/A')}
- Competidores con precio registrado: {ticket_data.get('n_con_precio',0)}"""
        if tier_key == "premium":
            datos_pro += f"""
- Ticket recomendado para el negocio: ${ticket_data.get('ticket_recomendado',0):,} MXN
- Posicionamiento sugerido: {ticket_data.get('posicionamiento','')}"""

    # Demografía edad/género para PRO y PREMIUM
    if demografia:
        dist_edad   = demografia.get('distribucion_edad', {})
        dist_genero = demografia.get('distribucion_genero', {})
        if dist_edad:
            datos_pro += f"""
PERFIL DEMOGRÁFICO DETALLADO (INEGI Censo 2020 · ENIGH 2022):
- Jóvenes 18-35 años: {dist_edad.get('18-35',0)}% del entorno
- Adultos 36-55 años: {dist_edad.get('36-55',0)}% del entorno
- Adultos mayores 56+: {dist_edad.get('56+',0)}% del entorno
- Género predominante: {'Femenino' if dist_genero.get('mujeres',50) > dist_genero.get('hombres',50) else 'Masculino'} ({max(dist_genero.get('mujeres',50), dist_genero.get('hombres',50))}%)"""

    datos_premium = ""
    if roi_data:
        datos_premium += f"""
ROI Y RECUPERACIÓN:
- ROI estimado 12 meses: {roi_data.get('roi_12m_pct',0)}%
- Meses recuperación: {roi_data.get('meses_recuperacion',0)}
- Utilidad mensual estimada: ${roi_data.get('utilidad_mensual_est',0):,}
- Clasificación ROI: {roi_data.get('clasificacion_roi','')}"""

    # ── Prompt según tier ──
    nombre_ub_principal = ubicacion[:60] if ubicacion else "ubicación principal"
    if tier_key in ["pro", "premium"]:
        contexto_datos = datos_base + datos_pro + datos_premium
        instrucciones = f"""Responde en {idioma} con este formato EXACTO. 
IMPORTANTE: Todas las referencias son EXCLUSIVAMENTE sobre la ubicación: "{nombre_ub_principal}"

**✅ RECOMENDACIÓN FINAL — {nombre_ub_principal}**

**Resumen de la oportunidad:**
[3-4 líneas que integren score, competencia, demografía, tráfico y mercado para ESTA ubicación específica]

**✅ Factores a favor:**
- [factor 1 — con dato específico de esta ubicación]
- [factor 2 — con dato específico de esta ubicación]
- [factor 3 — con dato específico de esta ubicación]

**⚠️ Factores de riesgo:**
- [riesgo 1 — con dato específico]
- [riesgo 2 — con dato específico]

**📊 Potencial comercial:**
[2 líneas sobre mercado y ventas estimadas para ESTA ubicación]

**🎫 Análisis de ticket y precio:**
[2 líneas: menciona el ticket promedio de la competencia en la zona, cómo se compara con el tipo de negocio evaluado, y qué estrategia de precio se recomienda según el score]

**👥 Perfil demográfico relevante:**
[2 líneas: describe el grupo de edad y género dominante en el entorno y explica específicamente por qué ese perfil es favorable o desfavorable para ESTE tipo de negocio]

**🛣️ Factibilidad vial:**
[2 líneas explicando el tipo de vialidad detectado, su impacto en puntos, y qué significa específicamente para ESTE tipo de negocio — si es negocio de paso o destino, si la vialidad favorece o perjudica la captación de clientes]

**🎯 VEREDICTO: [PROCEDER ✅ / PRECAUCIÓN 🟡 / NO PROCEDER ❌]**
[2-3 líneas de justificación integrando TODOS los datos incluyendo el factor vial]

**Nivel de confianza: [ALTO / MEDIO / BAJO]** — [motivo en 1 línea]"""
    else:
        contexto_datos = datos_base
        instrucciones = f"""Responde en {idioma} con este formato EXACTO.
IMPORTANTE: Todo el análisis es sobre: "{nombre_ub_principal}"

**✅ RECOMENDACIÓN FINAL — {nombre_ub_principal}**

**Resumen de la oportunidad:**
[2-3 líneas sobre el potencial de esta ubicación específica]

**✅ Factores a favor:**
- [factor 1 — menciona la vialidad si es favorable]
- [factor 2]
- [factor 3]

**⚠️ Factores de riesgo:**
- [riesgo 1 — menciona la vialidad si es desfavorable]
- [riesgo 2]

**🛣️ Factibilidad vial:**
[2 líneas explicando el tipo de vialidad, su impacto en puntos y qué significa para este tipo de negocio específico]

**🎯 VEREDICTO: [PROCEDER ✅ / PRECAUCIÓN 🟡 / NO PROCEDER ❌]**
[1-2 líneas de justificación incluyendo el factor vial]

*Actualiza al plan {'Básico ($99)' if tier_key == 'free' else 'PRO ($499)'} para incluir tráfico, mercado y forecast en tu recomendación.*"""

    prompt = f"""Eres un analista senior de inteligencia comercial y ubicaciones.
Tu tarea es emitir el Recomendación Final final de una ubicación, integrando TODOS los datos disponibles.
NO inventes datos que no estén en el contexto. Sé directo, preciso y accionable.

DATOS DEL ANÁLISIS:
{contexto_datos}

{instrucciones}"""

    try:
        msg = claude.messages.create(
            model="claude-sonnet-4-20250514",
            max_tokens=1200,
            messages=[{"role": "user", "content": prompt}]
        )
        return msg.content[0].text
    except Exception as e:
        return f"Error al generar diagnóstico: {e}"


# ============================================
# NARRATIVA INTERPRETATIVA POR SECCIÓN (PRO+)
# ============================================

def generar_narrativa_seccion(seccion, datos, tipo_negocio_nombre, idioma="es"):
    """
    Genera 2-3 líneas interpretando los datos de una sección específica
    en el contexto del negocio a validar. Solo se llama para PDFs PRO y PREMIUM.
    """
    prompts_por_seccion = {
        "competencia": f"""Eres analista comercial. En 2-3 líneas concisas, interpreta estos datos de competencia
para un negocio tipo "{tipo_negocio_nombre}":
{json.dumps(datos, ensure_ascii=False, indent=2)}
Explica qué significan estos números para el éxito o riesgo del negocio. 
Sé directo, usa los datos específicos. NO uses bullet points. Responde en {idioma}.""",

        "demografia": f"""Eres analista comercial. En 2-3 líneas concisas, interpreta estos datos demográficos
para un negocio tipo "{tipo_negocio_nombre}":
{json.dumps(datos, ensure_ascii=False, indent=2)}
Explica qué implica este perfil de área para la viabilidad del negocio (NSE, ingreso, densidad, crecimiento).
Sé directo, usa los datos. NO uses bullet points. Responde en {idioma}.""",

        "trafico": f"""Eres analista comercial. En 2-3 líneas concisas, interpreta estos datos de tráfico
para un negocio tipo "{tipo_negocio_nombre}":
{json.dumps(datos, ensure_ascii=False, indent=2)}
Explica qué significan las horas pico y el flujo semanal para la operación del negocio.
Sé directo, usa los datos. NO uses bullet points. Responde en {idioma}.""",

        "mercado": f"""Eres analista comercial. En 2-3 líneas concisas, interpreta estos datos de mercado
para un negocio tipo "{tipo_negocio_nombre}":
{json.dumps(datos, ensure_ascii=False, indent=2)}
Explica qué representa el mercado potencial capturables y los clientes estimados en términos prácticos.
Sé directo, usa los datos. NO uses bullet points. Responde en {idioma}.""",

        "forecast": f"""Eres analista comercial. En 2-3 líneas concisas, interpreta estos datos de forecast
para un negocio tipo "{tipo_negocio_nombre}":
{json.dumps(datos, ensure_ascii=False, indent=2)}
Explica qué significan los 3 escenarios y cuál es el rango realista de ventas esperado.
Sé directo, usa los datos. NO uses bullet points. Responde en {idioma}.""",

        "roi": f"""Eres analista comercial. En 2-3 líneas concisas, interpreta estos datos de ROI
para un negocio tipo "{tipo_negocio_nombre}":
{json.dumps(datos, ensure_ascii=False, indent=2)}
Explica qué significa el retorno estimado y el tiempo de recuperación para la decisión de inversión.
Sé directo, usa los datos. NO uses bullet points. Responde en {idioma}.""",

        "mapa": f"""Eres analista comercial. En 2-3 líneas concisas, interpreta estos datos de competidores
en el mapa para un negocio tipo "{tipo_negocio_nombre}":
{json.dumps(datos, ensure_ascii=False, indent=2)}
Explica qué implica la distribución y calidad de competidores en el radio para el posicionamiento del negocio.
Sé directo, usa los datos. NO uses bullet points. Responde en {idioma}.""",
    }

    prompt = prompts_por_seccion.get(seccion)
    if not prompt:
        return ""

    try:
        msg = claude.messages.create(
            model="claude-sonnet-4-20250514",
            max_tokens=200,
            messages=[{"role": "user", "content": prompt}]
        )
        return msg.content[0].text.strip()
    except:
        return ""


def _bloque_narrativa_pdf(narrativa, s):
    """Renderiza el bloque de narrativa interpretativa en el PDF"""
    from reportlab.platypus import Table as RLTable
    from reportlab.lib import colors as rcolors
    if not narrativa:
        return []
    elements = []
    elements.append(Spacer(1, 0.08*inch))
    narr_tbl = RLTable([[narrativa]], colWidths=[5.5*inch])
    narr_tbl.setStyle(TableStyle([
        ('BACKGROUND', (0,0), (-1,-1), colors.HexColor('#F0F7FF')),
        ('LEFTPADDING', (0,0), (-1,-1), 10),
        ('RIGHTPADDING', (0,0), (-1,-1), 10),
        ('TOPPADDING', (0,0), (-1,-1), 8),
        ('BOTTOMPADDING', (0,0), (-1,-1), 8),
        ('BOX', (0,0), (-1,-1), 0.5, colors.HexColor('#0D0D0D44')),
    ]))
    # Texto dentro de la caja
    narr_para = Paragraph(
        f'<font color="#0D0D0D">💬</font> <i>{narrativa}</i>',
        ParagraphStyle('narr', parent=s["base"]['Normal'],
            fontSize=9, textColor=colors.HexColor('#333333'),
            backColor=colors.HexColor('#F0F7FF'),
            borderPad=8, leading=13)
    )
    elements.append(narr_para)
    elements.append(Spacer(1, 0.08*inch))
    return elements


# ============================================
# MAPA ESTÁTICO
# ============================================

def _plaza_block_pdf(story, contexto, s):
    """Inserta seccion de plaza comercial en el PDF si aplica."""
    _ip = (contexto or {}).get("info_plaza", {})
    if not _ip or not _ip.get("tipo_plaza"):
        return
    _pl_tipo  = _ip["tipo_plaza"]
    _pl_label = "Mall / Plaza Cerrada" if _pl_tipo == "mall" else "Strip Center / Plaza Abierta"
    _pl_icono = "🏬" if _pl_tipo == "mall" else "🏪"
    _pl_nota  = _ip.get("nota", "").replace("\n", "  |  ")
    _pl_roi   = _ip.get("ajuste_roi_pct", 0)
    _pl_zona  = _ip.get("zona_mall", "")

    story.append(Spacer(1, 0.1*inch))
    story.extend(_seccion_pdf("Ubicacion en Plaza Comercial", _pl_icono, s))

    zona_labels = {
        "food_court":         "Food Court / zona de comida",
        "moda_accesorios":    "Moda / accesorios / lifestyle",
        "servicios":          "Servicios (salud, belleza, etc.)",
        "ancla_principal":    "Cerca de ancla principal (Liverpool, Walmart, Cine)",
        "planta_baja_entrada":"Planta baja / entrada principal",
        "piso_superior":      "Piso superior / menor trafico",
    }
    plaza_rows = [["Tipo de plaza:", f"{_pl_label}"]]
    if _pl_zona and _pl_tipo == "mall":
        plaza_rows.append(["Zona en mall:", zona_labels.get(_pl_zona, _pl_zona)])
    if _pl_roi > 0:
        plaza_rows.append(["Impacto ROI:", f"Reduccion estimada -{_pl_roi}% vs local en calle por renta de plaza"])
    if _pl_nota:
        plaza_rows.append(["Consideraciones:", _pl_nota[:350]])

    tbl_plaza = Table(plaza_rows, colWidths=[1.6*inch, 3.9*inch])
    tbl_plaza.setStyle(TableStyle([
        ('FONTNAME',  (0,0),(0,-1), 'Helvetica-Bold'),
        ('FONTSIZE',  (0,0),(-1,-1), 9),
        ('TEXTCOLOR', (0,0),(0,-1), colors.HexColor('#6B6B6B')),
        ('ROWBACKGROUNDS',(0,0),(-1,-1),[colors.HexColor('#FFF8F0'), colors.white]),
        ('GRID',(0,0),(-1,-1),0.3,colors.HexColor('#DDDDDD')),
        ('TOPPADDING',(0,0),(-1,-1),5),('BOTTOMPADDING',(0,0),(-1,-1),5),
        ('VALIGN',(0,0),(-1,-1),'TOP'),
    ]))
    story.append(tbl_plaza)
    story.append(Spacer(1, 0.06*inch))
    _warn_style = ParagraphStyle('plaza_warn_pdf', parent=s["base"]['Normal'],
        fontSize=8.5, textColor=colors.HexColor('#7B4500'),
        backColor=colors.HexColor('#FFF3CD'),
        borderPad=6, leading=13, leftIndent=4, rightIndent=4)
    if _pl_tipo == "mall":
        story.append(Paragraph(
            "El score ya incluye ajuste por trafico captivo y renta de mall. "
            "La Recomendacion Final considera las implicaciones de estar en un mall cerrado.",
            _warn_style))
    else:
        story.append(Paragraph(
            "El score ya incluye ajuste por visibilidad de calle y tienda ancla. "
            "La Recomendacion Final considera el impacto de la tienda ancla del strip center.",
            _warn_style))


def _render_recomendacion_pdf(story, analisis, s):
    """Renderiza el bloque de Recomendación Final en el PDF con estilo consistente"""
    story.extend(_seccion_pdf("Recomendación Final", "✅", s))
    story.append(Paragraph(
        "Análisis integrado por IA — combina todos los datos del análisis en una conclusión accionable.",
        ParagraphStyle('rec_intro', parent=s["base"]['Normal'],
        fontSize=9, textColor=colors.HexColor('#666666'), spaceAfter=8)))
    for linea in analisis.split('\n'):
        linea = linea.strip()
        if not linea:
            story.append(Spacer(1, 0.04*inch))
        elif any(kw in linea for kw in ['VEREDICTO','PROCEDER ✅','PRECAUCIÓN 🟡','NO PROCEDER ❌','🎯 VEREDICTO']):
            clean = linea.replace('**','')
            v_tbl = Table([[clean]], colWidths=[5.5*inch])
            v_tbl.setStyle(TableStyle([
                ('BACKGROUND',(0,0),(-1,-1),colors.HexColor('#E3F2FD')),
                ('TEXTCOLOR', (0,0),(-1,-1),colors.HexColor('#6B6B6B')),
                ('FONTNAME',  (0,0),(-1,-1),'Helvetica-Bold'),
                ('FONTSIZE',  (0,0),(-1,-1),12),
                ('TOPPADDING',(0,0),(-1,-1),10),
                ('BOTTOMPADDING',(0,0),(-1,-1),10),
                ('LEFTPADDING',(0,0),(-1,-1),12),
                ('BOX',(0,0),(-1,-1),1.5,colors.HexColor('#6B6B6B')),
            ]))
            story.append(v_tbl)
            story.append(Spacer(1, 0.08*inch))
        elif 'Nivel de confianza' in linea:
            story.append(Paragraph(linea.replace('**',''), ParagraphStyle('conf_r',
                parent=s["base"]['Normal'], fontSize=9,
                textColor=colors.HexColor('#555'), fontName='Helvetica-Oblique', spaceAfter=4)))
        elif linea.startswith('**') and linea.endswith('**'):
            story.append(Paragraph(linea.replace('**',''), ParagraphStyle('sec_r',
                parent=s["base"]['Normal'], fontSize=11, fontName='Helvetica-Bold',
                textColor=colors.HexColor('#6B6B6B'), spaceBefore=10, spaceAfter=4)))
        elif linea.startswith('- ') or linea.startswith('• '):
            story.append(Paragraph(f"  {linea}", ParagraphStyle('bul_r',
                parent=s["base"]['Normal'], fontSize=9.5, leftIndent=15, spaceAfter=3)))
        elif '**' in linea:
            import re as _re
            txt = _re.sub(r'\*\*(.+?)\*\*', r'<b></b>', linea)
            story.append(Paragraph(txt, s["normal"]))
        else:
            story.append(Paragraph(linea, s["normal"]))


def generar_mapa_estatico(lat, lng, competidores=None, con_competidores=False):
    """
    Genera mapa estático Google.
    Competidores coloreados: verde=rating>=4.3, naranja=3.5-4.2, rojo=<3.5 / sin rating.
    """
    base = "https://maps.googleapis.com/maps/api/staticmap"

    # Círculo de 500m
    circle_pts = []
    for i in range(0, 360, 10):
        angle = math.radians(i)
        p_lat = lat + (0.0045 * math.cos(angle))
        p_lng = lng + (0.006 * math.sin(angle))
        circle_pts.append(f"{p_lat},{p_lng}")

    params = [
        f"center={lat},{lng}",
        "zoom=16",
        "size=640x480",
        "maptype=roadmap",
        "style=feature:poi|element:labels|visibility:off",  # Quitar POIs para más limpio
        # Tu ubicación — marcador grande rojo con estrella
        f"markers=color:0xFF0000|size:large|label:★|{lat},{lng}",
        f"path=color:0x0047AB99|weight:2|fillcolor:0x0047AB15|{'|'.join(circle_pts)}",
        f"key={GOOGLE_API_KEY}"
    ]

    # Competidores en mapa con colores por rating
    if con_competidores and competidores:
        grupos = {
            "0x22C55E": [],  # verde — rating >= 4.3
            "0xF97316": [],  # naranja — rating 3.5–4.2
            "0xEF4444": [],  # rojo — rating < 3.5 o sin rating
        }
        for comp in competidores[:20]:
            loc = comp.get('location', {})
            c_lat = loc.get('latitude')
            c_lng = loc.get('longitude')
            if not c_lat or not c_lng:
                continue
            rating = comp.get('rating', 0)
            if rating >= 4.3:
                grupos["0x22C55E"].append(f"{c_lat},{c_lng}")
            elif rating >= 3.5:
                grupos["0xF97316"].append(f"{c_lat},{c_lng}")
            else:
                grupos["0xEF4444"].append(f"{c_lat},{c_lng}")

        for hex_color, locs in grupos.items():
            if locs:
                # Google Static Maps acepta múltiples marcadores en un solo param
                locs_str = "|".join(locs)
                params.append(f"markers=color:{hex_color}|size:small|{locs_str}")

    url = f"{base}?" + "&".join(params)
    try:
        resp = requests.get(url, timeout=10)
        if resp.status_code == 200:
            return BytesIO(resp.content)
    except:
        pass
    return None


def generar_texto_interpretacion_mapa(competidores, contexto, modo_mapa="basico", idioma="es"):
    """Genera texto de interpretación debajo del mapa"""
    n_comp = len(competidores) if competidores else 0

    # Ratings de competidores
    ratings = [c.get('rating', 0) for c in (competidores or []) if c.get('rating')]
    avg_rating = round(sum(ratings)/len(ratings), 1) if ratings else 0
    n_verdes   = sum(1 for r in ratings if r >= 4.3)
    n_naranjas = sum(1 for r in ratings if 3.5 <= r < 4.3)
    n_rojos    = sum(1 for r in ratings if r < 3.5)

    tipo_zona = contexto.get('tipo_zona', 'mixto') if contexto else 'mixto'
    zona_desc = {
        'paso': 'zona de alto tráfico vehicular',
        'comercial': 'zona comercial activa',
        'residencial': 'zona residencial',
        'mixto': 'zona mixta'
    }.get(tipo_zona, 'zona urbana')

    if modo_mapa == "basico":
        return (f"📍 La ubicación analizada (marcador rojo) se encuentra en {zona_desc}. "
                f"El círculo azul representa el radio de análisis de 500 metros. "
                f"Se identificaron {n_comp} negocios similares en este radio.")
    else:
        texto = (f"📍 Mapa de competencia en 500m. Tu ubicación: marcador rojo ★. "
                 f"Competidores: 🟢 {n_verdes} excelentes (≥4.3★) · "
                 f"🟠 {n_naranjas} buenos (3.5-4.2★) · "
                 f"🔴 {n_rojos} débiles o sin datos. ")
        if n_verdes > 3:
            texto += "⚠️ Alta concentración de competidores bien calificados — diferenciación clave."
        elif n_verdes == 0 and n_comp > 0:
            texto += "💡 Ningún competidor destaca con rating alto — oportunidad de posicionarte como el mejor."
        elif n_comp == 0:
            texto += "✅ Sin competidores directos identificados en el radio — mercado sin saturar."
        return texto

# ============================================
# GENERADORES DE PDF POR TIER
# ============================================

def _estilos_pdf():
    """Estilos PDF con identidad de marca 4SITE — Poppins + paleta negro/dorado/gris"""
    FONT_DIR = "/usr/share/fonts/truetype/google-fonts/"
    try:
        pdfmetrics.registerFont(TTFont("P-Bold",  FONT_DIR + "Poppins-Bold.ttf"))
        pdfmetrics.registerFont(TTFont("P-Reg",   FONT_DIR + "Poppins-Regular.ttf"))
        pdfmetrics.registerFont(TTFont("P-Med",   FONT_DIR + "Poppins-Medium.ttf"))
        pdfmetrics.registerFont(TTFont("P-Light", FONT_DIR + "Poppins-Light.ttf"))
        _FB, _FR, _FM, _FL = "P-Bold", "P-Reg", "P-Med", "P-Light"
    except Exception:
        _FB, _FR, _FM, _FL = "Helvetica-Bold", "Helvetica", "Helvetica", "Helvetica"

    styles = getSampleStyleSheet()
    _N = colors.HexColor("#0D0D0D")   # negro
    _D = colors.HexColor("#BB944F")   # dorado
    _G = colors.HexColor("#6B6B6B")   # gris texto
    _L = colors.HexColor("#EDEDED")   # gris claro
    _W = colors.HexColor("#FFFFFF")   # blanco

    title = ParagraphStyle("4Title", parent=styles["Normal"],
        fontName=_FB, fontSize=26, leading=32,
        textColor=_N, alignment=TA_CENTER, spaceAfter=4)
    subtitle = ParagraphStyle("4Sub", parent=styles["Normal"],
        fontName=_FR, fontSize=11, leading=16,
        textColor=_D, alignment=TA_CENTER, spaceAfter=16)
    h2 = ParagraphStyle("4H2", parent=styles["Normal"],
        fontName=_FB, fontSize=12, leading=16,
        textColor=_W, spaceBefore=0, spaceAfter=0)
    h3 = ParagraphStyle("4H3", parent=styles["Normal"],
        fontName=_FB, fontSize=10.5, leading=15,
        textColor=_N, spaceBefore=10, spaceAfter=4)
    normal = ParagraphStyle("4Normal", parent=styles["Normal"],
        fontName=_FR, fontSize=9.5, leading=14,
        textColor=_N, spaceAfter=4)
    small = ParagraphStyle("4Small", parent=styles["Normal"],
        fontName=_FL, fontSize=8, leading=11,
        textColor=_G, alignment=TA_CENTER, spaceAfter=2)
    italic = ParagraphStyle("4Italic", parent=styles["Normal"],
        fontName=_FL, fontSize=9, leading=13,
        textColor=_G, leftIndent=10, rightIndent=10, spaceAfter=3)
    watermark = ParagraphStyle("4WM", parent=styles["Normal"],
        fontName=_FL, fontSize=8,
        textColor=colors.HexColor("#AAAAAA"), alignment=TA_CENTER)
    kpi_val = ParagraphStyle("4KpiVal", parent=styles["Normal"],
        fontName=_FB, fontSize=20, leading=24,
        textColor=_N, alignment=TA_CENTER)
    kpi_val_gold = ParagraphStyle("4KpiGold", parent=styles["Normal"],
        fontName=_FB, fontSize=20, leading=24,
        textColor=_D, alignment=TA_CENTER)
    kpi_label = ParagraphStyle("4KpiLbl", parent=styles["Normal"],
        fontName=_FR, fontSize=8, leading=11,
        textColor=_G, alignment=TA_CENTER)

    return {"title": title, "subtitle": subtitle, "h2": h2, "h3": h3,
            "normal": normal, "small": small, "italic": italic,
            "watermark": watermark, "kpi_val": kpi_val,
            "kpi_val_gold": kpi_val_gold, "kpi_label": kpi_label,
            "base": styles,
            # Helpers de color para uso directo
            "_N": _N, "_D": _D, "_G": _G, "_L": _L, "_W": _W,
            "_FB": _FB, "_FR": _FR, "_FM": _FM, "_FL": _FL}

def _tabla_competidores(competidores, t, n_max, styles):
    """Genera tabla de competidores limitada según tier"""
    data = [[t["tabla_nombre"], t["tabla_rating"], t["tabla_reseñas"]]]
    for comp in competidores[:n_max]:
        nombre  = comp.get('displayName', {}).get('text', 'N/A')[:30]
        rating  = comp.get('rating', 'N/A')
        reviews = comp.get('userRatingCount', 0)
        data.append([nombre, f"{rating}★" if rating != 'N/A' else 'N/A', str(reviews)])
    tbl = Table(data, colWidths=[3.2*inch, 1*inch, 1*inch])
    tbl.setStyle(TableStyle([
        ('BACKGROUND', (0,0), (-1,0), colors.HexColor('#6B6B6B')),
        ('TEXTCOLOR',  (0,0), (-1,0), colors.white),
        ('FONTNAME',   (0,0), (-1,0), 'Helvetica-Bold'),
        ('FONTSIZE',   (0,0), (-1,-1), 9),
        ('ALIGN',      (0,0), (-1,-1), 'CENTER'),
        ('BACKGROUND', (0,1), (-1,-1), colors.HexColor('#F5F8FF')),
        ('ROWBACKGROUNDS', (0,1), (-1,-1), [colors.white, colors.HexColor('#F5F8FF')]),
        ('GRID',       (0,0), (-1,-1), 0.5, colors.HexColor('#CCCCCC')),
        ('TOPPADDING', (0,0), (-1,-1), 5),
        ('BOTTOMPADDING', (0,0), (-1,-1), 5),
    ]))
    return tbl

def _header_pdf(story, s, tier_key, idioma):
    """Header común para todos los PDFs"""
    t = TEXTOS[idioma]
    tier = TIERS[tier_key]
    story.append(Paragraph("4SITE", s["title"]))
    story.append(Paragraph("Don't guess. Foresee.", s["subtitle"]))
    # Badge de tier
    tier_color = tier["color"]
    story.append(Paragraph(
        f"<font color='{tier_color}'><b>[ {tier['nombre'].upper()} — {tier['precio']} ]</b></font>",
        ParagraphStyle('badge', parent=s["base"]['Normal'], fontSize=10, alignment=TA_CENTER, spaceAfter=10)
    ))
    story.append(Paragraph(f"{datetime.datetime.now().strftime('%d/%m/%Y')}", s["small"]))
    story.append(Spacer(1, 0.2*inch))

def _score_block(story, score, desglose, tier_key, t, s, idioma):
    """Bloque de score con desglose condicional"""
    nivel, color = nivel_score(score, idioma)
    score_data = [[f"{score}/100", nivel, TIERS[tier_key]['pdf_paginas']]]
    tbl = Table(score_data, colWidths=[1.5*inch, 1.8*inch, 2*inch])
    tbl.setStyle(TableStyle([
        ('ALIGN',      (0,0), (-1,-1), 'CENTER'),
        ('FONTSIZE',   (0,0), (0,0), 28),
        ('FONTSIZE',   (1,0), (-1,-1), 11),
        ('FONTNAME',   (0,0), (-1,-1), 'Helvetica-Bold'),
        ('TEXTCOLOR',  (0,0), (0,0), colors.HexColor(color)),
        ('TEXTCOLOR',  (1,0), (1,0), colors.HexColor(color)),
        ('TEXTCOLOR',  (2,0), (2,0), colors.HexColor('#6B6B6B')),
        ('BOX',        (0,0), (-1,-1), 1.5, colors.HexColor(color)),
        ('TOPPADDING', (0,0), (-1,-1), 12),
        ('BOTTOMPADDING', (0,0), (-1,-1), 12),
    ]))
    story.append(tbl)
    story.append(Spacer(1, 0.15*inch))
    
    # Desglose solo en BASIC+
    if TIERS[tier_key]["muestra_score_desglose"] and desglose:
        d_data = [
            [t["densidad_comp"], f"{desglose['densidad']}/100", "50%"],
            [t["calidad_comp"],  f"{desglose['calidad']}/100",  "30%"],
            [t["consolidacion"], f"{desglose['consolidacion']}/100", "20%"],
        ]
        d_tbl = Table(d_data, colWidths=[2.5*inch, 1*inch, 0.8*inch])
        d_tbl.setStyle(TableStyle([
            ('FONTSIZE',  (0,0), (-1,-1), 8),
            ('TEXTCOLOR', (0,0), (0,-1), colors.HexColor('#555555')),
            ('TEXTCOLOR', (1,0), (1,-1), colors.HexColor('#6B6B6B')),
            ('FONTNAME',  (1,0), (1,-1), 'Helvetica-Bold'),
            ('GRID',      (0,0), (-1,-1), 0.3, colors.HexColor('#DDDDDD')),
            ('TOPPADDING', (0,0), (-1,-1), 4),
            ('BOTTOMPADDING', (0,0), (-1,-1), 4),
        ]))
        story.append(Paragraph(t["desglose_score"], s["h2"]))
        story.append(d_tbl)
        story.append(Spacer(1, 0.15*inch))

def _cta_block(story, tier_key, t, s, idioma):
    """Bloque CTA al final del PDF — diferente según tier"""
    story.append(Spacer(1, 0.3*inch))
    
    if tier_key == "free":
        upgrades = [
            ("✅ Básico — $99 MXN", (
                "Recomendación Final completo · Análisis de competencia (10 negocios) · "
                "Perfil demográfico con NSE e ingreso · Desglose del score con explicación · "
                "Características de la zona · Mapa de ubicación · Análisis completo PDF"
            )),
            ("🚀 PRO — $499 MXN", (
                "Todo lo del Básico + Análisis de tráfico por hora y día · "
                "Tamaño de mercado potencial · Forecast de ventas 3 escenarios · "
                "Mapa de competidores coloreados · Heatmap de competencia · "
                "Mapa de isócronas (5-15 min) · Mapa de canibalización · "
                "Narrativa interpretativa en cada sección · Análisis PRO PDF"
            )),
            ("💎 PREMIUM — $999 MXN", (
                "Todo lo del PRO + ROI y punto de equilibrio · "
                "Dashboard ejecutivo en una página · Comparativa de hasta 3 ubicaciones · "
                "Curva de recuperación de inversión · Análisis PREMIUM PDF"
            )),
        ]
    elif tier_key == "basic":
        upgrades = [
            ("🚀 PRO — $499 MXN", (
                "Análisis de tráfico por hora y día · Tamaño de mercado potencial · "
                "Forecast 3 escenarios · 4 mapas avanzados (competidores, heatmap, isócronas, canibalización) · "
                "Narrativa interpretativa por sección · Análisis PRO PDF"
            )),
            ("💎 PREMIUM — $999 MXN", (
                "ROI + punto de equilibrio · Dashboard ejecutivo · "
                "Comparativa 3 ubicaciones · Curva de recuperación · Análisis PREMIUM PDF"
            )),
        ]
    elif tier_key == "pro":
        upgrades = [
            ("💎 PREMIUM — $999 MXN", (
                "ROI y punto de equilibrio · Dashboard ejecutivo completo · "
                "Comparativa simultánea de hasta 3 ubicaciones · "
                "Curva de recuperación de inversión · Análisis PREMIUM PDF"
            )),
        ]
    else:
        upgrades = []

    if upgrades:
        cta_style = ParagraphStyle('cta_h', parent=s["base"]['Heading2'], fontSize=12,
                                   textColor=colors.HexColor('#6B6B6B'), alignment=TA_CENTER)
        story.append(Paragraph(t["cta_titulo"], cta_style))
        for nombre, desc in upgrades:
            # Usar Paragraph para que el texto haga wrap automático
            from reportlab.platypus import Paragraph as _P
            nombre_p = _P(f"<b>{nombre}</b>",
                ParagraphStyle('cta_n', parent=s["base"]['Normal'],
                fontSize=9, textColor=colors.HexColor('#BB944F')))
            desc_p = _P(desc,
                ParagraphStyle('cta_d', parent=s["base"]['Normal'],
                fontSize=8.5, textColor=colors.HexColor('#444444')))
            row = Table([[nombre_p, desc_p]], colWidths=[1.8*inch, 3.7*inch])
            row.setStyle(TableStyle([
                ('VALIGN',  (0,0), (-1,-1), 'TOP'),
                ('TOPPADDING', (0,0), (-1,-1), 6),
                ('BOTTOMPADDING', (0,0), (-1,-1), 6),
                ('LEFTPADDING', (0,0), (-1,-1), 6),
                ('LINEBELOW', (0,0), (-1,-1), 0.3, colors.HexColor('#EEEEEE')),
            ]))
            story.append(row)
        story.append(Spacer(1, 0.15*inch))

    # Footer
    story.append(Spacer(1, 0.1*inch))
    footer_style = ParagraphStyle('footer', parent=s["base"]['Normal'], fontSize=8,
                                   textColor=colors.grey, alignment=TA_CENTER)
    story.append(Paragraph("📧 hola@4site.mx  |  🌐 4site.mx", footer_style))
    story.append(Paragraph(t["footer_derechos"], footer_style))

    if TIERS[tier_key]["watermark"]:
        story.append(Spacer(1, 0.2*inch))
        story.append(Paragraph("━" * 30, s["watermark"]))
        story.append(Paragraph(t["watermark"], s["watermark"]))
        story.append(Paragraph("━" * 30, s["watermark"]))


# ─────────────────────────────────────────────
# PDF FREE  (Análisis básico)
# ─────────────────────────────────────────────


def _encontrar_asset(base_dir: str, nombre_base: str) -> str:
    """Busca un archivo de imagen probando múltiples extensiones y ubicaciones."""
    extensiones = ["png", "jpg", "jpeg", "PNG", "JPG", "JPEG"]
    carpetas = [
        os.path.join(base_dir, "assets"),
        base_dir,
        os.path.join(base_dir, "static"),
        os.path.join(base_dir, "img"),
    ]
    for carpeta in carpetas:
        for ext in extensiones:
            ruta = os.path.join(carpeta, f"{nombre_base}.{ext}")
            if os.path.exists(ruta):
                return ruta
    return ""

# ════════════════════════════════════════════════════════════════
# PORTADA ESTILO 4SITE — foto de fondo + panel blanco
# ════════════════════════════════════════════════════════════════
def _portada_foto_canvas(output_canvas, ubicacion, score, tier_key,
                          fecha, tipo_negocio_nombre, logo_path, bg_path):
    """Portada estilo redes 4SITE: foto de calle + panel blanco redondeado."""
    import tempfile, os as _os
    from PIL import Image as _PIL, ImageFilter as _IFilter
    from reportlab.lib.pagesizes import letter as _letter
    W, H = _letter  # 612 x 792 pts

    _D  = colors.HexColor("#BB944F")   # dorado
    _GT = colors.HexColor("#6B6B6B")   # gris texto
    _GB = colors.HexColor("#EDEDED")   # gris bg
    _N  = colors.HexColor("#0D0D0D")   # negro
    _BL = colors.HexColor("#FFFFFF")   # blanco
    _SC = (colors.HexColor("#2ECC71") if score>=65 else
           colors.HexColor("#F39C12") if score>=40 else
           colors.HexColor("#E74C3C"))
    _LB = ("EXCELENTE" if score>=80 else "BUENO" if score>=65 else
           "REGULAR"   if score>=50 else "BAJO"  if score>=35 else "MUY BAJO")

    FD = "/usr/share/fonts/truetype/google-fonts/"
    try:
        pdfmetrics.registerFont(TTFont("P-Bold",  FD+"Poppins-Bold.ttf"))
        pdfmetrics.registerFont(TTFont("P-Reg",   FD+"Poppins-Regular.ttf"))
        pdfmetrics.registerFont(TTFont("P-Med",   FD+"Poppins-Medium.ttf"))
        pdfmetrics.registerFont(TTFont("P-Light", FD+"Poppins-Light.ttf"))
        FB2,FR2,FM2,FL2 = "P-Bold","P-Reg","P-Med","P-Light"
    except:
        FB2,FR2,FM2,FL2 = "Helvetica-Bold","Helvetica","Helvetica","Helvetica"

    _tmp = []
    def _t(img, fmt="JPEG", q=90):
        tf = tempfile.NamedTemporaryFile(suffix=f".{fmt.lower()}", delete=False)
        img.save(tf.name, fmt, quality=q); tf.close(); _tmp.append(tf.name)
        return tf.name

    # Fondo
    try:
        bg = _PIL.open(bg_path).convert("RGB")
        pw,ph = int(W*2),int(H*2)
        br=bg.width/bg.height; pr=pw/ph
        nw,nh = (pw,int(pw/br)) if br<=pr else (int(ph*br),ph)
        bg2 = bg.resize((nw,nh),_PIL.LANCZOS)
        l,t = (nw-pw)//2,(nh-ph)//2
        bg3 = bg2.crop((l,t,l+pw,t+ph))
        bg4 = bg3.filter(_IFilter.GaussianBlur(1.5))
        bg5 = _PIL.blend(bg4, _PIL.new("RGB",bg4.size,(255,255,255)), 0.38)
        output_canvas.drawImage(_t(bg5), 0,0,width=W,height=H)
    except:
        output_canvas.setFillColor(colors.HexColor("#CCCCCC"))
        output_canvas.rect(0,0,W,H,fill=1,stroke=0)

    # Panel blanco
    px,py,pw2,ph2,rd = 36,105,W-72,H-142,16
    output_canvas.setFillColorRGB(0,0,0,alpha=0.07)
    output_canvas.roundRect(px+3,py-4,pw2,ph2,rd,fill=1,stroke=0)
    output_canvas.setFillColor(_BL)
    output_canvas.roundRect(px,py,pw2,ph2,rd,fill=1,stroke=0)
    output_canvas.setStrokeColorRGB(0,0,0,alpha=0.05); output_canvas.setLineWidth(0.8)
    output_canvas.roundRect(px,py,pw2,ph2,rd,fill=0,stroke=1)

    padx=px+26; topy=py+ph2-26; cy=topy

    # Logo
    try:
        logo=_PIL.open(logo_path).convert("RGB")
        lw=150; lh=int(lw*logo.height/logo.width)
        output_canvas.drawImage(_t(logo.resize((lw*2,lh*2),_PIL.LANCZOS)),
                                padx,cy-lh,width=lw,height=lh)
        cy -= lh+12
    except:
        output_canvas.setFillColor(_N); output_canvas.setFont(FB2,28)
        output_canvas.drawString(padx,cy-36,"4SITE"); cy -= 48

    # Línea dorada
    output_canvas.setStrokeColor(_D); output_canvas.setLineWidth(1.5)
    output_canvas.line(padx,cy,px+pw2-26,cy); cy -= 18

    # Badge tier
    _tl = {"pro":"PLAN PRO · $499","premium":"PLAN PREMIUM · $999",
           "basico":"PLAN BÁSICO · $99","free":"GRATIS"}
    _tt = _tl.get(tier_key, f"PLAN {tier_key.upper()}")
    bw = len(_tt)*5.8+18
    output_canvas.setFillColor(_D)
    output_canvas.roundRect(padx,cy-17,bw,19,4,fill=1,stroke=0)
    output_canvas.setFillColor(_N); output_canvas.setFont(FB2,8)
    output_canvas.drawCentredString(padx+bw/2,cy-9,_tt); cy -= 30

    output_canvas.setFillColor(_GT); output_canvas.setFont(FM2,8.5)
    output_canvas.drawString(padx,cy,"REPORTE DE ANÁLISIS DE UBICACIÓN"); cy -= 16

    output_canvas.setFillColor(_N); output_canvas.setFont(FB2,11)
    mc=50
    if len(ubicacion)>mc:
        cut=ubicacion.rfind(",",0,mc); cut=cut if cut>0 else mc
        output_canvas.drawString(padx,cy,ubicacion[:cut+1].strip()); cy -= 15
        output_canvas.drawString(padx,cy,ubicacion[cut+1:].strip()); cy -= 15
    else:
        output_canvas.drawString(padx,cy,ubicacion); cy -= 15

    if tipo_negocio_nombre:
        output_canvas.setFillColor(_D); output_canvas.setFont(FM2,10)
        output_canvas.drawString(padx,cy,tipo_negocio_nombre); cy -= 26
    cy -= 8

    # Score
    cx2=px+pw2/2; scy=cy-80
    output_canvas.setFillColor(colors.HexColor("#F4F4F4"))
    output_canvas.circle(cx2,scy,74,fill=1,stroke=0)
    output_canvas.setStrokeColor(_SC); output_canvas.setLineWidth(6)
    output_canvas.circle(cx2,scy,68,fill=0,stroke=1)
    output_canvas.setFillColor(_SC); output_canvas.setFont(FB2,42)
    output_canvas.drawCentredString(cx2,scy-14,str(int(score)))
    output_canvas.setFillColor(_GT); output_canvas.setFont(FR2,8.5)
    output_canvas.drawCentredString(cx2,scy-27,"SCORE / 100")
    lw2=len(_LB)*7+22
    output_canvas.setFillColor(_D)
    output_canvas.roundRect(cx2-lw2/2,scy-58,lw2,18,4,fill=1,stroke=0)
    output_canvas.setFillColor(_N); output_canvas.setFont(FB2,8.5)
    output_canvas.drawCentredString(cx2,scy-48,_LB)
    cy = scy-96-10

    # KPIs 2 col
    kw2=(pw2-52)/2
    kpis=[("NSE PREDOMINANTE","—"),("COMPETIDORES","—"),
          ("DENSIDAD HAB/KM²","—"),("INGRESO PROM.","—")]
    for i,(lk,vk) in enumerate(kpis):
        col=i%2; row=i//2
        kx=padx+col*(kw2+8); ky=cy-row*36
        output_canvas.setStrokeColor(colors.HexColor("#EBEBEB"))
        output_canvas.setLineWidth(0.5)
        output_canvas.line(kx,ky+2,kx+kw2-8,ky+2)
        output_canvas.setFillColor(_GT); output_canvas.setFont(FR2,7)
        output_canvas.drawString(kx,ky-9,lk)
        output_canvas.setFillColor(_N); output_canvas.setFont(FB2,10.5)
        output_canvas.drawString(kx,ky-21,vk)

    output_canvas.setFillColor(_GT); output_canvas.setFont(FL2,7.5)
    output_canvas.drawCentredString(px+pw2/2,py+12,fecha)
    output_canvas.setFillColor(colors.HexColor("#555555")); output_canvas.setFont(FL2,8)
    output_canvas.drawCentredString(W/2,72,"4SITE · Don't guess. Foresee. · 4site.mx")

    output_canvas.showPage()
    for f in _tmp:
        try: _os.unlink(f)
        except: pass

# ════════════════════════════════════════════════════════════════
# BÚSQUEDA DE ZONA IDEAL — Opción A: cuadrícula de puntos
# ════════════════════════════════════════════════════════════════

def buscar_zona_ideal(lat_ref: float, lng_ref: float, radio_km: float,
                      tipo_negocio_key: str, idioma: str = "es",
                      progress_callback=None) -> list:
    """
    Muestrea puntos en cuadrícula dentro del radio y evalúa el score.
    Retorna lista ordenada por score desc.
    """
    import math

    # Validar que el tipo existe
    if tipo_negocio_key not in TIPOS_NEGOCIO:
        tipo_negocio_key = list(TIPOS_NEGOCIO.keys())[0]

    # Puntos en anillos concéntricos (más orgánico que cuadrícula)
    # Anillo 0: centro; anillo 1: radio/2; anillo 2: radio
    import random as _rnd
    _rnd.seed(int(lat_ref * 1000) % 999)  # seed reproducible por ubicación

    puntos = [(lat_ref, lng_ref)]  # siempre incluir el centro

    # Anillo interior (radio * 0.4)
    _r1 = radio_km * 0.4
    for angulo in range(0, 360, 60):  # 6 puntos
        _rad = math.radians(angulo + _rnd.uniform(-15, 15))
        _d   = _r1 * _rnd.uniform(0.8, 1.2)
        p_lat = lat_ref + (_d * math.cos(_rad)) / 111.32
        p_lng = lng_ref + (_d * math.sin(_rad)) / (111.32 * math.cos(math.radians(lat_ref)))
        puntos.append((p_lat, p_lng))

    # Anillo exterior (radio * 0.85)
    _r2 = radio_km * 0.85
    for angulo in range(0, 360, 45):  # 8 puntos
        _rad = math.radians(angulo + _rnd.uniform(-20, 20))
        _d   = _r2 * _rnd.uniform(0.75, 1.0)
        if _d <= radio_km:
            p_lat = lat_ref + (_d * math.cos(_rad)) / 111.32
            p_lng = lng_ref + (_d * math.sin(_rad)) / (111.32 * math.cos(math.radians(lat_ref)))
            puntos.append((p_lat, p_lng))

    # Limitar según radio para no hacer demasiadas llamadas
    _max_pts = {1: 7, 3: 10, 5: 13, 10: 18}
    _lim = next((v for k,v in sorted(_max_pts.items()) if radio_km <= k), 18)
    puntos = puntos[:_lim]

    resultados = []
    errores    = []

    for i, (p_lat, p_lng) in enumerate(puntos):
        if progress_callback:
            progress_callback(i + 1, len(puntos))
        try:
            # Geocoding inverso — extraer solo colonia/zona, no calle ni número
            try:
                geo_r = gmaps.reverse_geocode((p_lat, p_lng))
                if geo_r:
                    # Buscar componente tipo "sublocality" o "neighborhood" (colonia)
                    _addr_components = geo_r[0].get('address_components', [])
                    _colonia = ""
                    _delegacion = ""
                    _ciudad = ""
                    for _comp in _addr_components:
                        _types = _comp.get('types', [])
                        if any(t in _types for t in ['sublocality_level_1','sublocality','neighborhood']):
                            _colonia = _comp['long_name']
                        elif any(t in _types for t in ['political','administrative_area_level_3','locality']):
                            if not _delegacion: _delegacion = _comp['long_name']
                        elif 'administrative_area_level_1' in _types:
                            _ciudad = _comp['short_name']
                    # Construir dirección solo nivel colonia
                    if _colonia:
                        direccion = _colonia
                        if _delegacion: direccion += f", {_delegacion}"
                        if _ciudad: direccion += f", {_ciudad}"
                    else:
                        # Fallback: quitar el primer componente (calle+número)
                        _full = geo_r[0]['formatted_address']
                        _parts = [p.strip() for p in _full.split(',')]
                        direccion = ', '.join(_parts[1:4]) if len(_parts) > 2 else _full
                else:
                    direccion = f"{p_lat:.4f},{p_lng:.4f}"
            except Exception:
                direccion = f"{p_lat:.4f},{p_lng:.4f}"

            # Competencia
            try:
                comp = buscar_competencia_por_tipo(p_lat, p_lng, tipo_negocio_key)
            except Exception:
                comp = []

            # Score — calcular_score_competencia devuelve TUPLA (score, desglose)
            try:
                _res_score = calcular_score_competencia(comp, tipo_negocio_key)
                if isinstance(_res_score, tuple):
                    score_base, _desglose = _res_score
                else:
                    score_base = float(_res_score)
                    _desglose  = {}
            except Exception:
                score_base = 50.0
                _desglose  = {}

            # Demografía
            try:
                dem = obtener_demografia(p_lat, p_lng, direccion) if MODULOS_OK else {}
            except Exception:
                dem = {}

            # Ajustes de score
            try:
                score_dem = ajustar_score_por_demografia(score_base, tipo_negocio_key, dem)
            except Exception:
                score_dem = score_base
            try:
                ctx_p       = detectar_contexto_ubicacion(p_lat, p_lng, direccion, tipo_negocio_key)
                score_final = ajustar_score_por_contexto(score_dem, tipo_negocio_key, ctx_p)
            except Exception:
                score_final = score_dem

            score_final = max(0, min(100, float(score_final)))

            try:
                nivel, _ = nivel_score(score_final, idioma)
            except Exception:
                nivel = "Regular"

            nse      = dem.get("nse_predominante", "—") if dem else "—"
            densidad = dem.get("densidad_actual", dem.get("densidad_hab_km2", 0)) if dem else 0
            ingreso  = dem.get("ingreso_actual",  dem.get("ingreso_promedio_mensual", 0)) if dem else 0

            # Ticket promedio de la zona
            _ticket_zona = 0
            try:
                if MODULOS_OK and comp and "calcular_ticket_competencia" in dir():
                    _tk_z = calcular_ticket_competencia(comp, tipo_negocio_key, score_final)
                    _ticket_zona = _tk_z.get("ticket_competencia", 0)
            except Exception:
                pass

            # Tipo de vialidad
            _vialidad = ""
            try:
                _vialidad = ctx_p.get("analisis_vial", {}).get("badge", "") if ctx_p else ""
            except Exception:
                pass

            # ── Clasificar zona y calcular bonus/penalización de score ──
            # Competidor de referencia: el mejor calificado con ≥30 reseñas
            _mejor_comp_nombre = ""
            _mejor_comp_rating = 0.0
            _mejor_comp_reviews = 0
            for _c in comp:
                _cr = float(_c.get("rating", 0) or 0)
                _cv = int(_c.get("userRatingCount", 0) or 0)
                if _cr >= _mejor_comp_rating and _cv >= _mejor_comp_reviews:
                    _mejor_comp_rating  = _cr
                    _mejor_comp_reviews = _cv
                    _mejor_comp_nombre  = _c.get("displayName", {}).get("text", "")

            n_comp = len(comp)
            _tiene_ref = (_mejor_comp_rating >= 3.8 and _mejor_comp_reviews >= 30)

            # Categoría y ajuste de score
            if n_comp > 5:
                _zona_cat   = "saturada"
                score_final = max(0, score_final - 10)   # penalización saturación
            elif _tiene_ref and n_comp <= 5:
                _zona_cat   = "mercado_probado"
                score_final = min(100, score_final + 15)  # bonus mercado probado
            elif n_comp == 0:
                _zona_cat   = "virgen"
                score_final = min(100, score_final + 5)   # bonus zona virgen
            else:
                _zona_cat   = "moderada"

            score_final = max(0, min(100, float(score_final)))

            # ── Recomendaciones de ubicación ideal por tipo de negocio ──
            _recom_vialidad  = ""
            _recom_visib     = ""
            _recom_estac     = ""
            _recom_plaza     = ""
            _ref_comp_txt    = ""
            _tip_nse_txt     = ""

            _tipo_info_r = TIPOS_NEGOCIO.get(tipo_negocio_key, {})
            _tipo_nombre_r = _tipo_info_r.get("nombre", "")

            # Vialidad ideal — según dependencia de paso
            _dep_paso = ("alto" if tipo_negocio_key in NEGOCIOS_PASO["alto"]
                         else "medio" if tipo_negocio_key in NEGOCIOS_PASO["medio"] else "bajo")
            if tipo_negocio_key in ("cafe_premium", "cafe_casual"):
                _recom_vialidad = "Avenida de alto tráfico o vialidad principal. El café necesita flujo peatonal constante y visibilidad desde la calle."
            elif tipo_negocio_key == "medicina_estetica":
                _recom_vialidad = "Calle tranquila o avenida secundaria residencial. La medicina estética requiere ambiente discreto y aspiracional — evitar vialidades de alto ruido."
            elif tipo_negocio_key == "nutricion":
                _recom_vialidad = "Calle tranquila o avenida secundaria. Puede estar en planta alta o interior de plaza; el cliente va por recomendación o búsqueda activa."
            elif tipo_negocio_key in ("comida_rapida", "empanadas_antojitos"):
                _recom_vialidad = "Vialidad principal o avenida de alto tráfico. Máxima visibilidad y captación de tráfico de paso son esenciales para el modelo de alto volumen."
            elif tipo_negocio_key in ("gimnasio_boutique", "gimnasio_regular", "yoga_wellness"):
                _recom_vialidad = "Avenida secundaria con estacionamiento disponible. El cliente va a propósito — la accesibilidad vehicular importa más que la visibilidad de paso."
            elif tipo_negocio_key in ("acabados_hogar", "ferreteria"):
                _recom_vialidad = "Avenida secundaria o vialidad con acceso vehicular fácil y posibilidad de carga/descarga. El cliente lleva material — necesita espacio para maniobrar."
            elif tipo_negocio_key == "farmacia":
                _recom_vialidad = "Calle transitada o esquina de avenida. La farmacia necesita ser visible y accesible, idealmente con paso peatonal constante."
            elif tipo_negocio_key == "restaurante_fino":
                _recom_vialidad = "Calle tranquila o zona de restaurantes. El restaurante fino evita vialidades de ruido; el cliente lo busca, no lo descubre de paso."
            else:
                if _dep_paso == "alto":
                    _recom_vialidad = "Avenida o vialidad con tráfico constante. Este tipo de negocio depende del flujo de paso para captar clientes."
                elif _dep_paso == "bajo":
                    _recom_vialidad = "Calle tranquila o interior de zona comercial. El cliente llega a propósito — la ubicación discreta reduce renta y no penaliza ventas."
                else:
                    _recom_vialidad = "Avenida secundaria con buena visibilidad. Equilibrio entre accesibilidad y costo de renta."

            # Visibilidad
            if tipo_negocio_key in ("cafe_premium","cafe_casual","comida_rapida","empanadas_antojitos","panaderia","tienda_conveniencia","farmacia"):
                _recom_visib = "Rotulación exterior visible a 50m mínimo. Google Maps con fotos actualizadas y horarios correctos. Escaparate o fachada que invite a entrar desde la banqueta."
            elif tipo_negocio_key in ("medicina_estetica","nutricion","yoga_wellness","gimnasio_boutique"):
                _recom_visib = "El cliente llega por búsqueda en Google y recomendación — prioriza una ficha de Google My Business impecable con fotos, reseñas y respuestas. La rotulación exterior puede ser discreta."
            elif tipo_negocio_key in ("acabados_hogar","ferreteria"):
                _recom_visib = "Muestra de producto visible desde la calle o vitrina de fachada. Los clientes buscan por referencia o Google Maps — asegura fotos actualizadas del local."
            else:
                _recom_visib = "Letrero visible desde la calle con nombre claro del negocio. Google Maps con fotos y reseñas actualizadas para captar búsquedas locales."

            # Estacionamiento
            if tipo_negocio_key in ("medicina_estetica","acabados_hogar","ferreteria","gimnasio_boutique","gimnasio_regular"):
                _recom_estac = "Indispensable. El cliente lleva materiales, equipo o viene en auto desde lejos. Sin estacionamiento propio o público cercano, el negocio pierde clientes directamente."
            elif tipo_negocio_key in ("cafe_premium","cafe_casual","restaurante_casual","restaurante_fino","bar"):
                _recom_estac = "Conveniente. Un café o restaurante puede operar sin estacionamiento en zona de alto peatonal, pero perderá clientes en zonas suburbanas. Evalúa según el perfil de la zona."
            elif tipo_negocio_key in ("comida_rapida","empanadas_antojitos","panaderia"):
                _recom_estac = "No crítico si hay alto tráfico peatonal. En zona vehicular, un cajón o bahía de parada rápida mejora las ventas de paso."
            else:
                _recom_estac = "Verifica disponibilidad de cajones propios o estacionamiento público a menos de 100m. Es un factor que muchos emprendedores subestiman al firmar el contrato."

            # Plaza comercial
            if tipo_negocio_key in ("cafe_premium","cafe_casual","comida_rapida","empanadas_antojitos","panaderia","restaurante_casual"):
                _recom_plaza = "Planta baja obligatorio, idealmente en food court o corredor gastronómico. En mall cerrado, la ubicación cerca de acceso principal o ancla multiplica el tráfico. Evitar piso superior."
            elif tipo_negocio_key in ("medicina_estetica","nutricion","yoga_wellness","guarderia"):
                _recom_plaza = "Puede ir en 2do piso o interior de plaza si el negocio tiene presencia digital fuerte. Reduce renta significativamente sin penalizar ventas — el cliente va a propósito."
            elif tipo_negocio_key in ("farmacia","tienda_conveniencia","servicios"):
                _recom_plaza = "Planta baja con acceso directo desde la calle. En plaza, posición cerca de la entrada o ancla principal. El tráfico de paso es la fuente principal de clientes."
            elif tipo_negocio_key in ("acabados_hogar","ferreteria"):
                _recom_plaza = "Strip center o bodega comercial con espacio para exhibición y carga. En mall cerrado no es el formato adecuado — el cliente necesita espacio y privacidad para decidir."
            else:
                _recom_plaza = "Planta baja preferible para visibilidad. Evalúa el tráfico real del corredor antes de firmar — no todas las plazas tienen flujo constante."

            # Competidor de referencia
            if _mejor_comp_nombre:
                _ref_comp_txt = (f"'{_mejor_comp_nombre[:50]}' ({_mejor_comp_rating}★, {_mejor_comp_reviews} reseñas) "
                                 f"es el competidor mejor calificado en esta zona. Visítalo antes de abrir: "
                                 f"analiza su ticket, servicio, decoración y qué le falta.")

            # Tip por NSE
            _nse_z = nse
            if tipo_negocio_key == "medicina_estetica":
                if "A" in _nse_z or "B" in _nse_z:
                    _tip_nse_txt = f"NSE {_nse_z}: cliente de medicina estética espera instalaciones premium, marca personal del médico y discreción. Invierte en diseño de interiores y comunicación aspiracional."
                else:
                    _tip_nse_txt = f"NSE {_nse_z}: considera un modelo de precios accesibles en tratamientos de entrada (depilación, limpiezas) antes de ofrecer procedimientos de alto costo."
            elif tipo_negocio_key == "nutricion":
                if "A" in _nse_z or "B" in _nse_z:
                    _tip_nse_txt = f"NSE {_nse_z}: cliente de nutrición en zona alta espera enfoque integral (bienestar, no solo dieta), horarios flexibles y plataforma digital para seguimiento."
                elif "C+" in _nse_z:
                    _tip_nse_txt = f"NSE {_nse_z}: cliente C+ valora resultados concretos y precio accesible. Paquetes por sesiones con precio cerrado funcionan mejor que cobro por hora."
                else:
                    _tip_nse_txt = f"NSE {_nse_z}: mercado más sensible al precio — considera alianza con gimnasios locales o consultas a domicilio para reducir barrera de entrada."
            elif tipo_negocio_key == "empanadas_antojitos":
                if "C" in _nse_z or "D" in _nse_z:
                    _tip_nse_txt = f"NSE {_nse_z}: tu cliente principal. Ticket de $30-80 MXN, combos visibles en menú exterior, pago en efectivo esencial, rápido y sabroso."
                else:
                    _tip_nse_txt = f"NSE {_nse_z}: en zona de mayor ingreso, eleva la propuesta: ingredientes premium, nombre de marca, presentación cuidada y opciones vegetarianas/saludables."
            elif "A" in _nse_z or "B" in _nse_z:
                _tip_nse_txt = f"NSE {_nse_z}: zona de alto poder adquisitivo — prioriza calidad, experiencia y marca. El precio es secundario; la percepción de valor es lo que convierte."
            elif "C+" in _nse_z:
                _tip_nse_txt = f"NSE {_nse_z}: zona C+ — equilibrio precio-calidad. Propuesta clara, beneficios visibles y buen trato son los diferenciadores principales."
            elif "C" in _nse_z and "C+" not in _nse_z:
                _tip_nse_txt = f"NSE {_nse_z}: zona C estándar — el precio importa. Promos de lanzamiento, ticket accesible y presencia en WhatsApp Business aceleran la adopción inicial."
            else:
                _tip_nse_txt = f"NSE {_nse_z}: mercado sensible al precio. Volumen sobre margen. Evalúa si el ticket mínimo viable del negocio es compatible con el ingreso de la zona."

            _recom_ubicacion = {
                "vialidad_ideal":       _recom_vialidad,
                "visibilidad":          _recom_visib,
                "estacionamiento":      _recom_estac,
                "plaza_recomendacion":  _recom_plaza,
                "ref_competidor":       _ref_comp_txt,
                "tip_nse":              _tip_nse_txt,
            }

            resultados.append({
                "lat":              p_lat,
                "lng":              p_lng,
                "score":            round(score_final, 1),
                "direccion":        direccion,
                "nse":              nse,
                "densidad":         int(densidad),
                "ingreso":          int(ingreso),
                "n_competidores":   n_comp,
                "nivel_score":      nivel,
                "tipo_negocio_key": tipo_negocio_key,
                "desglose":         _desglose,
                "ticket_zona":      _ticket_zona,
                "vialidad":         _vialidad,
                "zona_categoria":   _zona_cat,
                "mejor_competidor": _mejor_comp_nombre,
                "mejor_rating":     _mejor_comp_rating,
                "mejor_reviews":    _mejor_comp_reviews,
                "recom_ubicacion":  _recom_ubicacion,
            })
        except Exception as _e:
            errores.append(str(_e))
            continue

    resultados.sort(key=lambda x: x["score"], reverse=True)

    # ── Asignar ranking_pos según zona_categoria ya calculada ──
    # Orden: mercado_probado y moderada primero (TOP 3), virgen y saturada después (TOP 4-6)
    _rank_prio = {"mercado_probado": 0, "moderada": 1, "virgen": 2, "saturada": 3}
    resultados.sort(key=lambda x: (_rank_prio.get(x.get("zona_categoria","moderada"), 1), -x["score"]))
    for _ri, _r in enumerate(resultados):
        _r["ranking_pos"] = _ri + 1

    return resultados[:10]


def generar_pptx_zona_potencial(resultados, ubicacion_ref, radio_km,
                                  tipo_negocio_nombre, tier_key="premium"):
    """
    Genera PPTX de zonas potenciales con el mismo formato que el reporte estándar 4SITE.
    Slides:
      1. Portada (misma foto bg_ciudad.jpg)
      2. Top 3 — Mercado Probado
      3. Top 4-6 — Zonas Vírgenes
      4. Detalle Zona #1 (mejor zona) con recomendaciones
      5. Guía de Ubicación Ideal
      6. Recomendación Final
    """
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN
        from io import BytesIO
        import datetime
        import os
    except ImportError:
        return None

    # _find_asset local — no depende de generar_pptx_4site
    def _find_asset_local(nombres):
        base_dirs = [
            os.path.dirname(os.path.abspath(__file__)) if '__file__' in dir() else ".",
            os.getcwd(),
            "/app", "/mount/src",
        ]
        for d in base_dirs:
            for n in nombres:
                p = os.path.join(d, n)
                if os.path.exists(p):
                    return p
        return None

    if not resultados:
        return None

    # Usar los colores y helpers del mismo generador principal
    C = {
        "dorado": "BB944F", "verde": "27AE60", "azul": "1A76D2",
        "amarillo": "F39C12", "rojo": "E74C3C", "blanco": "FFFFFF",
        "oscuro": "1A1A1A", "gris_t": "6B6B6B", "gris_c": "F5F5F5",
        "gris_b": "EBEBEB", "card_sc": "262626",
    }
    F = "Montserrat"

    def rgb(h):
        h = h.lstrip("#")
        return RGBColor(int(h[0:2],16), int(h[2:4],16), int(h[4:6],16))

    def bg_s(sl, col):
        sl.background.fill.solid()
        sl.background.fill.fore_color.rgb = rgb(col)

    def box(sl, x, y, w, h, fill, lc=None, lw=0.5):
        s = sl.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
        s.fill.solid(); s.fill.fore_color.rgb = rgb(fill)
        if lc: s.line.color.rgb = rgb(lc); s.line.width = Pt(lw)
        else: s.line.fill.background()
        return s

    def txt(sl, t, x, y, w, h, sz=10, bold=False, col=None, al=None, it=False):
        col = col or C["oscuro"]
        tb = sl.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
        tf = tb.text_frame; tf.word_wrap = True
        p = tf.paragraphs[0]; p.text = str(t)
        if al: p.alignment = al
        r = p.runs[0]; r.font.size = Pt(sz); r.font.bold = bold; r.font.italic = it
        r.font.color.rgb = rgb(col)
        try: r.font.name = F
        except: pass

    def hdr(sl, title, sub=None):
        box(sl, 0, 0, 13.33, 0.72, C["oscuro"])
        box(sl, 0, 0, 0.08, 0.72, C["dorado"])
        txt(sl, title, 0.22, 0.1, 12.8, 0.5, sz=15, bold=True, col=C["blanco"])
        if sub: txt(sl, sub, 0.22, 0.56, 12.8, 0.2, sz=8, col="888888", it=True)

    def kpi(sl, x, y, w, val, label, acc):
        box(sl, x, y, w, 0.95, C["blanco"], C["gris_b"], 0.4)
        box(sl, x, y, w, 0.055, acc)
        txt(sl, str(val), x+0.08, y+0.1, w-0.16, 0.48,
            sz=16, bold=True, col=acc, al=PP_ALIGN.CENTER)
        txt(sl, label, x+0.05, y+0.6, w-0.1, 0.3,
            sz=8, col=C["gris_t"], al=PP_ALIGN.CENTER)

    def interp(sl, y, text, h=0.55):
        box(sl, 0.3, y, 12.7, h, "FFF8EE", C["dorado"], 0.7)
        box(sl, 0.3, y, 0.06, h, C["dorado"])
        txt(sl, text, 0.46, y+0.06, 12.38, h-0.1, sz=9, col="3D3000", it=True)

    def zona_row(sl, pos, res, y, is_top3=True):
        sc_v = res.get("score", 0)
        sc_c = C["verde"] if sc_v>=65 else C["amarillo"] if sc_v>=40 else C["rojo"]
        bdr_c = C["verde"] if is_top3 else C["azul"]
        box(sl, 0.3, y, 12.7, 0.88, C["gris_c"], bdr_c, 0.4)
        box(sl, 0.3, y, 0.06, 0.88, bdr_c)
        # Posición
        txt(sl, f"#{pos}", 0.45, y+0.08, 0.4, 0.38, sz=12, bold=True, col=C["oscuro"])
        # Dirección
        txt(sl, res.get("direccion","")[:55], 0.92, y+0.06, 7.0, 0.3, sz=10, bold=True, col=C["oscuro"])
        # Categoria
        cat = {"mercado_probado":"Mercado probado","compartir_mercado":"Compartir mercado",
               "virgen":"Zona virgen","moderada":"Zona moderada"}.get(res.get("zona_categoria",""), "")
        cat_c = C["verde"] if is_top3 else C["azul"]
        txt(sl, cat, 0.92, y+0.42, 3.5, 0.28, sz=8, col=cat_c, bold=True)
        # Badges
        badges = f"NSE {res.get('nse','—')}  |  {res.get('densidad',0):,} hab/km2  |  {res.get('n_competidores',0)} competidores"
        if res.get("vialidad"): badges += f"  |  {res.get('vialidad','')}"
        txt(sl, badges, 4.0, y+0.42, 6.5, 0.28, sz=8, col=C["gris_t"])
        # Score
        txt(sl, f"{sc_v:.0f}", 11.5, y+0.06, 0.75, 0.44, sz=22, bold=True, col=sc_c, al=PP_ALIGN.RIGHT)
        txt(sl, "/100", 12.25, y+0.24, 0.7, 0.22, sz=8, col=C["gris_t"])
        # Competidor ref si aplica
        if res.get("mejor_competidor") and is_top3:
            ref = f"Ref: {res['mejor_competidor'][:35]} — {res.get('mejor_rating',0)} estrellas, {res.get('mejor_reviews',0)} resenas"
            txt(sl, ref, 0.92, y+0.62, 9.5, 0.22, sz=7.5, col="555555", it=True)

    def foot(sl, n, tot):
        box(sl, 0, 7.3, 13.33, 0.2, C["gris_c"])
        txt(sl, "4SITE · 4site.mx · hola@4site.mx", 0.3, 7.32, 8, 0.16, sz=7, col=C["gris_t"])
        txt(sl, f"{n} / {tot}", 12.0, 7.32, 1.2, 0.16, sz=7, col=C["gris_t"], al=PP_ALIGN.RIGHT)

    # ── Datos ───────────────────────────────────────────────
    top3   = [r for r in resultados if r.get("ranking_pos",99) <= 3][:3]
    top456 = [r for r in resultados if r.get("ranking_pos",99) > 3][:3]
    top1   = resultados[0] if resultados else {}
    bg_path = _find_asset_local(["bg_calle.jpg","bg_ciudad.jpg"])
    fecha = datetime.datetime.now().strftime("%d de %B de %Y")
    tl = {"premium":"PREMIUM $999","premium2":"ZONA ELITE $1,500","pro":"PRO ($499)"}.get(tier_key,"PREMIUM")
    total_slides = 6
    tmp_files = []

    prs = Presentation()
    prs.slide_width = Inches(13.33); prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]; sn = 0

    # ── S1: PORTADA (misma que reporte estándar) ────────────
    sn += 1
    sl = prs.slides.add_slide(blank)
    bg_s(sl, "1A1A1A")

    if bg_path and os.path.exists(bg_path):
        try:
            from PIL import Image as _PIL, ImageDraw as _IDraw
            im = _PIL.open(bg_path).convert("RGB")
            tw, th = 1333, 750
            ir = im.width/im.height; pr = tw/th
            nw,nh = (int(th*ir),th) if ir>pr else (tw,int(tw/ir))
            im2 = im.resize((nw,nh))
            lx,ly = (nw-tw)//2,(nh-th)//2
            im3 = im2.crop((lx,ly,lx+tw,ly+th)).convert("RGBA")
            overlay = _PIL.new("RGBA",(tw,th),(0,0,0,0))
            draw = _IDraw.Draw(overlay)
            draw.rectangle([743,161,743+431,161+429], fill=(26,26,26,195))
            im_final = _PIL.alpha_composite(im3,overlay).convert("RGB")
            import tempfile
            _tf = tempfile.NamedTemporaryFile(suffix=".jpg",delete=False)
            im_final.save(_tf.name,"JPEG",quality=92); _tf.close()
            tmp_files.append(_tf.name)
            sl.shapes.add_picture(_tf.name,Inches(0),Inches(-1.45),Inches(13.333),Inches(8.952))
        except: pass

    # Logo
    txt(sl, "4", 0.28, 0.12, 0.6, 0.72, sz=38, bold=True, col=C["dorado"])
    txt(sl, "SITE", 0.82, 0.15, 2.2, 0.65, sz=34, bold=True, col=C["blanco"])
    txt(sl, "Don't guess. Foresee.", 0.28, 0.9, 4.5, 0.32, sz=11, it=True, col=C["dorado"])

    # Panel derecho
    PX = 5.4
    box(sl, PX, 0.44, 2.6, 0.38, C["dorado"])
    txt(sl, tl, PX, 0.44, 2.6, 0.38, sz=11, bold=True, col=C["oscuro"], al=PP_ALIGN.CENTER)
    box(sl, PX, 0.9, 2.6, 0.025, C["dorado"])
    txt(sl, "REPORTE DE ZONAS POTENCIALES", PX, 1.0, 7.8, 0.26, sz=7.5, col="999999")
    txt(sl, f"Referencia: {ubicacion_ref}", PX, 1.36, 7.8, 0.38, sz=13, bold=True, col=C["blanco"])
    txt(sl, tipo_negocio_nombre, PX, 1.82, 7.8, 0.38, sz=12, col=C["dorado"])
    txt(sl, f"Radio: {radio_km} km  |  {len(resultados)} zonas evaluadas", PX, 2.28, 7.8, 0.28, sz=10, col="AAAAAA")
    # Top score
    sc_top = top1.get("score",0)
    sc_c_top = C["verde"] if sc_top>=65 else C["amarillo"] if sc_top>=40 else C["rojo"]
    txt(sl, f"{sc_top:.0f}/100", PX, 2.7, 4.0, 0.85, sz=48, bold=True, col=sc_c_top)
    txt(sl, "Score mejor zona", PX, 3.52, 4.0, 0.28, sz=9, col="888888")
    txt(sl, fecha, PX, 6.88, 4, 0.28, sz=9, col="777777")
    txt(sl, "4SITE · 4site.mx", PX+4.0, 6.88, 3.6, 0.28, sz=9, col=C["dorado"], al=PP_ALIGN.RIGHT)

    # ── S2: TOP 3 — Mercado Probado ─────────────────────────
    sn += 1
    sl = prs.slides.add_slide(blank)
    bg_s(sl, C["gris_c"])
    hdr(sl, "TOP 3 — Mercado Probado",
        "Zonas con 1 competidor bien calificado (3.8+, 30+ resenas) — la demanda esta comprobada")

    # Explicación metodología
    box(sl, 0.3, 0.78, 12.7, 0.52, "E8F8F0", C["verde"], 0.5)
    box(sl, 0.3, 0.78, 0.06, 0.52, C["verde"])
    txt(sl, "Logica: 1 competidor fuerte en la zona = mercado existente que puedes capturar con propuesta superior. "
           "El cliente ya tiene el habito — tu llegas con mejor calidad, precio o experiencia.",
        0.46, 0.84, 12.3, 0.4, sz=8.5, col="1A6B3C", it=True)

    if top3:
        for i, res in enumerate(top3):
            try:
                zona_row(sl, i+1, res, 1.42 + i*1.0, is_top3=True)
            except Exception:
                pass
    else:
        txt(sl, "No se encontraron zonas con competidor de referencia en el radio analizado. "
               "Considera aumentar el radio de busqueda o cambiar el tipo de negocio.",
            0.3, 1.5, 12.7, 0.6, sz=10, col=C["gris_t"], it=True)
    foot(sl, sn, total_slides)

    # ── S3: TOP 4-6 — Zonas Vírgenes ────────────────────────
    sn += 1
    sl = prs.slides.add_slide(blank)
    bg_s(sl, C["gris_c"])
    hdr(sl, "TOP 4-6 — Zonas Virgenes de Perfil Compatible",
        "Sin competencia directa — perfil demografico y vial similar al entorno del mejor competidor del radio")
    box(sl, 0.3, 0.78, 12.7, 0.52, "EBF5FF", C["azul"], 0.5)
    box(sl, 0.3, 0.78, 0.06, 0.52, C["azul"])
    txt(sl, "Logica: zonas sin competidores directos que comparten el perfil NSE, densidad y vialidad del entorno "
           "de los mejores competidores encontrados. Alto potencial si se valida demanda en campo antes de invertir.",
        0.46, 0.84, 12.3, 0.4, sz=8.5, col="1A4A7A", it=True)

    if top456:
        for i, res in enumerate(top456):
            try:
                zona_row(sl, len(top3)+i+1, res, 1.42 + i*1.0, is_top3=False)
            except Exception:
                pass
    else:
        txt(sl, "No se encontraron zonas virgenes con perfil compatible en el radio analizado.",
            0.3, 1.5, 12.7, 0.4, sz=10, col=C["gris_t"], it=True)
    foot(sl, sn, total_slides)

    # ── S4: Detalle Zona #1 ──────────────────────────────────
    sn += 1
    sl = prs.slides.add_slide(blank)
    bg_s(sl, C["gris_c"])
    dir1 = top1.get("direccion","—")
    sc1  = top1.get("score",0)
    sc_c1 = C["verde"] if sc1>=65 else C["amarillo"] if sc1>=40 else C["rojo"]
    hdr(sl, f"Detalle Zona #1 — {dir1[:55]}", f"Score: {sc1:.0f}/100  |  Categoria: {top1.get('zona_categoria','')}")

    # KPIs
    kpi(sl, 0.3,  0.82, 2.4, top1.get("nse","—"), "NSE predominante", C["azul"])
    kpi(sl, 2.88, 0.82, 2.4, f"{top1.get('densidad',0):,}", "Hab/km2", C["azul"])
    kpi(sl, 5.46, 0.82, 2.4, top1.get("n_competidores",0), "Competidores directos", C["dorado"])
    kpi(sl, 8.04, 0.82, 2.4, f"{top1.get('mejor_rating',0)}", "Rating mejor competidor", C["amarillo"])
    kpi(sl, 10.62, 0.82, 2.38, top1.get("mejor_reviews",0), "Resenas mejor competidor", C["amarillo"])

    # Recomendaciones de ubicación
    recom1 = top1.get("recom_ubicacion", {})
    txt(sl, "Recomendaciones de Ubicacion Ideal:", 0.3, 1.95, 9, 0.28, sz=11, bold=True, col=C["oscuro"])

    rec_items = [
        ("Vialidad ideal",     recom1.get("vialidad_ideal","—")),
        ("Visibilidad",        recom1.get("visibilidad","—")),
        ("Estacionamiento",    recom1.get("estacionamiento","—")),
        ("En plaza comercial", recom1.get("plaza_recomendacion","—")),
    ]
    for ri, (label, texto) in enumerate(rec_items):
        ry = 2.28 + ri * 0.88
        box(sl, 0.3, ry, 6.2, 0.78, C["blanco"], C["gris_b"], 0.4)
        box(sl, 0.3, ry, 0.06, 0.78, C["dorado"])
        txt(sl, label, 0.46, ry+0.05, 5.9, 0.22, sz=8.5, bold=True, col=C["gris_t"])
        txt(sl, texto[:120], 0.46, ry+0.28, 5.9, 0.44, sz=8.5, col=C["oscuro"])

        # Lado derecho: competidor ref y tip NSE
        if ri == 0 and recom1.get("ref_competidor"):
            box(sl, 6.75, ry, 6.28, 0.78, C["blanco"], C["gris_b"], 0.4)
            box(sl, 6.75, ry, 0.06, 0.78, C["azul"])
            txt(sl, "Competidor de referencia", 6.91, ry+0.05, 5.9, 0.22, sz=8.5, bold=True, col=C["gris_t"])
            txt(sl, recom1["ref_competidor"][:130], 6.91, ry+0.28, 5.9, 0.44, sz=8.5, col=C["oscuro"])
        elif ri == 1 and recom1.get("tip_nse"):
            box(sl, 6.75, ry, 6.28, 0.78, C["blanco"], C["gris_b"], 0.4)
            box(sl, 6.75, ry, 0.06, 0.78, "27AE60")
            txt(sl, "Tip para este NSE", 6.91, ry+0.05, 5.9, 0.22, sz=8.5, bold=True, col=C["gris_t"])
            txt(sl, recom1["tip_nse"][:130], 6.91, ry+0.28, 5.9, 0.44, sz=8.5, col=C["oscuro"])
    foot(sl, sn, total_slides)

    # ── S5: Guia de Ubicacion Ideal ──────────────────────────
    sn += 1
    sl = prs.slides.add_slide(blank)
    bg_s(sl, C["gris_c"])
    hdr(sl, "Guia de Ubicacion Ideal",
        f"Recomendaciones especificas para {tipo_negocio_nombre} segun el analisis del entorno")

    guia_items = [
        (C["dorado"],  "Tipo de Vialidad Recomendada",
         recom1.get("vialidad_ideal", "Avenida secundaria con buena visibilidad y acceso vehicular.")),
        (C["azul"],    "Visibilidad desde la Calle",
         recom1.get("visibilidad", "Rotulacion exterior visible a 50m. Google Maps con fotos actualizadas.")),
        (C["verde"],   "Estacionamiento y Accesibilidad",
         recom1.get("estacionamiento", "Verifica cajones propios o estacionamiento publico cercano.")),
        (C["amarillo"],"Ubicacion en Plaza Comercial",
         recom1.get("plaza_recomendacion", "Planta baja preferible para visibilidad. Evalua flujo real antes de firmar.")),
        ("9B59B6",     "Competidor de Referencia",
         recom1.get("ref_competidor", "Analiza el competidor mejor calificado de la zona antes de decidir.")),
        (C["dorado"],  "Tip por NSE del Area",
         recom1.get("tip_nse", "Adapta tu propuesta al nivel socioeconomico predominante.")),
    ]

    for gi, (gc, gtitle, gtext) in enumerate(guia_items):
        gx = 0.3 + (gi%2) * 6.52
        gy = 0.88 + (gi//2) * 1.12
        box(sl, gx, gy, 6.22, 1.02, C["blanco"], C["gris_b"], 0.4)
        box(sl, gx, gy, 0.06, 1.02, gc)
        txt(sl, gtitle, gx+0.16, gy+0.06, 5.9, 0.24, sz=9.5, bold=True, col=C["gris_t"])
        txt(sl, gtext[:140], gx+0.16, gy+0.32, 5.9, 0.64, sz=8.5, col=C["oscuro"])
    foot(sl, sn, total_slides)

    # ── S6: Recomendación Final ──────────────────────────────
    sn += 1
    sl = prs.slides.add_slide(blank)
    bg_s(sl, "3D3D3D")
    box(sl, 0, 0, 13.33, 0.72, C["dorado"])
    box(sl, 0, 0, 0.08, 0.72, "3D3D3D")
    txt(sl, "Recomendacion Final", 0.22, 0.1, 9, 0.5, sz=17, bold=True, col=C["oscuro"])
    txt(sl, f"{ubicacion_ref[:70]}  |  {tipo_negocio_nombre}",
        0.22, 0.56, 12, 0.22, sz=8.5, col="444444")

    # Zona recomendada
    box(sl, 0.3, 0.86, 12.7, 0.68, "252525")
    txt(sl, "ZONA RECOMENDADA", 0.46, 0.92, 4, 0.24, sz=8.5, bold=True, col=C["dorado"])
    txt(sl, dir1, 0.46, 1.16, 11, 0.3, sz=13, bold=True, col=C["blanco"])

    # Lógica explicada
    y = 1.66
    box(sl, 0.3, y, 12.7, 0.06, C["dorado"])
    y += 0.16

    txt(sl, "Por que esta es la mejor zona:", 0.3, y, 9, 0.28, sz=10, bold=True, col=C["blanco"])
    y += 0.38

    cat_top = top1.get("zona_categoria","")
    if cat_top in ["mercado_probado","compartir_mercado"]:
        logica = (f"Esta zona tiene {top1.get('n_competidores',0)} competidor(es) con {top1.get('mejor_rating',0)} estrellas "
                  f"y {top1.get('mejor_reviews',0)} resenas, lo que confirma que la demanda para {tipo_negocio_nombre} "
                  f"existe y esta activa en este mercado. No estas apostando en el vacio: hay clientes que buscan "
                  f"este servicio aqui. Tu ventaja es entrar con una propuesta mejor calibrada al NSE {top1.get('nse','—')} del area.")
    else:
        logica = (f"Esta zona no tiene competidores directos para {tipo_negocio_nombre}, pero su perfil demografico "
                  f"(NSE {top1.get('nse','—')}, {top1.get('densidad',0):,} hab/km2) es compatible con las zonas donde si "
                  f"hay demanda comprobada. Seras el primero en capturar este mercado — la clave es validar el trafico "
                  f"real en campo antes de firmar contrato.")
    txt(sl, logica, 0.3, y, 12.7, 0.8, sz=9.5, col="CCCCCC")
    y += 0.92

    # Tips finales
    box(sl, 0.3, y, 12.7, 0.06, "333333")
    y += 0.14
    txt(sl, "La ubicacion es el punto de partida — el exito tambien depende de:",
        0.3, y, 12.7, 0.28, sz=9.5, bold=True, col=C["dorado"])
    y += 0.38
    tips = [
        ("Ticket calibrado al NSE", "Ajusta tu precio al poder adquisitivo de la zona. Un ticket mal calibrado cancela una buena ubicacion."),
        ("Atencion al cliente",     "En zonas con competencia cercana, la experiencia es el diferenciador. Tu rating en Google impacta captacion."),
        ("Presencia digital",       "Google Maps, Instagram y WhatsApp Business son esenciales. Sin visibilidad digital, pierdes clientes de paso."),
    ]
    for tx2, (tip_t, tip_d) in enumerate(tips):
        bx = 0.3 + tx2 * 4.3
        box(sl, bx, y, 4.1, 1.0, "252525")
        txt(sl, tip_t, bx+0.12, y+0.08, 3.86, 0.26, sz=9, bold=True, col=C["dorado"])
        txt(sl, tip_d, bx+0.12, y+0.36, 3.86, 0.56, sz=8, col="AAAAAA")

    txt(sl, "4SITE · 4site.mx", 11.0, 7.15, 2.2, 0.24, sz=8, col="555555", al=PP_ALIGN.RIGHT)

    buf = BytesIO()
    try:
        prs.save(buf)
    except Exception:
        # Si falla al guardar, intentar una versión mínima de 1 slide
        try:
            prs2 = Presentation()
            prs2.slide_width = Inches(13.33); prs2.slide_height = Inches(7.5)
            sl_err = prs2.slides.add_slide(prs2.slide_layouts[6])
            from pptx.util import Pt as _Pt2
            tb = sl_err.shapes.add_textbox(Inches(1), Inches(3), Inches(11), Inches(1.5))
            tf = tb.text_frame
            tf.paragraphs[0].text = f"Reporte zona potencial — {ubicacion_ref[:60]}"
            buf2 = BytesIO(); prs2.save(buf2); buf2.seek(0)
            for f in tmp_files:
                try: os.unlink(f)
                except: pass
            return buf2.getvalue()
        except Exception:
            return None
    buf.seek(0)
    for f in tmp_files:
        try: os.unlink(f)
        except: pass
    return buf.getvalue()


def _render_busqueda_zona(idioma: str = "es"):
    """Renderiza la UI de búsqueda de zona ideal en Streamlit."""
    t = TEXTOS[idioma]
    st.markdown("### 🔍 Buscar zona ideal")
    st.caption("Encuentra las mejores zonas para tu negocio dentro de un radio de búsqueda.")

    col1, col2 = st.columns([2, 1])
    with col1:
        dir_ref = st.text_input(
            "📍 Punto de referencia",
            placeholder="Ej: Av. Insurgentes Sur 1458, CDMX",
            help="No tiene que ser la futura ubicación — puede ser tu zona de interés o tu casa."
        )
    with col2:
        tipo_busq = st.selectbox(
            "Tipo de negocio",
            options=list(TIPOS_NEGOCIO.keys()),
            format_func=lambda x: TIPOS_NEGOCIO[x]["nombre"],
            key="tipo_busqueda_zona"
        )

    radios_label = {"1 km (~10 min caminando)": 1.0, "3 km (recomendado)": 3.0,
                    "5 km (zona amplia)": 5.0,        "10 km (ciudad entera)": 10.0}
    radio_sel = st.radio(
        "Radio de búsqueda",
        options=list(radios_label.keys()),
        index=1, horizontal=True
    )
    radio_km = radios_label[radio_sel]

    n_puntos = {1.0: "~9", 3.0: "~12", 5.0: "~16", 10.0: "~25"}.get(radio_km, "~12")
    st.caption(f"Se evaluarán {n_puntos} puntos · Tiempo estimado: {'~20 seg' if radio_km<=3 else '~45 seg' if radio_km<=5 else '~2 min'}")

    if st.button("🗺️ Buscar mejores zonas", type="primary", use_container_width=True):
        if not dir_ref.strip():
            st.warning("Ingresa un punto de referencia para iniciar la búsqueda.")
            return

        with st.spinner(f"Analizando zona de {radio_km:.0f} km a la redonda..."):
            try:
                _geo_result = geocodificar_direccion(dir_ref)
                if not _geo_result or _geo_result == (None, None):
                    st.error("No se pudo geocodificar la dirección. Intenta con una más específica.")
                    return
                lat_r, lng_r = _geo_result if isinstance(_geo_result, tuple) else (_geo_result.get("lat"), _geo_result.get("lng"))
                if not lat_r or not lng_r:
                    st.error("No se encontraron coordenadas. Intenta con una dirección más específica.")
                    return

                # Barra de progreso
                _p_bar  = st.progress(0, text="Analizando zonas...")
                _p_txt  = st.empty()
                def _upd(a, t):
                    _p_bar.progress(int(a/t*100), text=f"Evaluando punto {a} de {t}...")
                try:
                    resultados = buscar_zona_ideal(lat_r, lng_r, radio_km, tipo_busq,
                                                   idioma, progress_callback=_upd)
                finally:
                    _p_bar.empty(); _p_txt.empty()

                if not resultados:
                    st.warning("⚠️ No se encontraron zonas evaluables. Prueba con radio mayor (5 km o 10 km) o verifica que la dirección esté en una ciudad con datos de Google Maps.")
                    return

                # ── Separar Top 3 (mercado probado) y Top 4-6 (zona virgen) ──
                _top3   = [r for r in resultados if r.get("ranking_pos",99) <= 3]
                _top456 = [r for r in resultados if r.get("ranking_pos",99) > 3]
                _tipo_nombre_z = TIPOS_NEGOCIO.get(tipo_busq,{}).get("nombre","tu negocio")

                st.success(f"Analisis completado: {len(resultados)} zonas evaluadas, {len(_top3)} con mercado probado, {len(_top456)} zonas virgenes")

                st.markdown(f"""
                <div style='background:#F8F9FA; border:1px solid #E0E0E0; border-radius:10px;
                     padding:14px 18px; margin:8px 0 16px;'>
                <p style='margin:0 0 6px; font-size:12px; font-weight:700; color:#333;'>
                    Metodologia de clasificacion para {_tipo_nombre_z}:</p>
                <div style='display:grid; grid-template-columns:1fr 1fr; gap:10px;'>
                    <div style='background:#E8F8F0; border-left:3px solid #27AE60;
                         border-radius:6px; padding:8px 12px;'>
                        <span style='font-size:11px; font-weight:700; color:#1A7A44;'>
                            TOP 3 - Mercado Probado</span><br>
                        <span style='font-size:11px; color:#2D6A4F;'>
                            Zonas con 1 competidor bien calificado (3.8+ estrellas, 30+ resenas).
                            La demanda esta comprobada. Entra con propuesta superior y captura ese mercado.</span>
                    </div>
                    <div style='background:#EBF5FF; border-left:3px solid #1A76D2;
                         border-radius:6px; padding:8px 12px;'>
                        <span style='font-size:11px; font-weight:700; color:#1A4A7A;'>
                            TOP 4-6 - Zona Virgen Compatible</span><br>
                        <span style='font-size:11px; color:#1A4A7A;'>
                            Sin competencia directa, con perfil demografico y vial similar
                            al entorno del mejor competidor del radio. Alto potencial si validas demanda en campo.</span>
                    </div>
                </div>
                </div>""", unsafe_allow_html=True)

                if _top3:
                    st.markdown("### TOP 3 - MERCADO PROBADO")
                    st.caption("Zonas donde ya existe demanda comprobada con poca competencia")

                for i, res in enumerate(_top3):
                    _sc = res["score"]
                    _ch = "#27AE60" if _sc>=65 else "#F39C12" if _sc>=40 else "#E74C3C"
                    _recom = res.get("recom_ubicacion", {})
                    _best_comp_n = res.get("mejor_competidor","")
                    _best_rat = res.get("mejor_rating", 0)
                    _best_rev = res.get("mejor_reviews", 0)
                    _cat_label = {"mercado_probado":"Mercado probado","compartir_mercado":"Mercado compartible"}.get(res.get("zona_categoria",""), "")
                    _ticket_badge = f"<span style='font-size:11px; background:#FFF3D6; color:#8B6914; padding:3px 9px; border-radius:20px;'>Ticket ${res.get('ticket_zona',0):,} MXN</span>" if res.get("ticket_zona") else ""
                    _via_badge = f"<span style='font-size:11px; background:#F0F0F0; padding:3px 9px; border-radius:20px;'>{res.get('vialidad','')}</span>" if res.get("vialidad") else ""
                    _comp_ref_html = f"<div style='margin-top:10px; padding:8px 12px; background:#FFFDE7; border-radius:6px; font-size:11px; color:#5D4037;'><b>Competidor referencia:</b> {_best_comp_n[:40]} — {_best_rat} estrellas, {_best_rev} resenas</div>" if _best_comp_n else ""
                    # Escapar dirección para evitar que comillas rompan el HTML
                    _dir_safe = res["direccion"].replace("'","&#39;").replace('"','&quot;').replace('<','&lt;').replace('>','&gt;')

                    st.markdown(f"""
                    <div style='border:1px solid #C8E6C9; border-left:5px solid {_ch};
                         border-radius:10px; padding:16px 18px; margin:6px 0 2px; background:#FAFFF8;'>
                      <div style='display:flex; justify-content:space-between; align-items:flex-start;'>
                        <div style='flex:1;'>
                          <div style='font-size:14px; font-weight:800; color:#1A1A1A; margin-bottom:4px;'>
                            #{i+1} — {_dir_safe}</div>
                          <span style='font-size:11px; background:#27AE6022; color:#1A7A44;
                              padding:3px 10px; border-radius:20px; font-weight:600;'>{_cat_label}</span>
                          {_ticket_badge}
                        </div>
                        <div style='text-align:right; min-width:70px;'>
                          <div style='font-size:28px; font-weight:900; color:{_ch}; line-height:1;'>{_sc:.0f}</div>
                          <div style='font-size:10px; color:#888;'>/100</div>
                        </div>
                      </div>
                      <div style='display:flex; gap:10px; margin-top:10px; flex-wrap:wrap;'>
                        <span style='font-size:11px; background:#F0F0F0; padding:3px 9px; border-radius:20px;'>NSE <b>{res.get("nse","—")}</b></span>
                        <span style='font-size:11px; background:#F0F0F0; padding:3px 9px; border-radius:20px;'>{res.get("densidad",0):,} hab/km2</span>
                        <span style='font-size:11px; background:#F0F0F0; padding:3px 9px; border-radius:20px;'>{res.get("n_competidores",0)} competidores</span>
                        {_via_badge}
                      </div>
                      {_comp_ref_html}
                    </div>""", unsafe_allow_html=True)

                    if _recom:
                        _key_exp_r = f"exp_recom_top3_{i}"
                        if _key_exp_r not in st.session_state: st.session_state[_key_exp_r] = (i==0)
                        _arrow_r = " ▲" if st.session_state[_key_exp_r] else " ▼"
                        if st.button(f"Recomendaciones de ubicacion ideal - Zona #{i+1}" + _arrow_r, key=f"btn_recom_t3_{i}"):
                            st.session_state[_key_exp_r] = not st.session_state[_key_exp_r]
                        if st.session_state[_key_exp_r]:
                            _c1, _c2 = st.columns(2)
                            with _c1:
                                if _recom.get("vialidad_ideal"):
                                    st.markdown("**Tipo de vialidad ideal**")
                                    st.info(_recom["vialidad_ideal"])
                                if _recom.get("visibilidad"):
                                    st.markdown("**Visibilidad requerida**")
                                    st.info(_recom["visibilidad"])
                                if _recom.get("estacionamiento"):
                                    st.markdown("**Estacionamiento**")
                                    st.info(_recom["estacionamiento"])
                            with _c2:
                                if _recom.get("plaza_recomendacion"):
                                    st.markdown("**En plaza comercial**")
                                    st.info(_recom["plaza_recomendacion"])
                                if _recom.get("ref_competidor"):
                                    st.markdown("**Sobre el competidor de referencia**")
                                    st.info(_recom["ref_competidor"])
                                if _recom.get("tip_nse"):
                                    st.markdown("**Tip para este NSE**")
                                    st.info(_recom["tip_nse"])

                if _top456:
                    st.markdown("")
                    st.markdown("### TOP 4-6 - ZONA VIRGEN COMPATIBLE")
                    st.caption("Sin competencia directa, perfil demografico compatible con el mercado probado")

                for i, res in enumerate(_top456):
                    _sc = res["score"]
                    _ch = "#1A76D2" if _sc>=50 else "#F39C12" if _sc>=35 else "#E74C3C"
                    _recom = res.get("recom_ubicacion", {})
                    _pos_global = i + len(_top3) + 1
                    _ticket_badge2 = f"<span style='font-size:11px; background:#FFF3D6; color:#8B6914; padding:3px 9px; border-radius:20px;'>Ticket ${res.get('ticket_zona',0):,} MXN</span>" if res.get("ticket_zona") else ""
                    _via_badge2 = f"<span style='font-size:11px; background:#F0F0F0; padding:3px 9px; border-radius:20px;'>{res.get('vialidad','')}</span>" if res.get("vialidad") else ""
                    # Escapar dirección para evitar que comillas rompan el HTML
                    _dir_safe2 = res["direccion"].replace("'","&#39;").replace('"','&quot;').replace('<','&lt;').replace('>','&gt;')

                    st.markdown(f"""
                    <div style='border:1px solid #BBDEFB; border-left:5px solid {_ch};
                         border-radius:10px; padding:16px 18px; margin:6px 0 2px; background:#F8FBFF;'>
                      <div style='display:flex; justify-content:space-between; align-items:flex-start;'>
                        <div style='flex:1;'>
                          <div style='font-size:14px; font-weight:800; color:#1A1A1A; margin-bottom:4px;'>
                            #{_pos_global} — {_dir_safe2}</div>
                          <span style='font-size:11px; background:#1A76D222; color:#1A4A7A;
                              padding:3px 10px; border-radius:20px; font-weight:600;'>Zona virgen</span>
                          {_ticket_badge2}
                        </div>
                        <div style='text-align:right; min-width:70px;'>
                          <div style='font-size:28px; font-weight:900; color:{_ch}; line-height:1;'>{_sc:.0f}</div>
                          <div style='font-size:10px; color:#888;'>/100</div>
                        </div>
                      </div>
                      <div style='display:flex; gap:10px; margin-top:10px; flex-wrap:wrap;'>
                        <span style='font-size:11px; background:#F0F0F0; padding:3px 9px; border-radius:20px;'>NSE <b>{res.get("nse","—")}</b></span>
                        <span style='font-size:11px; background:#F0F0F0; padding:3px 9px; border-radius:20px;'>{res.get("densidad",0):,} hab/km2</span>
                        <span style='font-size:11px; background:#F0F0F0; padding:3px 9px; border-radius:20px;'>{res.get("n_competidores",0)} competidores</span>
                        {_via_badge2}
                      </div>
                      <div style='margin-top:8px; font-size:11px; color:#555; font-style:italic;'>
                        Sin competencia directa detectada — valida demanda en campo antes de invertir.</div>
                    </div>""", unsafe_allow_html=True)

                    if _recom:
                        _key_exp_v = f"exp_recom_vir_{i}"
                        if _key_exp_v not in st.session_state: st.session_state[_key_exp_v] = False
                        _arrow_v = " ▲" if st.session_state[_key_exp_v] else " ▼"
                        if st.button(f"Recomendaciones de ubicacion ideal - Zona #{_pos_global}" + _arrow_v, key=f"btn_recom_v_{i}"):
                            st.session_state[_key_exp_v] = not st.session_state[_key_exp_v]
                        if st.session_state[_key_exp_v]:
                            _c1, _c2 = st.columns(2)
                            with _c1:
                                if _recom.get("vialidad_ideal"):
                                    st.markdown("**Tipo de vialidad ideal**")
                                    st.info(_recom["vialidad_ideal"])
                                if _recom.get("visibilidad"):
                                    st.markdown("**Visibilidad requerida**")
                                    st.info(_recom["visibilidad"])
                                if _recom.get("estacionamiento"):
                                    st.markdown("**Estacionamiento**")
                                    st.info(_recom["estacionamiento"])
                            with _c2:
                                if _recom.get("plaza_recomendacion"):
                                    st.markdown("**En plaza comercial**")
                                    st.info(_recom["plaza_recomendacion"])
                                if _recom.get("ref_competidor"):
                                    st.markdown("**Competidor de referencia**")
                                    st.info(_recom["ref_competidor"])
                                if _recom.get("tip_nse"):
                                    st.markdown("**Tip para este NSE**")
                                    st.info(_recom["tip_nse"])

                st.markdown("")

                # Mapa con los resultados si MAPAS_OK
                if MAPAS_OK and len(resultados) > 0:
                    try:
                        import folium
                        m_zona = folium.Map(location=[lat_r, lng_r], zoom_start=13, tiles="CartoDB positron")
                        for _ri, res in enumerate(resultados[:8]):
                            _sc2 = res["score"]
                            _col2 = "green" if _sc2>=65 else "orange" if _sc2>=40 else "red"
                            _rad2 = 16 if _ri == 0 else 13
                            _tk_str = f"<br>🎫 Ticket zona: <b>${res.get('ticket_zona',0):,} MXN</b>" if res.get('ticket_zona') else ""
                            _vl_str = f"<br>🛣️ {res.get('vialidad','')}" if res.get('vialidad') else ""
                            _dsg2   = res.get("desglose", {})
                            _dsg_str = ""
                            if _dsg2.get("densidad"):
                                _dsg_str = (f"<br><small>Competencia: {_dsg2.get('densidad',0)}/100 · "
                                            f"Calidad rival: {_dsg2.get('calidad',0)}/100</small>")
                            _popup_html = (
                                f"<div style='font-family:sans-serif; min-width:200px;'>"
                                f"<b style='font-size:15px; color:{'#2ECC71' if _sc2>=65 else '#F39C12' if _sc2>=40 else '#E74C3C'};'>"
                                f"#{_ri+1} Score: {_sc2:.0f}/100</b><br>"
                                f"<span style='font-size:11px; color:#555;'>{res['direccion'][:60]}</span><br>"
                                f"<br>🏘️ NSE: <b>{res['nse']}</b>"
                                f"<br>👥 Densidad: <b>{res['densidad']:,} hab/km²</b>"
                                f"<br>🏪 Competidores: <b>{res['n_competidores']}</b>"
                                f"{_tk_str}{_vl_str}{_dsg_str}"
                                f"</div>"
                            )
                            folium.CircleMarker(
                                location=[res["lat"], res["lng"]],
                                radius=_rad2, color=_col2, fill=True,
                                fill_color=_col2, fill_opacity=0.75,
                                tooltip=f"#{_ri+1} Score {_sc2:.0f} — {res['direccion'][:35]}",
                                popup=folium.Popup(_popup_html, max_width=260)
                            ).add_to(m_zona)
                        folium.Marker([lat_r, lng_r],
                            tooltip="📍 Tu punto de referencia",
                            icon=folium.Icon(color="blue", icon="star", prefix="fa")
                        ).add_to(m_zona)
                        from mapas_4site import render_mapa_con_interpretacion
                        render_mapa_con_interpretacion(m_zona,
                            f"Top {min(5,len(resultados))} zonas encontradas para {TIPOS_NEGOCIO[tipo_busq]['nombre']}",
                            altura=420)
                    except Exception:
                        pass

                # ── RESUMEN COMPARATIVO ──────────────────────────────────
                st.markdown("---")
                st.markdown("### 📊 Resumen Comparativo")

                # Top zona vs contexto general
                _top = resultados[0]
                _avg_score = sum(r["score"] for r in resultados) / len(resultados)
                _top_nse   = _top.get("nse", "—")
                _top_comp  = _top.get("n_competidores", 0)
                _top_tk    = _top.get("ticket_zona", 0)
                _tipo_nombre = TIPOS_NEGOCIO[tipo_busq]["nombre"]

                # Factores de la mejor zona
                _factores_favor = []
                _factores_riesgo = []

                if _top["score"] >= 70:
                    _factores_favor.append(f"Score de viabilidad **{_top['score']:.0f}/100** — zona con buenas condiciones generales")
                if _top_nse in ["A", "A/B", "B", "B/C+", "C+"]:
                    _factores_favor.append(f"NSE **{_top_nse}** — poder adquisitivo adecuado para {_tipo_nombre}")
                if _top_comp >= 1:
                    _factores_favor.append(f"**{_top_comp} competidores** en el área confirman que hay mercado activo para este giro")
                elif _top_comp == 0 and tipo_busq not in ["acabados_hogar","ferreteria","gimnasio_boutique"]:
                    _factores_favor.append("**Sin competencia directa** — oportunidad de ser el primero en la zona")
                if _top.get("densidad", 0) > 8000:
                    _factores_favor.append(f"**Alta densidad poblacional** ({_top['densidad']:,} hab/km²) — mercado potencial amplio")
                if _top.get("vialidad"):
                    _factores_favor.append(f"Vialidad: **{_top['vialidad']}** — accesibilidad y visibilidad")

                if _top_comp == 0 and tipo_busq in ["acabados_hogar","ferreteria","gimnasio_boutique"]:
                    _factores_riesgo.append("Sin competidores detectados — valida que haya demanda real antes de invertir")
                if _top.get("densidad", 0) < 5000:
                    _factores_riesgo.append("Densidad poblacional baja — mercado potencial limitado")
                if _top_nse in ["D+", "D/E"]:
                    _factores_riesgo.append(f"NSE {_top_nse} — puede limitar el ticket promedio del negocio")

                # Mostrar resumen
                _col_r1, _col_r2 = st.columns(2)
                with _col_r1:
                    st.markdown("**✅ Factores favorables en la mejor zona:**")
                    for _f in _factores_favor[:4]:
                        st.markdown(f"- {_f}")
                with _col_r2:
                    st.markdown("**⚠️ Factores de riesgo a considerar:**")
                    if _factores_riesgo:
                        for _r in _factores_riesgo[:3]:
                            st.markdown(f"- {_r}")
                    else:
                        st.markdown("- No se detectaron riesgos críticos en la zona líder")

                # Comparativa numérica de zonas
                st.markdown("**📊 Comparativa de las mejores zonas encontradas:**")
                _tabla_comp = []
                for _idx, _rz in enumerate(resultados[:5]):
                    _nivel, _ = nivel_score(_rz["score"], idioma)
                    _tabla_comp.append({
                        "Posición": f"#{_idx+1}",
                        "Score": f"{_rz['score']:.0f}/100",
                        "Nivel": _nivel,
                        "NSE": _rz.get("nse","—"),
                        "Competidores": _rz.get("n_competidores",0),
                        "Densidad": f"{_rz.get('densidad',0):,}",
                        "Ticket zona": f"${_rz.get('ticket_zona',0):,}" if _rz.get('ticket_zona') else "N/D",
                        "Dirección": _rz["direccion"][:35],
                    })
                st.dataframe(_tabla_comp, use_container_width=True, hide_index=True)

                # ── MENSAJE SIEMPRE PRESENTE: apoyo del negocio ──────────
                st.markdown("---")
                st.markdown("""
                <div style='background:#0D0D0D; border-radius:12px; padding:20px 24px; color:white;'>
                <p style='font-size:13px; font-weight:700; color:#BB944F; margin:0 0 10px;'>
                    💡 La ubicación es el punto de partida — el éxito depende también de:</p>
                <div style='display:grid; grid-template-columns:1fr 1fr 1fr; gap:12px; margin-top:8px;'>
                    <div style='background:#1A1A1A; border-radius:8px; padding:12px;'>
                        <div style='color:#BB944F; font-weight:700; font-size:12px; margin-bottom:4px;'>🎫 Ticket promedio</div>
                        <div style='color:#CCCCCC; font-size:11px;'>Ajusta tu precio al nivel socioeconómico y competencia de la zona. Un ticket mal calibrado puede cancelar una buena ubicación.</div>
                    </div>
                    <div style='background:#1A1A1A; border-radius:8px; padding:12px;'>
                        <div style='color:#BB944F; font-weight:700; font-size:12px; margin-bottom:4px;'>🤝 Atención al cliente</div>
                        <div style='color:#CCCCCC; font-size:11px;'>En zonas con competencia cercana, la experiencia del cliente es el diferenciador principal. El rating en Google impacta directamente en la captación.</div>
                    </div>
                    <div style='background:#1A1A1A; border-radius:8px; padding:12px;'>
                        <div style='color:#BB944F; font-weight:700; font-size:12px; margin-bottom:4px;'>📱 Presencia en redes</div>
                        <div style='color:#CCCCCC; font-size:11px;'>Una buena ubicación sin visibilidad digital pierde clientes. Google Maps, Instagram y WhatsApp Business son esenciales para convertir tráfico de paso en clientes recurrentes.</div>
                    </div>
                </div>
                </div>""", unsafe_allow_html=True)

            # ── DESCARGA PPTX ──────────────────────────────────────
                st.markdown("---")
                st.markdown("#### 📥 Descargar Reporte de Zona Potencial")
                _col_dl1, _col_dl2 = st.columns(2)
                with _col_dl1:
                    try:
                        _pptx_buf = generar_pptx_zona_potencial(
                            resultados=resultados,
                            ubicacion_ref=dir_ref,
                            radio_km=radio_km,
                            tipo_negocio_nombre=TIPOS_NEGOCIO.get(tipo_busq,{}).get("nombre",""),
                            tier_key=tier_key,
                        )
                        if _pptx_buf:
                            _fn_pptx = f"4site_zona_{dir_ref[:20].replace(' ','_').replace(',','')}.pptx"
                            st.download_button(
                                label="⬇️  Descargar PowerPoint (.pptx)",
                                data=_pptx_buf,
                                file_name=_fn_pptx,
                                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                                use_container_width=True,
                                type="primary",
                                key="btn_dl_pptx_zona"
                            )
                            st.caption("Presentación editable · 4 slides · Identidad 4SITE")
                    except Exception as _ep:
                        st.error(f"Error al generar PPTX: {_ep}")

            except Exception as ex:
                st.error(f"Error en búsqueda: {ex}")
                import traceback
                with st.expander("Detalle del error"):
                    st.code(traceback.format_exc())


def _seccion_pdf(titulo, icono, s, color=None):
    """Banda de color con título de sección — consistente en todos los PDFs."""
    _color = color or s.get("_N", colors.HexColor("#0D0D0D"))
    return [
        Spacer(1, 0.12 * inch),
        Table(
            [[Paragraph(f"{icono}  {titulo}", s["h2"])]],
            colWidths=[522],
            style=TableStyle([
                ("BACKGROUND",    (0,0),(-1,-1), _color),
                ("TOPPADDING",    (0,0),(-1,-1), 8),
                ("BOTTOMPADDING", (0,0),(-1,-1), 8),
                ("LEFTPADDING",   (0,0),(-1,-1), 14),
                ("RIGHTPADDING",  (0,0),(-1,-1), 14),
            ]),
        ),
        Spacer(1, 0.07 * inch),
    ]


def _kpi_pdf(kpis, s, col_w=None):
    """Fila de tarjetas KPI para el PDF."""
    n = len(kpis)
    if col_w is None:
        col_w = 522 / n
    row_val, row_lbl = [], []
    tstyle = [
        ("ALIGN",         (0,0),(-1,-1), "CENTER"),
        ("VALIGN",        (0,0),(-1,-1), "MIDDLE"),
        ("TOPPADDING",    (0,0),(-1,-1), 10),
        ("BOTTOMPADDING", (0,0),(-1,-1), 10),
        ("LEFTPADDING",   (0,0),(-1,-1), 6),
        ("RIGHTPADDING",  (0,0),(-1,-1), 6),
        ("GRID",          (0,0),(-1,-1), 0.4, colors.HexColor("#DDDDDD")),
    ]
    for i, k in enumerate(kpis):
        est = s["kpi_val_gold"] if k.get("gold") else s["kpi_val"]
        row_val.append(Paragraph(k["valor"], est))
        row_lbl.append(Paragraph(k["label"], s["kpi_label"]))
        bg = colors.HexColor("#FFF9F0") if k.get("gold") else colors.HexColor("#F5F5F5")
        tstyle.append(("BACKGROUND", (i,0), (i,1), bg))
    t = Table([row_val, row_lbl], colWidths=[col_w]*n)
    t.setStyle(TableStyle(tstyle))
    return t


def _tabla_pdf(filas, col_widths, s, header_color=None):
    """Tabla zebra con encabezado de color."""
    hc = header_color or s.get("_N", colors.HexColor("#0D0D0D"))
    t = Table(filas, colWidths=col_widths)
    ts = [
        ("BACKGROUND",    (0,0),(-1,0), hc),
        ("TEXTCOLOR",     (0,0),(-1,0), colors.white),
        ("FONTNAME",      (0,0),(-1,0), s.get("_FB","Helvetica-Bold")),
        ("FONTSIZE",      (0,0),(-1,0), 8.5),
        ("ALIGN",         (0,0),(-1,0), "CENTER"),
        ("TOPPADDING",    (0,0),(-1,-1), 5),
        ("BOTTOMPADDING", (0,0),(-1,-1), 5),
        ("LEFTPADDING",   (0,0),(-1,-1), 7),
        ("RIGHTPADDING",  (0,0),(-1,-1), 7),
        ("FONTNAME",      (0,1),(-1,-1), s.get("_FR","Helvetica")),
        ("FONTSIZE",      (0,1),(-1,-1), 8.5),
        ("GRID",          (0,0),(-1,-1), 0.4, colors.HexColor("#DDDDDD")),
    ]
    for i in range(1, len(filas)):
        if i % 2 == 0:
            ts.append(("BACKGROUND",(0,i),(-1,i), colors.HexColor("#F9F9F9")))
    t.setStyle(TableStyle(ts))
    return t


def _footer_pdf(canvas_obj, doc_obj):
    """Footer consistente en todas las páginas."""
    canvas_obj.saveState()
    canvas_obj.setFillColor(colors.HexColor("#AAAAAA"))
    canvas_obj.setFont("Helvetica", 7)
    canvas_obj.drawCentredString(306, 22,
        f"4SITE · Don't guess. Foresee. · 4site.mx  |  Página {doc_obj.page}")
    canvas_obj.setStrokeColor(colors.HexColor("#BB944F"))
    canvas_obj.setLineWidth(0.8)
    canvas_obj.line(45, 32, 567, 32)
    canvas_obj.restoreState()


def generar_pdf_free(ubicacion, score, desglose, analisis, competidores,
                      idioma, lat, lng, modo, tipo_negocio=None, recomendaciones=None,
                      demografia=None, contexto=None):
    buffer = BytesIO()
    doc    = SimpleDocTemplate(buffer, pagesize=letter,
                               leftMargin=50, rightMargin=50, topMargin=50, bottomMargin=50,
                               title=f"4SITE — Reporte: {ubicacion[:60]}",
                               author="4SITE · Don't guess. Foresee.",
                               subject="Análisis de viabilidad comercial · 4site.mx",
                               creator="4SITE App", producer="4SITE")
    s     = _estilos_pdf()
    t     = TEXTOS[idioma]
    story = []
    tier  = "free"

    _header_pdf(story, s, tier, idioma)
    story.append(Paragraph(f"<b>{t['input_ubicacion']}:</b> {ubicacion}", s["normal"]))
    if modo == "validar" and tipo_negocio:
        story.append(Paragraph(f"<b>Tipo:</b> {TIPOS_NEGOCIO[tipo_negocio]['nombre']}", s["normal"]))
    story.append(Spacer(1, 0.2*inch))

    _score_block(story, score, desglose, tier, t, s, idioma)

    # Mapa (sin competidores)
    mapa = generar_mapa_estatico(lat, lng, competidores=None, con_competidores=False)
    if mapa:
        story.extend(_seccion_pdf("Ubicación", "📍", s))
        story.append(RLImage(mapa, width=4.5*inch, height=3*inch))
        story.append(Spacer(1, 0.1*inch))

    # Análisis vial básico en Free
    av_f = contexto.get("analisis_vial") if contexto else None
    if av_f:
        ajuste_f = av_f.get("ajuste_score", 0)
        color_f = '#4CAF50' if ajuste_f > 0 else '#F44336' if ajuste_f < -5 else '#FF9800'
        story.append(Paragraph(
            f"🛣️ Vialidad: {av_f.get('badge','')}  |  "
            f"Impacto: <b><font color='{color_f}'>{'+'if ajuste_f>=0 else ''}{ajuste_f} pts</font></b>",
            ParagraphStyle('vial_f', parent=s["base"]['Normal'], fontSize=9,
                           textColor=colors.HexColor('#333333'), spaceAfter=4)))
        story.append(Paragraph(
            "🔒 Análisis vial completo con narrativa interpretativa disponible en plan Básico ($99)",
            ParagraphStyle('vial_lock', parent=s["base"]['Normal'], fontSize=8,
                           textColor=colors.HexColor('#888888'), fontName='Helvetica-Oblique')))
        story.append(Spacer(1, 0.1*inch))

    # Competidores — solo 3 filas
    if competidores:
        story.append(Paragraph(t["competencia_titulo"], s["h2"]))
        story.append(Paragraph(f"Mostrando 3 de {len(competidores)} — Actualiza para ver todos", s["small"]))
        story.append(Spacer(1, 0.1*inch))
        story.append(_tabla_competidores(competidores, t, 3, s))
        story.append(Spacer(1, 0.1*inch))

    # Demografía básica (free): solo población, viviendas y densidad
    if demografia:
        densidad_info = formatear_densidad(demografia['densidad_hab_km2'])
        story.append(Paragraph("📍 Datos del Área (500m)" if idioma == "es" else "📍 Area Data (500m)", s["h2"]))
        dem_data = [
            ["Población estimada" if idioma == "es" else "Estimated population",
             f"{demografia['poblacion_estimada']:,} hab"],
            ["Viviendas habitadas" if idioma == "es" else "Occupied housing",
             f"{demografia['viviendas_habitadas']:,}"],
            ["Densidad" if idioma == "es" else "Density",
             densidad_info['clasificacion']],
            ["Aprox. por manzana" if idioma == "es" else "Approx. per block",
             f"~{densidad_info['personas_manzana']} personas"],
        ]
        dem_tbl = Table(dem_data, colWidths=[2.8*inch, 2.2*inch])
        dem_tbl.setStyle(TableStyle([
            ('FONTSIZE',   (0,0), (-1,-1), 9),
            ('FONTNAME',   (0,0), (0,-1), 'Helvetica-Bold'),
            ('TEXTCOLOR',  (0,0), (0,-1), colors.HexColor('#6B6B6B')),
            ('ROWBACKGROUNDS', (0,0), (-1,-1), [colors.white, colors.HexColor('#F5F8FF')]),
            ('GRID',       (0,0), (-1,-1), 0.3, colors.HexColor('#DDDDDD')),
            ('TOPPADDING', (0,0), (-1,-1), 5),
            ('BOTTOMPADDING', (0,0), (-1,-1), 5),
        ]))
        story.append(dem_tbl)
        story.append(Paragraph(densidad_info['descripcion'], s["small"]))
        story.append(Paragraph(
            "🔒 NSE, ingreso promedio y datos INEGI reales disponibles en plan Básico ($99)" if idioma == "es"
            else "🔒 NSE, average income and real INEGI data available in Basic plan ($99)",
            ParagraphStyle('lock', parent=s["base"]['Normal'], fontSize=8,
                           textColor=colors.HexColor('#888888'), fontName='Helvetica-Oblique', spaceAfter=6)
        ))
        story.append(Spacer(1, 0.1*inch))

    # PageBreak → Resumen ejecutivo simple + CTA
    story.append(PageBreak())
    story.extend(_seccion_pdf("Resumen Ejecutivo", "📋", s))
    lineas = analisis.split('\n')
    for linea in lineas:
        linea = linea.strip()
        if not linea:
            story.append(Spacer(1, 0.04*inch))
        elif linea.startswith('**') and linea.endswith('**'):
            story.append(Paragraph(linea.replace('**',''), ParagraphStyle('re_h',
                parent=s["base"]['Normal'], fontSize=10, fontName='Helvetica-Bold',
                textColor=colors.HexColor('#6B6B6B'), spaceBefore=8, spaceAfter=3)))
        elif linea.startswith('- ') or linea.startswith('• '):
            story.append(Paragraph(f"  {linea}", ParagraphStyle('re_bl',
                parent=s["base"]['Normal'], fontSize=9, leftIndent=12, spaceAfter=2)))
        elif '**' in linea:
            txt = linea
            for _ in range(linea.count('**') // 2):
                txt = txt.replace('**', '<b>', 1).replace('**', '</b>', 1)
            story.append(Paragraph(txt, s["normal"]))
        else:
            story.append(Paragraph(linea, s["normal"]))
    story.append(Spacer(1, 0.1*inch))
    story.append(Paragraph(
        "[ Para análisis completo con Recomendación Final, tráfico y forecast actualiza al plan Básico ($99) ]",
        ParagraphStyle('upsell_f', parent=s["base"]['Normal'], fontSize=8.5,
            textColor=colors.HexColor('#BB944F'), fontName='Helvetica-Oblique')))

    _cta_block(story, tier, t, s, idioma)
    doc.build(story, onFirstPage=_footer_pdf, onLaterPages=_footer_pdf)
    buffer.seek(0)
    return buffer


# ─────────────────────────────────────────────
# PDF BASIC  (Análisis completo)
# ─────────────────────────────────────────────
def generar_pdf_basic(ubicacion, score, desglose, analisis, competidores,
                       idioma, lat, lng, modo, tipo_negocio=None,
                       recomendaciones=None, demografia=None, contexto=None):
    buffer = BytesIO()
    doc    = SimpleDocTemplate(buffer, pagesize=letter,
                               leftMargin=50, rightMargin=50, topMargin=50, bottomMargin=50,
                               title=f"4SITE — Reporte: {ubicacion[:60]}",
                               author="4SITE · Don't guess. Foresee.",
                               subject="Análisis de viabilidad comercial · 4site.mx",
                               creator="4SITE App", producer="4SITE")
    s     = _estilos_pdf()
    t     = TEXTOS[idioma]
    story = []
    tier  = "basic"

    # ── P1: Portada ──
    _header_pdf(story, s, tier, idioma)
    story.append(Paragraph(f"<b>{t['input_ubicacion']}:</b> {ubicacion}", s["normal"]))
    if modo == "validar" and tipo_negocio:
        tipo_info = TIPOS_NEGOCIO[tipo_negocio]
        story.append(Paragraph(f"<b>Tipo:</b> {tipo_info['nombre']}", s["normal"]))
        story.append(Paragraph(f"<b>Inversión estimada:</b> ${tipo_info['inversion_min']:,}–${tipo_info['inversion_max']:,} MXN", s["normal"]))
    story.append(Spacer(1, 0.2*inch))
    _score_block(story, score, desglose, tier, t, s, idioma)

    # Badges
    if contexto and contexto.get("badges"):
        # Filtrar badges de vialidad — ya se muestran en el bloque de análisis vial
        badges_vial = {"🛣️", "🚦", "🏘️", "📍"}
        badges_filtrados = [b for b in contexto["badges"]
                           if not any(b.startswith(v) for v in badges_vial)]
        for badge in badges_filtrados:
            story.append(Paragraph(f"  {badge}", s["normal"]))
        story.append(Spacer(1, 0.1*inch))

    # ── P2: Mapa + Análisis Vial ──
    story.append(PageBreak())
    story.extend(_seccion_pdf("Mapa de Ubicación", "📍", s))
    mapa = generar_mapa_estatico(lat, lng, competidores=None, con_competidores=False)
    if mapa:
        story.append(RLImage(mapa, width=5*inch, height=3.3*inch))
    story.append(Paragraph("Radio de análisis: 500m", s["small"]))

    # Tabla análisis vial con narrativa
    av_b = contexto.get("analisis_vial") if contexto else None
    if av_b and av_b.get("descripcion"):
        story.append(Spacer(1, 0.15*inch))
        story.extend(_seccion_pdf("Análisis de Vialidad", "🛣️", s))
        ajuste_b = av_b.get("ajuste_score", 0)
        color_b  = '#4CAF50' if ajuste_b > 0 else '#F44336' if ajuste_b < -5 else '#FF9800'
        dep_b = {"alto":"Negocio de paso (crítico)","medio":"Negocio mixto","bajo":"Negocio de destino"}.get(
            av_b.get("dependencia_paso","medio"),"")
        vial_rows_b = [
            ["Tipo de vialidad:",  av_b.get("badge","")],
            ["Dependencia al tráfico:", dep_b],
            ["Impacto en score:", f"{'+'if ajuste_b>=0 else ''}{ajuste_b} puntos"],
            ["Análisis:", av_b.get("descripcion","")],
        ]
        tbl_b = Table(vial_rows_b, colWidths=[1.8*inch, 3.7*inch])
        tbl_b.setStyle(TableStyle([
            ('FONTNAME',(0,0),(0,-1),'Helvetica-Bold'),
            ('FONTSIZE',(0,0),(-1,-1),9),
            ('TEXTCOLOR',(0,0),(0,-1),colors.HexColor('#6B6B6B')),
            ('TEXTCOLOR',(1,2),(1,2),colors.HexColor(color_b)),
            ('FONTNAME',(1,2),(1,2),'Helvetica-Bold'),
            ('ROWBACKGROUNDS',(0,0),(-1,-1),[colors.HexColor('#F0F4FF'),colors.white,
                                              colors.HexColor('#F0F4FF'),colors.white]),
            ('GRID',(0,0),(-1,-1),0.3,colors.HexColor('#DDDDDD')),
            ('TOPPADDING',(0,0),(-1,-1),6),('BOTTOMPADDING',(0,0),(-1,-1),6),
            ('VALIGN',(0,0),(-1,-1),'TOP'),
        ]))
        story.append(tbl_b)
        # Narrativa interpretativa vial
        tipo_nom_b = TIPOS_NEGOCIO.get(tipo_negocio, {}).get('nombre','') if tipo_negocio else (
            recomendaciones[0]['nombre'] if recomendaciones else '')
        narr_vial_b = generar_narrativa_seccion("mapa", {
            "tipo_vialidad": av_b.get("tipo_vialidad",""),
            "ajuste_score": ajuste_b,
            "descripcion": av_b.get("descripcion",""),
            "dependencia_paso": av_b.get("dependencia_paso","medio"),
        }, tipo_nom_b, idioma)
        for el in _bloque_narrativa_pdf(narr_vial_b, s):
            story.append(el)

    story.append(Spacer(1, 0.2*inch))

    # ── P3: Competidores (10 filas) — movido, Recomendación va al final ──
    story.append(PageBreak())
    if competidores:
        story.append(Paragraph(t["competencia_titulo"], s["h2"]))
        story.append(Paragraph(f"{len(competidores)} negocios similares encontrados en 500m", s["small"]))
        story.append(Spacer(1, 0.1*inch))
        story.append(_tabla_competidores(competidores, t, 10, s))
        # Narrativa interpretativa de competencia (Basic)
        ratings_b  = [c.get('rating',0) for c in (competidores or []) if c.get('rating')]
        avg_r_b    = round(sum(ratings_b)/len(ratings_b), 1) if ratings_b else 0
        n_verde_b  = sum(1 for r in ratings_b if r >= 4.3)
        tipo_nom_b = TIPOS_NEGOCIO.get(tipo_negocio, {}).get('nombre','') if tipo_negocio else (
            recomendaciones[0]['nombre'] if recomendaciones else '')
        narr_b = generar_narrativa_seccion("competencia", {
            "total_competidores": len(competidores),
            "rating_promedio": avg_r_b,
            "competidores_excelentes": n_verde_b,
        }, tipo_nom_b, idioma)
        for el in _bloque_narrativa_pdf(narr_b, s):
            story.append(el)

    # ── P5: Demografía estimada ──
    story.append(PageBreak())
    story.append(Paragraph(t["demografía"], s["h2"]))
    story.append(Paragraph(t["datos_estimados"], s["small"]))
    story.append(Spacer(1, 0.1*inch))
    if demografia:
        densidad_info = formatear_densidad(demografia['densidad_hab_km2'])
        dem_data = [
            [t["poblacion"],    f"{demografia['poblacion_estimada']:,} hab"],
            [t["nse"],          demografia['nse_predominante']],
            [t["densidad_hab"], densidad_info['clasificacion']],
            ["Viviendas",       f"{demografia['viviendas_habitadas']:,}"],
            ["Ingreso prom. mensual", f"${demografia['ingreso_promedio_mensual']:,}"],
        ]
        dem_tbl = Table(dem_data, colWidths=[2.5*inch, 2.5*inch])
        dem_tbl.setStyle(TableStyle([
            ('FONTSIZE',  (0,0), (-1,-1), 10),
            ('FONTNAME',  (0,0), (0,-1), 'Helvetica-Bold'),
            ('TEXTCOLOR', (0,0), (0,-1), colors.HexColor('#6B6B6B')),
            ('ROWBACKGROUNDS', (0,0), (-1,-1), [colors.white, colors.HexColor('#F5F8FF')]),
            ('GRID', (0,0), (-1,-1), 0.3, colors.HexColor('#DDDDDD')),
            ('TOPPADDING', (0,0), (-1,-1), 6),
            ('BOTTOMPADDING', (0,0), (-1,-1), 6),
        ]))
        story.append(dem_tbl)
        story.append(Spacer(1, 0.1*inch))
        story.append(Paragraph(densidad_info['descripcion'], s["small"]))

    # Nota PRO
    story.append(Spacer(1, 0.2*inch))
    story.append(Paragraph(
        "🔜 PRO: Datos INEGI reales + tamaño de mercado + forecast de ventas + mapa con competidores",
        ParagraphStyle('pro_nota', parent=s["base"]['Normal'], fontSize=9,
                       textColor=colors.HexColor('#BB944F'), fontName='Helvetica-Oblique')
    ))

    # ── P6: Recomendación Final ──
    story.append(PageBreak())
    _plaza_block_pdf(story, contexto, s)
    _render_recomendacion_pdf(story, analisis, s)

    # ── P7: CTA ──
    story.append(PageBreak())
    _cta_block(story, tier, t, s, idioma)

    doc.build(story, onFirstPage=_footer_pdf, onLaterPages=_footer_pdf)
    buffer.seek(0)
    return buffer


# ─────────────────────────────────────────────
# PDF PRO  (Análisis PRO — con gráficas reales)
# ─────────────────────────────────────────────
def generar_pdf_pro(ubicacion, score, desglose, analisis, competidores,
                     idioma, lat, lng, modo, tipo_negocio=None,
                     recomendaciones=None, demografia=None, contexto=None):
    """PRO: portada + mapa + análisis + competidores + demografía + tráfico + mercado + forecast + CTA"""
    buffer = BytesIO()
    doc    = SimpleDocTemplate(buffer, pagesize=letter,
                               leftMargin=45, rightMargin=45, topMargin=45, bottomMargin=45,
                               title=f"4SITE — Reporte de Análisis: {ubicacion[:60]}",
                               author="4SITE · Don't guess. Foresee.",
                               subject="Análisis de viabilidad comercial",
                               creator="4SITE App · 4site.mx",
                               producer="4SITE")
    s     = _estilos_pdf()
    t     = TEXTOS[idioma]
    story = []
    tier  = "pro"
    tipo_info = TIPOS_NEGOCIO.get(tipo_negocio, {}) if tipo_negocio else {}

    # Resolver tipo de negocio para modo recomendar (usar top1)
    tipo_negocio_efectivo = tipo_negocio
    tipo_info_efectivo    = tipo_info
    if not tipo_negocio_efectivo and modo == "recomendar" and recomendaciones:
        tipo_negocio_efectivo = recomendaciones[0]["tipo_key"]
        tipo_info_efectivo    = TIPOS_NEGOCIO.get(tipo_negocio_efectivo, {})
        competidores = recomendaciones[0]["competidores"] or competidores

    # Calcular datos PRO si módulos disponibles
    trafico_data = None
    mercado_data = None
    forecast_data = None
    roi_data = None
    ticket_data = None
    if MODULOS_OK and demografia and tipo_negocio_efectivo:
        try:
            tipo_zona = contexto.get("tipo_zona", "mixto") if contexto else "mixto"
            densidad_a = demografia.get("densidad_actual", demografia.get("densidad_hab_km2", 8000))
            trafico_data  = generar_reporte_trafico(tipo_zona, tipo_negocio_efectivo, densidad_a)
            mercado_data  = calcular_mercado_potencial(demografia, tipo_negocio_efectivo, len(competidores or []))
            forecast_data = generar_forecast(mercado_data, tipo_negocio_efectivo, score,
                                             tipo_info_efectivo.get("inversion_min", 300000))
            roi_data      = calcular_roi(forecast_data, tipo_info_efectivo.get("inversion_min", 300000),
                                         tipo_info_efectivo.get("inversion_max", 600000), tipo_negocio_efectivo)
        except:
            pass
    if MODULOS_OK and competidores:
        try: ticket_data = calcular_ticket_competencia(competidores, tipo_negocio_efectivo or "default", score)
        except: pass

    # ── P1: PORTADA ESTILO 4SITE (foto + panel blanco) ──
    _portada_buf = BytesIO()
    _portada_cv  = canvas.Canvas(_portada_buf, pagesize=letter)
    _tn_portada  = tipo_info_efectivo.get('nombre', tipo_info.get('nombre',''))
    if not _tn_portada and recomendaciones:
        _tn_portada = recomendaciones[0].get('nombre','')
    _base_dir   = os.path.dirname(os.path.abspath(__file__))
    _logo_path  = _encontrar_asset(_base_dir, "logo_4site")
    _bg_path    = _encontrar_asset(_base_dir, "bg_calle")
    _portada_ok = bool(_logo_path) and bool(_bg_path)
    if _portada_ok:
        _portada_foto_canvas(_portada_cv, ubicacion, score, "pro",
            datetime.datetime.now().strftime("%d de %B de %Y"),
            _tn_portada, _logo_path, _bg_path)
        _portada_cv.save(); _portada_buf.seek(0)

    # Score gauge embebido
    if GRAFICAS_OK:
        nivel_txt, _ = nivel_score(score, idioma)
        gauge = grafica_score_gauge(score, nivel_txt)
        story.append(RLImage(gauge, width=3.5*inch, height=2*inch))
    else:
        _score_block(story, score, desglose, tier, t, s, idioma)

    # Análisis vial en PDF
    av_pdf = contexto.get("analisis_vial") if contexto else None
    if av_pdf and av_pdf.get("descripcion"):
        story.append(Spacer(1, 0.1*inch))
        ajuste_v = av_pdf.get("ajuste_score", 0)
        color_v  = '#4CAF50' if ajuste_v > 0 else '#F44336' if ajuste_v < -5 else '#FF9800'
        vial_rows = [
            ["Tipo de vialidad:", av_pdf.get("badge","")],
            ["Impacto en score:", f"{'+'if ajuste_v>=0 else ''}{ajuste_v} puntos"],
            ["Análisis:", av_pdf.get("descripcion","")],
        ]
        tbl_vial = Table(vial_rows, colWidths=[1.5*inch, 4*inch])
        tbl_vial.setStyle(TableStyle([
            ('FONTNAME',  (0,0),(0,-1),'Helvetica-Bold'),
            ('FONTSIZE',  (0,0),(-1,-1),9),
            ('TEXTCOLOR', (0,0),(0,-1),colors.HexColor('#6B6B6B')),
            ('TEXTCOLOR', (1,1),(1,1),colors.HexColor(color_v)),
            ('FONTNAME',  (1,1),(1,1),'Helvetica-Bold'),
            ('ROWBACKGROUNDS',(0,0),(-1,-1),[colors.HexColor('#F0F4FF'),colors.white,colors.HexColor('#F0F4FF')]),
            ('GRID',(0,0),(-1,-1),0.3,colors.HexColor('#DDDDDD')),
            ('TOPPADDING',(0,0),(-1,-1),5),('BOTTOMPADDING',(0,0),(-1,-1),5),
            ('VALIGN',(0,0),(-1,-1),'TOP'),
        ]))
        story.extend(_seccion_pdf("Análisis de Vialidad:", "🛣️", s))
        story.append(tbl_vial)

    # Sección de plaza comercial en PDF — si aplica
    _ip_pdf = contexto.get("info_plaza", {}) if contexto else {}
    if _ip_pdf and _ip_pdf.get("tipo_plaza"):
        _pl_tipo_pdf  = _ip_pdf["tipo_plaza"]
        _pl_label_pdf = "🏬 Mall / Plaza Cerrada" if _pl_tipo_pdf == "mall" else "🏪 Strip Center / Plaza Abierta"
        _pl_nota_pdf  = _ip_pdf.get("nota", "").replace("\n", "  |  ")
        _pl_roi_pdf   = _ip_pdf.get("ajuste_roi_pct", 0)
        _pl_zona_pdf  = _ip_pdf.get("zona_mall", "")
        story.append(Spacer(1, 0.08*inch))
        story.extend(_seccion_pdf("Tipo de Ubicación — Plaza Comercial:", "🏢", s))
        plaza_rows = [["Tipo de plaza:", _pl_label_pdf]]
        if _pl_zona_pdf and _pl_tipo_pdf == "mall":
            plaza_rows.append(["Zona en mall:", _pl_zona_pdf])
        if _pl_roi_pdf > 0:
            plaza_rows.append(["Impacto ROI:", f"⚠️ -{_pl_roi_pdf}% vs local en calle equivalente por renta de plaza"])
        plaza_rows.append(["Notas:", _pl_nota_pdf[:300]])
        tbl_plaza = Table(plaza_rows, colWidths=[1.5*inch, 4*inch])
        tbl_plaza.setStyle(TableStyle([
            ('FONTNAME',  (0,0),(0,-1), 'Helvetica-Bold'),
            ('FONTSIZE',  (0,0),(-1,-1), 9),
            ('TEXTCOLOR', (0,0),(0,-1), colors.HexColor('#6B6B6B')),
            ('ROWBACKGROUNDS',(0,0),(-1,-1),[colors.HexColor('#FFF8F0'), colors.white]),
            ('GRID',(0,0),(-1,-1),0.3,colors.HexColor('#DDDDDD')),
            ('TOPPADDING',(0,0),(-1,-1),5),('BOTTOMPADDING',(0,0),(-1,-1),5),
            ('VALIGN',(0,0),(-1,-1),'TOP'),
        ]))
        story.append(tbl_plaza)
        # Narrativa interpretativa vial para PRO/PREMIUM
        _tipo_nom_vial = (tipo_info_efectivo.get('nombre','') if 'tipo_info_efectivo' in dir()
                          else tipo_info.get('nombre',''))
        if not _tipo_nom_vial and recomendaciones:
            _tipo_nom_vial = recomendaciones[0].get('nombre','')
        narr_vial_p = generar_narrativa_seccion("mapa", {
            "tipo_vialidad": av_pdf.get("tipo_vialidad",""),
            "ajuste_score": av_pdf.get("ajuste_score", 0),
            "descripcion": av_pdf.get("descripcion",""),
            "dependencia_paso": av_pdf.get("dependencia_paso","medio"),
        }, _tipo_nom_vial, idioma)
        for el in _bloque_narrativa_pdf(narr_vial_p, s):
            story.append(el)
        story.append(Spacer(1, 0.05*inch))

    # Características de la zona / Generadores de tráfico
    if contexto and contexto.get("badges"):
        story.append(Spacer(1, 0.12*inch))
        story.append(Paragraph("📍 Características de la Zona y Generadores de Tráfico:",
            ParagraphStyle('badges_titulo', parent=s["base"]['Normal'],
            fontSize=10, fontName='Helvetica-Bold',
            textColor=colors.HexColor('#6B6B6B'), spaceAfter=4)))
        BADGE_DESCRIPCIONES = {
            "🚗 Alto tráfico vehicular": "Vialidad de alto flujo — ideal para negocios de paso y visibilidad",
            "⛽ Gasolinera cercana":      "Gasolinera en 500m — genera tráfico vehicular constante",
            "🏬 Plaza comercial":         "Centro comercial cercano — concentrador de shoppers y tráfico",
            "🏥 Hospital cercano":        "Hospital/clínica en 500m — clientela cautiva de alta frecuencia",
            "🎓 Escuela/Universidad":     "Institución educativa — alto flujo estudiantil y docente",
            "🚌 Transporte público":      "Estación de transporte — alta accesibilidad peatonal",
        }
        # Filtrar badges de vialidad — ya se muestran en el bloque de análisis vial
        badges_vial = {"🛣️", "🚦", "🏘️", "📍"}
        badges_filtrados = [b for b in contexto["badges"]
                           if not any(b.startswith(v) for v in badges_vial)]
        for badge in badges_filtrados:
            desc = BADGE_DESCRIPCIONES.get(badge, "Característica detectada en la zona")
            story.append(Paragraph(f"  {badge}  —  {desc}",
                ParagraphStyle('badge_item', parent=s["base"]['Normal'],
                fontSize=8.5, textColor=colors.HexColor('#1565C0'),
                leftIndent=8, spaceAfter=3)))

    story.append(Paragraph(f"<i>Reporte generado por 4SITE · hola@4site.mx · 4site.mx</i>",
                            ParagraphStyle('footer_p', parent=s["base"]['Normal'],
                            fontSize=8, textColor=colors.grey, alignment=TA_CENTER)))

    # ── P2: DESGLOSE SCORE ──
    # ── P5: COMPETIDORES ──
    story.append(PageBreak())
    story.append(Paragraph(t["competencia_titulo"], s["h2"]))
    story.append(Paragraph(f"{len(competidores or [])} negocios similares encontrados en radio de 500m", s["small"]))
    story.append(Spacer(1, 0.1*inch))
    if competidores:
        story.append(_tabla_competidores(competidores, t, 20, s))
        # Narrativa interpretativa de competencia
        ratings_c  = [c.get('rating',0) for c in (competidores or []) if c.get('rating')]
        n_verdes_c = sum(1 for r in ratings_c if r >= 4.3)
        avg_r_c    = round(sum(ratings_c)/len(ratings_c),1) if ratings_c else 0
        datos_comp_narr = {
            "total_competidores": len(competidores or []),
            "rating_promedio": avg_r_c,
            "competidores_excelentes": n_verdes_c,
            "tipo_negocio": tipo_negocio_efectivo or tipo_negocio or ""
        }
        narr_comp = generar_narrativa_seccion("competencia", datos_comp_narr,
                                               tipo_info_efectivo.get('nombre', tipo_info.get('nombre','')), idioma)
        for el in _bloque_narrativa_pdf(narr_comp, s):
            story.append(el)

    # ── P5b: TICKET PROMEDIO (PRO) ──
    if ticket_data:
        story.append(PageBreak())
        story.extend(_seccion_pdf("Ticket Promedio de la Competencia", "🎫", s))

        # Detectar si el tipo necesita nota de confiabilidad
        _tipo_eff_tk = tipo_negocio_efectivo or "default"
        _tk_override = _get_ticket_override(_tipo_eff_tk)
        _tk_nota     = _get_ticket_nota(_tipo_eff_tk)

        # Override de valores si aplica
        if _tk_override:
            ticket_data = dict(ticket_data)  # copia para no mutar el original
            ticket_data["ticket_competencia"] = _tk_override["ticket_tipico"]
            ticket_data["rango_zona"]          = _tk_override["rango_txt"][:40]
            ticket_data["ticket_recomendado"]  = int(_tk_override["ticket_tipico"] * 1.1)
            story.append(Paragraph(
                f"Referencia de mercado · {_tk_override['rango_txt']} · "
                "Google Places price_level no es confiable para este tipo de negocio",
                s["small"]
            ))
        else:
            story.append(Paragraph("Fuente: price_level Google Places + ENIGH 2022", s["small"]))

        story.append(Spacer(1, 0.08*inch))
        _kpi_t = ParagraphStyle('kpit',parent=s["base"]['Normal'],fontSize=24,fontName='Helvetica-Bold',textColor=colors.HexColor('#BB944F'),alignment=TA_CENTER)
        _kpi_l = ParagraphStyle('kpil',parent=s["base"]['Normal'],fontSize=8,textColor=colors.HexColor('#6B6B6B'),alignment=TA_CENTER)
        _label_ticket = "Ticket referencial zona" if _tk_override else "Ticket promedio competencia"
        _kt = Table([
            [Paragraph(f"${ticket_data['ticket_competencia']:,} MXN",_kpi_t), Paragraph(ticket_data['rango_zona'][:30],_kpi_t)],
            [Paragraph(_label_ticket,_kpi_l), Paragraph("Nivel de precio zona",_kpi_l)],
        ], colWidths=[3*inch, 3*inch])
        _kt.setStyle(TableStyle([
            ('ALIGN',(0,0),(-1,-1),'CENTER'),
            ('BACKGROUND',(0,0),(-1,-1),colors.HexColor('#F9F9F9')),
            ('GRID',(0,0),(-1,-1),0.3,colors.HexColor('#DDDDDD')),
            ('TOPPADDING',(0,0),(-1,-1),10),('BOTTOMPADDING',(0,0),(-1,-1),10),
        ]))
        story.append(_kt)
        story.append(Spacer(1,0.08*inch))

        # Nota de advertencia para tipos con datos no confiables
        if _tk_nota:
            _nota_style = ParagraphStyle('tk_warn', parent=s["base"]['Normal'],
                fontSize=8.5, textColor=colors.HexColor('#7B4500'),
                backColor=colors.HexColor('#FFF3CD'),
                borderPad=8, leading=13, leftIndent=8, rightIndent=8)
            story.append(Paragraph(f"⚠️ {_tk_nota.replace('⚠️','').replace('ℹ️','').strip()}", _nota_style))
            if _tk_override and _tk_override.get("ejemplos"):
                story.append(Spacer(1,0.04*inch))
                story.append(Paragraph(
                    f"Rangos de referencia MX: {_tk_override['ejemplos']}", s["small"]))
        else:
            story.append(Paragraph(ticket_data['interpretacion_pro'], s["normal"]))
            if ticket_data.get('detalle_competidores'):
                _dr = [["Competidor","Precio","Ticket est."]]
                for _d in ticket_data['detalle_competidores'][:10]:
                    _dr.append([_d['nombre'][:32], _d['nivel'], f"${_d['ticket_estimado']:,}"])
                _dt = Table(_dr, colWidths=[3*inch,1.2*inch,1.8*inch])
                _dt.setStyle(TableStyle([
                    ('BACKGROUND',(0,0),(-1,0),colors.HexColor('#6B6B6B')),
                    ('TEXTCOLOR',(0,0),(-1,0),colors.white),
                    ('FONTNAME',(0,0),(-1,0),'Helvetica-Bold'),
                    ('FONTSIZE',(0,0),(-1,-1),8.5),
                    ('ALIGN',(1,0),(-1,-1),'CENTER'),
                    ('GRID',(0,0),(-1,-1),0.3,colors.HexColor('#DDDDDD')),
                    ('ROWBACKGROUNDS',(0,1),(-1,-1),[colors.white,colors.HexColor('#F5F5F5')]),
                    ('TOPPADDING',(0,0),(-1,-1),5),('BOTTOMPADDING',(0,0),(-1,-1),5),
                    ('LEFTPADDING',(0,0),(-1,-1),6),
                ]))
                story.append(Spacer(1,0.06*inch))
                story.append(_dt)

    # ── P6: DEMOGRAFÍA ──
    story.append(PageBreak())
    story.extend(_seccion_pdf("Perfil Demográfico del Área (500m)", "👥", s))
    fuente = demografia.get('fuente', 'Estimación por zona') if demografia else ''
    story.append(Paragraph(f"📊 Fuente: {fuente}", s["small"]))
    story.append(Spacer(1, 0.1*inch))
    if demografia and GRAFICAS_OK:
        graf_dem = grafica_demografia(demografia)
        story.append(RLImage(graf_dem, width=5.5*inch, height=2.6*inch))
    if demografia:
        año_act = datetime.datetime.now().year
        densidad_info = clasificar_densidad(demografia.get('densidad_actual', 8000)) if MODULOS_OK else formatear_densidad(demografia.get('densidad_hab_km2', 8000))
        dem_rows = [
            [f"Población {año_act}", f"{demografia.get('poblacion_actual', demografia.get('poblacion_estimada',0)):,} hab",
             f"(Censo 2020: {demografia.get('poblacion_2020','-'):,})" if demografia.get('poblacion_2020') else ""],
            ["Viviendas habitadas", f"{demografia.get('viviendas_actual', demografia.get('viviendas_habitadas',0)):,}", ""],
            ["Densidad urbana", densidad_info.get('nivel', densidad_info.get('clasificacion','')),
             f"~{demografia.get('personas_manzana', densidad_info.get('personas_manzana',0))} personas/manzana"],
            ["NSE predominante", demografia.get('nse_predominante','C'), "Nivel Socioeconómico AMAI"],
            ["Ingreso prom. mensual", f"${demografia.get('ingreso_actual', demografia.get('ingreso_promedio_mensual',0)):,}", "Proyectado 2026 (INPC)"],
            ["Gasto prom. mensual", f"${demografia.get('gasto_actual', demografia.get('gasto_promedio_mensual',0)):,}", "~78% del ingreso (ENIGH)"],
            ["Crecimiento anual zona", f"{demografia.get('tasa_crecimiento_pct',0.55):.2f}%", "Tasa CONAPO 2020-2030"],
        ]
        # KPIs visuales de demografía
        _pob_act = demografia.get('poblacion_actual', demografia.get('poblacion_estimada',0))
        _nse_d   = demografia.get('nse_predominante','C')
        _ing_d   = demografia.get('ingreso_actual', demografia.get('ingreso_promedio_mensual',0))
        _cre_d   = demografia.get('tasa_crecimiento_pct',0)
        story.append(_kpi_pdf([
            {"label": f"Población {año_act}",    "valor": f"{_pob_act:,} hab"},
            {"label": "NSE predominante",         "valor": _nse_d, "gold": True},
            {"label": "Ingreso prom./mes",         "valor": f"${_ing_d:,}"},
            {"label": "Crecimiento anual",         "valor": f"{_cre_d:.2f}%"},
        ], s))
        story.append(Spacer(1, 0.1*inch))
        # Tabla detallada con _tabla_pdf
        _hdr_dem = ["Indicador", "Valor", "Referencia"]
        _dem_rows_full = [_hdr_dem] + dem_rows
        story.append(_tabla_pdf(_dem_rows_full, [2*inch, 1.8*inch, 2.2*inch], s))
        # Gráfica edad/género (datos reales INEGI/ENIGH)
        if GRAFICAS_OK and "grafica_edad_genero" in dir():
            try:
                _geg = grafica_edad_genero(demografia, tipo_info_efectivo.get('nombre', tipo_info.get('nombre','')))
                story.append(Spacer(1,0.08*inch))
                story.append(RLImage(_geg, width=5.3*inch, height=2.6*inch))
            except: pass
        # Interpretación demográfica específica al negocio
        if MODULOS_OK and "interpretar_demografia_negocio" in dir():
            try:
                _idem = interpretar_demografia_negocio(demografia,
                    tipo_info_efectivo.get('nombre', tipo_info.get('nombre','')),
                    tipo_negocio_efectivo or tipo_negocio or "")
                story.append(Spacer(1,0.08*inch))
                for _ln in _idem.split("\n\n"):
                    if _ln.strip():
                        _lf = _ln.strip()
                        import re as _re
                        _lf = _re.sub(r'\*\*(.+?)\*\*', r'<b></b>', _lf)
                        story.append(Paragraph(_lf, s["normal"]))
                        story.append(Spacer(1,0.04*inch))
            except: pass

    # ── P7: TRÁFICO ──
    story.append(PageBreak())
    story.extend(_seccion_pdf("Análisis de Tráfico Estimado", "🚗", s))
    story.append(Paragraph(
        "Perfil de flujo peatonal y vehicular estimado para este tipo de negocio "
        "basado en el tipo de zona detectada y densidad poblacional.",
        s["normal"]))
    if trafico_data and GRAFICAS_OK:
        graf_trafico_h = grafica_trafico_horario(
            trafico_data["trafico_horario"],
            tipo_info.get('nombre','') if tipo_negocio else '')
        story.append(RLImage(graf_trafico_h, width=5.5*inch, height=2.8*inch))
        story.append(Spacer(1, 0.1*inch))

        graf_trafico_s = grafica_trafico_semanal(trafico_data["trafico_semanal"])
        story.append(RLImage(graf_trafico_s, width=5.5*inch, height=1.8*inch))

        # Tabla horas pico
        story.append(Spacer(1, 0.1*inch))
        story.append(Paragraph("🕐 Horas pico recomendadas para operar / mayor flujo:", s["normal"]))
        pico_rows = [["Horario", "Nivel de flujo", "Recomendación"]]
        for pico in trafico_data["horas_pico"]:
            rec = ("Máxima atención al cliente" if pico["flujo"] >= 80
                   else "Alta actividad" if pico["flujo"] >= 60
                   else "Flujo moderado")
            pico_rows.append([pico["hora_str"], f"{pico['nivel']} ({pico['flujo']}%)", rec])
        story.append(_tabla_pdf(pico_rows, [1.8*inch, 1.8*inch, 2.4*inch], s))
        # Narrativa interpretativa de tráfico
        _tipo_nom_narr = (tipo_info_efectivo.get('nombre','') if 'tipo_info_efectivo' in dir()
                         else tipo_info.get('nombre',''))
        if trafico_data:
            datos_traf_narr = {
                "nivel_general": trafico_data.get('nivel_general',''),
                "horas_pico": [p['hora_str'] for p in trafico_data.get('horas_pico',[])[:2]],
                "mejor_dia": trafico_data.get('dia_pico',''),
                "dia_bajo": trafico_data.get('dia_bajo',''),
            }
            narr_traf = generar_narrativa_seccion("trafico", datos_traf_narr,
                                                   _tipo_nom_narr, idioma)
            for el in _bloque_narrativa_pdf(narr_traf, s):
                story.append(el)
    else:
        story.append(Paragraph("🔜 Datos de tráfico no disponibles — verifica que modulos_4site.py esté instalado.", s["small"]))

    # ── P8: MERCADO + FORECAST ──
    story.append(PageBreak())
    story.extend(_seccion_pdf("Tamaño de Mercado Potencial", "📊", s))
    if mercado_data and GRAFICAS_OK:
        graf_mercado = grafica_mercado_donut(mercado_data)
        story.append(RLImage(graf_mercado, width=5*inch, height=3*inch))
        story.append(Paragraph(f"📐 Metodología: {mercado_data['metodologia']}", s["small"]))
        # Narrativa interpretativa de mercado
        narr_mercado = generar_narrativa_seccion("mercado", {
            "mercado_total_mensual": mercado_data.get('mercado_total_mensual',0),
            "mercado_captura_mensual": mercado_data.get('mercado_captura_mensual',0),
            "factor_captura_pct": mercado_data.get('factor_captura_pct',0),
            "clientes_dia": mercado_data.get('clientes_dia_estimados',0),
        }, _tipo_nom_narr if '_tipo_nom_narr' in dir() else tipo_info.get('nombre',''), idioma)
        for el in _bloque_narrativa_pdf(narr_mercado, s):
            story.append(el)

    story.append(Spacer(1, 0.15*inch))
    story.extend(_seccion_pdf("Forecast de Ventas — 12 Meses (3 Escenarios)", "📈", s))
    if forecast_data and GRAFICAS_OK:
        graf_forecast = grafica_forecast(forecast_data)
        fc_rows = [["Escenario", "Ventas mes 1", "Ventas mes 9", "Total año 1", "Prom. mensual"]]
        story.append(RLImage(graf_forecast, width=5.5*inch, height=3*inch))
        for key, emoji in [("pesimista","🔴"), ("base","🔵"), ("optimista","🟢")]:
            e = forecast_data["escenarios"][key]
            fc_rows.append([
                f"{emoji} {key.capitalize()}",
                f"${e['ventas_mensuales'][0]/1000:.0f}K",
                f"${e['ventas_mensuales'][8]/1000:.0f}K",
                f"${e['total_anual']/1000:.0f}K",
                f"${e['promedio_mensual']/1000:.0f}K",
            ])
        tbl_fc = Table(fc_rows, colWidths=[1.2*inch, 1*inch, 1*inch, 1.2*inch, 1.2*inch])
        tbl_fc.setStyle(TableStyle([
            ('BACKGROUND', (0,0), (-1,0), colors.HexColor('#6B6B6B')),
            ('TEXTCOLOR',  (0,0), (-1,0), colors.white),
            ('FONTNAME',   (0,0), (-1,0), 'Helvetica-Bold'),
            ('FONTSIZE',   (0,0), (-1,-1), 9),
            ('ALIGN',      (1,0), (-1,-1), 'CENTER'),
            ('GRID', (0,0), (-1,-1), 0.3, colors.HexColor('#DDDDDD')),
            ('ROWBACKGROUNDS', (0,1), (-1,-1), [
                colors.HexColor('#FFEBEE'),
                colors.HexColor('#E3F2FD'),
                colors.HexColor('#E8F5E9'),
            ]),
            ('TOPPADDING', (0,0), (-1,-1), 6),
            ('BOTTOMPADDING', (0,0), (-1,-1), 6),
        ]))
        story.append(tbl_fc)
        story.append(Spacer(1, 0.08*inch))
        story.append(Paragraph("📋 Supuestos: " + " · ".join(forecast_data.get("supuestos",[])), s["small"]))
        # Narrativa interpretativa de forecast
        esc = forecast_data.get('escenarios',{})
        narr_fc = generar_narrativa_seccion("forecast", {
            "total_pesimista": esc.get('pesimista',{}).get('total_anual',0),
            "total_base": esc.get('base',{}).get('total_anual',0),
            "total_optimista": esc.get('optimista',{}).get('total_anual',0),
            "promedio_mensual_base": esc.get('base',{}).get('promedio_mensual',0),
        }, _tipo_nom_narr if '_tipo_nom_narr' in dir() else tipo_info.get('nombre',''), idioma)
        for el in _bloque_narrativa_pdf(narr_fc, s):
            story.append(el)

    # ── P9: Recomendación Final ──
    story.append(PageBreak())
    _plaza_block_pdf(story, contexto, s)
    _render_recomendacion_pdf(story, analisis, s)


    # ── P11: CTA ──
    story.append(PageBreak())
    _cta_block(story, tier, t, s, idioma)

    doc.build(story, onFirstPage=_footer_pdf, onLaterPages=_footer_pdf)
    buffer.seek(0)
    return buffer


# ─────────────────────────────────────────────
# PDF PREMIUM  (Análisis PREMIUM — dashboard + comparativa)
# ─────────────────────────────────────────────
def _dsg_get(dsg, key, default=0):
    """Accede a desglose de forma segura (dict, tupla o None)."""
    if isinstance(dsg, dict): return dsg.get(key, default)
    if isinstance(dsg, (list, tuple)):
        for item in dsg:
            if isinstance(item, dict): return item.get(key, default)
    return default


def generar_pdf_premium(ubicacion, score, desglose, analisis, competidores,
                          idioma, lat, lng, modo, tipo_negocio=None,
                          recomendaciones=None, demografia=None, contexto=None):
    """PREMIUM: todo lo de PRO + dashboard página completa + ROI + comparativa"""
    buffer = BytesIO()
    doc    = SimpleDocTemplate(buffer, pagesize=letter,
                               leftMargin=45, rightMargin=45, topMargin=45, bottomMargin=45,
                               title=f"4SITE — Reporte de Análisis: {ubicacion[:60]}",
                               author="4SITE · Don't guess. Foresee.",
                               subject="Análisis de viabilidad comercial",
                               creator="4SITE App · 4site.mx",
                               producer="4SITE")
    s     = _estilos_pdf()
    t     = TEXTOS[idioma]
    story = []
    tier  = "premium"
    tipo_info = TIPOS_NEGOCIO.get(tipo_negocio, {}) if tipo_negocio else {}

    # Resolver tipo de negocio para modo recomendar (usar top1)
    tipo_negocio_efectivo = tipo_negocio
    tipo_info_efectivo    = tipo_info
    if not tipo_negocio_efectivo and modo == "recomendar" and recomendaciones:
        tipo_negocio_efectivo = recomendaciones[0]["tipo_key"]
        tipo_info_efectivo    = TIPOS_NEGOCIO.get(tipo_negocio_efectivo, {})
        if not competidores and recomendaciones[0].get("competidores"):
            competidores = recomendaciones[0]["competidores"]

    # Calcular datos
    trafico_data = None
    mercado_data = None
    forecast_data = None
    roi_data = None
    ticket_data = None
    if MODULOS_OK and demografia and tipo_negocio_efectivo:
        try:
            tipo_zona = contexto.get("tipo_zona","mixto") if contexto else "mixto"
            densidad_a = demografia.get("densidad_actual", demografia.get("densidad_hab_km2", 8000))
            trafico_data  = generar_reporte_trafico(tipo_zona, tipo_negocio_efectivo, densidad_a)
            mercado_data  = calcular_mercado_potencial(demografia, tipo_negocio_efectivo, len(competidores or []))
            forecast_data = generar_forecast(mercado_data, tipo_negocio_efectivo, score,
                                             tipo_info_efectivo.get("inversion_min", 300000))
            roi_data      = calcular_roi(forecast_data, tipo_info_efectivo.get("inversion_min", 300000),
                                         tipo_info_efectivo.get("inversion_max", 600000), tipo_negocio_efectivo)
        except:
            pass
    if MODULOS_OK and competidores:
        try: ticket_data = calcular_ticket_competencia(competidores, tipo_negocio_efectivo or "default", score)
        except: pass

    # ── P1: PORTADA ESTILO 4SITE (foto + panel blanco) ──
    _portada_buf_p = BytesIO()
    _portada_cv_p  = canvas.Canvas(_portada_buf_p, pagesize=letter)
    _tn_p  = tipo_info_efectivo.get('nombre', tipo_info.get('nombre',''))
    if not _tn_p and recomendaciones: _tn_p = recomendaciones[0].get('nombre','')
    _bd2   = os.path.dirname(os.path.abspath(__file__))
    _lp2   = _encontrar_asset(_bd2, "logo_4site")
    _bp2   = _encontrar_asset(_bd2, "bg_calle")
    _ok_p  = bool(_lp2) and bool(_bp2)
    if _ok_p:
        _portada_foto_canvas(_portada_cv_p, ubicacion, score, "premium",
            datetime.datetime.now().strftime("%d de %B de %Y"),
            _tn_p, _lp2, _bp2)
        _portada_cv_p.save(); _portada_buf_p.seek(0)
    # ── P2: MAPA CON COMPETIDORES ──
    story.append(PageBreak())
    story.extend(_seccion_pdf("Mapa de Ubicación y Competidores", "📍", s))
    story.append(Paragraph("🟢 Excelente (≥4.3★)   🟠 Bueno (3.5-4.2★)   🔴 Débil/sin datos   ★ Tu ubicación",
        ParagraphStyle('ley', parent=s["base"]['Normal'], fontSize=8.5,
        textColor=colors.HexColor('#555'), spaceAfter=8)))
    mapa = generar_mapa_estatico(lat, lng, competidores=competidores, con_competidores=True)
    if mapa:
        story.append(RLImage(mapa, width=5.5*inch, height=3.7*inch))
    interp = generar_texto_interpretacion_mapa(competidores or [], contexto or {}, "competidores", idioma)
    story.append(Paragraph(interp, ParagraphStyle('interp2', parent=s["base"]['Normal'],
        fontSize=8.5, textColor=colors.HexColor('#444'), spaceAfter=6,
        backColor=colors.HexColor('#F5F8FF'), borderPad=6)))

    # ── P3b: ANÁLISIS VIAL ──
    av_prem = contexto.get("analisis_vial") if contexto else None
    if av_prem and av_prem.get("descripcion"):
        story.append(PageBreak())
        story.extend(_seccion_pdf("Análisis de Vialidad", "🛣️", s))
        ajuste_prem = av_prem.get("ajuste_score", 0)
        color_prem  = '#4CAF50' if ajuste_prem > 0 else '#F44336' if ajuste_prem < -5 else '#FF9800'
        dep_prem = {"alto":"Negocio de paso (crítico)","medio":"Negocio mixto",
                    "bajo":"Negocio de destino"}.get(av_prem.get("dependencia_paso","medio"),"")
        vial_rows_prem = [
            ["Tipo de vialidad:",       av_prem.get("badge","")],
            ["Dependencia al tráfico:", dep_prem],
            ["Impacto en score:",       f"{'+'if ajuste_prem>=0 else ''}{ajuste_prem} puntos"],
            ["Análisis:",               av_prem.get("descripcion","")],
        ]
        tbl_vial_prem = Table(vial_rows_prem, colWidths=[1.8*inch, 3.7*inch])
        tbl_vial_prem.setStyle(TableStyle([
            ('FONTNAME',  (0,0),(0,-1),'Helvetica-Bold'),
            ('FONTSIZE',  (0,0),(-1,-1),9),
            ('TEXTCOLOR', (0,0),(0,-1),colors.HexColor('#6B6B6B')),
            ('TEXTCOLOR', (1,2),(1,2),colors.HexColor(color_prem)),
            ('FONTNAME',  (1,2),(1,2),'Helvetica-Bold'),
            ('ROWBACKGROUNDS',(0,0),(-1,-1),[colors.HexColor('#F0F4FF'),colors.white,
                                             colors.HexColor('#F0F4FF'),colors.white]),
            ('GRID',(0,0),(-1,-1),0.3,colors.HexColor('#DDDDDD')),
            ('TOPPADDING',(0,0),(-1,-1),6),('BOTTOMPADDING',(0,0),(-1,-1),6),
            ('VALIGN',(0,0),(-1,-1),'TOP'),
        ]))
        story.append(tbl_vial_prem)
        # Narrativa interpretativa vial
        _tipo_nom_vp = tipo_info_efectivo.get('nombre', tipo_info.get('nombre',''))
        if not _tipo_nom_vp and recomendaciones:
            _tipo_nom_vp = recomendaciones[0].get('nombre','')
        narr_vial_prem = generar_narrativa_seccion("mapa", {
            "tipo_vialidad":    av_prem.get("tipo_vialidad",""),
            "ajuste_score":     ajuste_prem,
            "descripcion":      av_prem.get("descripcion",""),
            "dependencia_paso": av_prem.get("dependencia_paso","medio"),
        }, _tipo_nom_vp, idioma)
        for el in _bloque_narrativa_pdf(narr_vial_prem, s):
            story.append(el)

    # ── P4: COMPETIDORES ──
    story.append(PageBreak())
    story.append(Paragraph(t["competencia_titulo"], s["h2"]))
    if competidores:
        story.append(_tabla_competidores(competidores, t, 20, s))

    # ── P4b: TICKET COMPETENCIA + RECOMENDADO (PREMIUM) ──
    if ticket_data:
        story.append(PageBreak())
        story.extend(_seccion_pdf("Análisis de Ticket Promedio", "🎫", s))

        # Detectar si el tipo necesita nota de confiabilidad
        _tipo_eff_tk2 = tipo_negocio_efectivo or "default"
        _tk_override2 = _get_ticket_override(_tipo_eff_tk2)
        _tk_nota2     = _get_ticket_nota(_tipo_eff_tk2)

        # Override de valores si aplica
        if _tk_override2:
            ticket_data = dict(ticket_data)
            ticket_data["ticket_competencia"] = _tk_override2["ticket_tipico"]
            ticket_data["rango_zona"]          = _tk_override2["rango_txt"][:40]
            ticket_data["ticket_recomendado"]  = int(_tk_override2["ticket_tipico"] * 1.1)
            story.append(Paragraph(
                f"Referencia de mercado · {_tk_override2['rango_txt']} · "
                "Google Places price_level no es confiable para este tipo de negocio",
                s["small"]
            ))
        else:
            story.append(Paragraph("Fuente: price_level Google Places + ENIGH 2022", s["small"]))

        story.append(Spacer(1,0.08*inch))
        _kpi_t2 = ParagraphStyle('kpit2',parent=s["base"]['Normal'],fontSize=22,fontName='Helvetica-Bold',textColor=colors.HexColor('#BB944F'),alignment=TA_CENTER)
        _kpi_r2 = ParagraphStyle('kpir2',parent=s["base"]['Normal'],fontSize=22,fontName='Helvetica-Bold',textColor=colors.HexColor('#1D9E75'),alignment=TA_CENTER)
        _kpi_l2 = ParagraphStyle('kpil2',parent=s["base"]['Normal'],fontSize=8,textColor=colors.HexColor('#6B6B6B'),alignment=TA_CENTER)
        _label_tk2 = "Ticket referencial zona" if _tk_override2 else "Ticket promedio competencia"
        _kt2 = Table([
            [Paragraph(f"${ticket_data['ticket_competencia']:,} MXN",_kpi_t2),
             Paragraph(ticket_data['rango_zona'][:30],_kpi_t2),
             Paragraph(f"${ticket_data['ticket_recomendado']:,} MXN",_kpi_r2)],
            [Paragraph(_label_tk2,_kpi_l2),
             Paragraph("Nivel de precio zona",_kpi_l2),
             Paragraph("Ticket recomendado para tu negocio",_kpi_l2)],
        ], colWidths=[2.1*inch,1.6*inch,2.3*inch])
        _kt2.setStyle(TableStyle([
            ('ALIGN',(0,0),(-1,-1),'CENTER'),
            ('BACKGROUND',(0,0),(1,-1),colors.HexColor('#F9F9F9')),
            ('BACKGROUND',(2,0),(2,-1),colors.HexColor('#F0FBF6')),
            ('GRID',(0,0),(-1,-1),0.3,colors.HexColor('#DDDDDD')),
            ('TOPPADDING',(0,0),(-1,-1),10),('BOTTOMPADDING',(0,0),(-1,-1),10),
        ]))
        story.append(_kt2)
        story.append(Spacer(1,0.08*inch))

        # Nota de advertencia o interpretación normal
        if _tk_nota2:
            _nota_style2 = ParagraphStyle('tk_warn2', parent=s["base"]['Normal'],
                fontSize=8.5, textColor=colors.HexColor('#7B4500'),
                backColor=colors.HexColor('#FFF3CD'),
                borderPad=8, leading=13, leftIndent=8, rightIndent=8)
            story.append(Paragraph(f"⚠️ {_tk_nota2.replace('⚠️','').replace('ℹ️','').strip()}", _nota_style2))
            if _tk_override2 and _tk_override2.get("ejemplos"):
                story.append(Spacer(1,0.04*inch))
                story.append(Paragraph(
                    f"Rangos de referencia MX: {_tk_override2['ejemplos']}", s["small"]))
        else:
            story.append(Paragraph(ticket_data['interpretacion_prem'], s["normal"]))
            story.append(Spacer(1,0.04*inch))
            story.append(Paragraph(f"<b>Estrategia sugerida:</b> {ticket_data['posicionamiento'].capitalize()}.", s["normal"]))

    # ── P5: DEMOGRAFÍA ──
    story.append(PageBreak())
    story.extend(_seccion_pdf("Perfil Demográfico del Área", "👥", s))
    if demografia:
        fuente = demografia.get('fuente','Estimación por zona')
        story.append(Paragraph(f"📊 {fuente}", s["small"]))
        if GRAFICAS_OK:
            try:
                story.append(RLImage(grafica_demografia(demografia), width=5.5*inch, height=2.6*inch))
            except:
                pass
        año_act = datetime.datetime.now().year
        densidad_info = clasificar_densidad(demografia.get('densidad_actual',8000)) if MODULOS_OK else formatear_densidad(demografia.get('densidad_hab_km2',8000))
        dem_rows = [
            [f"Población {año_act}", f"{demografia.get('poblacion_actual', demografia.get('poblacion_estimada',0)):,} hab",
             f"Censo 2020: {demografia.get('poblacion_2020','-'):,}" if demografia.get('poblacion_2020') else ""],
            ["Viviendas habitadas", f"{demografia.get('viviendas_actual', demografia.get('viviendas_habitadas',0)):,}", ""],
            ["Densidad urbana", densidad_info.get('nivel', densidad_info.get('clasificacion','')),
             f"~{demografia.get('personas_manzana', densidad_info.get('personas_manzana',0))} personas/manzana"],
            ["NSE predominante", demografia.get('nse_predominante','C'), "AMAI"],
            ["Ingreso prom/mes", f"${demografia.get('ingreso_actual', demografia.get('ingreso_promedio_mensual',0)):,}", "Proyectado 2026"],
            ["Gasto prom/mes", f"${demografia.get('gasto_actual', demografia.get('gasto_promedio_mensual',0)):,}", "78% del ingreso (ENIGH)"],
            ["Crecimiento anual", f"{demografia.get('tasa_crecimiento_pct',0.55):.2f}%", "CONAPO 2020-2030"],
        ]
        story.append(_tabla_pdf(dem_rows, [2*inch, 1.8*inch, 2.2*inch], s))
        # Gráfica edad/género PREMIUM
        if GRAFICAS_OK and "grafica_edad_genero" in dir():
            try:
                _geg_p = grafica_edad_genero(demografia, tipo_info_efectivo.get('nombre', tipo_info.get('nombre','')))
                story.append(Spacer(1,0.08*inch))
                story.append(RLImage(_geg_p, width=5.3*inch, height=2.6*inch))
            except: pass
        if MODULOS_OK and "interpretar_demografia_negocio" in dir():
            try:
                _idem_p = interpretar_demografia_negocio(demografia,
                    tipo_info_efectivo.get('nombre', tipo_info.get('nombre','')),
                    tipo_negocio_efectivo or tipo_negocio or "")
                story.append(Spacer(1,0.08*inch))
                for _ln2 in _idem_p.split("\n\n"):
                    if _ln2.strip():
                        _lf2 = _ln2.strip()
                        import re as _re2
                        _lf2 = _re2.sub(r'\*\*(.+?)\*\*', r'<b></b>', _lf2)
                        story.append(Paragraph(_lf2, s["normal"]))
                        story.append(Spacer(1,0.04*inch))
            except: pass

    # ── P6: TRÁFICO ──
    story.append(PageBreak())
    story.extend(_seccion_pdf("Análisis de Tráfico Estimado", "🚗", s))
    if trafico_data and GRAFICAS_OK:
        story.append(RLImage(grafica_trafico_horario(
            trafico_data["trafico_horario"], tipo_info.get('nombre','')),
            width=5.5*inch, height=2.8*inch))
        story.append(RLImage(grafica_trafico_semanal(trafico_data["trafico_semanal"]),
            width=5.5*inch, height=1.8*inch))
        pico_rows = [["Horario","Nivel","Recomendación"]]
        for p in trafico_data["horas_pico"]:
            pico_rows.append([p["hora_str"], f"{p['nivel']} ({p['flujo']}%)",
                "Máx. atención" if p["flujo"]>=80 else "Alta actividad"])
        story.append(_tabla_pdf(pico_rows, [1.8*inch,1.8*inch,2.4*inch], s))
        # Narrativa tráfico PREMIUM
        _t_nom_p = tipo_info_efectivo.get('nombre', tipo_info.get('nombre','')) if 'tipo_info_efectivo' in dir() else tipo_info.get('nombre','')
        if trafico_data:
            narr_traf_p = generar_narrativa_seccion("trafico", {
                "nivel_general": trafico_data.get('nivel_general',''),
                "horas_pico": [p['hora_str'] for p in trafico_data.get('horas_pico',[])[:2]],
                "mejor_dia": trafico_data.get('dia_pico',''),
                "dia_bajo": trafico_data.get('dia_bajo',''),
            }, _t_nom_p, idioma)
            for el in _bloque_narrativa_pdf(narr_traf_p, s):
                story.append(el)

    # ── P7: MERCADO + FORECAST ──
    story.append(PageBreak())
    story.extend(_seccion_pdf("Tamaño de Mercado Potencial", "📊", s))
    if mercado_data and GRAFICAS_OK:
        story.append(RLImage(grafica_mercado_donut(mercado_data), width=5*inch, height=3*inch))
        story.append(Paragraph(f"📐 {mercado_data['metodologia']}", s["small"]))
        # Narrativa mercado PREMIUM
        narr_merc_p = generar_narrativa_seccion("mercado", {
            "mercado_total_mensual": mercado_data.get('mercado_total_mensual',0),
            "mercado_captura_mensual": mercado_data.get('mercado_captura_mensual',0),
            "factor_captura_pct": mercado_data.get('factor_captura_pct',0),
            "clientes_dia": mercado_data.get('clientes_dia_estimados',0),
        }, _t_nom_p if '_t_nom_p' in dir() else '', idioma)
        for el in _bloque_narrativa_pdf(narr_merc_p, s):
            story.append(el)

    story.append(Spacer(1,0.15*inch))
    story.extend(_seccion_pdf("Forecast de Ventas — 12 Meses", "📈", s))
    if forecast_data and GRAFICAS_OK:
        story.append(RLImage(grafica_forecast(forecast_data), width=5.5*inch, height=3*inch))
        fc_rows = [["Escenario","Mes 1","Mes 9","Total año 1","Prom/mes"]]
        for key, emoji in [("pesimista","🔴"),("base","🔵"),("optimista","🟢")]:
            e = forecast_data["escenarios"][key]
            fc_rows.append([f"{emoji} {key.capitalize()}",
                f"${e['ventas_mensuales'][0]/1000:.0f}K",
                f"${e['ventas_mensuales'][8]/1000:.0f}K",
                f"${e['total_anual']/1000:.0f}K",
                f"${e['promedio_mensual']/1000:.0f}K"])
        tbl_fc = Table(fc_rows, colWidths=[1.2*inch,1*inch,1*inch,1.2*inch,1.2*inch])
        tbl_fc.setStyle(TableStyle([
            ('BACKGROUND',(0,0),(-1,0),colors.HexColor('#6B6B6B')),
            ('TEXTCOLOR', (0,0),(-1,0),colors.white),
            ('FONTNAME',  (0,0),(-1,0),'Helvetica-Bold'),
            ('FONTSIZE',  (0,0),(-1,-1),9),
            ('ALIGN',     (1,0),(-1,-1),'CENTER'),
            ('GRID',(0,0),(-1,-1),0.3,colors.HexColor('#DDDDDD')),
            ('ROWBACKGROUNDS',(0,1),(-1,-1),[
                colors.HexColor('#FFEBEE'),
                colors.HexColor('#E3F2FD'),
                colors.HexColor('#E8F5E9')]),
            ('TOPPADDING',(0,0),(-1,-1),6),('BOTTOMPADDING',(0,0),(-1,-1),6),
        ]))
        story.append(tbl_fc)
        # Narrativa forecast PREMIUM
        esc_p = forecast_data.get('escenarios',{})
        narr_fc_p = generar_narrativa_seccion("forecast", {
            "total_pesimista": esc_p.get('pesimista',{}).get('total_anual',0),
            "total_base": esc_p.get('base',{}).get('total_anual',0),
            "total_optimista": esc_p.get('optimista',{}).get('total_anual',0),
            "promedio_mensual_base": esc_p.get('base',{}).get('promedio_mensual',0),
        }, _t_nom_p if '_t_nom_p' in dir() else '', idioma)
        for el in _bloque_narrativa_pdf(narr_fc_p, s):
            story.append(el)

    # ── P8: ROI ──
    story.append(PageBreak())
    story.extend(_seccion_pdf("ROI + Punto de Equilibrio", "💰", s))
    story.append(Paragraph(
        "Análisis de retorno sobre inversión basado en el forecast de ventas "
        "y estructura de costos estándar para el tipo de negocio seleccionado.", s["normal"]))
    if roi_data and forecast_data and GRAFICAS_OK:
        story.append(RLImage(grafica_roi_dashboard(roi_data, forecast_data),
            width=6*inch, height=3*inch))
        roi_rows = [
            ["ROI estimado 12 meses", f"{roi_data['roi_12m_pct']}%", roi_data['clasificacion_roi']],
            ["Meses de recuperación est.", f"{roi_data['meses_recuperacion']} meses", "Escenario base"],
            ["Utilidad mensual estimada", f"${roi_data['utilidad_mensual_est']:,}", "35% margen bruto"],
            ["Punto de equilibrio", f"${roi_data['punto_eq_ventas_mes']:,}/mes", "Ventas mínimas para no perder"],
            ["Inversión mínima", f"${roi_data['inversion_min']:,}", tipo_info.get('nombre','')],
        ]
        story.append(_tabla_pdf(roi_rows, [2.2*inch,1.5*inch,2.3*inch], s))
        story.append(Paragraph("⚠️ Estimaciones con supuestos estándar. Los resultados dependen de gestión, producto y mercado real.", s["small"]))
        # Narrativa ROI PREMIUM
        if roi_data:
            narr_roi = generar_narrativa_seccion("roi", {
                "roi_12m_pct": roi_data.get('roi_12m_pct',0),
                "meses_recuperacion": roi_data.get('meses_recuperacion',0),
                "utilidad_mensual": roi_data.get('utilidad_mensual_est',0),
                "clasificacion": roi_data.get('clasificacion_roi',''),
            }, tipo_info_efectivo.get('nombre', tipo_info.get('nombre','')), idioma)
            for el in _bloque_narrativa_pdf(narr_roi, s):
                story.append(el)

    # ── P9: COMPARATIVA DE UBICACIONES ──
    story.append(PageBreak())
    story.extend(_seccion_pdf("Comparativa de Ubicaciones", "🗺️", s))

    # Datos de la ubicación principal
    ubicaciones_comp = [{
        "nombre": (ubicacion or "Ubicación principal")[:40],
        "score": score,
        "densidad_comp": _dsg_get(desglose, 'densidad', 0) if desglose else 0,
        "calidad_comp":  _dsg_get(desglose, 'calidad', 0) if desglose else 0,
        "consolidacion": _dsg_get(desglose, 'consolidacion', 0) if desglose else 0,
        "poblacion":     (demografia or {}).get('poblacion_actual', (demografia or {}).get('poblacion_estimada', 0)),
        "nse":           (demografia or {}).get('nse_predominante', 'C'),
        "ingreso":       (demografia or {}).get('ingreso_actual', (demografia or {}).get('ingreso_promedio_mensual', 0)),
        "num_competidores": len(competidores or []),
    }]
    # Agregar ubicaciones comparadas si existen en session (se pasan como parámetro extra)
    # Por ahora mostrar tabla con la ubicación principal y un espacio para las comparadas
    comp_rows = [["Indicador", (ubicacion or "Ubicación principal")[:25]]]
    comp_rows += [
        ["Score viabilidad",    f"{score}/100"],
        ["NSE predominante",    (demografia or {}).get('nse_predominante','C')],
        ["Competidores 500m",   str(len(competidores or []))],
        ["Densidad comp.",      f"{_dsg_get(desglose,'densidad')}/100" if desglose else "N/A"],
        ["Calidad comp.",       f"{_dsg_get(desglose,'calidad')}/100" if desglose else "N/A"],
        ["Consolidación",       f"{_dsg_get(desglose,'consolidacion')}/100" if desglose else "N/A"],
        ["Ingreso prom. área",  f"${(demografia or {}).get('ingreso_actual', (demografia or {}).get('ingreso_promedio_mensual',0)):,}"],
    ]
    tbl_comp = Table(comp_rows, colWidths=[2.2*inch, 3.3*inch])
    tbl_comp.setStyle(TableStyle([
        ('BACKGROUND', (0,0),(-1,0), colors.HexColor('#6B6B6B')),
        ('TEXTCOLOR',  (0,0),(-1,0), colors.white),
        ('FONTNAME',   (0,0),(-1,0), 'Helvetica-Bold'),
        ('FONTSIZE',   (0,0),(-1,-1), 9),
        ('FONTNAME',   (0,0),(0,-1), 'Helvetica-Bold'),
        ('TEXTCOLOR',  (0,1),(0,-1), colors.HexColor('#6B6B6B')),
        ('ROWBACKGROUNDS',(0,1),(-1,-1),[colors.white, colors.HexColor('#F0F4FF')]),
        ('GRID',(0,0),(-1,-1),0.3,colors.HexColor('#DDDDDD')),
        ('TOPPADDING',(0,0),(-1,-1),6),('BOTTOMPADDING',(0,0),(-1,-1),6),
    ]))
    story.append(tbl_comp)
    story.append(Spacer(1, 0.1*inch))
    story.append(Paragraph(
        "💡 Para comparar ubicaciones adicionales, usa la función de comparativa en la aplicación. "
        "El análisis compara score, demografía y competencia entre todas las ubicaciones seleccionadas.",
        ParagraphStyle('comp_nota', parent=s["base"]['Normal'], fontSize=8.5,
        textColor=colors.HexColor('#555'), fontName='Helvetica-Oblique')))

    if GRAFICAS_OK and len(ubicaciones_comp) >= 1:
        try:
            graf_comp_pdf = grafica_comparativa(ubicaciones_comp)
            if graf_comp_pdf:
                story.append(Spacer(1, 0.1*inch))
                story.append(RLImage(graf_comp_pdf, width=5.5*inch, height=3.5*inch))
        except:
            pass

    # ── P10: Recomendación Final ──
    story.append(PageBreak())
    _plaza_block_pdf(story, contexto, s)
    _render_recomendacion_pdf(story, analisis, s)

    # ── P11: CTA ──
    story.append(PageBreak())
    _cta_block(story, tier, t, s, idioma)

    doc.build(story, onFirstPage=_footer_pdf, onLaterPages=_footer_pdf)
    buffer.seek(0)
    if _ok_p:
        try:
            from pypdf import PdfWriter as _PW2, PdfReader as _PR2
            _w2=_PW2(); _w2.append(_PR2(_portada_buf_p)); _w2.append(_PR2(buffer))
            _o2=BytesIO(); _w2.write(_o2); _o2.seek(0); return _o2
        except: pass
    return buffer



def generar_pdf_por_tier(tier_key, **kwargs):
    """Router central de generación de PDF"""
    if tier_key == "basic":
        return generar_pdf_basic(**kwargs)
    elif tier_key == "pro":
        return generar_pdf_pro(**kwargs)
    elif tier_key == "premium":
        return generar_pdf_premium(**kwargs)
    else:
        return generar_pdf_free(**kwargs)

# ============================================
# INTERFAZ STREAMLIT
# ============================================

# Selector idioma
idioma = st.selectbox(
    TEXTOS["es"]["selector_idioma"],
    options=["es", "en"],
    format_func=lambda x: "Español 🇲🇽" if x == "es" else "English 🇺🇸"
)
t = TEXTOS[idioma]

# Header
st.markdown(f"""
<div style='text-align:center; padding:20px;'>
    <h1 style='color:#0D0D0D; font-size:52px; margin:0; letter-spacing:-2px;'>4SITE</h1>
    <p style='color:#BB944F; font-size:18px; margin:5px 0; font-style:italic;'>Don't guess. Foresee.</p>
</div>
""", unsafe_allow_html=True)

st.markdown("---")

# ── SISTEMA DE CÓDIGOS ──────────────────────
st.markdown(f"### {t['ingresa_codigo']}")

col_code, col_btn = st.columns([3, 1])
with col_code:
    codigo_input = st.text_input("", placeholder=t["placeholder_codigo"],
                                  key="codigo_input", label_visibility="collapsed")
with col_btn:
    aplicar = st.button(t["validar_codigo"], use_container_width=True)

# Estado del tier en session_state
if "tier_activo" not in st.session_state:
    st.session_state.tier_activo = "free"

if aplicar and codigo_input:
    resultado = validar_codigo(codigo_input)
    if resultado is None:
        st.error(t["codigo_invalido"])
    elif resultado == "agotado":
        st.error("❌ Este código ya alcanzó el límite de 3 análisis. Adquiere un nuevo código en hola@4site.mx")
    else:
        st.session_state.tier_activo = resultado
        st.session_state.codigo_activo = codigo_input.strip().upper()
        activar_codigo(codigo_input)
        tier_info = TIERS[resultado]
        _usos_rest = usos_restantes(codigo_input)
        _usos_txt = f"· {_usos_rest} análisis restante{'s' if _usos_rest != 1 else ''}" if not (codigo_input.strip().upper().endswith("-TEST") or codigo_input.strip().upper() == "TEST") else ""
        st.success(f"{t['codigo_valido']} **{tier_info['nombre']}** — {tier_info['precio']} · {tier_info['pdf_paginas']} {_usos_txt}")

# Botón para continuar sin código
if st.button(t["sin_codigo"], use_container_width=False):
    st.session_state.tier_activo = "free"

# Mostrar tier activo
tier_key   = st.session_state.tier_activo
tier_info  = TIERS[tier_key]
tier_color = tier_info["color"]

st.markdown(f"""
<div style='padding:10px 16px; background:{tier_color}18; border-left:4px solid {tier_color};
     border-radius:6px; margin:8px 0;'>
    <b style='color:{tier_color};'>{t['tier_activo']} {tier_info['nombre']} — {tier_info['precio']}</b>
    &nbsp;·&nbsp; {tier_info['pdf_paginas']}
    &nbsp;·&nbsp; Top {tier_info['top_negocios']} negocios
    &nbsp;·&nbsp; {tier_info['competidores_mostrados']} competidores
    {'&nbsp;·&nbsp; ✅ Demografía' if tier_info['muestra_demografia'] else ''}
    {'&nbsp;·&nbsp; ✅ Desglose score' if tier_info['muestra_score_desglose'] else ''}
</div>
""", unsafe_allow_html=True)

st.markdown("---")

# ── INPUT DE UBICACIÓN ──────────────────────
st.markdown(f"### {t['input_ubicacion']}")
direccion_texto = st.text_input("", placeholder=t["placeholder_direccion"],
                                 key="direccion_input", label_visibility="collapsed")

# ── Detección automática de tipo de plaza ─────────────────────
_tipo_plaza_detectado = detectar_tipo_plaza(direccion_texto or "")
tipo_plaza = None  # se definirá abajo

# PRO y PREMIUM: selector siempre visible (con default inteligente de auto-detección)
if tier_key in ["pro", "premium", "premium2"]:

    _default_plaza = {
        "mall":          "mall",
        "strip_center":  "strip_center",
        "plaza_ambigua": "strip_center",
    }.get(_tipo_plaza_detectado, "calle")

    if _tipo_plaza_detectado in ["mall", "strip_center", "plaza_ambigua"]:
        _det_labels = {
            "mall":          "🏬 Detectamos que tu dirección corresponde a un mall o plaza cerrada.",
            "strip_center":  "🏪 Detectamos que tu dirección corresponde a un strip center o plaza abierta.",
            "plaza_ambigua": "🏢 Detectamos que tu dirección puede estar dentro de una plaza — confirma el tipo.",
        }
        st.markdown(f"""<div style='background:#1A1A1A; border-left:4px solid #F39C12;
            border-radius:6px; padding:8px 14px; margin:2px 0 6px;'>
            <span style='color:#F39C12; font-size:12px;'>{_det_labels[_tipo_plaza_detectado]}
            Confirma o ajusta abajo.</span></div>""", unsafe_allow_html=True)

    st.markdown("**¿Esta ubicación está dentro de una plaza comercial?**")
    _col_pz1, _col_pz2 = st.columns([3, 2])
    with _col_pz1:
        _tipo_plaza_sel = st.radio(
            "",
            options=["calle", "strip_center", "mall"],
            format_func=lambda x: {
                "calle":        "🚶 No — está en calle, avenida o edificio normal",
                "strip_center": "🏪 Sí — Plaza abierta / Strip center (locales con cajones al frente)",
                "mall":         "🏬 Sí — Mall / Plaza cerrada (tipo Galerías, Antara, Perisur)",
            }.get(x, x),
            index=["calle", "strip_center", "mall"].index(_default_plaza),
            key="tipo_plaza_radio_pro",
            horizontal=False,
            label_visibility="collapsed",
        )
    tipo_plaza = None if _tipo_plaza_sel == "calle" else _tipo_plaza_sel

    with _col_pz2:
        if tipo_plaza == "mall":
            st.markdown("""<div style='background:#1A1A1A; border-left:3px solid #BB944F;
                border-radius:6px; padding:8px 12px; font-size:11px; color:#AAAAAA;'>
                <b style='color:#BB944F;'>Mall cerrado</b><br>
                Amplía radio de competencia a 1km, aplica lógica de tráfico captivo
                y ajusta ROI por renta elevada.</div>""", unsafe_allow_html=True)
        elif tipo_plaza == "strip_center":
            st.markdown("""<div style='background:#1A1A1A; border-left:3px solid #1A76D2;
                border-radius:6px; padding:8px 12px; font-size:11px; color:#AAAAAA;'>
                <b style='color:#1A76D2;'>Strip center / plaza abierta</b><br>
                Visibilidad desde la calle, tráfico mixto paso + destino.
                La tienda ancla tiene alto impacto.</div>""", unsafe_allow_html=True)
        else:
            st.markdown("""<div style='background:#1A1A1A; border-left:3px solid #555;
                border-radius:6px; padding:8px 12px; font-size:11px; color:#AAAAAA;'>
                <b style='color:#CCC;'>Calle / avenida</b><br>
                Análisis estándar: visibilidad de paso, vialidad y demografía del entorno.
                </div>""", unsafe_allow_html=True)

else:
    # FREE y BÁSICO: solo auto-detección
    if _tipo_plaza_detectado == "mall":
        tipo_plaza = "mall"
        st.markdown("""<div style='background:#1A1A1A; border-left:4px solid #BB944F;
            border-radius:6px; padding:10px 14px; margin:4px 0 8px;'>
            <span style='color:#BB944F; font-weight:700;'>🏬 Mall / Plaza Cerrada detectado</span><br>
            <span style='color:#AAAAAA; font-size:12px;'>
            El análisis aplicará lógica especial para plazas cerradas.</span>
            </div>""", unsafe_allow_html=True)
    elif _tipo_plaza_detectado == "strip_center":
        tipo_plaza = "strip_center"
        st.markdown("""<div style='background:#1A1A1A; border-left:4px solid #1A76D2;
            border-radius:6px; padding:10px 14px; margin:4px 0 8px;'>
            <span style='color:#1A76D2; font-weight:700;'>🏪 Strip Center / Plaza Abierta detectada</span><br>
            <span style='color:#AAAAAA; font-size:12px;'>
            Análisis con lógica de plaza abierta.</span>
            </div>""", unsafe_allow_html=True)
    elif _tipo_plaza_detectado == "plaza_ambigua" and direccion_texto:
        st.markdown("""<div style='background:#1A1A1A; border-left:4px solid #F39C12;
            border-radius:6px; padding:10px 14px; margin:4px 0 8px;'>
            <span style='color:#F39C12; font-weight:700;'>🏢 ¿Qué tipo de plaza es?</span><br>
            <span style='color:#AAAAAA; font-size:12px;'>
            Detectamos que tu ubicación está dentro de una plaza.</span>
            </div>""", unsafe_allow_html=True)
        tipo_plaza = st.radio(
            "",
            options=["strip_center", "mall", "no_plaza"],
            format_func=lambda x: {
                "strip_center": "🏪 Plaza abierta / Strip center",
                "mall":         "🏬 Mall / Plaza cerrada",
                "no_plaza":     "🚶 No es plaza — calle normal",
            }.get(x, x),
            key="tipo_plaza_radio",
            horizontal=False,
        )
        if tipo_plaza == "no_plaza":
            tipo_plaza = None

# ── Info adicional para MALL (todos los tiers) ──────────────────────────────
if tipo_plaza == "mall":
    st.markdown("<div style='margin-top:6px'></div>", unsafe_allow_html=True)
    col_m1, col_m2 = st.columns(2)
    with col_m1:
        zona_mall = st.selectbox(
            "📍 Zona dentro del mall:",
            options=["no_especificado","food_court","moda_accesorios","servicios",
                     "ancla_principal","planta_baja_entrada","piso_superior"],
            format_func=lambda x: {
                "no_especificado":    "No sé / no especificado",
                "food_court":         "🍽️ Food court / zona de comida",
                "moda_accesorios":    "👗 Moda / accesorios / lifestyle",
                "servicios":          "💈 Servicios (salud, belleza, etc.)",
                "ancla_principal":    "⚓ Cerca de ancla (Liverpool, Walmart, Cine)",
                "planta_baja_entrada":"🚪 Planta baja / entrada principal",
                "piso_superior":      "⬆️ Piso superior / menor tráfico",
            }.get(x, x),
            key="zona_mall_select",
        )
    with col_m2:
        nivel_renta_mall = st.selectbox(
            "💰 Nivel de renta estimada:",
            options=["desconocido","bajo","medio","alto","muy_alto"],
            format_func=lambda x: {
                "desconocido": "No sé / no tengo dato",
                "bajo":        "💚 Baja — < $400/m² mes",
                "medio":       "🟡 Media — $400–$800/m² mes",
                "alto":        "🟠 Alta — $800–$1,500/m² mes",
                "muy_alto":    "🔴 Muy alta — > $1,500/m² mes",
            }.get(x, x),
            key="nivel_renta_mall",
        )
    # zona_mall y nivel_renta_mall ya están en session_state via sus widget keys

elif tipo_plaza == "strip_center":
    st.markdown("<div style='margin-top:6px'></div>", unsafe_allow_html=True)
    _col_sc1, _col_sc2 = st.columns(2)
    with _col_sc1:
        ancla_strip = st.selectbox(
            "⚓ ¿Hay una tienda ancla en la plaza?",
            options=["sin_ancla","walmart","chedraui","soriana","oxxo","farmacias","banco","otra"],
            format_func=lambda x: {
                "sin_ancla":  "Sin ancla identificable",
                "walmart":    "🛒 Walmart / Bodega Aurrerá / Sam's",
                "chedraui":   "🛒 Chedraui / Superama",
                "soriana":    "🛒 Soriana / Comercial Mexicana",
                "oxxo":       "🏪 OXXO / tienda de conveniencia",
                "farmacias":  "💊 Farmacia cadena (Similares, Benavides, etc.)",
                "banco":      "🏦 Banco / cajero",
                "otra":       "Otra tienda ancla",
            }.get(x, x),
            key="ancla_strip_select",
        )
    with _col_sc2:
        _tips_ancla = {
            "walmart":   "🛒 Flujo familiar constante — ideal para comida rápida, farmacia y servicios.",
            "chedraui":  "🛒 NSE ligeramente más alto — funciona bien para café y papelería.",
            "soriana":   "🛒 Buen flujo de compra semanal, similar a Chedraui.",
            "oxxo":      "🏪 Muy alto tráfico de paso — ideal para comida y servicios rápidos.",
            "farmacias": "💊 Clientela recurrente y de confianza — excelente para salud y nutrición.",
            "banco":     "🏦 Genera cola de espera — clientes disponibles frente al local.",
            "sin_ancla": "Sin ancla: el tráfico depende 100% de visibilidad exterior.",
            "otra":      "Evalúa qué clientela atrae el ancla y si es compatible con tu negocio.",
        }
        if ancla_strip in _tips_ancla:
            st.markdown(f"""<div style='background:#1A1A1A; border-left:3px solid #1A76D2;
                border-radius:6px; padding:8px 12px; font-size:11px; color:#AAAAAA; margin-top:22px;'>
                {_tips_ancla[ancla_strip]}</div>""", unsafe_allow_html=True)
    # ancla_strip ya está en session_state via su widget key

# Guardar tipo_plaza en session_state para usarlo en el análisis
st.session_state["tipo_plaza"] = tipo_plaza

# Mapa interactivo Folium
st.markdown(f"**{'O haz click en el mapa:' if idioma == 'es' else 'Or click on the map:'}**")
try:
    import folium
    from streamlit_folium import st_folium

    m = folium.Map(location=[19.432608, -99.133209], zoom_start=12)
    m.add_child(folium.LatLngPopup())
    mapa_data = st_folium(m, width=700, height=350, key="mapa_folium")

    lat_mapa, lng_mapa, direccion_mapa = None, None, None
    if mapa_data and mapa_data.get("last_clicked"):
        lat_mapa = mapa_data["last_clicked"]["lat"]
        lng_mapa = mapa_data["last_clicked"]["lng"]
        direccion_mapa = geocodificar_inversa(lat_mapa, lng_mapa)
        st.success(f"📍 {'Seleccionado' if idioma == 'es' else 'Selected'}: {direccion_mapa}")
except ImportError:
    st.caption("💡 `pip install folium streamlit-folium` para habilitar el mapa interactivo")
    lat_mapa, lng_mapa, direccion_mapa = None, None, None

st.markdown("---")

# Modo de análisis — "Descubrir negocio" solo desde Basic+
st.markdown(f"### {t['modo_analisis']}")
if tier_key == "free":
    modo = "validar"
    st.radio("", options=["validar"],
             format_func=lambda x: t["modo_validar"],
             horizontal=True)
    st.caption("🔒 Modo **Descubrir qué negocio es mejor** disponible desde plan Básico ($99 MXN)")
else:
    if tier_key == "premium2":
        # ZONA ELITE: validar + descubrir + identificar zona potencial
        modo = st.radio("", options=["validar", "recomendar", "zona_potencial"],
                        format_func=lambda x: {
                            "validar":        t["modo_validar"],
                            "recomendar":     t["modo_recomendar"],
                            "zona_potencial": "🔍 Identificar zona potencial",
                        }.get(x, x),
                        horizontal=True,
                        key="modo_elite")
        if modo == "zona_potencial":
            # En zona_potencial se ocultan validar/descubrir — banner informativo
            st.markdown("""<div style='background:#0D0D0D; border-radius:10px;
                padding:12px 18px; margin:6px 0;'>
                <span style='color:#BB944F; font-weight:700;'>🔍 Identificar Zona Potencial</span><br>
                <span style='color:#AAAAAA; font-size:12px;'>
                Ingresa un punto de referencia y encontraremos las mejores zonas
                para tu negocio en el radio que elijas.</span>
                </div>""", unsafe_allow_html=True)
    elif tier_key == "premium":
        # PREMIUM: validar + descubrir + descubrir zona (nueva función)
        modo = st.radio("", options=["validar", "recomendar", "zona_potencial"],
                        format_func=lambda x: {
                            "validar":        t["modo_validar"],
                            "recomendar":     t["modo_recomendar"],
                            "zona_potencial": "🔍 Descubrir zona potencial",
                        }.get(x, x),
                        horizontal=True,
                        key="modo_premium")
        if modo == "zona_potencial":
            st.markdown("""<div style='background:#0D0D0D; border-radius:10px;
                padding:12px 18px; margin:6px 0;'>
                <span style='color:#D4AF37; font-weight:700;'>🔍 Descubrir Zona Potencial — PREMIUM 💎</span><br>
                <span style='color:#AAAAAA; font-size:12px;'>
                Ingresa un punto de referencia y encontraremos las mejores zonas
                para tu negocio en el radio que elijas.</span>
                </div>""", unsafe_allow_html=True)
    elif tier_key == "pro":
        modo = st.radio("", options=["validar", "recomendar"],
                        format_func=lambda x: t["modo_validar"] if x == "validar" else t["modo_recomendar"],
                        horizontal=True)
    else:
        modo = st.radio("", options=["validar", "recomendar"],
                        format_func=lambda x: t["modo_validar"] if x == "validar" else t["modo_recomendar"],
                        horizontal=True)

tipo_negocio_seleccionado = None
if modo == "zona_potencial":
    # ── ZONA POTENCIAL TEMPORALMENTE BLOQUEADA ──
    st.markdown("---")
    st.markdown("""
    <div style='background:#1A1A1A; border:2px solid #BB944F; border-radius:12px;
         padding:28px 32px; margin:12px 0; text-align:center;'>
        <div style='font-size:32px; margin-bottom:10px;'>🔧</div>
        <div style='color:#BB944F; font-weight:700; font-size:18px; margin-bottom:8px;'>
            Zona Potencial — En mantenimiento</div>
        <div style='color:#CCCCCC; font-size:13px; line-height:1.6; margin-bottom:16px;'>
            Estamos mejorando la detección de zonas y el reporte PowerPoint.<br>
            Esta función estará disponible muy pronto con mayor precisión y estabilidad.
        </div>
        <div style='color:#888; font-size:11px;'>
            Mientras tanto usa los modos <b style="color:#CCC">Validar negocio</b> o
            <b style="color:#CCC">Descubrir negocio</b> para analizar ubicaciones específicas.
        </div>
    </div>""", unsafe_allow_html=True)
    st.stop()

if modo == "validar":
    st.markdown(f"**{t['selecciona_tipo']}**")
    tipo_negocio_seleccionado = st.selectbox("", options=list(TIPOS_NEGOCIO.keys()),
                                              format_func=lambda x: TIPOS_NEGOCIO[x]["nombre"])

st.markdown("---")

# ── BOTÓN ANALIZAR ──────────────────────────
if st.button(t["boton_analizar"], type="primary", use_container_width=True):

    # Prioridad: click en mapa > texto
    if 'lat_mapa' in dir() and lat_mapa and lng_mapa:
        lat, lng  = lat_mapa, lng_mapa
        ubicacion = direccion_mapa or f"{lat:.6f}, {lng:.6f}"
    elif direccion_texto:
        lat, lng  = geocodificar_direccion(direccion_texto)
        ubicacion = direccion_texto
    else:
        st.error("Por favor ingresa una ubicación o haz click en el mapa." if idioma == "es"
                 else "Please enter a location or click on the map.")
        st.stop()

    if not lat or not lng:
        st.error("No se pudo geocodificar la ubicación." if idioma == "es" else "Could not geocode the location.")
        st.stop()

    # Recuperar datos de tipo de plaza del session_state
    _tipo_plaza  = st.session_state.get("tipo_plaza", None)
    _zona_mall   = st.session_state.get("zona_mall_select", "no_especificado")
    _renta_mall  = st.session_state.get("nivel_renta_mall", "desconocido")
    _ancla_strip = st.session_state.get("ancla_strip_select", "sin_ancla")

    with st.spinner("🔍 Analizando ubicación..."):

        # Obtener demografía primero (no depende del tipo de negocio)
        demografia = obtener_demografia(lat, lng, ubicacion)

        # Para mall: usar NSE del mall si es conocido (override)
        if _tipo_plaza == "mall":
            # Los malls de alto nivel elevan el NSE efectivo del análisis
            nse_actual = demografia.get("nse_predominante","C")
            if _renta_mall in ["alto","muy_alto"] and nse_actual not in ["A","A/B","B"]:
                demografia["nse_predominante"] = "B/C+"
                demografia["nse_override_mall"] = True

        if modo == "validar":
            competidores = buscar_competencia_por_tipo(
                lat, lng, tipo_negocio_seleccionado,
                tipo_plaza=_tipo_plaza
            )
            # Detectar contexto con tipo de negocio para motor vial preciso
            contexto = detectar_contexto_ubicacion(lat, lng, ubicacion, tipo_negocio_seleccionado)

            # Inyectar info de plaza en el contexto
            if _tipo_plaza:
                contexto["tipo_plaza"]  = _tipo_plaza
                contexto["zona_mall"]   = _zona_mall
                contexto["ancla_strip"] = _ancla_strip
                # Para mall: tipo_zona siempre es "comercial"
                if _tipo_plaza == "mall":
                    contexto["tipo_zona"] = "comercial"
                    contexto["badges"].append("🏬 Mall / Plaza Cerrada")
                elif _tipo_plaza == "strip_center":
                    contexto["tipo_zona"] = "comercial"
                    contexto["badges"].append("🏪 Strip Center")

            score_base, desglose = calcular_score_competencia(competidores, tipo_negocio_seleccionado)
            score_ctx  = ajustar_score_por_contexto(score_base, tipo_negocio_seleccionado, contexto)
            score_dem  = ajustar_score_por_demografia(score_ctx, tipo_negocio_seleccionado, demografia)

            # Ajuste final por tipo de plaza
            score, info_plaza = ajustar_score_por_plaza(
                score_dem, tipo_negocio_seleccionado,
                _tipo_plaza, _zona_mall, _renta_mall, _ancla_strip
            )
            # Guardar info plaza en contexto para reportes
            if _tipo_plaza:
                contexto["info_plaza"] = info_plaza

            analisis    = generar_analisis_claude(ubicacion, competidores, idioma, "validar",
                                                  tipo_negocio=tipo_negocio_seleccionado,
                                                  tipo_plaza=_tipo_plaza,
                                                  info_plaza=info_plaza if _tipo_plaza else None)
            recomendaciones = None
            tipo_negocio_top1 = tipo_negocio_seleccionado
        else:
            recomendaciones, contexto, demografia = recomendar_tipos_negocio(lat, lng, ubicacion)
            competidores = recomendaciones[0]["competidores"] if recomendaciones else []
            score        = recomendaciones[0]["score"] if recomendaciones else 0
            desglose     = recomendaciones[0]["desglose"] if recomendaciones else {}
            tipo_negocio_top1 = recomendaciones[0]["tipo_key"] if recomendaciones else None
            # Inyectar plaza en contexto para recomendar también
            if _tipo_plaza:
                contexto["tipo_plaza"] = _tipo_plaza
                if _tipo_plaza == "mall":
                    contexto["tipo_zona"] = "comercial"
                    contexto["badges"].append("🏬 Mall / Plaza Cerrada")
                    # Ajustar scores de recomendaciones también
                    for rec in (recomendaciones or []):
                        rec["score"], _ip = ajustar_score_por_plaza(
                            rec["score"], rec.get("tipo_key",""),
                            _tipo_plaza, _zona_mall, _renta_mall, "sin_ancla"
                        )
                        rec["score"] = max(0, min(100, rec["score"]))
                elif _tipo_plaza == "strip_center":
                    contexto["tipo_zona"] = "comercial"
                    contexto["badges"].append("🏪 Strip Center")
            # Re-detectar contexto con el tipo de negocio top1 para motor vial preciso
            if tipo_negocio_top1:
                contexto = detectar_contexto_ubicacion(lat, lng, ubicacion, tipo_negocio_top1)
                if _tipo_plaza:
                    contexto["tipo_plaza"] = _tipo_plaza
            analisis = generar_analisis_claude(ubicacion, [], idioma, "recomendar",
                                               recomendaciones=recomendaciones,
                                               tipo_plaza=_tipo_plaza)

    # Guardar resultados en session_state para evitar re-run al interactuar con mapas
    st.session_state.resultados = {
        "ubicacion": ubicacion, "lat": lat, "lng": lng,
        "score": score, "desglose": desglose, "analisis": analisis,
        "competidores": competidores, "contexto": contexto,
        "demografia": demografia, "recomendaciones": recomendaciones,
        "modo": modo, "tipo_negocio": tipo_negocio_seleccionado,
        "tipo_negocio_top1": tipo_negocio_top1, "idioma": idioma,
    }
    # META CAPI — ViewContent
    _meta_capi("ViewContent", content_name=f"{ubicacion[:50]} | score:{score:.0f}")

# ── Recuperar resultados de session_state si existen ──────────────
if "resultados" in st.session_state:
    _r = st.session_state.resultados
    ubicacion            = _r["ubicacion"]
    lat                  = _r["lat"]
    lng                  = _r["lng"]
    score                = _r["score"]
    desglose             = _r["desglose"]
    analisis             = _r["analisis"]
    competidores         = _r["competidores"]
    contexto             = _r["contexto"]
    demografia           = _r["demografia"]
    recomendaciones      = _r["recomendaciones"]
    modo                 = _r["modo"]
    tipo_negocio_seleccionado = _r["tipo_negocio"]
    tipo_negocio_top1    = _r.get("tipo_negocio_top1")
    idioma               = _r["idioma"]
    t                    = TEXTOS[idioma]
    tier_key             = st.session_state.tier_activo
    tier_info            = TIERS[tier_key]

    # ──────────────────────────────────────────
    # RESULTADOS EN PANTALLA
    # ──────────────────────────────────────────

    # Definir tipo_para_pro tempranamente para evitar NameError en secciones desplegables
    _r_safe = _r if '_r' in dir() and _r else {}
    tipo_para_pro = tipo_negocio_seleccionado if modo == "validar" else (
        tipo_negocio_top1 if 'tipo_negocio_top1' in dir() and tipo_negocio_top1 else
        (recomendaciones[0].get("tipo_key") if recomendaciones else None))

    st.markdown(f"## {t['titulo_analisis']}")
    if modo == "validar":
        st.markdown(f"### {TIPOS_NEGOCIO[tipo_negocio_seleccionado]['nombre']}")

    # ── ANÁLISIS VIAL — visible en todos los tiers ──
    av_data = contexto.get("analisis_vial")
    if av_data and av_data.get("descripcion"):
        dep_label = {"alto": "negocio de paso", "medio": "negocio mixto", "bajo": "negocio de destino"}.get(
            av_data.get("dependencia_paso","medio"), "negocio")
        color_vial = "#4CAF50" if av_data["ajuste_score"] > 0 else "#F44336" if av_data["ajuste_score"] < -5 else "#FFC107"
        st.markdown(f"""
        <div style='background:{color_vial}11; border-left:4px solid {color_vial};
             padding:10px 14px; border-radius:0 8px 8px 0; margin:8px 0;'>
            <b style='color:{color_vial};'>{av_data["badge"]}</b>
            &nbsp;·&nbsp; <span style='font-size:12px; color:#555;'>
            Impacto en score: <b style='color:{color_vial};'>
            {"+" if av_data["ajuste_score"] >= 0 else ""}{av_data["ajuste_score"]} pts</b>
            ({dep_label})</span><br>
            <span style='font-size:12px; color:#444;'>{av_data["descripcion"]}</span>
        </div>""", unsafe_allow_html=True)

    # Características de la zona (Basic+)
    if tier_info["muestra_badges"] and contexto.get("badges"):
        st.markdown("**📍 Características de la zona:**")
        BADGE_DESC = {
            "🚗 Alto tráfico vehicular": "La ubicación está sobre una vialidad de alto flujo — ideal para negocios de paso",
            "⛽ Gasolinera cercana":      "Presencia de gasolinera en 500m — atrae tráfico vehicular constante",
            "🏬 Plaza comercial":         "Plaza o centro comercial cercano — genera tráfico de shoppers",
            "🏥 Hospital cercano":        "Hospital o clínica en 500m — clientela cautiva y tráfico constante",
            "🎓 Escuela/Universidad":     "Institución educativa en 500m — alto flujo estudiantil",
            "🚌 Transporte público":      "Estación de transporte cercana — alta accesibilidad peatonal",
        }
        chips_html = "<div style='display:flex; flex-wrap:wrap; gap:8px; margin:8px 0;'>"
        # Filtrar badges de vialidad — ya se muestran en el bloque de análisis vial
        badges_vial = {"🛣️", "🚦", "🏘️", "📍"}
        badges_filtrados = [b for b in contexto["badges"]
                           if not any(b.startswith(v) for v in badges_vial)]
        for badge in badges_filtrados:
            desc = BADGE_DESC.get(badge, "Característica detectada en la zona")
            chips_html += f"""
            <div title='{desc}' style='
                background:#E3F2FD; border:1px solid #1976D2; border-radius:20px;
                padding:6px 14px; font-size:13px; color:#1565C0; cursor:help;
                font-weight:500;'>{badge}</div>"""
        chips_html += "</div>"
        chips_html += "<p style='font-size:11px; color:#888; margin:0;'>💡 Pasa el cursor sobre cada chip para ver su significado</p>"
        st.markdown(chips_html, unsafe_allow_html=True)
        st.markdown("") 

    # Score visual
    nivel, color = nivel_score(score, idioma)
    col1, col2, col3 = st.columns([1, 2, 1])
    with col2:
        st.markdown(f"""
        <div style='text-align:center; padding:25px; background:linear-gradient(135deg,{color}22 0%,{color}44 100%);
             border-radius:15px; border:3px solid {color};'>
            <h1 style='color:{color}; font-size:72px; margin:0;'>{score}</h1>
            <p style='color:#555; font-size:16px; margin:4px 0;'>{t['score_viabilidad']}</p>
            <b style='color:{color}; font-size:18px;'>{nivel}</b>
        </div>""", unsafe_allow_html=True)

    # Desglose score (Basic+)
    if tier_info["muestra_score_desglose"] and desglose:

        # ── Banner de plaza si aplica ──
        _ip = contexto.get("info_plaza", {}) if contexto else {}
        if _ip and _ip.get("tipo_plaza"):
            _pl_label = "🏬 Mall / Plaza Cerrada" if _ip["tipo_plaza"] == "mall" else "🏪 Strip Center / Plaza Abierta"
            _pl_color = "#BB944F" if _ip["tipo_plaza"] == "mall" else "#1A76D2"
            _pl_notas = _ip.get("nota", "").replace("\n", "<br>")
            _pl_roi   = _ip.get("ajuste_roi_pct", 0)
            st.markdown(f"""
            <div style='background:#1A1A1A; border-left:4px solid {_pl_color};
                 border-radius:8px; padding:12px 16px; margin:10px 0;'>
                <div style='color:{_pl_color}; font-weight:700; font-size:13px; margin-bottom:6px;'>
                    {_pl_label} — Impacto en el análisis</div>
                <div style='color:#CCCCCC; font-size:12px; line-height:1.6;'>{_pl_notas}</div>
                {"<div style='color:#F39C12; font-size:11px; margin-top:6px;'>⚠️ El ROI estimado ya incluye una reducción de ~" + str(_pl_roi) + "% por renta de plaza.</div>" if _pl_roi > 0 else ""}
            </div>""", unsafe_allow_html=True)

        st.markdown(f"#### {t['desglose_score']}")

        # Definir interpretaciones dinámicas
        def _interp_factor(valor, tipo):
            if tipo == "densidad":
                if valor >= 80: return "✅ Poca competencia — mercado libre"
                elif valor >= 60: return "🟡 Competencia moderada"
                elif valor >= 40: return "🟠 Zona saturada"
                else: return "🔴 Alta saturación"
            elif tipo == "calidad":
                if valor >= 80: return "✅ Competidores débiles — oportunidad"
                elif valor >= 60: return "🟡 Calidad media en la zona"
                elif valor >= 40: return "🟠 Competidores bien calificados"
                else: return "🔴 Competencia muy fuerte"
            elif tipo == "consolidacion":
                if valor >= 80: return "✅ Mercado joven — fácil de ganar"
                elif valor >= 60: return "🟡 Mercado en desarrollo"
                elif valor >= 40: return "🟠 Mercado establecido"
                else: return "🔴 Mercado muy consolidado"

        c1, c2, c3 = st.columns(3)
        with c1:
            st.metric(
                "Densidad de competencia (50%)",
                f"{desglose['densidad']}/100",
                help="Cuántos negocios similares hay en 500m. Mayor puntaje = menos competidores = más oportunidad."
            )
            st.caption(_interp_factor(desglose['densidad'], 'densidad'))
        with c2:
            st.metric(
                "Calidad de competencia (30%)",
                f"{desglose['calidad']}/100",
                help="Rating promedio de los competidores. Mayor puntaje = competidores con rating bajo = mercado insatisfecho."
            )
            st.caption(_interp_factor(desglose['calidad'], 'calidad'))
        with c3:
            st.metric(
                "Consolidación del mercado (20%)",
                f"{desglose['consolidacion']}/100",
                help="Volumen de reseñas de los competidores. Mayor puntaje = mercado joven con pocos reviews = más fácil de entrar."
            )
            st.caption(_interp_factor(desglose['consolidacion'], 'consolidacion'))
    elif not tier_info["muestra_score_desglose"]:
        st.caption("🔒 Desglose del score disponible en plan Básico ($99)")

    st.markdown("---")

    # ── RESUMEN / DIAGNÓSTICO según tier ──────────────────────────
    if tier_key == "free":
        # Free: resumen ejecutivo simple
        st.markdown("### 📋 Resumen Ejecutivo")
        if modo == "validar":
            st.markdown(analisis)
    else:
        # Basic+: Recomendación Final completo al FINAL (se agrega después)
        pass

    if modo == "validar" and tier_key != "free":
        _key_comp = "exp_comp_open"
        if _key_comp not in st.session_state: st.session_state[_key_comp] = False
        if st.button("Ver analisis de competencia" + (" ▲" if st.session_state[_key_comp] else " ▼"), key="btn_comp_exp", use_container_width=True):
            st.session_state[_key_comp] = not st.session_state[_key_comp]
        if st.session_state[_key_comp]:
            st.markdown(analisis)
    elif modo != "validar":
        # Recomendador — mostrar top N según tier
        top_n = tier_info["top_negocios"]
        st.markdown(f"#### {t['top_recomendaciones']}")
        for i, rec in enumerate(recomendaciones[:top_n]):
            nivel_r, color_r = nivel_score(rec['score'], idioma)
            _exp_label = re.sub(r'[^\x00-\x7F]','',rec['nombre']).strip()
            _key_rec = f"exp_rec_{i}"
            if _key_rec not in st.session_state:
                st.session_state[_key_rec] = (i == 0)
            _arrow = " ▲" if st.session_state[_key_rec] else " ▼"
            if st.button(f"#{i+1} {_exp_label} - Score: {rec['score']}/100" + _arrow,
                         key=f"btn_rec_{i}", use_container_width=True):
                st.session_state[_key_rec] = not st.session_state[_key_rec]
            if st.session_state[_key_rec]:
                c1, c2 = st.columns([2, 1])
                with c1:
                    st.markdown(f"""
                    <div style='padding:12px; background:{color_r}18; border-left:4px solid {color_r}; border-radius:5px;'>
                        <h2 style='color:{color_r}; margin:0;'>{rec['score']}/100</h2>
                        <p style='margin:4px 0;'>{rec['descripcion']}</p>
                    </div>""", unsafe_allow_html=True)
                    if i == 0: st.success(f"⭐ {t['mejor_opcion']}")
                    st.markdown(f"**Competidores encontrados:** {rec['num_competidores']}")
                with c2:
                    st.markdown(f"**{t['inversion_estimada']}:**")
                    st.markdown(f"${rec['inversion_min']//1000}K - ${rec['inversion_max']//1000}K MXN")

        if len(recomendaciones) > top_n:
            st.caption(f"🔒 Actualiza para ver los {len(recomendaciones) - top_n} tipos de negocio restantes")

        st.markdown("---")
        st.markdown(f"### 🤖 {t['resumen']}")
        st.markdown(analisis)

        # Negocios a evitar
        st.markdown(f"### {t['evitar']}")
        evitar_data = [{"Tipo": r['nombre'], "Score": f"{r['score']}/100"} for r in recomendaciones[-3:]]
        st.dataframe(evitar_data, use_container_width=True)

    st.markdown("---")

    # Competidores
    if competidores:
        n_mostrados = tier_info["competidores_mostrados"]
        st.markdown(f"### {t['competencia_titulo']}")
        comp_data = []
        for comp in competidores[:n_mostrados]:
            comp_data.append({
                t["tabla_nombre"]:  comp.get('displayName', {}).get('text', 'N/A'),
                t["tabla_rating"]:  f"{comp.get('rating','N/A')}★" if comp.get('rating') else 'N/A',
                t["tabla_reseñas"]: comp.get('userRatingCount', 0)
            })
        st.dataframe(comp_data, use_container_width=True)
        if len(competidores) > n_mostrados:
            st.caption(f"🔒 Mostrando {n_mostrados} de {len(competidores)} — actualiza tu plan para ver todos")

    st.markdown("---")

    # Demografía — free básico, Basic+ completo
    st.markdown(f"### {t['demografía']}")
    st.caption(t["datos_estimados"])
    densidad_info = formatear_densidad(demografia['densidad_hab_km2'])

    if tier_info["muestra_demografia"]:
        # Basic+: datos completos con INEGI
        fuente = demografia.get('fuente', 'Estimación por zona')
        st.caption(f"📊 Fuente: {fuente}")

        c1, c2, c3, c4 = st.columns(4)
        pob_actual = demografia.get('poblacion_actual', demografia.get('poblacion_estimada', 0))
        pob_2020   = demografia.get('poblacion_2020', 0)
        delta_pob  = f"+{pob_actual - pob_2020:,}" if pob_2020 else None

        c1.metric("👥 Población " + str(datetime.datetime.now().year),
                  f"{pob_actual:,} hab", delta_pob)
        c2.metric("🏠 Viviendas habitadas",
                  f"{demografia.get('viviendas_actual', demografia.get('viviendas_habitadas',0)):,}")
        c3.metric("🏙️ Densidad",
                  densidad_info.get('nivel', densidad_info.get('clasificacion','N/A')))
        c4.metric("📐 Aprox/manzana",
                  f"~{demografia.get('personas_manzana', densidad_info.get('personas_manzana',0))} personas")

        st.caption(densidad_info.get('descripcion',''))

        # NSE + ingreso + gasto — tarjetas visuales
        st.markdown("---")
        st.markdown("**💰 Perfil Económico del Área**")
        e1, e2, e3, e4 = st.columns(4)
        nse = demografia.get('nse_predominante','C')
        nse_colores = {"A":"#7B1FA2","A/B":"#9C27B0","B":"#1976D2","B/C+":"#0097A7",
                       "C+":"#2E7D32","C":"#558B2F","C/D+":"#F57F17","D+":"#E65100","D/E":"#B71C1C"}
        nse_color = nse_colores.get(nse, "#555555")
        e1.markdown(f"""<div style='background:{nse_color}18; border-left:3px solid {nse_color};
            padding:10px; border-radius:6px;'>
            <b style='color:{nse_color}; font-size:22px;'>{nse}</b><br>
            <span style='font-size:11px; color:#666;'>NSE Predominante</span></div>""",
            unsafe_allow_html=True)
        e2.metric("💵 Ingreso prom/mes",
                  f"${demografia.get('ingreso_actual', demografia.get('ingreso_promedio_mensual',0)):,}")
        e3.metric("🛒 Gasto prom/mes",
                  f"${demografia.get('gasto_actual', demografia.get('gasto_promedio_mensual',0)):,}")
        tasa = demografia.get('tasa_crecimiento_pct', 0.55)
        e4.metric("📈 Crec. anual zona", f"{tasa:.2f}%")

        if tier_key in ["pro", "premium"]:
            fuente_str = demografia.get('fuente', 'INEGI Censo 2020 · Proyección CONAPO') if demografia else ''
            if fuente_str:
                st.caption(f"📊 Fuente: {fuente_str}")

            # ── GRÁFICA EDAD Y GÉNERO ──
            st.markdown("---")
            st.markdown("**👥 Distribución de edad y género del entorno**")
            st.caption("Fuente: INEGI Censo 2020 · ENIGH 2022")
            if GRAFICAS_OK and "grafica_edad_genero" in dir():
                try:
                    import matplotlib.pyplot as plt
                    _fig_eg = grafica_edad_genero(
                        demografia,
                        TIPOS_NEGOCIO.get(tipo_para_pro if 'tipo_para_pro' in dir() else tipo_negocio_seleccionado or '', {}).get('nombre','')
                    )
                    st.image(_fig_eg, use_container_width=True)
                except Exception as _e:
                    st.caption(f"Gráfica no disponible: {_e}")

            # ── INTERPRETACIÓN DEMOGRÁFICA ──
            if MODULOS_OK and "interpretar_demografia_negocio" in dir():
                try:
                    _tipo_nom = TIPOS_NEGOCIO.get(
                        tipo_para_pro if 'tipo_para_pro' in dir() and tipo_para_pro else tipo_negocio_seleccionado or '', {}
                    ).get('nombre','')
                    _interp = interpretar_demografia_negocio(
                        demografia, _tipo_nom,
                        tipo_para_pro if 'tipo_para_pro' in dir() and tipo_para_pro else tipo_negocio_seleccionado or ''
                    )
                    for _p in _interp.split("\n\n"):
                        if _p.strip():
                            import re as _re
                            _p2 = _re.sub(r'\*\*(.+?)\*\*', r'**\1**', _p.strip())
                            st.markdown(_p2)
                except Exception:
                    pass

            # ── TICKET PROMEDIO COMPETENCIA (PRO) ──
            st.markdown("---")
            st.markdown("### 🎫 Ticket Promedio de la Competencia")

            # Determinar tipo activo para el ticket
            _tipo_tk_activo = (tipo_para_pro if 'tipo_para_pro' in dir() and tipo_para_pro
                               else tipo_negocio_seleccionado or 'default')

            # Caption base — más honesta
            _ticket_override = _get_ticket_override(_tipo_tk_activo)
            if _ticket_override:
                st.caption(
                    f"Referencia de mercado · {_ticket_override['rango_txt']} · "
                    f"Google Places price_level no es confiable para este tipo de negocio"
                )
            else:
                st.caption("Fuente: price_level Google Places + rangos ENIGH 2022 por categoría")

            if MODULOS_OK and "calcular_ticket_competencia" in dir() and competidores:
                try:
                    _tk = calcular_ticket_competencia(
                        competidores,
                        _tipo_tk_activo,
                        score
                    )

                    # Override de ticket para tipos con datos no confiables
                    if _ticket_override:
                        _tk["ticket_competencia"] = _ticket_override["ticket_tipico"]
                        _tk["rango_zona"]          = _ticket_override["rango_txt"]
                        _tk["ticket_recomendado"]  = int(_ticket_override["ticket_tipico"] * 1.1)

                    tk1, tk2, tk3 = st.columns(3)
                    tk1.metric("💰 Ticket referencial zona", f"${_tk['ticket_competencia']:,} MXN")
                    tk2.metric("📊 Nivel de precio zona", _tk['rango_zona'][:30])
                    if tier_key == "premium":
                        tk3.metric("✅ Ticket recomendado", f"${_tk['ticket_recomendado']:,} MXN",
                                   help="Calculado según el score de tu ubicación")
                    else:
                        tk3.metric("🔒 Ticket recomendado", "Solo PREMIUM",
                                   help="Actualiza a PREMIUM para ver el ticket recomendado para tu negocio")

                    # Nota de confiabilidad — amarilla para tipos referenciales, gris para confiables
                    _nota_tk = _get_ticket_nota(_tipo_tk_activo)
                    if _nota_tk:
                        st.warning(_nota_tk)
                        if _ticket_override.get("ejemplos"):
                            st.caption(f"📋 Rangos de referencia MX: {_ticket_override['ejemplos']}")
                    else:
                        st.info(_tk['interpretacion_pro'])

                    _key_tk = "exp_ticket_open"
                    if _key_tk not in st.session_state: st.session_state[_key_tk] = False
                    if not _ticket_override:  # Solo mostrar detalle por competidor si los datos son confiables
                        if st.button("Ver detalle por competidor" + (" ▲" if st.session_state[_key_tk] else " ▼"), key="btn_ticket_exp", use_container_width=False):
                            st.session_state[_key_tk] = not st.session_state[_key_tk]
                        if st.session_state[_key_tk]:
                            for _d in _tk.get('detalle_competidores', [])[:8]:
                                st.markdown(f"**{_d['nombre'][:35]}** — Nivel: {_d['nivel']} → Ticket est. ${_d['ticket_estimado']:,} MXN")
                except Exception as _te:
                    st.caption(f"Ticket no disponible: {_te}")
    else:
        # Free: solo básico
        c1, c2, c3, c4 = st.columns(4)
        c1.metric("Población",     f"{demografia.get('poblacion_estimada',0):,} hab")
        c2.metric("Viviendas",     f"{demografia.get('viviendas_habitadas',0):,}")
        c3.metric("Densidad",      densidad_info.get('nivel', densidad_info.get('clasificacion','N/A')))
        c4.metric("Aprox/manzana", f"~{demografia.get('personas_manzana', densidad_info.get('personas_manzana',0))} personas")
        st.caption(densidad_info.get('descripcion',''))
        st.caption("🔒 NSE, ingreso promedio y análisis completo disponibles en plan Básico ($99 MXN)")

    # ── PRO: Tráfico + Mercado + Forecast ──────────────────────────
    # Se muestra tanto en "validar" como en "recomendar" (usando el #1 recomendado)
    # tipo_para_pro ya fue definido arriba al inicio de la sección de resultados
    if tier_key in ["pro", "premium"] and MODULOS_OK and tipo_para_pro:
        st.markdown("---")

        tipo_zona = contexto.get("tipo_zona", "mixto")
        densidad_a = demografia.get("densidad_actual", demografia.get("densidad_hab_km2", 8000))

        # Mostrar qué negocio se usa como base si es modo recomendar
        if modo == "recomendar" and recomendaciones:
            top1 = recomendaciones[0]
            st.info(f"📊 Análisis PRO basado en tu negocio con mayor oportunidad: **{top1['nombre']}** (Score: {top1['score']}/100)")

        # ── TRÁFICO ──
        st.markdown("### 🚗 Análisis de Tráfico Estimado")
        trafico_data = generar_reporte_trafico(tipo_zona, tipo_para_pro, densidad_a)
        # Formato plano de tráfico para PPTX
        _trafico_flat = {
            "nivel_general": trafico_data.get("nivel_general",""),
            "dia_pico": trafico_data.get("dia_pico",""),
            "dia_bajo": trafico_data.get("dia_bajo",""),
            "horas_pico": trafico_data.get("horas_pico",[]),
            "trafico_horario": trafico_data.get("trafico_horario",[]),
        }
        st.session_state["_trafico_data_pptx"] = _trafico_flat

        t1, t2, t3 = st.columns(3)
        t1.metric("📊 Flujo general",   trafico_data["nivel_general"])
        t2.metric("📅 Mejor día",        trafico_data["dia_pico"])
        t3.metric("📅 Día más bajo",     trafico_data["dia_bajo"])

        st.markdown("**🕐 Horas pico para tu tipo de negocio:**")
        for pico in trafico_data["horas_pico"]:
            color_p = "#4CAF50" if pico["flujo"] >= 80 else "#FFC107" if pico["flujo"] >= 60 else "#9E9E9E"
            st.markdown(f"""<div style='display:inline-block; margin:4px; padding:8px 16px;
                background:{color_p}22; border:1px solid {color_p}; border-radius:8px; font-size:13px;'>
                <b>{pico["hora_str"]}</b> — {pico["nivel"]} ({pico["flujo"]}%)
                </div>""", unsafe_allow_html=True)
        st.markdown("")

        # Gráfica horaria con barras ASCII → visual con HTML
        st.markdown("**📈 Perfil de flujo horario (0h – 23h):**")
        horas_html = "<div style='display:flex; align-items:flex-end; gap:2px; height:80px; margin:8px 0;'>"
        for hora, val in enumerate(trafico_data["trafico_horario"]):
            color_b = "#0D0D0D" if val >= 70 else "#BB944F" if val >= 45 else "#E0E0E0"
            tooltip = f"{hora:02d}h: {val}%"
            horas_html += f"""<div title='{tooltip}' style='
                flex:1; background:{color_b}; height:{val}%;
                border-radius:2px 2px 0 0; min-height:3px;'></div>"""
        horas_html += "</div>"
        horas_html += "<div style='display:flex; justify-content:space-between; font-size:10px; color:#888;'><span>0h</span><span>6h</span><span>12h</span><span>18h</span><span>23h</span></div>"
        st.markdown(horas_html, unsafe_allow_html=True)

        # Tráfico semanal
        st.markdown("**📅 Flujo por día de la semana:**")
        dias_html = "<div style='display:flex; gap:6px; margin:8px 0;'>"
        for i, (dia, val) in enumerate(zip(DIAS_SEMANA, trafico_data["trafico_semanal"])):
            color_d = "#0D0D0D" if val >= 85 else "#BB944F" if val >= 65 else "#E0E0E0"
            txt_color = "white" if val >= 65 else "#555"
            dias_html += f"""<div style='flex:1; background:{color_d}; color:{txt_color};
                text-align:center; padding:8px 4px; border-radius:6px; font-size:11px; font-weight:bold;'>
                {dia[:3]}<br><span style='font-size:13px;'>{val}%</span></div>"""
        dias_html += "</div>"
        st.markdown(dias_html, unsafe_allow_html=True)
        st.caption("⚠️ Estimación basada en perfil de zona + tipo de negocio. No refleja datos de tráfico en tiempo real.")

        st.markdown("---")

        # ── MERCADO POTENCIAL ──
        st.markdown("### 📊 Tamaño de Mercado Potencial")
        mercado = calcular_mercado_potencial(demografia, tipo_para_pro, len(competidores))
        _mercado_flat = {
            "mercado_total": mercado.get("mercado_total_mensual",0),
            "mercado_capturable": mercado.get("mercado_captura_mensual",0),
            "clientes_potenciales_dia": mercado.get("clientes_dia_estimados",0),
            "penetracion_estimada_pct": mercado.get("factor_captura_pct",15),
        }
        st.session_state["_mercado_data_pptx"] = _mercado_flat

        m1, m2, m3 = st.columns(3)
        m1.metric("🏪 Mercado total área/mes", f"${mercado['mercado_total_mensual']:,}")
        m2.metric(f"🎯 Tu captura estimada ({mercado['factor_captura_pct']:.0f}%)",
                  f"${mercado['mercado_captura_mensual']:,}/mes")
        m3.metric("👤 Clientes/día estimados", str(mercado["clientes_dia_estimados"]))
        st.caption(f"📐 Metodología: {mercado['metodologia']}")

        st.markdown("---")

        # ── FORECAST ──
        tipo_info_fc = TIPOS_NEGOCIO[tipo_para_pro]
        forecast = generar_forecast(mercado, tipo_para_pro, score,
                                    tipo_info_fc["inversion_min"])
        roi_data = calcular_roi(forecast, tipo_info_fc["inversion_min"],
                                tipo_info_fc["inversion_max"], tipo_para_pro)
        # Convertir estructura anidada a formato plano para el PPTX
        _forecast_flat = {
            "pesimista": forecast["escenarios"]["pesimista"]["promedio_mensual"],
            "realista":  forecast["escenarios"]["base"]["promedio_mensual"],
            "optimista": forecast["escenarios"]["optimista"]["promedio_mensual"],
            "nota": ("; ".join(forecast.get("supuestos", [])[:2])),
        }
        _roi_flat = {
            "inversion_min": roi_data.get("inversion_min", 0),
            "inversion_max": roi_data.get("inversion_max", 0),
            "retorno_mensual": roi_data.get("utilidad_mensual_est", 0),
            "meses_recuperacion_realista": roi_data.get("meses_recuperacion", "—"),
            "roi_anual_pct": roi_data.get("roi_12m_pct", "—"),
        }
        st.session_state["_forecast_data_pptx"] = _forecast_flat
        st.session_state["_roi_data_pptx"] = _roi_flat

        st.markdown("### 📈 Forecast de Ventas — 12 Meses (3 Escenarios)")
        f1, f2, f3 = st.columns(3)

        for col, escenario, color_e, emoji_e in [
            (f1, "pesimista", "#F44336", "🔴"),
            (f2, "base",      "#BB944F", "🟡"),
            (f3, "optimista", "#4CAF50", "🟢"),
        ]:
            datos_e = forecast["escenarios"][escenario]
            col.markdown(f"""<div style='background:{color_e}11; border:1px solid {color_e};
                border-radius:10px; padding:14px; text-align:center;'>
                <div style='font-size:13px; color:{color_e}; font-weight:bold;'>{emoji_e} {escenario.upper()}</div>
                <div style='font-size:22px; font-weight:bold; color:{color_e}; margin:6px 0;'>
                    ${datos_e["total_anual"]:,}</div>
                <div style='font-size:11px; color:#666;'>año 1 total</div>
                <div style='font-size:13px; margin-top:6px;'>
                    ${datos_e["promedio_mensual"]:,}/mes prom.</div>
                </div>""", unsafe_allow_html=True)

        # Gráfica de líneas 12 meses con matplotlib (colores 4SITE)
        st.markdown("**Ventas mensuales — 3 escenarios:**")
        if GRAFICAS_OK:
            try:
                _buf_fc = grafica_forecast(forecast)
                st.image(_buf_fc, use_container_width=True)
            except Exception as _efc:
                # Fallback: barras HTML
                _meses_fc = forecast["escenarios"]["base"]["ventas_mensuales"]
                _max_fc = max(_meses_fc) or 1
                _html_fc = "<div style='display:flex;align-items:flex-end;gap:3px;height:80px;margin:8px 0;'>"
                for _i, _v in enumerate(_meses_fc):
                    _pct = int(_v / _max_fc * 100)
                    _html_fc += (f"<div style='flex:1;display:flex;flex-direction:column;align-items:center;'>"
                                 f"<div style='background:#BB944F;width:100%;height:{_pct}%;border-radius:3px 3px 0 0;min-height:3px;'></div>"
                                 f"<div style='font-size:9px;color:#888;margin-top:2px;'>M{_i+1}</div></div>")
                _html_fc += "</div>"
                st.markdown(_html_fc, unsafe_allow_html=True)

        _key_sup = "exp_supuestos_open"
        if _key_sup not in st.session_state: st.session_state[_key_sup] = False
        if st.button("Supuestos del forecast" + (" ▲" if st.session_state[_key_sup] else " ▼"), key="btn_sup_exp"):
            st.session_state[_key_sup] = not st.session_state[_key_sup]
        if st.session_state[_key_sup]:
            for sup in forecast["supuestos"]:
                st.caption(f"• {sup}")

        st.markdown("---")

        # ── ROI ── (PREMIUM)
        if tier_key == "premium":
            st.markdown("### 💰 ROI + Punto de Equilibrio")
            r1, r2, r3, r4 = st.columns(4)
            r1.metric("📊 ROI estimado 12m", f"{roi_data['roi_12m_pct']}%")
            r2.metric("⏱️ Recuperación est.", f"{roi_data['meses_recuperacion']} meses")
            r3.metric("💵 Utilidad mensual est.", f"${roi_data['utilidad_mensual_est']:,}")
            r4.metric("⚖️ Punto de equilibrio", f"${roi_data['punto_eq_ventas_mes']:,}/mes")
            st.markdown(f"**Clasificación:** {roi_data['clasificacion_roi']}")
            st.caption("⚠️ Estimaciones con supuestos estándar. Los resultados reales dependen de gestión, producto y mercado.")

    # ── BÚSQUEDA DE ZONA IDEAL — temporalmente bloqueada ──
    _lat_zona = lat if 'lat' in dir() and lat else None
    _lng_zona = lng if 'lng' in dir() and lng else None
    if tier_key == "pro" and _lat_zona and _lng_zona:
        st.markdown("---")
        st.info("🔧 **¿Dónde más podría funcionar este negocio?** — Esta función está en mantenimiento y estará disponible muy pronto.")

    # ── MAPAS AVANZADOS (PRO y PREMIUM) — fuera del bloque de cálculo ──
    if tier_key in ["pro", "premium"]:
        st.markdown("---")
        st.markdown("### 🗺️ Mapas Avanzados")
        # Asegurar que competidores y coordenadas están disponibles
        _comp_mapas = competidores or []
        _tipo_mapas = (tipo_para_pro if 'tipo_para_pro' in dir() and tipo_para_pro
                       else tipo_negocio_seleccionado or
                       (recomendaciones[0]['tipo_key'] if recomendaciones else None))

        if MAPAS_OK and lat and lng:
            tab_comp, tab_heat, tab_iso, tab_cani, tab_vial = st.tabs([
                "📍 Competidores", "🔥 Heatmap", "⏱️ Isócronas", "⚠️ Canibalización", "🛣️ Vialidades"
            ])

            with tab_comp:
                st.caption("Competidores coloreados por rating en tu radio de 500m")
                try:
                    tipo_nombre_m = TIPOS_NEGOCIO.get(_tipo_mapas, {}).get('nombre','') if _tipo_mapas else ''
                    mapa_c, interp_c = crear_mapa_competidores(
                        lat, lng, _comp_mapas, tipo_nombre_m, idioma)
                    render_mapa_con_interpretacion(mapa_c, interp_c, altura=420)
                except Exception as e_mc:
                    st.error(f"Error en mapa de competidores: {e_mc}")

            with tab_heat:
                st.caption("Densidad de competencia ponderada por rating y número de reseñas")
                try:
                    mapa_h, interp_h = crear_heatmap_competencia(lat, lng, _comp_mapas, idioma)
                    render_mapa_con_interpretacion(mapa_h, interp_h, altura=420)
                except Exception as e_mh:
                    st.error(f"Error en heatmap: {e_mh}")

            with tab_iso:
                perfil_iso = st.radio("Modo de desplazamiento:",
                    ["foot-walking", "driving-car"],
                    format_func=lambda x: "🚶 A pie" if x == "foot-walking" else "🚗 En auto",
                    horizontal=True, key="radio_iso")
                if not ORS_API_KEY:
                    st.info("💡 Mostrando isócronas estimadas por distancia. Para isócronas reales "
                            "por red vial agrega `ORS_API_KEY` a tu `.env` (gratis en openrouteservice.org).")
                try:
                    mapa_iso, interp_iso = crear_mapa_isocronas(
                        lat, lng, ORS_API_KEY, [5, 10, 15], perfil_iso, _comp_mapas, idioma)
                    render_mapa_con_interpretacion(mapa_iso, interp_iso, altura=420)
                except Exception as e_iso:
                    st.error(f"Error en isócronas: {e_iso}")

            with tab_cani:
                st.caption("Identifica competidores que comparten tu área de influencia inmediata")
                try:
                    mapa_cani, interp_cani, datos_cani = crear_mapa_canibalizacion(
                        lat, lng, _comp_mapas, 500, idioma)
                    render_mapa_con_interpretacion(mapa_cani, interp_cani, altura=400)
                    if datos_cani["muy_cercanos"]:
                        st.markdown("**⚠️ Competidores críticos (≤150m) — riesgo alto de canibalización:**")
                        crit_data = [{
                            "Nombre": c.get('displayName',{}).get('text','N/A')[:25],
                            "Distancia": f"{c.get('_distancia_m',0)}m",
                            "Rating": c.get('rating','N/A'),
                            "Reseñas": c.get('userRatingCount',0)
                        } for c in datos_cani["muy_cercanos"]]
                        st.dataframe(crit_data, use_container_width=True)
                    elif not datos_cani["muy_cercanos"] and not datos_cani["cercanos"]:
                        st.success("✅ Sin riesgo de canibalización — no hay competidores directos en los primeros 300m")
                except Exception as e_cani:
                    st.error(f"Error en mapa de canibalización: {e_cani}")
            with tab_vial:
                st.caption("Vialidades cercanas clasificadas por nivel de flujo vehicular")
                st.info("🛣️ Clasifica las calles del entorno en primarias, comerciales y locales. Útil para evaluar visibilidad y tráfico de paso.")
                try:
                    mapa_vial, interp_vial = crear_mapa_vialidades(lat, lng, GOOGLE_API_KEY, _comp_mapas, idioma)
                    render_mapa_con_interpretacion(mapa_vial, interp_vial, altura=440, titulo="🛣️ Vialidades y movilidad del entorno")
                except Exception as e_vial:
                    st.error(f"Error en mapa de vialidades: {e_vial}")

        elif not MAPAS_OK:
            st.warning("Instala `streamlit-folium` para ver mapas avanzados: `pip install streamlit-folium folium`")

    # ── COMPARATIVA (PREMIUM en modo validar/recomendar, NO zona_potencial) ──
    if tier_key == "premium" and modo in ["validar", "recomendar"]:
        st.markdown("---")
        st.markdown("### 🗺️ Comparativa de Ubicaciones")
        st.caption("Analiza hasta 3 ubicaciones en paralelo y encuentra la mejor opción.")

        if "ubicaciones_comparativa" not in st.session_state:
            st.session_state.ubicaciones_comparativa = []

        _key_comp2 = "exp_comp2_open"
        if _key_comp2 not in st.session_state: st.session_state[_key_comp2] = False
        if st.button("+ Agregar ubicaciones para comparar" + (" ▲" if st.session_state[_key_comp2] else " ▼"), key="btn_comp2_exp", use_container_width=True):
            st.session_state[_key_comp2] = not st.session_state[_key_comp2]
        if st.session_state[_key_comp2]:
            pass
        if st.session_state.get(_key_comp2, False):
            col_a, col_b = st.columns([3,1])
            with col_a:
                nueva_ub = st.text_input("Dirección adicional a comparar:",
                    placeholder="Ej: Av. Reforma 222, CDMX", key="nueva_ub_input")
            with col_b:
                if st.button("Agregar", key="btn_agregar_ub"):
                    if not nueva_ub:
                        st.warning("Escribe una dirección primero")
                    elif len(st.session_state.ubicaciones_comparativa) >= 2:
                        st.warning("Máximo 2 ubicaciones adicionales (3 en total)")
                    else:
                        with st.spinner(f"Analizando {nueva_ub[:30]}..."):
                            try:
                                lat2, lng2 = geocodificar_direccion(nueva_ub)
                                if not lat2 or not lng2:
                                    st.error("No se pudo geocodificar la dirección")
                                else:
                                    # Tipo de negocio para comparar
                                    tipo_comp = tipo_para_pro if 'tipo_para_pro' in dir() else tipo_negocio_seleccionado
                                    dem2   = obtener_demografia(lat2, lng2, nueva_ub)
                                    comps2 = buscar_competencia_por_tipo(lat2, lng2, tipo_comp) if tipo_comp else []
                                    score2_b, desglose2 = calcular_score_competencia(comps2, tipo_comp) if (tipo_comp and comps2) else (0, {})
                                    ctx2   = detectar_contexto_ubicacion(lat2, lng2, nueva_ub)
                                    score2_c = ajustar_score_por_contexto(score2_b, tipo_comp, ctx2) if tipo_comp else score2_b
                                    score2   = ajustar_score_por_demografia(score2_c, tipo_comp, dem2) if tipo_comp else score2_c
                                    st.session_state.ubicaciones_comparativa.append({
                                        "nombre": nueva_ub[:25],
                                        "score": score2,
                                        "densidad_comp": desglose2.get("densidad", 0),
                                        "calidad_comp":  desglose2.get("calidad", 0),
                                        "consolidacion": desglose2.get("consolidacion", 0),
                                        "poblacion":     dem2.get("poblacion_actual", dem2.get("poblacion_estimada", 0)),
                                        "nse":           dem2.get("nse_predominante", "C"),
                                        "ingreso":       dem2.get("ingreso_actual", dem2.get("ingreso_promedio_mensual", 0)),
                                        "num_competidores": len(comps2),
                                    })
                                    st.success(f"✅ Agregada: {nueva_ub[:30]} — Score: {score2}/100")
                                    st.rerun()
                            except Exception as e_comp:
                                st.error(f"Error al analizar: {e_comp}")

        # Mostrar comparativa si hay ubicaciones
        # Fix: desglose puede ser dict, lista de tuplas, o None
        def _get_desglose_val(dsg, key, default=0):
            if isinstance(dsg, dict): return dsg.get(key, default)
            if isinstance(dsg, (list, tuple)):
                for item in dsg:
                    if isinstance(item, (list,tuple)) and len(item)>=2:
                        if str(item[0]).lower() == key: return item[1]
            return default

        todas_ubs = [{
            "nombre": ubicacion[:25],
            "score": float(score),
            "densidad_comp": _get_desglose_val(desglose, "densidad"),
            "calidad_comp":  _get_desglose_val(desglose, "calidad"),
            "consolidacion": _get_desglose_val(desglose, "consolidacion"),
            "poblacion": (demografia or {}).get("poblacion_actual",
                          (demografia or {}).get("poblacion_estimada",0)),
            "nse":    (demografia or {}).get("nse_predominante","C"),
            "ingreso":(demografia or {}).get("ingreso_actual",
                       (demografia or {}).get("ingreso_promedio_mensual",0)),
            "num_competidores": len(competidores or []),
        }] + st.session_state.ubicaciones_comparativa

        if len(todas_ubs) > 1 and GRAFICAS_OK:
            try:
                graf_comp = grafica_comparativa(todas_ubs)
                if graf_comp:
                    st.image(graf_comp, use_container_width=True,
                             caption="Análisis comparativo — radar y barras por ubicación")
            except Exception as e_graf:
                st.caption(f"Gráfica comparativa no disponible: {e_graf}")

            # Tabla resumen
            comp_table_data = []
            for ub in todas_ubs:
                nivel_c, _ = nivel_score(ub['score'], idioma)
                comp_table_data.append({
                    "Ubicación": ub['nombre'],
                    "Score": f"{ub['score']}/100",
                    "Nivel": nivel_c,
                    "Competidores": ub['num_competidores'],
                    "NSE": ub['nse'],
                    "Ingreso prom": f"${ub['ingreso']:,}",
                })
            st.dataframe(comp_table_data, use_container_width=True)

            # Recomendación automática
            mejor = max(todas_ubs, key=lambda x: x['score'])
            st.success(f"✅ **Recomendación:** La mejor ubicación es **{mejor['nombre']}** con score {mejor['score']}/100")

        if st.button("🗑️ Limpiar comparativa", key="btn_limpiar"):
            st.session_state.ubicaciones_comparativa = []
            st.rerun()

    st.markdown("---")

    # ── DIAGNÓSTICO COMERCIAL AL FINAL (Basic, PRO, PREMIUM) ──────
    if tier_key != "free":
        st.markdown("### ✅ Recomendación Final")

        _trafico_diag = _mercado_diag = _forecast_diag = _roi_diag = None
        if tier_key in ["pro", "premium"] and MODULOS_OK:
            _tipo_diag = tipo_para_pro if 'tipo_para_pro' in dir() and tipo_para_pro else tipo_negocio_seleccionado
            if _tipo_diag:
                try:
                    _tz  = contexto.get("tipo_zona", "mixto")
                    _den = demografia.get("densidad_actual", demografia.get("densidad_hab_km2", 8000))
                    _trafico_diag  = generar_reporte_trafico(_tz, _tipo_diag, _den)
                    _mercado_diag  = calcular_mercado_potencial(demografia, _tipo_diag, len(competidores))
                    _tipo_info_d   = TIPOS_NEGOCIO.get(_tipo_diag, {})
                    _forecast_diag = generar_forecast(_mercado_diag, _tipo_diag, score,
                                                       _tipo_info_d.get("inversion_min", 300000))
                    _roi_diag      = calcular_roi(_forecast_diag,
                                                   _tipo_info_d.get("inversion_min", 300000),
                                                   _tipo_info_d.get("inversion_max", 600000), _tipo_diag)
                except:
                    pass

        with st.spinner("🏆 Generando Recomendación Final..."):
            # Calcular ticket_data para la recomendación final
            _ticket_diag = None
            if MODULOS_OK and competidores:
                try:
                    _tipo_t = tipo_para_pro if 'tipo_para_pro' in dir() and tipo_para_pro else tipo_negocio_seleccionado
                    if _tipo_t:
                        _ticket_diag = calcular_ticket_competencia(competidores, _tipo_t, score)
                except Exception:
                    pass

            diagnostico = generar_recomendacion_final(
                ubicacion=ubicacion, score=score, desglose=desglose,
                competidores=competidores, demografia=demografia,
                contexto=contexto, idioma=idioma, tier_key=tier_key, modo=modo,
                tipo_negocio=tipo_negocio_seleccionado if modo == "validar" else None,
                recomendaciones=recomendaciones,
                trafico_data=_trafico_diag, mercado_data=_mercado_diag,
                forecast_data=_forecast_diag, roi_data=_roi_diag,
                ticket_data=_ticket_diag,
            )

        # Renderizar con estilo destacado
        import re as _re
        def md_to_html(txt):
            # Convierte **texto** a <b>texto</b> y saltos de línea a <br>
            txt = _re.sub(r'\*\*(.+?)\*\*', r'<b>\1</b>', txt)
            txt = txt.replace('\n', '<br>')
            return txt

        st.session_state["_diagnostico_pptx"] = diagnostico
        st.markdown(f"""
        <div style='background:linear-gradient(135deg,#0D0D0D08 0%,#BB944F15 100%);
             border:2px solid #0D0D0D55; border-radius:12px;
             padding:22px 26px; margin:10px 0;'>
        {md_to_html(diagnostico)}
        </div>""", unsafe_allow_html=True)

        # Mensaje de apoyo — siempre presente
        st.markdown("""
        <div style='background:#0D0D0D; border-radius:12px; padding:18px 22px;
             color:white; margin-top:12px;'>
        <p style='font-size:12px; font-weight:700; color:#BB944F; margin:0 0 10px;'>
            💡 La ubicación es el punto de partida — el éxito depende también de:</p>
        <div style='display:flex; gap:12px; flex-wrap:wrap;'>
            <div style='flex:1; min-width:160px; background:#1A1A1A; border-radius:8px; padding:10px 12px;'>
                <div style='color:#BB944F; font-weight:700; font-size:11px;'>🎫 Ticket promedio</div>
                <div style='color:#AAAAAA; font-size:10px; margin-top:4px;'>Calibra tu precio al NSE y competencia de la zona.</div>
            </div>
            <div style='flex:1; min-width:160px; background:#1A1A1A; border-radius:8px; padding:10px 12px;'>
                <div style='color:#BB944F; font-weight:700; font-size:11px;'>🤝 Atención al cliente</div>
                <div style='color:#AAAAAA; font-size:10px; margin-top:4px;'>Tu rating en Google es tu segundo score de viabilidad.</div>
            </div>
            <div style='flex:1; min-width:160px; background:#1A1A1A; border-radius:8px; padding:10px 12px;'>
                <div style='color:#BB944F; font-weight:700; font-size:11px;'>📱 Presencia en redes</div>
                <div style='color:#AAAAAA; font-size:10px; margin-top:4px;'>Google Maps, Instagram y WhatsApp Business son esenciales para convertir tráfico en clientes.</div>
            </div>
        </div>
        </div>""", unsafe_allow_html=True)

        st.markdown("---")

    # ── DESCARGA REPORTE (PowerPoint) ──
    st.markdown("---")
    st.markdown("#### 📊 Descargar Reporte")
    if PPTX_OK:
        try:
            # Obtener nombre del tipo de negocio
            _tipo_nombre_dl = ""
            if tipo_negocio_seleccionado and tipo_negocio_seleccionado in TIPOS_NEGOCIO:
                _tipo_nombre_dl = TIPOS_NEGOCIO[tipo_negocio_seleccionado].get("nombre","")
            elif recomendaciones:
                _tipo_nombre_dl = recomendaciones[0].get("nombre","") if recomendaciones else ""

            # Calcular ticket_data si no existe
            _ticket_dl = None
            _tipo_dl = tipo_negocio_seleccionado or (
                tipo_para_pro if "tipo_para_pro" in dir() and tipo_para_pro else None)
            if MODULOS_OK and competidores and _tipo_dl:
                try: _ticket_dl = calcular_ticket_competencia(competidores, _tipo_dl, score)
                except: pass

            _analisis_dl = (st.session_state.get("_diagnostico_pptx") or
                            diagnostico if 'diagnostico' in dir() and diagnostico else analisis)

            # Sanitizar recomendaciones — asegurar tipos numéricos correctos
            # para evitar 'int vs str' en el generador PPTX
            _recom_dl = []
            for _r in (recomendaciones or []):
                _recom_dl.append({
                    **_r,
                    "score":        float(_r.get("score") or 0),
                    "inversion_min": int(_r.get("inversion_min") or 0),
                    "inversion_max": int(_r.get("inversion_max") or 0),
                    "num_competidores": int(_r.get("num_competidores") or 0),
                })

            # Sanitizar desglose — asegurar que todos los valores numéricos son float/int
            # rating_promedio y otros valores pueden llegar como str desde calcular_score
            _desglose_dl = {}
            for _k, _v in (desglose or {}).items():
                try:
                    _desglose_dl[_k] = float(_v) if isinstance(_v, str) else _v
                except (ValueError, TypeError):
                    _desglose_dl[_k] = _v

            pptx_buf = _generar_pptx_por_tier(
                tier_key,
                ubicacion=ubicacion, score=float(score or 0), desglose=_desglose_dl,
                analisis=_analisis_dl, competidores=competidores,
                idioma=idioma, lat=lat, lng=lng, modo=modo,
                tipo_negocio=_tipo_nombre_dl or tipo_negocio_seleccionado or "",
                recomendaciones=_recom_dl,
                demografia=demografia, contexto=contexto,
                ticket_data=_ticket_dl,
                trafico_data=st.session_state.get("_trafico_data_pptx"),
                mercado_data=st.session_state.get("_mercado_data_pptx"),
                forecast_data=st.session_state.get("_forecast_data_pptx"),
                roi_data=st.session_state.get("_roi_data_pptx"),
            )

            if pptx_buf:
                _dir_clean = ubicacion[:25].replace(" ","_").replace(",","").replace(".","")
                _fname = f"4site_{tier_key}_{_dir_clean}.pptx"
                tier_nombres = {"free":"Gratis","basico":"Basico $99","pro":"PRO ($499)",
                                "premium":"PREMIUM $999","premium2":"ZONA ELITE $1,500"}
                st.download_button(
                    label=f"⬇️  Descargar presentación PowerPoint",
                    data=pptx_buf,
                    file_name=_fname,
                    mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                    use_container_width=True,
                    type="primary",
                    key="btn_descarga_pptx_main"
                )
                _slides = {"free":"2 slides","basico":"5 slides","pro":"7 slides",
                           "premium":"10 slides","premium2":"4 slides"}
                st.caption(f"📊 {tier_nombres.get(tier_key,'')} · {_slides.get(tier_key,'')} · "
                           f"Editable en PowerPoint · {ubicacion[:40]}...")
            else:
                st.error("No se pudo generar el reporte. Intenta de nuevo.")
        except Exception as e_pptx:
            st.error(f"Error al generar reporte: {e_pptx}")
            _key_err = "exp_err_open"
            if _key_err not in st.session_state: st.session_state[_key_err] = False
            if st.button("Ver detalle del error", key="btn_err_exp"):
                st.session_state[_key_err] = not st.session_state[_key_err]
            if st.session_state[_key_err]:
                import traceback
                st.code(traceback.format_exc())
    else:
        st.warning("Módulo PPTX no disponible. Instala: `pip install python-pptx`")

# ── CTA BOTTOM ──────────────────────────────
STRIPE_BASICO  = "https://buy.stripe.com/4gM28r66M5sD3Gqfmk5Rm02"
STRIPE_PRO     = "https://buy.stripe.com/9B6eVdcvadZ990K7TS5Rm04"
STRIPE_PREMIUM = "https://buy.stripe.com/3cI28r66M6wHgtca205Rm00"

st.markdown("---")

_cta_titulo = t.get('cta_titulo', '🚀 ¿Quieres el análisis completo?')
_cta_intro  = t.get('cta_intro',  'Cada análisis incluye un reporte PowerPoint descargable y editable.')

import streamlit.components.v1 as _components
_components.html(f"""
<!DOCTYPE html>
<html>
<head>
<link href="https://fonts.googleapis.com/css2?family=Montserrat:wght@400;500;600;700&display=swap" rel="stylesheet">
<style>
  * {{ box-sizing: border-box; margin: 0; padding: 0; }}
  body {{ background: transparent; font-family: 'Montserrat', sans-serif; padding: 8px 0; }}
  .wrap {{ text-align: center; padding: 36px 32px 28px; background: #0D0D0D; border-radius: 16px; color: white; }}
  .titulo {{ font-size: 22px; font-weight: 700; margin-bottom: 6px; }}
  .intro  {{ font-size: 12px; color: #888888; margin-bottom: 24px; }}
  .cards  {{ display: flex; justify-content: center; gap: 14px; flex-wrap: wrap; }}

  .card {{ background: #1A1A1A; border-radius: 12px; padding: 20px 18px;
           width: 100%; max-width: 215px; min-width: 180px;
           text-align: left; display: flex; flex-direction: column; justify-content: space-between; }}
  .card-basico  {{ border: 1px solid #333; }}
  .card-pro     {{ border: 2px solid #BB944F; position: relative; }}
  .card-premium {{ border: 1px solid #D4AF37; }}

  .badge-popular {{ position: absolute; top: -11px; left: 50%; transform: translateX(-50%);
    background: linear-gradient(135deg,#BB944F,#D4A85A); color: #0D0D0D;
    font-size: 9px; font-weight: 700; letter-spacing: .08em;
    padding: 3px 10px; border-radius: 20px; white-space: nowrap; }}

  .tag {{ font-size: 10px; font-weight: 700; letter-spacing: .1em; margin-bottom: 8px; text-transform: uppercase; }}
  .precio {{ font-size: 22px; font-weight: 700; margin-bottom: 4px; }}
  .precio span {{ font-size: 11px; color: #888; font-weight: 400; }}
  .usos {{ font-size: 10px; color: #666; margin-bottom: 12px; }}

  .features {{ font-size: 10.5px; color: #BBBBBB; line-height: 1.9; margin-bottom: 16px; }}
  .features li {{ list-style: none; padding-left: 0; }}
  .features li::before {{ content: "✓ "; color: #BB944F; font-weight: 700; }}
  .features li.dim {{ color: #555; }}
  .features li.dim::before {{ content: "— "; color: #444; }}

  .btn {{ display: block; text-align: center; font-size: 11px; font-weight: 700; padding: 9px 12px;
    border-radius: 8px; text-decoration: none; letter-spacing: .05em; cursor: pointer; transition: all .2s; }}
  .btn-basico  {{ background: #2A2A2A; color: #FFFFFF; }}
  .btn-basico:hover {{ background: #3A3A3A; }}
  .btn-pro     {{ background: linear-gradient(135deg,#BB944F,#D4A85A); color: #0D0D0D; }}
  .btn-pro:hover {{ background: linear-gradient(135deg,#D4A85A,#E8BC6A); }}
  .btn-premium {{ background: #1E1A00; color: #D4AF37; border: 1px solid #D4AF37; }}
  .btn-premium:hover {{ background: #2A2500; }}

  .footer {{ margin-top: 22px; font-size: 11px; color: #555; }}
  .footer a {{ color: #BB944F; text-decoration: none; }}

  @media (max-width: 600px) {{
    .wrap {{ padding: 24px 14px 20px; border-radius: 12px; }}
    .titulo {{ font-size: 18px; }}
    .cards {{ flex-direction: column; align-items: stretch; gap: 20px; }}
    .card {{ max-width: 100%; min-width: unset; width: 100%; }}
    .card-pro {{ margin-top: 10px; }}
  }}
</style>
</head>
<body>
<div class="wrap">
  <div class="titulo">{_cta_titulo}</div>
  <div class="intro">{_cta_intro}</div>
  <div class="cards">

    <!-- BÁSICO -->
    <div class="card card-basico">
      <div>
        <div class="tag" style="color:#AAAAAA;">BÁSICO</div>
        <div class="precio" style="color:#FFFFFF;">$99 <span>MXN</span></div>
        <div class="usos">1 análisis incluido</div>
        <ul class="features">
          <li>Score de viabilidad con desglose por factor</li>
          <li>Perfil demográfico de la zona (edad, NSE, hogares)</li>
          <li>Hasta 10 competidores cercanos identificados</li>
          <li>Top 5 tipos de negocio recomendados para ese punto</li>
          <li>Reporte PowerPoint de 5 slides descargable</li>
          <li class="dim">Mapa interactivo de competencia</li>
          <li class="dim">Forecast de ventas</li>
          <li class="dim">Análisis de mercado potencial</li>
        </ul>
      </div>
      <a class="btn btn-basico" href="{STRIPE_BASICO}" target="_blank"
         onclick="fetch('https://graph.facebook.com/v18.0/{META_PIXEL_ID}/events?access_token={META_CAPI_TOKEN}',{{method:'POST',headers:{{'Content-Type':'application/json'}},body:JSON.stringify({{test_event_code:'{META_TEST_CODE}',data:[{{event_name:'InitiateCheckout',event_time:Math.floor(Date.now()/1000),action_source:'website',custom_data:{{value:99,currency:'MXN',content_name:'Basico'}}}}]}})}})">Comprar Básico →</a>
    </div>

    <!-- PRO -->
    <div class="card card-pro">
      <div class="badge-popular">⭐ MÁS POPULAR</div>
      <div>
        <div class="tag" style="color:#BB944F;">PRO</div>
        <div class="precio" style="color:#BB944F;">$499 <span>MXN</span></div>
        <div class="usos">2 análisis incluidos</div>
        <ul class="features">
          <li>Todo lo de Básico</li>
          <li>Mapa interactivo con pins de cada competidor</li>
          <li>Ticket promedio estimado de la competencia cercana</li>
          <li>Tamaño de mercado potencial en la zona ($MXN/mes)</li>
          <li>Forecast de ventas proyectado a 12 meses</li>
          <li>Patrones de tráfico por hora y día de la semana</li>
          <li>Hasta 20 competidores analizados</li>
          <li>Reporte PowerPoint de 7 slides descargable</li>
          <li class="dim">ROI y payback estimado</li>
          <li class="dim">Análisis comparativo multi-ubicación</li>
        </ul>
      </div>
      <a class="btn btn-pro" href="{STRIPE_PRO}" target="_blank"
         onclick="fetch('https://graph.facebook.com/v18.0/{META_PIXEL_ID}/events?access_token={META_CAPI_TOKEN}',{{method:'POST',headers:{{'Content-Type':'application/json'}},body:JSON.stringify({{test_event_code:'{META_TEST_CODE}',data:[{{event_name:'InitiateCheckout',event_time:Math.floor(Date.now()/1000),action_source:'website',custom_data:{{value:499,currency:'MXN',content_name:'PRO'}}}}]}})}})">Comprar PRO →</a>
    </div>

    <!-- PREMIUM -->
    <div class="card card-premium">
      <div>
        <div class="tag" style="color:#D4AF37;">PREMIUM 💎</div>
        <div class="precio" style="color:#D4AF37;">$999 <span>MXN</span></div>
        <div class="usos">3 análisis incluidos</div>
        <ul class="features">
          <li>Todo lo de PRO</li>
          <li>ROI estimado a 12 meses en %</li>
          <li>Meses para recuperar tu inversión inicial</li>
          <li>Punto de equilibrio: cuánto necesitas vender al mes para no perder</li>
          <li>Ticket recomendado para tu negocio específico</li>
          <li>Comparativa radar entre hasta 3 ubicaciones candidatas</li>
          <li>Mapa de canibalización con competidores críticos ≤150m</li>
          <li>Reporte PowerPoint de 10 slides con gráficas ejecutivas</li>
        </ul>
      </div>
      <a class="btn btn-premium" href="{STRIPE_PREMIUM}" target="_blank"
         onclick="fetch('https://graph.facebook.com/v18.0/{META_PIXEL_ID}/events?access_token={META_CAPI_TOKEN}',{{method:'POST',headers:{{'Content-Type':'application/json'}},body:JSON.stringify({{test_event_code:'{META_TEST_CODE}',data:[{{event_name:'InitiateCheckout',event_time:Math.floor(Date.now()/1000),action_source:'website',custom_data:{{value:999,currency:'MXN',content_name:'Premium'}}}}]}})}})">Comprar Premium →</a>
    </div>

  </div>
  <div class="footer">3 análisis por código · Reporte incluido · <a href="mailto:hola@4site.mx">hola@4site.mx</a></div>
</div>
</body>
</html>
""", height=1800, scrolling=True)

st.markdown("---")
st.markdown(f"""
<div style='text-align:center; color:#999; font-size:13px;'>
    {t['footer_derechos']}<br>{t['footer_contacto']}
</div>
""", unsafe_allow_html=True)